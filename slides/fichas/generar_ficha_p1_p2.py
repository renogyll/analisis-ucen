# -*- coding: utf-8 -*-
"""
generar_ficha_p1_p2.py
Genera FICHA_Ampliar_P1_P2.pptx en formato dark v3:
  Slide 1 - Portada
  Slide 2 - Cascada P1: 917 -> 915 -> 493 -> sub-cascadas
  Slide 3 - El gran salto: 424 docentes sin dotacion (Nomina vs Dotacion)
  Slide 4 - Sub-cascadas dentro de los 493 (brechas menores)
  Slide 5 - Cascada P2: 917 -> 357 formados -> 615 iniciativas
  Slide 6 - Datos adicionales solicitados para P1 y P2
"""
import sys; sys.stdout.reconfigure(encoding="utf-8")
import os, zipfile
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.patheffects as pe
import pandas as pd
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─────────────────────────────────────────────────────────────────────────────
# Rutas
# ─────────────────────────────────────────────────────────────────────────────
BASE      = os.path.dirname(os.path.abspath(__file__))
PROC918   = os.path.join(BASE, "..", "PROCESADO")
ROOT      = os.path.normpath(os.path.join(BASE, "..", ".."))
CONS      = os.path.join(ROOT, "CONSOLIDADO DOCENTES 3-05-2026")
FONDOTIPO = os.path.join(ROOT, "Fondotipop.pptx")
OUT_DIR   = os.path.join(BASE, "dark_slides_v3")
OUT_PPTX  = os.path.join(ROOT, "FICHA_Ampliar_P1_P2.pptx")
os.makedirs(OUT_DIR, exist_ok=True)

# ─────────────────────────────────────────────────────────────────────────────
# Datos
# ─────────────────────────────────────────────────────────────────────────────
doc  = pd.read_csv(os.path.join(PROC918, "docente_918.csv"), encoding="utf-8-sig")
p2   = pd.read_csv(os.path.join(PROC918, "participacion_p2_918.csv"), encoding="utf-8-sig")

# P1 cascada principal
N_TOTAL        = len(doc)                                           # 917
N_CON_JER      = int(doc["jerarquia"].notna().sum())               # 915
dotacion_mask  = doc["fuente"].isin(["NOMINA_DOTACION","SOLO_DOTACION"])
N_DOTACION     = int(dotacion_mask.sum())                          # 493
N_SOLO_NOMINA  = int((doc["fuente"] == "NOMINA").sum())            # 424
N_HON_SIN_DOT  = int(((doc["fuente"] == "NOMINA") &
                       (doc["tipo_contrato"] == "Honorario")).sum())  # 364
N_JOR_SIN_DOT  = N_SOLO_NOMINA - N_HON_SIN_DOT                    # 60

# Sub-cascadas dentro de los 493 (desde CASCADAS CON N 7-7-2026)
SUBCASC = [
    (493, "Total en dotacion",       "Edad y Sexo · Facultad · Carga Academica"),
    (491, "Con antiguedad",          "Antiguedad"),
    (482, "Con edad y jerarquia",    "Edad x Jerarquia"),
    (478, "Con nivel de formacion",  "Nivel Formacion · Institucion · Pais"),
    (477, "Con grado clasificado",   "GRADOREC · Jerarquia x Nivel Formacion"),
    (433, "Con fechas completas",    "Anos hasta Jerarquizacion"),
]

# P2 cascada
N_FORMADOS   = int(p2["rut_key"].nunique())    # 357
N_INICIATIVAS= len(p2)                          # 615
N_TALL       = int((p2["tipo_formacion"] == "TALLER").sum())     # 376
N_DIP        = int((p2["tipo_formacion"] == "DIPLOMADO").sum())  # 201
N_PROY       = int((p2["tipo_formacion"] == "PROYECTO").sum())   # 38
N_SIN_FORM   = N_TOTAL - N_FORMADOS             # 560

# Archivo fuente
F_DOC  = "CONSOLIDADO DOCENTES 3-05-2026.xlsx"

print(f"P1: {N_TOTAL} total | {N_CON_JER} con jerarquia | {N_DOTACION} dotacion | "
      f"{N_SOLO_NOMINA} solo nomina ({N_HON_SIN_DOT} Hon + {N_JOR_SIN_DOT} Jor)")
print(f"P2: {N_FORMADOS} formados | {N_INICIATIVAS} iniciativas | "
      f"Tall={N_TALL} Dip={N_DIP} Proy={N_PROY}")

# ─────────────────────────────────────────────────────────────────────────────
# Assets
# ─────────────────────────────────────────────────────────────────────────────
BG_PATH   = os.path.join(OUT_DIR, "_fondotipo_bg.jpg")
LOGO_PATH = os.path.join(OUT_DIR, "_fondotipo_logo.png")
for path, zname in [(BG_PATH, "ppt/media/image1.jpg"),
                    (LOGO_PATH, "ppt/media/image2.png")]:
    if not os.path.exists(path):
        with zipfile.ZipFile(FONDOTIPO) as z:
            with open(path, "wb") as f:
                f.write(z.read(zname))

with PILImage.open(BG_PATH) as _im:
    _rgb = _im.convert("RGB"); _iw, _ih = _rgb.size
    _nh  = int(_iw / (16/9)); _y0 = min(int(_ih * 0.12), _ih - _nh)
    bg_arr = np.array(_rgb.crop((0, _y0, _iw, _y0 + _nh)))

with PILImage.open(LOGO_PATH) as _logo:
    logo_arr = np.array(_logo.convert("RGBA")).astype(np.float32) / 255.0

H_GRAD = 600; grad = np.zeros((H_GRAD, 1, 4), dtype=np.float32)
for _r in range(H_GRAD):
    _t = _r / (H_GRAD - 1)
    _stops = [(0.00, (0, 33, 71)), (0.54, (0, 70, 128)), (1.00, (144, 171, 196))]
    for _i in range(len(_stops) - 1):
        _t0, _c0 = _stops[_i]; _t1, _c1 = _stops[_i + 1]
        if _t0 <= _t <= _t1:
            _s = (_t - _t0) / (_t1 - _t0)
            grad[_r, 0] = [(_c0[0]+_s*(_c1[0]-_c0[0]))/255,
                           (_c0[1]+_s*(_c1[1]-_c0[1]))/255,
                           (_c0[2]+_s*(_c1[2]-_c0[2]))/255, 0.82]; break

# ─────────────────────────────────────────────────────────────────────────────
# Layout constants
# ─────────────────────────────────────────────────────────────────────────────
SW, SH         = 13.333, 7.5
SW_EMU, SH_EMU = 12192000, 6858000
PIC_L, PIC_T, PIC_W, PIC_H         = 786581, 1125000, 10599174, 3720000
BUL_L, BUL_T, BUL_W, BUL_H         = 786581, 4870000, 10599174, 1870000
LOGO_L, LOGO_T, LOGO_W, LOGO_H     = 9813773, 656354, 1756626, 697725
TITLE_L, TITLE_T, TITLE_W, TITLE_H = PIC_L, 185000, PIC_W, 710000
POP_L,   POP_T,   POP_W,   POP_H   = PIC_L, 845000, 9000000, 255000

def _ex(e): return e / SW_EMU
def _ey(e): return e / SH_EMU
def _fig_rect(l, t, w, h): return (l, 1-t-h, w, h)
PIC_RECT  = _fig_rect(_ex(PIC_L), _ey(PIC_T), _ex(PIC_W), _ey(PIC_H))
LOGO_RECT = _fig_rect(_ex(LOGO_L), _ey(LOGO_T), _ex(LOGO_W), _ey(LOGO_H))

# ─────────────────────────────────────────────────────────────────────────────
# Helpers matplotlib
# ─────────────────────────────────────────────────────────────────────────────
def _bg_fig():
    fig = plt.figure(figsize=(SW, SH), facecolor="#101820")
    fig.patch.set_facecolor("#101820")
    for z, arr in [(0, bg_arr), (1, grad)]:
        ax = fig.add_axes([0, 0, 1, 1], zorder=z)
        ax.imshow(arr, extent=[0, 1, 0, 1], aspect="auto", origin="upper")
        ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis("off")
    al = fig.add_axes([LOGO_RECT[0], LOGO_RECT[1], LOGO_RECT[2], LOGO_RECT[3]],
                      zorder=10, facecolor="none")
    al.imshow(logo_arr, aspect="auto"); al.axis("off"); al.patch.set_visible(False)
    return fig

def _tr_fig():
    fig = plt.figure(figsize=(SW, SH), facecolor="none")
    fig.patch.set_facecolor("none"); return fig

SHARED_BG = os.path.join(OUT_DIR, "_background.png")

def _ensure_bg():
    if not os.path.exists(SHARED_BG):
        fig = _bg_fig()
        plt.savefig(SHARED_BG, dpi=150, facecolor=fig.get_facecolor()); plt.close()
        print("  bg generado")

def _save_ch(fig, name):
    path = os.path.join(OUT_DIR, name)
    plt.savefig(path, dpi=150, facecolor="none", transparent=True)
    plt.close(); return path

# ─────────────────────────────────────────────────────────────────────────────
# Helpers python-pptx
# ─────────────────────────────────────────────────────────────────────────────
def _new_sl(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _pic(sl, path, prs):
    sl.shapes.add_picture(path, Emu(0), Emu(0), prs.slide_width, prs.slide_height)

def _txt(sl, text, left, top, width, height,
         fs=12, bold=False, italic=False, color="#FFFFFF",
         align=PP_ALIGN.LEFT, wrap=True, lspc=0):
    txb = sl.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf  = txb.text_frame; tf.word_wrap = wrap
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if lspc > 0 and i > 0: p.space_before = Pt(lspc)
        run = p.add_run(); run.text = line
        run.font.size = Pt(fs); run.font.bold = bold; run.font.italic = italic
        r, g, b = int(color[1:3],16), int(color[3:5],16), int(color[5:7],16)
        run.font.color.rgb = RGBColor(r, g, b)

def _T(sl, text, fs=20):
    _txt(sl, text, TITLE_L, TITLE_T, TITLE_W, TITLE_H,
         fs=fs, bold=True, color="#FFFFFF", align=PP_ALIGN.CENTER)

def _POP(sl, text):
    _txt(sl, text, POP_L, POP_T, POP_W, POP_H,
         fs=7.5, italic=True, color="#C8DCF0")

def _BUL(sl, items, fs=11.5):
    txb = sl.shapes.add_textbox(Emu(BUL_L), Emu(BUL_T), Emu(BUL_W), Emu(BUL_H))
    tf  = txb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4); p.alignment = PP_ALIGN.LEFT
        run = p.add_run(); run.text = "•  " + item
        run.font.size = Pt(fs)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Portada
# ─────────────────────────────────────────────────────────────────────────────
def slide_portada(prs):
    sl = _new_sl(prs)
    _pic(sl, SHARED_BG, prs)
    _txt(sl, "Ficha Tecnica", TITLE_L, 1_700_000, TITLE_W, 600_000,
         fs=30, bold=True, color="#FFFFFF", align=PP_ALIGN.CENTER)
    _txt(sl, "Ampliar los Datos de P1 y P2",
         TITLE_L, 2_350_000, TITLE_W, 700_000,
         fs=26, bold=True, color="#7EC8E3", align=PP_ALIGN.CENTER)
    _txt(sl, "Cascadas de datos, brechas de cobertura y solicitudes para mayor completitud",
         TITLE_L, 3_150_000, TITLE_W, 500_000,
         fs=13, italic=True, color="#C8DCF0", align=PP_ALIGN.CENTER)
    _txt(sl,
         f"P1 — Caracterizacion: {N_TOTAL} jerarquizados  →  {N_DOTACION} con perfil completo"
         f"  |  {N_SOLO_NOMINA} sin dotacion",
         TITLE_L, 3_850_000, TITLE_W, 380_000,
         fs=11.5, color="#FFD580", align=PP_ALIGN.CENTER)
    _txt(sl,
         f"P2 — Formacion: {N_FORMADOS} formados de {N_TOTAL}  ({N_INICIATIVAS} iniciativas:"
         f" {N_TALL} Talleres · {N_DIP} Diplomados · {N_PROY} Proyectos)",
         TITLE_L, 4_300_000, TITLE_W, 380_000,
         fs=11.5, color="#A8E8A0", align=PP_ALIGN.CENTER)
    _txt(sl, f"Archivo fuente: {F_DOC}  (hojas: NOMINA, DOTACION, DIPLOMADO, TALLERES, PROYECTOS)",
         TITLE_L, 4_800_000, TITLE_W, 350_000,
         fs=9.5, italic=True, color="#A0B8D0", align=PP_ALIGN.CENTER)
    print("  slide 1 - Portada OK")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Cascada P1: 917 -> 915 -> 493 -> sub-cascadas
# ─────────────────────────────────────────────────────────────────────────────
def slide_cascada_p1(prs):
    fig = _tr_fig()
    gx = PIC_RECT[0] + 0.01
    gy = PIC_RECT[1] + 0.00
    gw = PIC_RECT[2] - 0.02
    gh = PIC_RECT[3] - 0.01
    ax = fig.add_axes([gx, gy, gw, gh], facecolor="none", zorder=5)
    ax.set_xlim(0, 10); ax.set_ylim(0, 10); ax.axis("off")

    # Columna izquierda: cascada principal (3 niveles)
    # Nodo 1: 917
    def _box(ax, x0, y0, w, h, n, label, sub, col, fs_n=22):
        rect = mpatches.FancyBboxPatch((x0, y0), w, h,
                                        boxstyle="round,pad=0.10",
                                        facecolor=col, edgecolor="white",
                                        linewidth=0.8, alpha=0.90, zorder=3)
        ax.add_patch(rect)
        ax.text(x0 + 0.60, y0 + h*0.62, n,
                ha="center", va="center", fontsize=fs_n, fontweight="bold",
                color="white", zorder=5,
                path_effects=[pe.withStroke(linewidth=3, foreground="#050D1A")])
        ax.text(x0 + 1.35, y0 + h*0.67, label,
                ha="left", va="center", fontsize=9.5, fontweight="bold",
                color="white", zorder=5)
        ax.text(x0 + 1.35, y0 + h*0.28, sub,
                ha="left", va="center", fontsize=7.8, color="#B8D0E8", zorder=5)

    def _arrow(ax, x, y_from, y_to, lbl, col="#FFD580"):
        mid = (y_from + y_to) / 2
        ax.annotate("", xy=(x, y_to + 0.05), xytext=(x, y_from - 0.05),
                    arrowprops=dict(arrowstyle="->", color=col, lw=1.4,
                                   mutation_scale=16), zorder=6)
        ax.text(x - 0.10, mid, lbl, ha="right", va="center",
                fontsize=7.5, color=col, fontstyle="italic", zorder=5)

    # Nodos principales (columna izquierda)
    BW, BH = 3.70, 1.35
    BX = 0.20

    _box(ax, BX, 8.40, BW, BH, "917", "docentes jerarquizados",
         f"Hoja NOMINA  ({F_DOC})", "#1B4E8F")
    _box(ax, BX, 6.60, BW, BH, "915", "con jerarquia informada",
         "2 sin jerarquia registrada", "#2E6AAD")
    _box(ax, BX, 4.50, BW, BH, "493", "con perfil completo",
         f"Hoja DOTACION  ({F_DOC})", "#1B7A4A", fs_n=22)

    _arrow(ax, BX + BW/2 - 0.2, 8.40, 6.60 + BH, "915 de 917")
    _arrow(ax, BX + BW/2 - 0.2, 6.60, 4.50 + BH,
           f"-{N_SOLO_NOMINA} sin dotacion", col="#E05A3A")

    # Etiqueta del gran salto
    ax.text(BX + BW/2 - 0.2 - 0.15, 5.85,
            f"  424 SOLO EN NOMINA  ",
            ha="right", va="center", fontsize=8.0, fontweight="bold",
            color="#E05A3A", zorder=5,
            bbox=dict(boxstyle="round,pad=0.3", facecolor="#2A0808",
                      edgecolor="#E05A3A", alpha=0.85))

    # Sub-cascadas (columna derecha) — rama de los 493
    # Linea horizontal desde 493 hacia la derecha
    midbox = 4.50 + BH / 2
    ax.plot([BX + BW, BX + BW + 0.45], [midbox, midbox],
            "-", color="white", linewidth=0.7, alpha=0.45, zorder=4)
    ax.plot([BX + BW + 0.45, BX + BW + 0.45], [midbox, 9.10],
            "-", color="white", linewidth=0.7, alpha=0.45, zorder=4)

    sx0 = BX + BW + 0.55
    sub_ys = [9.10, 7.85, 6.60, 5.35, 4.10, 2.85]
    sub_colors = ["#1B7A4A", "#27824A", "#348A52", "#42925A", "#529A62", "#646A36"]

    for (n, lbl, analisis), y, col in zip(SUBCASC, sub_ys, sub_colors):
        sh = 1.05
        srect = mpatches.FancyBboxPatch(
            (sx0, y - sh/2), 5.85, sh,
            boxstyle="round,pad=0.07",
            facecolor=col, edgecolor="white",
            linewidth=0.5, alpha=0.82, zorder=3)
        ax.add_patch(srect)
        ax.text(sx0 + 0.55, y + 0.08, str(n),
                ha="center", va="center", fontsize=16, fontweight="bold",
                color="white", zorder=5,
                path_effects=[pe.withStroke(linewidth=2.5, foreground="#050D1A")])
        ax.text(sx0 + 1.10, y + 0.20, lbl,
                ha="left", va="center", fontsize=8.5, fontweight="bold",
                color="white", zorder=5)
        ax.text(sx0 + 1.10, y - 0.22, f"-> {analisis}",
                ha="left", va="center", fontsize=7.5,
                color="#C8E8C8", zorder=5)
        # linea horizontal
        ax.plot([BX + BW + 0.45, sx0 - 0.02], [y, y],
                "-", color="white", linewidth=0.5, alpha=0.35, zorder=4)

    path = _save_ch(fig, "ficha_p1_cascada.png")

    sl = _new_sl(prs)
    _pic(sl, SHARED_BG, prs)
    _pic(sl, path, prs)
    _T(sl, "P1 — Cascada de Datos: de 917 a 493 y Sub-niveles de Analisis", fs=13)
    _POP(sl, f"Fuente: {F_DOC}  |  Hoja NOMINA (917) cruzada con hoja DOTACION (493 jerarquizados con perfil)")
    _BUL(sl, [
        f"El salto critico es de 917 a 493: los {N_SOLO_NOMINA} docentes "
        f"({N_HON_SIN_DOT} Honorarios + {N_JOR_SIN_DOT} Jornada) estan en NOMINA "
        f"pero NO en DOTACION, por lo que carecen de edad, antiguedad, nivel de "
        f"formacion, unidad/facultad y carga academica.",
        f"Dentro de los 493 con dotacion, las sub-cascadas reflejan brechas menores: "
        f"2 sin antiguedad, 11 con jerarquia 'NO INFORMA', 15 sin nivel de formacion "
        f"declarado, y 60 sin fecha de jerarquizacion (necesaria para 'Anos hasta Jerarquizacion').",
    ])
    print("  slide 2 - Cascada P1 OK")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — El gran salto: 424 sin dotacion (Nomina vs Dotacion)
# ─────────────────────────────────────────────────────────────────────────────
def slide_nomina_vs_dotacion(prs):
    fig = _tr_fig()
    gx = PIC_RECT[0] + 0.01
    gy = PIC_RECT[1] + 0.00
    gw = PIC_RECT[2] - 0.02
    gh = PIC_RECT[3] - 0.01
    ax = fig.add_axes([gx, gy, gw, gh], facecolor="none", zorder=5)
    ax.set_xlim(0, 10); ax.set_ylim(0, 10); ax.axis("off")

    # --- Panel izquierdo: 493 NOMINA_DOTACION ---
    def _panel(ax, x0, y0, w, h, n, title, col, items_ok, items_no=None):
        rect = mpatches.FancyBboxPatch((x0, y0), w, h,
                                        boxstyle="round,pad=0.12",
                                        facecolor=col, edgecolor="white",
                                        linewidth=0.8, alpha=0.28, zorder=2)
        ax.add_patch(rect)
        # barra superior de color
        bar = mpatches.FancyBboxPatch(
            (x0, y0 + h - 0.62), w, 0.62,
            boxstyle="round,pad=0.05",
            facecolor=col, edgecolor="none", alpha=0.88, zorder=3)
        ax.add_patch(bar)
        ax.text(x0 + w/2, y0 + h - 0.31, f"{n}  {title}",
                ha="center", va="center", fontsize=11, fontweight="bold",
                color="white", zorder=5)
        y_item = y0 + h - 0.90
        for it in items_ok:
            ax.text(x0 + 0.25, y_item, it,
                    ha="left", va="top", fontsize=8.0, color="white",
                    zorder=5, linespacing=1.35)
            y_item -= 0.60
        if items_no:
            ax.text(x0 + 0.25, y_item - 0.10, "NO disponible:",
                    ha="left", va="top", fontsize=8.0, color="#E08080",
                    fontweight="bold", zorder=5)
            y_item -= 0.52
            for it in items_no:
                ax.text(x0 + 0.25, y_item, it,
                        ha="left", va="top", fontsize=8.0, color="#E08080",
                        zorder=5)
                y_item -= 0.55

    ok_vars_dot = [
        "[OK]  Sexo (de NOMINA)",
        "[OK]  Jerarquia academica",
        "[OK]  Tipo contrato",
        "[OK]  Edad / fecha nacimiento",
        "[OK]  Antiguedad / fecha ingreso",
        "[OK]  Unidad / Facultad",
        "[OK]  Nivel formacion (grado)",
        "[OK]  Institucion y pais grado",
        "[OK]  Carga horaria (jornada hrs)",
        "[OK]  Cargo especifico",
    ]
    no_vars_nom = [
        "[---]  Edad / fecha nacimiento",
        "[---]  Antiguedad / fecha ingreso",
        "[---]  Unidad / Facultad detallada",
        "[---]  Nivel de formacion",
        "[---]  Nombre grado / titulo",
        "[---]  Institucion y pais grado",
        "[---]  Carga horaria (horas semanales)",
        "[---]  Cargo especifico",
    ]
    ok_vars_nom = [
        "[OK]  Sexo",
        "[OK]  Jerarquia academica",
        "[OK]  Tipo contrato (Jornada/Honor.)",
        "[OK]  Funcion principal academica",
        "[OK]  Fecha jerarquizacion (60% tiene)",
    ]

    _panel(ax, 0.10, 0.30, 4.15, 9.35,
           N_DOTACION, "docentes con DOTACION",
           "#1B7A4A", ok_vars_dot)

    # Divisor central
    ax.axvline(4.60, ymin=0.03, ymax=0.97,
               color="white", linewidth=0.6, alpha=0.30, zorder=4)

    # Panel derecho: 424 SOLO NOMINA
    _panel(ax, 4.80, 0.30, 4.85, 9.35,
           N_SOLO_NOMINA, "docentes SOLO en NOMINA",
           "#8B2020", ok_vars_nom, no_vars_nom)

    # Anotacion del desglose
    ax.text(4.83, 0.12,
            f"Desglose: {N_HON_SIN_DOT} Honorarios + {N_JOR_SIN_DOT} Jornada  "
            f"(la hoja DOTACION no los incluye)",
            ha="left", va="center", fontsize=8.0, color="#E08080",
            fontstyle="italic", zorder=5)

    path = _save_ch(fig, "ficha_p1_brecha.png")

    sl = _new_sl(prs)
    _pic(sl, SHARED_BG, prs)
    _pic(sl, path, prs)
    _T(sl, f"El Gran Salto — {N_SOLO_NOMINA} Docentes Sin Dotacion", fs=15)
    _POP(sl,
         f"NOMINA ({F_DOC}): {N_TOTAL} jerarquizados  |  "
         f"DOTACION: {N_DOTACION} con perfil completo  |  "
         f"Gap: {N_SOLO_NOMINA} sin datos de dotacion ({N_HON_SIN_DOT} Honorarios + "
         f"{N_JOR_SIN_DOT} Jornada)")
    _BUL(sl, [
        f"{N_HON_SIN_DOT} Honorarios jerarquizados ({100*N_HON_SIN_DOT/N_SOLO_NOMINA:.0f}% del gap) "
        f"estan en NOMINA pero NO en la hoja DOTACION: la dotacion actual cubre "
        f"esencialmente docentes de Jornada. Para analizar P1 con cobertura total "
        f"se necesita dotacion para los Honorarios.",
        f"{N_JOR_SIN_DOT} docentes Jornada tambien estan fuera de dotacion. "
        f"Podrian ser nuevos ingresos posteriores al corte del archivo, o docentes "
        f"con algun error de cruce por RUT o nombre.",
    ])
    print("  slide 3 - Nomina vs Dotacion OK")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Sub-cascadas dentro de los 493
# ─────────────────────────────────────────────────────────────────────────────
def slide_subcascadas(prs):
    fig = _tr_fig()
    gx = PIC_RECT[0] + 0.04
    gy = PIC_RECT[1] + 0.02
    gw = PIC_RECT[2] * 0.72
    gh = PIC_RECT[3] - 0.04
    ax = fig.add_axes([gx, gy, gw, gh], facecolor="none", zorder=5)

    ns     = [sc[0] for sc in SUBCASC]
    labels = [sc[1] for sc in SUBCASC]
    anals  = [sc[2] for sc in SUBCASC]
    cols   = ["#1B7A4A","#27924A","#35A452","#44B660","#54C870","#4A9040"]

    bars = ax.barh(range(len(ns)), ns, color=cols, alpha=0.88,
                   edgecolor="white", linewidth=0.6, height=0.58, zorder=3)

    for bar, n, lbl, ana, col in zip(bars, ns, labels, anals, cols):
        # N al final de la barra
        ax.text(n + 3, bar.get_y() + bar.get_height()/2,
                str(n), va="center", ha="left", fontsize=16, fontweight="bold",
                color="white",
                path_effects=[pe.withStroke(linewidth=2, foreground="#050D1A")])
        # Diferencia vs 493
        diff = 493 - n
        if diff > 0:
            ax.text(n - 5, bar.get_y() + bar.get_height()/2,
                    f"-{diff}", va="center", ha="right", fontsize=9,
                    color="#FFD580", fontweight="bold",
                    path_effects=[pe.withStroke(linewidth=1.5, foreground="#050D1A")])
        # analisis habilitado (panel derecho, via ax2)

    ax.set_xlim(360, 510)
    ax.set_xlabel("N docentes disponibles para cada analisis", color="#AAAAAA", fontsize=9)
    ax.set_yticks(range(len(ns)))
    ax.set_yticklabels(labels, fontsize=9.5, color="white")
    ax.tick_params(axis="x", colors="#AAAAAA", labelsize=8.5)
    for sp in ax.spines.values():
        sp.set_edgecolor("white"); sp.set_alpha(0.25); sp.set_linewidth(0.7)
    ax.xaxis.grid(True, color="white", alpha=0.08, linewidth=0.5)
    ax.set_axisbelow(True)

    # Panel derecho: que analisis habilita cada nivel
    ax2 = fig.add_axes([gx + gw + 0.018, gy + gh*0.04, 0.20, gh*0.93],
                       facecolor="none", zorder=5)
    ax2.axis("off")
    ax2.text(0.0, 1.0, "ANALISIS HABILITADO",
             transform=ax2.transAxes, va="top", ha="left",
             fontsize=9.5, fontweight="bold", color="#FFD580", zorder=5)
    ax2.text(0.0, 0.92, "(columna derecha del grafico)",
             transform=ax2.transAxes, va="top", ha="left",
             fontsize=7.5, color="#A0B8D0", fontstyle="italic", zorder=5)

    y_pos = 0.82
    for n, lbl, ana, col in zip(ns, labels, anals, cols):
        circ = mpatches.Circle((0.05, y_pos), 0.025,
                                facecolor=col, edgecolor="white",
                                linewidth=0.5, alpha=0.85,
                                transform=ax2.transAxes, zorder=3)
        ax2.add_patch(circ)
        ax2.text(0.12, y_pos, ana,
                 transform=ax2.transAxes, va="center", ha="left",
                 fontsize=8.0, color="white", zorder=5)
        y_pos -= 0.145

    # Notas sobre brechas menores
    ax2.text(0.0, 0.08,
             "Causas de los drops:\n"
             "- 2 con fecha ingreso anomala\n"
             "- 11 jerarquia='NO INFORMA'\n"
             "- 15 nivel formacion='NO INFORMA'\n"
             "- 60 sin fecha jerarquizacion\n"
             "  en hoja NOMINA",
             transform=ax2.transAxes, va="bottom", ha="left",
             fontsize=7.5, color="#E8A080", zorder=5, linespacing=1.4)

    path = _save_ch(fig, "ficha_p1_subcascadas.png")

    sl = _new_sl(prs)
    _pic(sl, SHARED_BG, prs)
    _pic(sl, path, prs)
    _T(sl, "P1 — Sub-Cascadas Dentro de los 493 con Dotacion", fs=15)
    _POP(sl, f"Fuente: {F_DOC}  (hoja DOTACION + NOMINA)  |  "
             f"Brechas menores de datos dentro de los {N_DOTACION} con perfil completo")
    _BUL(sl, [
        "La mayor sub-caida es de 493 a 433 (60 sin fecha de jerarquizacion), "
        "imprescindible para calcular 'Anos hasta jerarquizacion'. "
        "La hoja NOMINA tiene 45 nulos en FECHA JERARQUIZACION — "
        "solicitando ese campo completado se recuperan esos casos.",
        "Las otras brechas (11 con jerarquia 'NO INFORMA', 15 con nivel formacion "
        "'NO INFORMA') son correcciones en el propio archivo DOTACION: "
        "UCEN puede confirmar esos valores o reclasificarlos.",
    ])
    print("  slide 4 - Sub-cascadas OK")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Cascada P2: 917 -> 357 formados -> 615 iniciativas
# ─────────────────────────────────────────────────────────────────────────────
def slide_cascada_p2(prs):
    fig = _tr_fig()
    gx = PIC_RECT[0] + 0.02
    gy = PIC_RECT[1] + 0.00
    gw = PIC_RECT[2] - 0.04
    gh = PIC_RECT[3] - 0.01
    ax = fig.add_axes([gx, gy, gw, gh], facecolor="none", zorder=5)
    ax.set_xlim(0, 10); ax.set_ylim(0, 10); ax.axis("off")

    BW, BH, BX = 3.70, 1.45, 0.20

    def _box(ax, x0, y0, n, label, sub, col, fs_n=22):
        rect = mpatches.FancyBboxPatch((x0, y0), BW, BH,
                                        boxstyle="round,pad=0.10",
                                        facecolor=col, edgecolor="white",
                                        linewidth=0.8, alpha=0.90, zorder=3)
        ax.add_patch(rect)
        ax.text(x0 + 0.60, y0 + BH*0.62, str(n),
                ha="center", va="center", fontsize=fs_n, fontweight="bold",
                color="white", zorder=5,
                path_effects=[pe.withStroke(linewidth=3, foreground="#050D1A")])
        ax.text(x0 + 1.35, y0 + BH*0.68, label,
                ha="left", va="center", fontsize=9.5, fontweight="bold",
                color="white", zorder=5)
        ax.text(x0 + 1.35, y0 + BH*0.22, sub,
                ha="left", va="center", fontsize=7.8, color="#B8D0E8", zorder=5)

    def _arrow(ax, x, y_from, y_to, lbl, col="#FFD580"):
        mid = (y_from + y_to) / 2
        ax.annotate("", xy=(x, y_to + 0.05), xytext=(x, y_from - 0.05),
                    arrowprops=dict(arrowstyle="->", color=col, lw=1.4,
                                   mutation_scale=16), zorder=6)
        ax.text(x - 0.10, mid, lbl, ha="right", va="center",
                fontsize=7.5, color=col, fontstyle="italic", zorder=5)

    # Nodos principales
    _box(ax, BX, 8.20, N_TOTAL, "docentes jerarquizados",
         f"Hoja NOMINA  ({F_DOC})", "#1B4E8F")

    # 560 sin formacion (burbuja a la derecha)
    bub_x, bub_y = BX + BW + 0.60, 7.90
    bub = mpatches.FancyBboxPatch((bub_x, bub_y), 5.10, 1.30,
                                   boxstyle="round,pad=0.10",
                                   facecolor="#5A2020", edgecolor="#E05A3A",
                                   linewidth=0.8, alpha=0.75, zorder=3)
    ax.add_patch(bub)
    ax.text(bub_x + 0.55, bub_y + 0.65, str(N_SIN_FORM),
            ha="center", va="center", fontsize=20, fontweight="bold",
            color="white", zorder=5,
            path_effects=[pe.withStroke(linewidth=2.5, foreground="#050D1A")])
    ax.text(bub_x + 1.25, bub_y + 0.82, "sin formacion registrada",
            ha="left", va="center", fontsize=9, fontweight="bold",
            color="white", zorder=5)
    ax.text(bub_x + 1.25, bub_y + 0.40, "Brecha de cobertura, no de datos",
            ha="left", va="center", fontsize=7.8,
            color="#E08080", fontstyle="italic", zorder=5)
    ax.text(bub_x + 1.25, bub_y + 0.10,
            f"({100*N_SIN_FORM/N_TOTAL:.0f}% del universo jerarquizado no ha participado)",
            ha="left", va="center", fontsize=7.5, color="#B8D0E8", zorder=5)
    # flecha horizontal desde 917 a burbuja
    ax.annotate("", xy=(bub_x - 0.02, bub_y + 0.65),
                xytext=(BX + BW + 0.05, 8.20 + BH*0.44),
                arrowprops=dict(arrowstyle="->", color="#E05A3A", lw=1.2,
                                connectionstyle="arc3,rad=0.0"), zorder=6)

    _box(ax, BX, 5.95, N_FORMADOS, "docentes formados (al menos 1 iniciativa)",
         f"Hojas DIPLOMADO 2022-25 · TALLERES · PROYECTOS  ({F_DOC})",
         "#2E6AAD")

    _arrow(ax, BX + BW/2 - 0.2, 8.20, 5.95 + BH,
           f"{N_FORMADOS} de {N_TOTAL}  ({100*N_FORMADOS/N_TOTAL:.0f}%)")

    # Desglose de 615 iniciativas
    _box(ax, BX, 3.65, N_INICIATIVAS, "iniciativas de formacion totales",
         f"{N_FORMADOS} docentes unicos (uno puede tener varias)", "#3A8BC4")

    _arrow(ax, BX + BW/2 - 0.2, 5.95, 3.65 + BH, f"{N_INICIATIVAS} eventos")

    # Breakdown por tipo (barras horizontales a la derecha de la caja 615)
    tipos = [
        (N_TALL,  "Talleres",    "#4B9CD3", "(hojas TALLERES 2023_2 / 2024_1 / 2024_2)"),
        (N_DIP,   "Diplomados",  "#3A7ABF", "(hojas DIPLOMADO 2022 / 2023 / 2024 / 2025)"),
        (N_PROY,  "Proyectos",   "#2A5EA0", "(hoja PROYECTOS DE INVESTIGACION)"),
    ]
    ty0 = 3.65
    for n_t, lbl_t, col_t, file_t in tipos:
        bx_t = BX + BW + 0.40
        bar_w = (n_t / N_TALL) * 5.0
        brect = mpatches.FancyBboxPatch(
            (bx_t, ty0 + 0.12), bar_w, 0.75,
            boxstyle="round,pad=0.05",
            facecolor=col_t, edgecolor="white",
            linewidth=0.4, alpha=0.85, zorder=3)
        ax.add_patch(brect)
        ax.text(bx_t + bar_w + 0.12, ty0 + 0.50,
                f"{n_t}  {lbl_t}", ha="left", va="center",
                fontsize=9, fontweight="bold", color="white", zorder=5)
        ax.text(bx_t + bar_w + 0.12, ty0 + 0.16,
                file_t, ha="left", va="center",
                fontsize=7.5, color="#B8D0E8", fontstyle="italic", zorder=5)
        ty0 += 1.00

    path = _save_ch(fig, "ficha_p2_cascada.png")

    sl = _new_sl(prs)
    _pic(sl, SHARED_BG, prs)
    _pic(sl, path, prs)
    _T(sl, "P2 — Cascada de Participacion en Formacion: 917 a 357 Formados", fs=13)
    _POP(sl, f"Fuente: {F_DOC}  |  Hojas: DIPLOMADO 2022-2025 · TALLERES 2023_2 / 2024_1 / 2024_2 · PROYECTOS")
    _BUL(sl, [
        f"De los {N_TOTAL} jerarquizados, {N_SIN_FORM} ({100*N_SIN_FORM/N_TOTAL:.0f}%) "
        f"no tienen ninguna iniciativa de formacion registrada. "
        f"Esto NO es un problema de datos faltantes — es una brecha de cobertura "
        f"de la politica de formacion docente.",
        f"Los {N_FORMADOS} formados generan {N_INICIATIVAS} iniciativas "
        f"en total: Talleres ({N_TALL}), Diplomados ({N_DIP}) y "
        f"Proyectos de Investigacion ({N_PROY}). Un mismo docente puede "
        f"haber participado en varias iniciativas de distinto tipo.",
    ])
    print("  slide 5 - Cascada P2 OK")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Solicitudes de datos para P1 y P2
# ─────────────────────────────────────────────────────────────────────────────
def slide_solicitudes(prs):
    fig = _tr_fig()
    gx = PIC_RECT[0] + 0.01
    gy = PIC_RECT[1] + 0.01
    gw = PIC_RECT[2] - 0.02
    gh = PIC_RECT[3] - 0.03
    ax = fig.add_axes([gx, gy, gw, gh], facecolor="none", zorder=5)
    ax.set_xlim(0, 10); ax.set_ylim(0, 10.5); ax.axis("off")

    BOX_H = 2.35
    solicitudes = [
        {
            "num": "1", "col": "#1B6A40",
            "title": "P1  —  Dotacion para los 364 Honorarios jerarquizados",
            "lines": [
                f"La hoja DOTACION cubre solo Jornada. Los {N_HON_SIN_DOT} Honorarios jerarquizados",
                "tienen jerarquia academica asignada pero carecen de:",
                "  edad / fecha nacimiento  |  antiguedad  |  unidad/facultad detallada",
                "  nivel de formacion  |  nombre y pais del grado  |  carga horaria",
                f"Solicitar: hoja DOTACION ampliada que incluya Honorarios (o archivo equivalente).",
            ],
            "impacto": f"Recupera {N_HON_SIN_DOT} docentes para todas las variables de P1",
            "yc": 9.45,
        },
        {
            "num": "2", "col": "#1B4E8F",
            "title": f"P1  —  Dotacion para los {N_JOR_SIN_DOT} Jornada sin cruce",
            "lines": [
                f"{N_JOR_SIN_DOT} docentes Jornada estan en NOMINA pero no en DOTACION.",
                "Posibles causas:",
                "  a) Ingresaron despues del corte del archivo DOTACION (nuevos 2026)",
                "  b) Error de cruce por RUT distinto o nombre con diferencias",
                "Solicitar: verificacion de esos RUTs en DOTACION o complemento de datos.",
            ],
            "impacto": f"Recupera hasta {N_JOR_SIN_DOT} docentes Jornada con perfil completo",
            "yc": 6.75,
        },
        {
            "num": "3", "col": "#3A5080",
            "title": "P1  —  Fecha de jerarquizacion para 60 sin dato",
            "lines": [
                "La hoja NOMINA tiene 45 nulos en FECHA JERARQUIZACION.",
                "Sumado a otros sin cruce, son ~60 docentes sin ese campo.",
                "Necesario para: calcular 'Anos hasta jerarquizacion' y trayectoria academica.",
                "Solicitar: completar FECHA JERARQUIZACION en NOMINA (o listado complementario).",
            ],
            "impacto": "Sube de 433 a 493 los disponibles para analisis de trayectoria",
            "yc": 4.05,
        },
        {
            "num": "4", "col": "#4A7030",
            "title": "P2  —  Confirmar completitud de registros de formacion 2025",
            "lines": [
                f"La base P2 contiene {N_INICIATIVAS} iniciativas ({N_FORMADOS} docentes unicos).",
                "Las hojas de TALLERES cubren solo 2023_2, 2024_1 y 2024_2.",
                "No hay hoja de TALLERES 2025 ni TALLERES 2023_1.",
                "Solicitar: completar con talleres 2025 y cualquier periodo faltante.",
            ],
            "impacto": f"Amplia el universo de {N_FORMADOS} formados y las {N_INICIATIVAS} iniciativas",
            "yc": 1.55,
        },
    ]

    for s in solicitudes:
        yc = s["yc"]; y0 = yc - BOX_H/2
        ax.add_patch(mpatches.FancyBboxPatch(
            (-0.05, y0), 10.1, BOX_H,
            boxstyle="round,pad=0.10", facecolor=s["col"],
            edgecolor="white", linewidth=0.6, alpha=0.32, zorder=2))
        circ = mpatches.Circle((0.42, yc + 0.65), 0.35,
                                facecolor=s["col"], edgecolor="white",
                                linewidth=0.8, alpha=0.95, zorder=3)
        ax.add_patch(circ)
        ax.text(0.42, yc + 0.65, s["num"],
                ha="center", va="center", fontsize=15, fontweight="bold",
                color="white", zorder=4)
        ax.text(1.00, yc + 0.83, s["title"],
                ha="left", va="center", fontsize=10, fontweight="bold",
                color="#FFD580", zorder=4)
        y_txt = yc + 0.40
        for line in s["lines"]:
            ax.text(1.00, y_txt, line,
                    ha="left", va="top", fontsize=7.8, color="#D0E8FF", zorder=4)
            y_txt -= 0.37
        ax.text(0.42, y0 + 0.18, "-> " + s["impacto"],
                ha="left", va="center", fontsize=7.8,
                color="#A8E8A0", zorder=4, fontstyle="italic")

    path = _save_ch(fig, "ficha_p1p2_solicitudes.png")

    sl = _new_sl(prs)
    _pic(sl, SHARED_BG, prs)
    _pic(sl, path, prs)
    _T(sl, "Datos Adicionales — Solicitudes para Ampliar P1 y P2", fs=15)
    _POP(sl, f"Todas las solicitudes apuntan al mismo archivo: {F_DOC}  "
             f"(ampliacion de hojas existentes o correcciones de campo)")
    _BUL(sl, [
        f"La solicitud mas impactante (N1) es incluir a los {N_HON_SIN_DOT} Honorarios "
        f"en la hoja DOTACION: ampliar ese archivo de Jornada a todos los jerarquizados "
        f"sube el universo de analisis P1 de {N_DOTACION} a {N_DOTACION + N_HON_SIN_DOT} "
        f"docentes — un aumento del {100*N_HON_SIN_DOT/N_DOTACION:.0f}%.",
        f"Las solicitudes N2 y N3 son correcciones menores que elevan la "
        f"completitud interna de los {N_DOTACION} ya cubiertos.",
    ])
    print("  slide 6 - Solicitudes OK")


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("Generando FICHA_Ampliar_P1_P2.pptx ...")
    _ensure_bg()
    prs = Presentation()
    prs.slide_width  = Emu(SW_EMU)
    prs.slide_height = Emu(SH_EMU)
    slide_portada(prs)
    slide_cascada_p1(prs)
    slide_nomina_vs_dotacion(prs)
    slide_subcascadas(prs)
    slide_cascada_p2(prs)
    slide_solicitudes(prs)
    prs.save(OUT_PPTX)
    print(f"\n  Guardado: {OUT_PPTX}")
    print("  6 slides: [1] Portada, [2] Cascada P1, [3] Nomina vs Dotacion, "
          "[4] Sub-cascadas, [5] Cascada P2, [6] Solicitudes")

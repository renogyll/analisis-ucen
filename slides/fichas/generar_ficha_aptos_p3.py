# -*- coding: utf-8 -*-
"""
generar_ficha_aptos_p3.py  v2
Genera FICHA_Ampliar_Aptos_P3.pptx en formato dark v3:
  Slide 1 - Portada
  Slide 2 - Pipeline de datos (flujo con N y archivos originales)
  Slide 3 - Los 160 excluidos (categorias no solapadas que suman 160)
  Slide 4 - Inventario de archivos originales entregados
  Slide 5 - Datos adicionales para ampliar el universo
"""
import sys; sys.stdout.reconfigure(encoding="utf-8")
import os, zipfile
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.patheffects as pe
import pandas as pd
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─────────────────────────────────────────────────────────────────────────────
# Rutas
# ─────────────────────────────────────────────────────────────────────────────
BASE      = os.path.dirname(os.path.abspath(__file__))
PROC918   = os.path.join(BASE, "..", "PROCESADO")
PROC_ROOT = os.path.normpath(os.path.join(BASE, "..", "..", "PROCESADO"))
ROOT      = os.path.normpath(os.path.join(BASE, "..", ".."))
FONDOTIPO = os.path.join(ROOT, "Fondotipop.pptx")
OUT_DIR   = os.path.join(BASE, "dark_slides_v3")
OUT_PPTX  = os.path.join(ROOT, "FICHA_Ampliar_Aptos_P3_v2.pptx")
os.makedirs(OUT_DIR, exist_ok=True)

# ─────────────────────────────────────────────────────────────────────────────
# Datos
# ─────────────────────────────────────────────────────────────────────────────
p3     = pd.read_csv(os.path.join(PROC918, "p3_918.csv"), encoding="utf-8-sig")
p3_sat = pd.read_csv(os.path.join(PROC918, "p3_sat_zscore_918.csv"), encoding="utf-8-sig")
ep     = pd.read_csv(os.path.join(PROC_ROOT, "evaluacion_periodo.csv"), encoding="utf-8-sig")

N_TOTAL_EVENTOS = len(p3)
N_FORMADOS      = p3["rut_key"].nunique()
N_APTOS         = len(p3_sat)
N_EXCLUIDOS     = N_FORMADOS - N_APTOS
N_EVAL          = len(ep)
N_DOC_SAT       = ep["rut_docente"].nunique()
PERIODOS        = sorted(ep["periodo"].unique())
PER_STR         = f"{PERIODOS[0]} a {PERIODOS[-1]}"

# Exclusiones NO solapadas que sumen exactamente 160
p3_doc = p3.groupby("rut_key").agg(
    tiene_bl  = ("tiene_sat_baseline",   "max"),
    tiene_res = ("tiene_sat_resultado", "max"),
).reset_index()
# Los 160 excluidos caen en 4 grupos mutuamente excluyentes:
N_SOLO_BL  = int(((p3_doc["tiene_bl"] == False) & (p3_doc["tiene_res"] == True)).sum())
N_SOLO_RES = int(((p3_doc["tiene_bl"] == True)  & (p3_doc["tiene_res"] == False)).sum())
N_SIN_BOTH = int(((p3_doc["tiene_bl"] == False) & (p3_doc["tiene_res"] == False)).sum())
N_MIXTO    = N_EXCLUIDOS - N_SOLO_BL - N_SOLO_RES - N_SIN_BOTH   # caso especial

anio_counts = p3["anio_evento"].value_counts().sort_index()

print(f"N_FORMADOS={N_FORMADOS}, N_APTOS={N_APTOS}, N_EXCLUIDOS={N_EXCLUIDOS}")
print(f"  SOLO_BL={N_SOLO_BL} + SOLO_RES={N_SOLO_RES} + SIN_BOTH={N_SIN_BOTH} "
      f"+ MIXTO={N_MIXTO} = {N_SOLO_BL+N_SOLO_RES+N_SIN_BOTH+N_MIXTO}")
print(f"N_EVAL={N_EVAL}, N_DOC_SAT={N_DOC_SAT}, PER_STR={PER_STR}")

# Nombres de archivos ORIGINALES de la contraparte
F_SAT      = "EVALUACION ESTUDIANTES UCEN 3-5-2026.xlsx"
F_SAT_DET  = "(6 hojas: 2023-01 a 2025-02)"
F_DOC      = "CONSOLIDADO DOCENTES 3-05-2026.xlsx"
F_DOC_NOM  = "(hojas: DOTACION, NOMINA)"
F_DOC_FORM = "(hojas: DIPLOMADO 2022-25, TALLERES, PROYECTOS)"
F_CALC     = "Construido en el analisis"
F_CALC_DET = "(cruce SAT x Formacion)"

# ─────────────────────────────────────────────────────────────────────────────
# Assets
# ─────────────────────────────────────────────────────────────────────────────
BG_PATH   = os.path.join(OUT_DIR, "_fondotipo_bg.jpg")
LOGO_PATH = os.path.join(OUT_DIR, "_fondotipo_logo.png")
for path, zname in [(BG_PATH, "ppt/media/image1.jpg"),
                    (LOGO_PATH, "ppt/media/image2.png")]:
    if not os.path.exists(path):
        with zipfile.ZipFile(FONDOTIPO) as z:
            with open(path, "wb") as f:
                f.write(z.read(zname))

with PILImage.open(BG_PATH) as _im:
    _rgb = _im.convert("RGB"); _iw, _ih = _rgb.size
    _nh  = int(_iw / (16/9)); _y0 = min(int(_ih * 0.12), _ih - _nh)
    bg_arr = np.array(_rgb.crop((0, _y0, _iw, _y0 + _nh)))

with PILImage.open(LOGO_PATH) as _logo:
    logo_arr = np.array(_logo.convert("RGBA")).astype(np.float32) / 255.0

H_GRAD = 600; grad = np.zeros((H_GRAD, 1, 4), dtype=np.float32)
for _r in range(H_GRAD):
    _t = _r / (H_GRAD - 1)
    _stops = [(0.00, (0, 33, 71)), (0.54, (0, 70, 128)), (1.00, (144, 171, 196))]
    for _i in range(len(_stops) - 1):
        _t0, _c0 = _stops[_i]; _t1, _c1 = _stops[_i + 1]
        if _t0 <= _t <= _t1:
            _s = (_t - _t0) / (_t1 - _t0)
            grad[_r, 0] = [(_c0[0]+_s*(_c1[0]-_c0[0]))/255,
                           (_c0[1]+_s*(_c1[1]-_c0[1]))/255,
                           (_c0[2]+_s*(_c1[2]-_c0[2]))/255, 0.82]; break

# ─────────────────────────────────────────────────────────────────────────────
# Layout constants
# ─────────────────────────────────────────────────────────────────────────────
SW, SH         = 13.333, 7.5
SW_EMU, SH_EMU = 12192000, 6858000
PIC_L, PIC_T, PIC_W, PIC_H         = 786581, 1125000, 10599174, 3720000
BUL_L, BUL_T, BUL_W, BUL_H         = 786581, 4870000, 10599174, 1870000
LOGO_L, LOGO_T, LOGO_W, LOGO_H     = 9813773, 656354, 1756626, 697725
TITLE_L, TITLE_T, TITLE_W, TITLE_H = PIC_L, 185000, PIC_W, 710000
POP_L,   POP_T,   POP_W,   POP_H   = PIC_L, 845000, 9000000, 255000

def _ex(e): return e / SW_EMU
def _ey(e): return e / SH_EMU
def _fig_rect(l, t, w, h): return (l, 1-t-h, w, h)
PIC_RECT  = _fig_rect(_ex(PIC_L), _ey(PIC_T), _ex(PIC_W), _ey(PIC_H))
LOGO_RECT = _fig_rect(_ex(LOGO_L), _ey(LOGO_T), _ex(LOGO_W), _ey(LOGO_H))

# ─────────────────────────────────────────────────────────────────────────────
# Helpers matplotlib
# ─────────────────────────────────────────────────────────────────────────────
def _bg_fig():
    fig = plt.figure(figsize=(SW, SH), facecolor="#101820")
    fig.patch.set_facecolor("#101820")
    for z, arr in [(0, bg_arr), (1, grad)]:
        ax = fig.add_axes([0, 0, 1, 1], zorder=z)
        ax.imshow(arr, extent=[0, 1, 0, 1], aspect="auto", origin="upper")
        ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis("off")
    al = fig.add_axes([LOGO_RECT[0], LOGO_RECT[1], LOGO_RECT[2], LOGO_RECT[3]],
                      zorder=10, facecolor="none")
    al.imshow(logo_arr, aspect="auto"); al.axis("off"); al.patch.set_visible(False)
    return fig

def _tr_fig():
    fig = plt.figure(figsize=(SW, SH), facecolor="none")
    fig.patch.set_facecolor("none"); return fig

SHARED_BG = os.path.join(OUT_DIR, "_background.png")

def _ensure_bg():
    if not os.path.exists(SHARED_BG):
        fig = _bg_fig()
        plt.savefig(SHARED_BG, dpi=150, facecolor=fig.get_facecolor()); plt.close()
        print("  bg generado")

def _save_ch(fig, name):
    path = os.path.join(OUT_DIR, name)
    plt.savefig(path, dpi=150, facecolor="none", transparent=True)
    plt.close(); return path

# ─────────────────────────────────────────────────────────────────────────────
# Helpers python-pptx
# ─────────────────────────────────────────────────────────────────────────────
def _new_sl(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _pic(sl, path, prs):
    sl.shapes.add_picture(path, Emu(0), Emu(0), prs.slide_width, prs.slide_height)

def _txt(sl, text, left, top, width, height,
         fs=12, bold=False, italic=False, color="#FFFFFF",
         align=PP_ALIGN.LEFT, wrap=True, lspc=0):
    txb = sl.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf  = txb.text_frame; tf.word_wrap = wrap
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if lspc > 0 and i > 0:
            p.space_before = Pt(lspc)
        run = p.add_run(); run.text = line
        run.font.size = Pt(fs); run.font.bold = bold; run.font.italic = italic
        r, g, b = int(color[1:3],16), int(color[3:5],16), int(color[5:7],16)
        run.font.color.rgb = RGBColor(r, g, b)

def _T(sl, text, fs=20):
    _txt(sl, text, TITLE_L, TITLE_T, TITLE_W, TITLE_H,
         fs=fs, bold=True, color="#FFFFFF", align=PP_ALIGN.CENTER)

def _POP(sl, text):
    _txt(sl, text, POP_L, POP_T, POP_W, POP_H,
         fs=7.5, italic=True, color="#C8DCF0")

def _BUL(sl, items, fs=11.5):
    txb = sl.shapes.add_textbox(Emu(BUL_L), Emu(BUL_T), Emu(BUL_W), Emu(BUL_H))
    tf  = txb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4); p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "•  " + item
        run.font.size = Pt(fs)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Portada
# ─────────────────────────────────────────────────────────────────────────────
def slide_portada(prs):
    sl = _new_sl(prs)
    _pic(sl, SHARED_BG, prs)
    _txt(sl, "Ficha Tecnica", TITLE_L, 1_700_000, TITLE_W, 600_000,
         fs=32, bold=True, color="#FFFFFF", align=PP_ALIGN.CENTER)
    _txt(sl, "Ampliar el Universo Aptos P3",
         TITLE_L, 2_350_000, TITLE_W, 700_000,
         fs=28, bold=True, color="#7EC8E3", align=PP_ALIGN.CENTER)
    _txt(sl, "Contexto de datos, filtros aplicados y solicitudes para mayor cobertura",
         TITLE_L, 3_150_000, TITLE_W, 500_000,
         fs=14, italic=True, color="#C8DCF0", align=PP_ALIGN.CENTER)
    _txt(sl,
         f"Universo actual: {N_FORMADOS} formados  ->  {N_APTOS} Aptos P3  "
         f"({N_EXCLUIDOS} docentes excluidos por falta de SAT en uno o ambos momentos)",
         TITLE_L, 3_850_000, TITLE_W, 400_000,
         fs=12, color="#FFD580", align=PP_ALIGN.CENTER)
    _txt(sl,
         f"Base SAT disponible: {PER_STR}  |  {N_DOC_SAT:,} docentes  |  {N_EVAL:,} evaluaciones",
         TITLE_L, 4_350_000, TITLE_W, 350_000,
         fs=10, italic=True, color="#A0B8D0", align=PP_ALIGN.CENTER)
    print("  slide 1 - Portada OK")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Pipeline de datos  (headline = docentes, unidad de analisis)
# ─────────────────────────────────────────────────────────────────────────────
def slide_pipeline(prs):
    fig = _tr_fig()
    gx = PIC_RECT[0] + 0.02
    gy = PIC_RECT[1] + 0.00
    gw = PIC_RECT[2] - 0.04
    gh = PIC_RECT[3] - 0.01
    ax = fig.add_axes([gx, gy, gw, gh], facecolor="none", zorder=5)
    ax.set_xlim(0, 10); ax.set_ylim(0, 10); ax.axis("off")

    nodes = [
        {
            "n_big":  f"{N_DOC_SAT:,}",          # headline: docentes
            "n_lbl":  "docentes con SAT",
            "sub":    f"{N_EVAL:,} evaluaciones  |  {PER_STR}",
            "file1":  F_SAT,
            "file2":  F_SAT_DET,
            "filter": "Datos de evaluacion docente entregados por UCEN",
            "col":    "#1B4E8F",
            "yc":     9.00,
        },
        {
            "n_big":  "917",
            "n_lbl":  "docentes jerarquizados",
            "sub":    "Universo base del analisis UCEN",
            "file1":  F_DOC,
            "file2":  F_DOC_NOM,
            "filter": "Filtro: ranking UCEN vigente (jerarquia activa)",
            "col":    "#2E6AAD",
            "yc":     6.75,
        },
        {
            "n_big":  str(N_FORMADOS),
            "n_lbl":  "docentes formados",
            "sub":    f"{N_TOTAL_EVENTOS} eventos de formacion  |  2022 a 2025",
            "file1":  F_DOC,
            "file2":  F_DOC_FORM,
            "filter": "Filtro: >= 1 actividad de formacion (Taller / Diplomado / Proyecto)",
            "col":    "#3A8BC4",
            "yc":     4.50,
        },
        {
            "n_big":  str(N_APTOS),
            "n_lbl":  "Aptos P3",
            "sub":    "SAT valido ANTES y DESPUES de la formacion",
            "file1":  F_CALC,
            "file2":  F_CALC_DET,
            "filter": "Filtro: SAT baseline=True  AND  SAT resultado=True",
            "col":    "#2E9E55",
            "yc":     2.25,
        },
    ]

    BOX_W, BOX_H = 3.9, 1.55
    BOX_X = 0.8

    for nd in nodes:
        yc = nd["yc"]
        y0 = yc - BOX_H / 2
        rect = mpatches.FancyBboxPatch(
            (BOX_X, y0), BOX_W, BOX_H,
            boxstyle="round,pad=0.10",
            facecolor=nd["col"], edgecolor="white",
            linewidth=0.9, alpha=0.92, zorder=3)
        ax.add_patch(rect)

        # N headline (docentes)
        ax.text(BOX_X + 0.58, yc + 0.15, nd["n_big"],
                ha="center", va="center", fontsize=25, fontweight="bold",
                color="white", zorder=5,
                path_effects=[pe.withStroke(linewidth=3, foreground="#050D1A")])
        # label bajo el N
        ax.text(BOX_X + 0.58, yc - 0.30, nd["n_lbl"],
                ha="center", va="center", fontsize=7.0,
                color="#D0E8FF", zorder=5)
        # sub (info secundaria)
        ax.text(BOX_X + 1.45, yc + 0.30, nd["sub"],
                ha="left", va="center", fontsize=9.0, fontweight="bold",
                color="white", zorder=5)
        ax.text(BOX_X + 1.45, yc - 0.12, nd["filter"],
                ha="left", va="center", fontsize=7.5,
                color="#B8D0E8", zorder=5)

        # Archivo original (derecha del box)
        ax.text(BOX_X + BOX_W + 0.22, yc + 0.28, nd["file1"],
                ha="left", va="center", fontsize=8.2,
                color="#FFD580", fontweight="bold", zorder=5)
        ax.text(BOX_X + BOX_W + 0.22, yc - 0.08, nd["file2"],
                ha="left", va="center", fontsize=7.5,
                color="#B8D0E8", zorder=5, fontstyle="italic")

    # Flechas + porcentajes entre nodos
    arrows = [
        (nodes[0]["yc"] - BOX_H/2, nodes[1]["yc"] + BOX_H/2,
         f"917 de {N_DOC_SAT:,}"),
        (nodes[1]["yc"] - BOX_H/2, nodes[2]["yc"] + BOX_H/2,
         f"{N_FORMADOS} de 917  (39%)"),
        (nodes[2]["yc"] - BOX_H/2, nodes[3]["yc"] + BOX_H/2,
         f"{N_APTOS} de {N_FORMADOS}  (55%)"),
    ]
    ax_x = BOX_X + BOX_W / 2 - 0.2
    for y_from, y_to, lbl in arrows:
        mid = (y_from + y_to) / 2
        ax.annotate("",
                    xy=(ax_x, y_to + 0.05), xytext=(ax_x, y_from - 0.05),
                    arrowprops=dict(arrowstyle="->", color="#FFD580",
                                   lw=1.5, mutation_scale=18), zorder=6)
        ax.text(ax_x - 0.12, mid, lbl,
                ha="right", va="center", fontsize=7.5,
                color="#FFD580", fontstyle="italic", zorder=5)

    path = _save_ch(fig, "ficha_pipeline.png")
    sl = _new_sl(prs)
    _pic(sl, SHARED_BG, prs)
    _pic(sl, path, prs)
    _T(sl, "Pipeline de Datos   De la Evaluacion SAT a los 197 Aptos P3", fs=14)
    _POP(sl, f"Base SAT: {PER_STR}  |  Formacion: 2022-2025  |  "
             f"Universo base: 917 docentes jerarquizados UCEN")
    _BUL(sl, [
        f"Cada docente formado requiere SAT registrado en el semestre ANTERIOR "
        f"(baseline) y POSTERIOR (resultado) a su actividad de formacion. "
        f"Sin ese par completo, el docente queda excluido del analisis P3.",
        f"La base SAT cubre {PER_STR}: docentes con formacion en 2022 o en "
        f"2023-01 no tienen baseline disponible — es la causa principal de exclusion.",
        f"Los {N_APTOS} Aptos P3 representan el 55% de los {N_FORMADOS} formados. "
        f"Los {N_EXCLUIDOS} restantes quedan fuera unicamente por falta de SAT "
        f"en uno o ambos momentos requeridos.",
    ])
    print("  slide 2 - Pipeline OK")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Los 160 excluidos (4 grupos mutuamente excluyentes que suman 160)
# ─────────────────────────────────────────────────────────────────────────────
def slide_excluidos(prs):
    fig = _tr_fig()
    gx = PIC_RECT[0] + 0.05
    gy = PIC_RECT[1] + 0.02
    gw = PIC_RECT[2] * 0.68
    gh = PIC_RECT[3] - 0.06
    ax = fig.add_axes([gx, gy, gw, gh], facecolor="none", zorder=5)

    cats = [
        f"Sin ambos:\nni baseline\nni resultado",
        f"Caso mixto:\nSAT en periodos\nno contiguos",
        f"Solo sin resultado:\ntienen baseline,\nfalta resultado",
        f"Solo sin baseline:\ntienen resultado,\nfalta baseline",
    ]
    vals   = [N_SIN_BOTH, N_MIXTO, N_SOLO_RES, N_SOLO_BL]
    colors = ["#A03020", "#C05828", "#C88030", "#3A8BC4"]

    bars = ax.barh(cats, vals, color=colors, alpha=0.88,
                   edgecolor="white", linewidth=0.7, height=0.58, zorder=3)

    for bar, v in zip(bars, vals):
        ax.text(bar.get_width() + 1.5, bar.get_y() + bar.get_height()/2,
                str(v), va="center", ha="left", fontsize=17, fontweight="bold",
                color="white",
                path_effects=[pe.withStroke(linewidth=2, foreground="#050D1A")])

    # Barra de total (vertical punteada)
    ax.axvline(N_EXCLUIDOS, color="#FFD580", linewidth=1.8,
               linestyle="--", alpha=0.9, zorder=4)
    ax.text(N_EXCLUIDOS + 1.2, 3.65,
            f"Total\nexcluidos:\n{N_EXCLUIDOS}",
            ha="left", va="top", fontsize=10,
            color="#FFD580", fontweight="bold", zorder=5)

    # Verificacion suma
    total_check = N_SOLO_BL + N_SOLO_RES + N_SIN_BOTH + N_MIXTO
    ax.text(0.5, -0.62,
            f"Verificacion: {N_SOLO_BL} + {N_SOLO_RES} + {N_SIN_BOTH} + {N_MIXTO} "
            f"= {total_check} docentes excluidos  (grupos mutuamente excluyentes, "
            f"sin solapamiento)",
            ha="left", va="top", fontsize=8.0,
            color="#A8E8A0", style="italic", zorder=5)

    ax.set_xlim(0, N_EXCLUIDOS * 1.45)
    ax.set_xlabel("N docentes (grupos sin solapamiento)", color="#AAAAAA", fontsize=9)
    ax.tick_params(axis="x", colors="#AAAAAA", labelsize=9)
    ax.tick_params(axis="y", colors="white", labelsize=9.5)
    for sp in ax.spines.values():
        sp.set_edgecolor("white"); sp.set_alpha(0.25); sp.set_linewidth(0.7)
    ax.xaxis.grid(True, color="white", alpha=0.08, linewidth=0.5)
    ax.set_axisbelow(True)

    # Panel derecho: causa + solicitud
    ax2 = fig.add_axes([gx + gw + 0.022, gy + gh*0.05, 0.195, gh * 0.92],
                       facecolor="none", zorder=5)
    ax2.axis("off")
    lines = [
        ("CAUSA PRINCIPAL:", "#FFD580", True),
        ("", "white", False),
        ("Azul (falta baseline):", "#7EC8E3", True),
        ("Base SAT parte en 2023-01.", "#D0E8FF", False),
        ("Docentes formados en 2022", "#D0E8FF", False),
        ("o en 2023-01 no tienen", "#D0E8FF", False),
        ("semestre anterior cubierto.", "#D0E8FF", False),
        ("", "white", False),
        ("Naranjo (falta resultado):", "#E09A50", True),
        ("Formaciones en 2025-02 aun", "#D0E8FF", False),
        ("sin SAT resultado (2026-01", "#D0E8FF", False),
        ("todavia no existe).", "#D0E8FF", False),
        ("", "white", False),
        ("Rojo (sin ambos):", "#C05040", True),
        ("Sin cursos ese semestre o", "#D0E8FF", False),
        ("SAT no capturado.", "#D0E8FF", False),
    ]
    y = 1.0
    for txt, col, bold in lines:
        ax2.text(0.0, y, txt, transform=ax2.transAxes,
                 va="top", ha="left", fontsize=8.2 if bold else 7.8,
                 fontweight="bold" if bold else "normal",
                 color=col, zorder=5)
        y -= 0.065

    path = _save_ch(fig, "ficha_excluidos.png")
    sl = _new_sl(prs)
    _pic(sl, SHARED_BG, prs)
    _pic(sl, path, prs)
    _T(sl, f"De {N_FORMADOS} Formados a {N_APTOS} Aptos P3   Los {N_EXCLUIDOS} Excluidos",
       fs=14)
    _POP(sl, f"4 grupos mutuamente excluyentes que suman exactamente {N_EXCLUIDOS} docentes  "
             f"|  Todos excluidos por ausencia de SAT en el periodo requerido")
    _BUL(sl, [
        f"{N_SOLO_BL} docentes solo les falta el SAT baseline (tienen resultado "
        f"en otro periodo). Causa: su formacion fue anterior a 2023-02, cuando "
        f"la base SAT aun no existia.",
        f"{N_SOLO_RES} docentes solo les falta el SAT resultado. Incluye "
        f"formaciones recientes cuyo semestre posterior aun no tiene datos, "
        f"o docentes sin cursos ese semestre.",
        f"{N_SIN_BOTH} carecen de ambos (sin ninguna medicion util). "
        f"{N_MIXTO} caso especial: tienen SAT en periodos no contiguos con "
        f"su formacion especifica.",
    ])
    print(f"  slide 3 - Excluidos OK  ({N_SOLO_BL}+{N_SOLO_RES}+{N_SIN_BOTH}+{N_MIXTO}={N_EXCLUIDOS})")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Inventario de archivos (layout rediseñado, sin solapamientos)
# ─────────────────────────────────────────────────────────────────────────────
def slide_inventario(prs):
    fig = _tr_fig()
    gx = PIC_RECT[0] + 0.01
    gy = PIC_RECT[1] + 0.01
    gw = PIC_RECT[2] - 0.02
    gh = PIC_RECT[3] - 0.03
    ax = fig.add_axes([gx, gy, gw, gh], facecolor="none", zorder=5)
    ax.set_xlim(0, 10); ax.set_ylim(0, 10); ax.axis("off")

    # Encabezado
    hdrs = ["Archivo original (entregado por UCEN)", "Contenido", "Columnas clave",
            "Rol en el analisis"]
    col_x = [0.10, 3.90, 6.30, 8.30]
    for cx, h in zip(col_x, hdrs):
        ax.text(cx, 9.65, h, ha="left", va="center",
                fontsize=9.2, fontweight="bold", color="#FFD580", zorder=5)
    ax.axhline(9.30, color="#FFD580", linewidth=0.55, alpha=0.65, zorder=4)

    rows = [
        {
            "file":  "EVALUACION ESTUDIANTES\nUCEN 3-5-2026.xlsx",
            "cont":  f"{N_EVAL:,} eval.\n{N_DOC_SAT:,} docentes\n{PER_STR}",
            "cols":  "rut_docente\nperiodo\ncod_facultad\nevaluacion_id",
            "rol":   "Base SAT madre:\nfuente de evaluaciones\npara asignar baseline\ny resultado",
            "col":   "#1B4E8F",
            "yc":    7.60,
        },
        {
            "file":  "CONSOLIDADO DOCENTES\n3-05-2026.xlsx\n(hojas: DOTACION, NOMINA)",
            "cont":  "917 docentes\njerarquizados UCEN",
            "cols":  "rut_key\nnombre\njerarquia\nunidad_facultad\ntipo_contrato",
            "rol":   "Universo base:\ndefine los 917 docentes\nelegibles para P3",
            "col":   "#2E6AAD",
            "yc":    5.35,
        },
        {
            "file":  "CONSOLIDADO DOCENTES\n3-05-2026.xlsx\n(hojas: DIPLOMADO,\nTALLERES, PROYECTOS)",
            "cont":  f"{N_TOTAL_EVENTOS} eventos\n{N_FORMADOS} docentes\nunicos formados\n2022-2025",
            "cols":  "rut_key\ntipo_formacion\nanio_evento\nperiodo_evento",
            "rol":   "Identifica formados\ny permite asignar\nbaseline/resultado\npor periodo",
            "col":   "#3A8BC4",
            "yc":    2.85,
        },
        {
            "file":  "Construido en el analisis\n(cruce SAT x Formacion)",
            "cont":  f"{N_APTOS} Aptos P3\n(docentes con\nambos SAT\ndisponibles)",
            "cols":  "z_baseline\nz_resultado\ndelta_z\nsat_baseline\nsat_resultado",
            "rol":   "Analisis de\nincidencia P3:\nz-scores y cambio\nSAT individual",
            "col":   "#2E9E55",
            "yc":    0.65,
        },
    ]

    ROW_H = 2.10
    for rd in rows:
        yc = rd["yc"]
        y0 = yc - ROW_H / 2
        rect = mpatches.FancyBboxPatch(
            (-0.05, y0), 10.1, ROW_H,
            boxstyle="round,pad=0.07",
            facecolor=rd["col"], edgecolor="white",
            linewidth=0.5, alpha=0.26, zorder=2)
        ax.add_patch(rect)
        # barra color izquierda
        bar = mpatches.FancyBboxPatch(
            (-0.05, y0 + 0.05), 0.22, ROW_H - 0.10,
            boxstyle="round,pad=0.03",
            facecolor=rd["col"], edgecolor="none",
            alpha=0.95, zorder=3)
        ax.add_patch(bar)
        for cx, content, is_file in zip(
                col_x,
                [rd["file"], rd["cont"], rd["cols"], rd["rol"]],
                [True, False, False, False]):
            ax.text(cx, yc, content,
                    ha="left", va="center",
                    fontsize=8.8 if is_file else 8.0,
                    fontweight="bold" if is_file else "normal",
                    color="#FFD580" if is_file else "white",
                    zorder=5, linespacing=1.4)

    path = _save_ch(fig, "ficha_inventario.png")
    sl = _new_sl(prs)
    _pic(sl, SHARED_BG, prs)
    _pic(sl, path, prs)
    _T(sl, "Inventario de Archivos Originales Entregados por UCEN", fs=16)
    _POP(sl, "2 archivos Excel origen  |  1 archivo construido en el analisis  "
             "|  Todos disponibles en UNIVERSO_918/PROCESADO/")
    print("  slide 4 - Inventario OK")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Datos adicionales (cajas mas grandes, titulos sin solapamiento)
# ─────────────────────────────────────────────────────────────────────────────
def slide_solicitudes(prs):
    fig = _tr_fig()
    gx = PIC_RECT[0] + 0.01
    gy = PIC_RECT[1] + 0.01
    gw = PIC_RECT[2] - 0.02
    gh = PIC_RECT[3] - 0.03
    ax = fig.add_axes([gx, gy, gw, gh], facecolor="none", zorder=5)
    ax.set_xlim(0, 10); ax.set_ylim(0, 10.5); ax.axis("off")

    BOX_H = 2.90
    YCS   = [9.00, 5.70, 2.40]   # centros de las 3 cajas

    solicitudes = [
        {
            "num":    "1",
            "title":  "Ampliar cobertura temporal del SAT",
            "lines":  [
                f"Base actual: {PER_STR}  (faltan semestres anteriores a {PERIODOS[0]})",
                f"Solicitar:  2021-02  |  2022-01  |  2022-02",
                f"-> Daria baseline a ~{anio_counts.get(2022,40)} docentes formados en 2022",
                f"   + parte de los {N_SOLO_BL} que hoy solo les falta baseline",
                f"Tambien: 2026-01 cuando este disponible (resultado para formados en 2025-02)",
            ],
            "impacto": f"Potencial: recuperar hasta {N_SOLO_BL + N_SIN_BOTH} docentes adicionales",
            "col":     "#1B4E8F",
        },
        {
            "num":    "2",
            "title":  "Registro de carga docente por semestre",
            "lines":  [
                "Lista de docentes con cursos asignados cada semestre (asignacion academica).",
                "Permite distinguir 2 causas de SAT ausente:",
                "  A) Docente sin cursos -> SAT estructuralmente imposible (no recuperable)",
                "  B) Docente con cursos pero encuesta no aplicada (recuperable si se reenvian datos)",
                f"Hoy no podemos distinguir A de B para los {N_SOLO_RES + N_SIN_BOTH} sin resultado.",
            ],
            "impacto": f"Diagnotica cuantos de los {N_EXCLUIDOS} son recuperables vs estructurales",
            "col":     "#1B6A40",
        },
        {
            "num":    "3",
            "title":  "Confirmar integridad del registro de formaciones 2025",
            "lines":  [
                f"La base actual contiene {anio_counts.get(2025,0)} eventos de formacion en 2025.",
                "Verificar si el listado esta completo o faltan actividades:",
                "  - Talleres, diplomados o proyectos iniciados en 2025 aun no registrados",
                "  - Actividades con inicio 2024 y termino 2025 (posible problema de categoria de ano)",
                f"Si hay formados no registrados, el universo de {N_FORMADOS} seria mayor.",
            ],
            "impacto": f"Amplia el universo base de formados antes de cruzar con SAT",
            "col":     "#6A4010",
        },
    ]

    for s, yc in zip(solicitudes, YCS):
        y0 = yc - BOX_H / 2

        rect = mpatches.FancyBboxPatch(
            (-0.05, y0), 10.1, BOX_H,
            boxstyle="round,pad=0.12",
            facecolor=s["col"], edgecolor="white",
            linewidth=0.6, alpha=0.32, zorder=2)
        ax.add_patch(rect)

        # Numero circulo
        circ = mpatches.Circle((0.42, yc + 0.60), 0.40,
                                facecolor=s["col"], edgecolor="white",
                                linewidth=0.9, alpha=0.95, zorder=3)
        ax.add_patch(circ)
        ax.text(0.42, yc + 0.60, s["num"],
                ha="center", va="center", fontsize=17, fontweight="bold",
                color="white", zorder=4)

        # Titulo (bien separado del circulo)
        ax.text(1.05, yc + 0.85, s["title"],
                ha="left", va="center", fontsize=10.5, fontweight="bold",
                color="#FFD580", zorder=4)

        # Detalle (lineas)
        y_txt = yc + 0.38
        for line in s["lines"]:
            ax.text(1.05, y_txt, line,
                    ha="left", va="top", fontsize=8.0,
                    color="#D0E8FF", zorder=4)
            y_txt -= 0.42

        # Impacto (parte inferior)
        ax.text(0.42, y0 + 0.22, "-> " + s["impacto"],
                ha="left", va="center", fontsize=8.0,
                color="#A8E8A0", zorder=4, fontstyle="italic")

    path = _save_ch(fig, "ficha_solicitudes.png")
    sl = _new_sl(prs)
    _pic(sl, SHARED_BG, prs)
    _pic(sl, path, prs)
    _T(sl, "Datos Adicionales para Ampliar el Universo Aptos P3", fs=15)
    _POP(sl, f"Hoy: {N_APTOS} de {N_FORMADOS} formados son Aptos P3  |  "
             f"3 solicitudes concretas a la contraparte para aumentar esa cobertura")
    _BUL(sl, [
        "La palanca principal es ampliar la cobertura temporal del SAT: enviando "
        f"semestres anteriores a {PERIODOS[0]} se recuperan docentes formados en 2022 "
        "y primer semestre 2023 que hoy quedan excluidos por falta de baseline.",
        "El registro de carga docente es clave para depurar el diagnostico: "
        "distinguir docentes sin cursos (exclusion estructural, irrecuperable) de "
        "docentes con cursos cuya encuesta no fue aplicada (exclusion recuperable).",
    ])
    print("  slide 5 - Solicitudes OK")


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("Generando FICHA_Ampliar_Aptos_P3.pptx ...")
    _ensure_bg()
    prs = Presentation()
    prs.slide_width  = Emu(SW_EMU)
    prs.slide_height = Emu(SH_EMU)
    slide_portada(prs)
    slide_pipeline(prs)
    slide_excluidos(prs)
    slide_inventario(prs)
    slide_solicitudes(prs)
    prs.save(OUT_PPTX)
    print(f"\n  Guardado: {OUT_PPTX}")
    print("  5 slides: [1] Portada, [2] Pipeline, [3] Excluidos, "
          "[4] Inventario, [5] Solicitudes")

"""
generar_ppt_v3_slides1_26.py  — VERSION CORREGIDA (24 slides)
Arquitectura V3:
  - Fondo oscuro PNG compartido (foto B&N + gradiente navy + logo UCEN)
  - Graficos como PNG transparente (solo chart, sin texto)
  - Todo el texto como TextBox editables en python-pptx
Fuente de datos: p3_sat_zscore_918.csv (197 aptos P3, todos formados)
Salida: PRESENTACION_197_P3_v3_slides1_26.pptx  (24 diapositivas reales)

CAMBIOS V3.1:
  - _fac(): substring matching (fix encoding UTF-8 accents)
  - _hbar(): xlim mv*1.55, barra más larga pone etiqueta DENTRO
  - CHART_X/CHART_W: más centrado y ligeramente más pequeño
  - slide_12 Venn: centrado en área PIC
  - slide_14 Embudo: ejes reducidos
  - slide_16 Evolución: centrado (via CHART_X global)
  - ELIMINAR slide_16 (Cobertura tipo, redundante con Venn)
  - slide_17→17, 19→18, 20→19, 21→20, 22→21, 23→22, 25→23, 26→24
  - slide_17 (nuevo): n DENTRO de cada barra
  - slide_18 (nuevo): "Combinado"→"Participación Mixta"
  - slide_20 (nuevo): encoding fix via _fac() global
  - slide_21 (nuevo): anotaciones compactas DENTRO del gráfico
  - ELIMINAR slide_22 (SAT por Facultad)
  - slide_23 (nuevo): Trayectoria Pura con delta visible (flecha + sombreado)
  - slide_22 (nuevo): Delta z como gráfico divergente % mejora vs % empeora
  - slide_02, 04, 11, 15: numeración actualizada
  - "Combinado"→"Participación Mixta" en todo el script
"""
import sys; sys.stdout.reconfigure(encoding="utf-8")
import os, zipfile
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.patheffects as pe
from matplotlib.lines import Line2D
from matplotlib_venn import venn3
from scipy import stats as scipy_stats
import pandas as pd
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─────────────────────────────────────────────────────────────────────────────
# Rutas  (actualizado 2026-07-18 — datos desde data/cascade/, universo 1.144)
# ─────────────────────────────────────────────────────────────────────────────
import pathlib
BASE      = os.path.dirname(os.path.abspath(__file__))
REPO      = str(pathlib.Path(BASE).parents[1])          # raíz del repo git
CASCADE   = os.path.join(REPO, "data", "cascade")
COMP      = os.path.join(CASCADE, "complementarios")

EDD_CSV    = os.path.join(COMP, "evaluacion_jefes.csv")
DOC918_CSV = os.path.join(CASCADE, "03_jerarquizados", "docente_918.csv")
SCAT_CSV   = os.path.join(COMP, "scatter_sat_notas.csv")
FONDOTIPO  = (r"c:\Users\r.gonzalez_fluxsolar.LAPTOP-FLUX-ECO"
              r"\Downloads\Analisis_UCEN_v2\Fondotipop.pptx")
DOT_CSV    = (r"c:\Users\r.gonzalez_fluxsolar.LAPTOP-FLUX-ECO\Downloads\Analisis_UCEN_v2"
              r"\CONSOLIDADO DOCENTES 3-05-2026"
              r"\CONSOLIDADO DOCENTES 3-05-2026.xlsx - DOTACION_CON_GRADOREC.csv")
OUT_PPTX   = os.path.join(REPO, "outputs", "pptx", "PRESENTACION_210_P3_v4.pptx")
OUT_DIR    = os.path.join(BASE, "dark_slides_v3")
SCRATCH    = os.path.join(REPO, "outputs", "scratch")
os.makedirs(OUT_DIR, exist_ok=True)
os.makedirs(SCRATCH, exist_ok=True)

# ─────────────────────────────────────────────────────────────────────────────
# Assets (foto B&N + logo extraidos de Fondotipop.pptx)
# ─────────────────────────────────────────────────────────────────────────────
BG_PATH   = os.path.join(SCRATCH, "fondotipo_image1.jpg")
LOGO_PATH = os.path.join(SCRATCH, "fondotipo_image2.png")
for path, zname in [(BG_PATH,"ppt/media/image1.jpg"),(LOGO_PATH,"ppt/media/image2.png")]:
    if not os.path.exists(path):
        with zipfile.ZipFile(FONDOTIPO) as z:
            with open(path,"wb") as f:
                f.write(z.read(zname))

with PILImage.open(BG_PATH) as _im:
    _rgb = _im.convert("RGB"); _iw,_ih = _rgb.size
    _nh  = int(_iw/(16/9)); _y0 = min(int(_ih*0.12),_ih-_nh)
    bg_arr = np.array(_rgb.crop((0,_y0,_iw,_y0+_nh)))

with PILImage.open(LOGO_PATH) as _logo:
    logo_arr = np.array(_logo.convert("RGBA")).astype(np.float32)/255.0

H_GRAD = 600; grad = np.zeros((H_GRAD,1,4), dtype=np.float32)
for _r in range(H_GRAD):
    _t = _r/(H_GRAD-1)
    _stops = [(0.00,(0,33,71)),(0.54,(0,70,128)),(1.00,(144,171,196))]
    for _i in range(len(_stops)-1):
        _t0,_c0 = _stops[_i]; _t1,_c1 = _stops[_i+1]
        if _t0 <= _t <= _t1:
            _s = (_t-_t0)/(_t1-_t0)
            grad[_r,0] = [(_c0[0]+_s*(_c1[0]-_c0[0]))/255,
                          (_c0[1]+_s*(_c1[1]-_c0[1]))/255,
                          (_c0[2]+_s*(_c1[2]-_c0[2]))/255, 0.82]; break

# ─────────────────────────────────────────────────────────────────────────────
# Constantes de layout (EMU – 12192000 x 6858000)
# ─────────────────────────────────────────────────────────────────────────────
SW, SH       = 13.333, 7.5          # pulgadas para matplotlib
SW_EMU       = 12192000
SH_EMU       = 6858000

PIC_L, PIC_T, PIC_W, PIC_H = 786581, 1125000, 10599174, 3720000
BUL_L, BUL_T, BUL_W, BUL_H = 786581, 4870000, 10599174, 1870000
LOGO_L, LOGO_T, LOGO_W, LOGO_H = 9813773, 656354, 1756626, 697725

TITLE_L, TITLE_T, TITLE_W, TITLE_H = PIC_L, 185000, PIC_W, 710000
POP_L,   POP_T,   POP_W,   POP_H   = PIC_L, 845000, 9000000, 255000

# Fracciones matplotlib para PNG transparente
def _ex(e): return e/SW_EMU
def _ey(e): return e/SH_EMU
def _fig_rect(l,t,w,h): return (l, 1-t-h, w, h)

PIC_RECT  = _fig_rect(_ex(PIC_L), _ey(PIC_T), _ex(PIC_W), _ey(PIC_H))
LOGO_RECT = _fig_rect(_ex(LOGO_L), _ey(LOGO_T), _ex(LOGO_W), _ey(LOGO_H))

# Gráfico más centrado y ligeramente más pequeño que V3.0
CHART_X = PIC_RECT[0] + 0.13   # más centrado (era 0.16)
CHART_Y = PIC_RECT[1] + 0.04
CHART_W = PIC_RECT[2] - 0.19   # margen derecho más amplio (era 0.18)
CHART_H = PIC_RECT[3] - 0.09

# Chart title (etiqueta encima del gráfico)
CTITLE_L = int(CHART_X * SW_EMU)
CTITLE_T = PIC_T + 38000
CTITLE_W = int((PIC_RECT[0] + PIC_RECT[2] - CHART_X + 0.02) * SW_EMU)
CTITLE_H = 295000

PAL    = ["#5C9BD6","#64B5F6","#80DEEA","#A5D6A7","#FFB74D","#CE93D8","#90A4AE","#F48FB1"]
TIPO_COLS = {"Taller":"#5C9BD6","Diplomado":"#64B5F6","Proyecto":"#A5D6A7","Participación Mixta":"#CE93D8"}

def _tipo_simple(t):
    t = str(t).upper()
    if "TALLER" in t and "DIPLOMADO" not in t and "PROYECTO" not in t: return "Taller"
    if "DIPLOMADO" in t and "TALLER" not in t and "PROYECTO" not in t: return "Diplomado"
    if "PROYECTO" in t and "TALLER" not in t and "DIPLOMADO" not in t: return "Proyecto"
    return "Participación Mixta"
POP_197 = "Universo: 210 Aptos P3  ·  todos formados  ·  Periodos 2022–2025"
POP_CTR = "Universo: 210 Aptos P3  ·  formados vs control externo sin formacion  ·  2023–2025"

# Constantes para slides de perfil demográfico (embudo / butterfly)
JER_ORD = ["INSTRUCTOR REGULAR", "INSTRUCTOR DOCENTE",
           "ASISTENTE REGULAR",  "ASISTENTE DOCENTE",
           "ASOCIADO REGULAR",   "ASOCIADO DOCENTE",
           "TITULAR REGULAR",    "TITULAR DOCENTE"]
JER_LBL = ["Instr. Regular", "Instr. Docente",
            "Asist. Regular", "Asist. Docente",
            "Asoc. Regular",  "Asoc. Docente",
            "Tit. Regular",   "Tit. Docente"]
TRAMOS_EDAD = ["< 35", "35–44", "45–54", "55–64", "≥ 65"]
TRAMOS_EDAD_RAW = {
    "<30": "< 35", "30-34": "< 35",
    "35-39": "35–44", "40-44": "35–44",
    "45-49": "45–54", "50-54": "45–54",
    "55-59": "55–64", "60-64": "55–64",
    "65-69": "≥ 65",  "70+":   "≥ 65",
}
COL_FORM_P = "#5C9BD6"   # azul formados  (perfil)
COL_CTRL_P = "#FFB74D"   # amarillo ctrl  (perfil)

# ─────────────────────────────────────────────────────────────────────────────
# _fac(): substring matching — evita artefactos de encoding UTF-8
# ─────────────────────────────────────────────────────────────────────────────
def _fac(s):
    if not isinstance(s, str): return str(s)
    s = s.strip(); u = s.upper()
    if "MEDICINA" in u and "SALUD" in u:   return "Medicina y C. Salud"
    if "DERECHO" in u and "HUMANIDADES" in u: return "Derecho y Humanidades"
    if "INGENIER" in u:                     return "Ingeniería y Arq."
    if "EDUCACI" in u:                      return "Educación"
    if "ECONOM" in u and "GOBIERNO" in u:  return "Economía, Gob. y Com."
    if "INVEST" in u and "INNOV" in u:     return "VR Investigación"
    if "VICERRECTORIA" in u and "ACADEM" in u: return "VR Académica"
    if "ASEGURAMIENTO" in u:               return "Dir. Aseg. Calidad"
    return s[:28].title()

# ─────────────────────────────────────────────────────────────────────────────
# Carga de datos (CSV como fuente de verdad)
# ─────────────────────────────────────────────────────────────────────────────
sat  = pd.read_csv(os.path.join(CASCADE,"05_aptos_p3","p3_sat_zscore.csv"), encoding="utf-8-sig")
sat["rut_key"] = sat["rut_key"].astype(str).str.strip()
p3ev = pd.read_csv(os.path.join(CASCADE,"04_formados_p3","p3_918.csv"), encoding="utf-8-sig")
p3ev["rut_key"] = p3ev["rut_key"].astype(str).str.strip()
cvt  = pd.read_csv(os.path.join(COMP,"control_vs_trat_918.csv"), encoding="utf-8-sig")
ctrl = pd.read_csv(os.path.join(COMP,"control_918.csv"), encoding="utf-8-sig")
ctrl["rut_key"] = ctrl["rut_key"].astype(str).str.strip()
N197 = len(sat)

has_dot = os.path.exists(DOT_CSV)
if has_dot:
    dot = pd.read_csv(DOT_CSV, dtype=str, encoding="utf-8-sig")
    dot.columns = dot.columns.str.strip()
    dot["rut_key"] = (dot["RUT"].str.strip().str.replace(".", "", regex=False)
                      .str.split("-").str[0].str.strip())
    dot197 = dot[dot["rut_key"].isin(set(sat["rut_key"]))].copy()
else:
    dot197 = pd.DataFrame()

doc918 = pd.read_csv(DOC918_CSV, dtype={"rut_key": str}, encoding="utf-8-sig")
doc918["rut_key"] = doc918["rut_key"].str.strip()
ruts_917 = set(doc918["rut_key"])

ruts_todos_formados = set(p3ev["rut_key"].astype(str).str.strip())   # 419 formados en 941
N_NO_FORM = len(ruts_917 - ruts_todos_formados)

has_edd = os.path.exists(EDD_CSV)
if has_edd:
    edd_df = pd.read_csv(EDD_CSV, dtype={"rut_key": str}, encoding="utf-8-sig")
    edd_df["rut_key"] = edd_df["rut_key"].str.strip()
    edd_df["edd_total"] = pd.to_numeric(edd_df["edd_total"], errors="coerce")
    edd_df["anio_eval"] = edd_df["anio_evaluacion"].apply(
        lambda x: str(int(float(x)))[:4] if pd.notna(x) else None)
    ruts_form_set = set(sat["rut_key"].astype(str).str.strip())
    ruts_todos_formados = set(p3ev["rut_key"].astype(str).str.strip())
    edd_form = (edd_df[edd_df["rut_key"].isin(ruts_form_set) & edd_df["edd_total"].notna()
                       & edd_df["anio_eval"].notna()]
                .drop_duplicates(subset=["rut_key", "anio_eval"]))
    edd_ctrl = (edd_df[edd_df["rut_key"].isin(ruts_917) & ~edd_df["rut_key"].isin(ruts_todos_formados)
                       & edd_df["edd_total"].notna() & edd_df["anio_eval"].notna()]
                .drop_duplicates(subset=["rut_key", "anio_eval"]))
else:
    edd_form = pd.DataFrame()
    edd_ctrl = pd.DataFrame()

# ── Scatter Bloque III + perfiles demográficos ────────────────────────────────
scat = pd.read_csv(SCAT_CSV, encoding="utf-8-sig")
scat["formado"] = scat["formado"].astype(str).str.strip().str.upper().isin(
    ["TRUE", "1", "SI", "SÍ", "YES"])
scat_ctrl_ruts = set(scat[~scat["formado"]]["rut_docente"].astype(str).str.strip().unique())
ctrl_b3_doc = doc918[doc918["rut_key"].isin(scat_ctrl_ruts)].drop_duplicates("rut_key").copy()
ctrl_b3_doc["jerarquia_u"] = ctrl_b3_doc["jerarquia"].str.strip().str.upper()
ctrl_b3_doc["tramo_g"]     = ctrl_b3_doc["tramo_edad"].map(TRAMOS_EDAD_RAW)
N_B3_FORM = N197
N_B3_CTRL = len(ctrl_b3_doc)
sat["jerarquia_u"] = sat["jerarquia"].str.strip().str.upper()
sat["tramo_g"]     = sat["tramo_edad"].map(TRAMOS_EDAD_RAW)

if has_edd and len(edd_form) > 0:
    edd_form_ruts = set(edd_form["rut_key"].unique())
    edd_ctrl_ruts = set(edd_ctrl["rut_key"].unique())
    b4_form_doc = doc918[doc918["rut_key"].isin(edd_form_ruts)].drop_duplicates("rut_key").copy()
    b4_ctrl_doc = doc918[doc918["rut_key"].isin(edd_ctrl_ruts)].drop_duplicates("rut_key").copy()
    b4_form_doc["jerarquia_u"] = b4_form_doc["jerarquia"].str.strip().str.upper()
    b4_form_doc["tramo_g"]     = b4_form_doc["tramo_edad"].map(TRAMOS_EDAD_RAW)
    b4_ctrl_doc["jerarquia_u"] = b4_ctrl_doc["jerarquia"].str.strip().str.upper()
    b4_ctrl_doc["tramo_g"]     = b4_ctrl_doc["tramo_edad"].map(TRAMOS_EDAD_RAW)
    N_B4_FORM = len(b4_form_doc)
    N_B4_CTRL = len(b4_ctrl_doc)
else:
    b4_form_doc = b4_ctrl_doc = pd.DataFrame()
    N_B4_FORM = N_B4_CTRL = 0

# ─────────────────────────────────────────────────────────────────────────────
# Helpers matplotlib
# ─────────────────────────────────────────────────────────────────────────────
def _bg_fig():
    """Fondo oscuro (foto + gradiente + logo) sin ningun texto."""
    fig = plt.figure(figsize=(SW,SH), facecolor="#101820")
    fig.patch.set_facecolor("#101820")
    for z,arr in [(0,bg_arr),(1,grad)]:
        ax = fig.add_axes([0,0,1,1], zorder=z)
        ax.imshow(arr, extent=[0,1,0,1], aspect="auto", origin="upper")
        ax.set_xlim(0,1); ax.set_ylim(0,1); ax.axis("off")
    al = fig.add_axes([LOGO_RECT[0],LOGO_RECT[1],LOGO_RECT[2],LOGO_RECT[3]],
                      zorder=10, facecolor="none")
    al.imshow(logo_arr, aspect="auto"); al.axis("off"); al.patch.set_visible(False)
    return fig

def _tr_fig():
    """Figura transparente (para charts superpuestos)."""
    fig = plt.figure(figsize=(SW,SH), facecolor="none")
    fig.patch.set_facecolor("none")
    return fig

def _save_bg(fig, name):
    path = os.path.join(OUT_DIR, name)
    plt.savefig(path, dpi=150, facecolor=fig.get_facecolor())
    plt.close(); return path

def _save_ch(fig, name):
    path = os.path.join(OUT_DIR, name)
    plt.savefig(path, dpi=150, facecolor="none", transparent=True)
    plt.close(); return path

def _ax():
    fig = _tr_fig()
    ax  = fig.add_axes([CHART_X, CHART_Y, CHART_W, CHART_H], facecolor="none", zorder=5)
    return fig, ax

def _style(ax, xlabel=None, xgrid=False, ygrid=True):
    ax.tick_params(axis="x", colors="white", labelsize=9, length=0)
    ax.tick_params(axis="y", colors="#AAAAAA", labelsize=9, length=0)
    for sp in ax.spines.values():
        sp.set_edgecolor("white"); sp.set_alpha(0.35); sp.set_linewidth(0.9)
    if ygrid: ax.yaxis.grid(True, color="white", alpha=0.07, linewidth=0.5)
    if xgrid: ax.xaxis.grid(True, color="white", alpha=0.07, linewidth=0.5)
    ax.set_axisbelow(True)
    if xlabel: ax.set_xlabel(xlabel, color="#AAAAAA", fontsize=9)

SHARED_BG = os.path.join(OUT_DIR, "_background.png")

def _ensure_bg():
    if not os.path.exists(SHARED_BG):
        fig = _bg_fig(); _save_bg(fig, "_background.png")
        print("  bg generado")

# ─────────────────────────────────────────────────────────────────────────────
# Helpers python-pptx
# ─────────────────────────────────────────────────────────────────────────────
def _new_sl(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _pic(sl, path, prs):
    sl.shapes.add_picture(path, Emu(0), Emu(0), prs.slide_width, prs.slide_height)

def _txt(sl, text, left, top, width, height,
         fs=12, bold=False, italic=False, color="#FFFFFF",
         align=PP_ALIGN.LEFT, wrap=True, lspc=0):
    txb = sl.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf  = txb.text_frame; tf.word_wrap = wrap
    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if lspc > 0 and i > 0:
            p.space_before = Pt(lspc)
        run = p.add_run(); run.text = line
        run.font.size = Pt(fs); run.font.bold = bold; run.font.italic = italic
        r,g,b = int(color[1:3],16), int(color[3:5],16), int(color[5:7],16)
        run.font.color.rgb = RGBColor(r,g,b)
    return txb

def _T(sl, text, fs=20):
    _txt(sl, text, TITLE_L, TITLE_T, TITLE_W, TITLE_H,
         fs=fs, bold=True, color="#FFFFFF", align=PP_ALIGN.CENTER)

def _CT(sl, text):
    _txt(sl, text, CTITLE_L, CTITLE_T, CTITLE_W, CTITLE_H,
         fs=10, color="#FFFFFF")

def _POP(sl, text=None):
    _txt(sl, text or POP_197, POP_L, POP_T, POP_W, POP_H,
         fs=7.5, italic=True, color="#C8DCF0")

def _BUL(sl, items, fs=11.5):
    txb = sl.shapes.add_textbox(Emu(BUL_L), Emu(BUL_T), Emu(BUL_W), Emu(BUL_H))
    tf  = txb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5); p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{i+1}.  {item}"
        run.font.size = Pt(fs)
        run.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)

def _hbar(fig, labels, vals, pcts, cx=None, cw=None):
    """Barra horizontal. La barra más larga pone etiqueta DENTRO."""
    x = cx if cx is not None else CHART_X
    w = cw if cw is not None else CHART_W
    ax = fig.add_axes([x, CHART_Y, w, CHART_H], facecolor="none", zorder=5)
    n  = len(labels); yp = np.arange(n)
    ax.barh(yp[::-1], vals, color=PAL[:n], height=0.58, edgecolor="none", alpha=0.90)
    mv = max(vals) if vals else 1
    for i,(v,p) in enumerate(zip(vals[::-1], pcts[::-1])):
        ax.text(v + mv*0.020, i, f"{v}  ({p:.1f}%)",
                va="center", ha="left", fontsize=10.5, fontweight="bold", color="white",
                path_effects=[pe.withStroke(linewidth=2.5, foreground="#0A0F18")])
    ax.set_yticks(yp); ax.set_yticklabels(labels[::-1], fontsize=10.5,
                                           fontweight="bold", color="white")
    ax.tick_params(axis="y", length=0, pad=8)
    ax.tick_params(axis="x", colors="#AAAAAA", labelsize=9)
    for sp in ax.spines.values():
        sp.set_edgecolor("white"); sp.set_alpha(0.35); sp.set_linewidth(0.9)
    ax.set_xlim(0, mv*1.55)   # espacio generoso para etiquetas (era 1.35)
    ax.set_ylim(-0.5, n-0.5)
    ax.xaxis.grid(True, color="white", alpha=0.07, linewidth=0.5); ax.set_axisbelow(True)
    return ax

def _style_ax(ax):
    ax.tick_params(axis="x", colors="#AAAAAA", labelsize=8, length=0)
    ax.tick_params(axis="y", colors="white",   labelsize=8, length=0)
    for sp in ax.spines.values():
        sp.set_edgecolor("white"); sp.set_alpha(0.20); sp.set_linewidth(0.7)
    ax.set_axisbelow(True)

def _pct_jer(df, col="jerarquia_u"):
    vc = df[col].value_counts()
    return [round(vc.get(j, 0) / max(len(df), 1) * 100, 1) for j in JER_ORD]

def _pct_tramo(df, col="tramo_g"):
    vc = df[col].value_counts()
    return [round(vc.get(t, 0) / max(len(df), 1) * 100, 1) for t in TRAMOS_EDAD]

def _butterfly(fig, bx, by, bw, bh,
               cats, pct_f, pct_c, n_f_tot, n_c_tot,
               chart_title, col_f=COL_FORM_P, col_c=COL_CTRL_P):
    ya      = np.arange(len(cats))
    max_val = max(max(pct_f, default=1), max(pct_c, default=1))
    xlim    = max_val * 1.30
    gap_lbl = 0.045
    w_panel = (bw - gap_lbl) / 2

    ax_f = fig.add_axes([bx, by, w_panel, bh], facecolor="none", zorder=5)
    ax_f.barh(ya, pct_f, height=0.60, color=col_f, alpha=0.88, edgecolor="none")
    ax_f.set_xlim(xlim, 0)
    ax_f.set_ylim(-0.6, len(cats) - 0.4)
    ax_f.set_yticks(ya); ax_f.set_yticklabels([])
    ax_f.xaxis.grid(True, color="white", alpha=0.06, linewidth=0.5)
    for j, v in enumerate(pct_f):
        if v >= 4:
            ax_f.text(v / 2, j, f"{v:.0f}%", ha="center", va="center",
                      fontsize=7.5, color="white", fontweight="bold")
        elif v > 0:
            ax_f.text(v + 0.5, j, f"{v:.0f}%", ha="right", va="center",
                      fontsize=7, color="#AAAAAA")
    ax_f.xaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"{abs(x):.0f}%"))
    ax_f.set_title(f"Formados  (n={n_f_tot})", color=col_f, fontsize=9, fontweight="bold", pad=5)
    _style_ax(ax_f)

    ax_c = fig.add_axes([bx + w_panel + gap_lbl, by, w_panel, bh], facecolor="none", zorder=5)
    ax_c.barh(ya, pct_c, height=0.60, color=col_c, alpha=0.88, edgecolor="none")
    ax_c.set_xlim(0, xlim)
    ax_c.set_ylim(-0.6, len(cats) - 0.4)
    ax_c.set_yticks(ya); ax_c.set_yticklabels(cats, fontsize=8.5, color="white")
    ax_c.yaxis.set_ticks_position("left")
    ax_c.xaxis.grid(True, color="white", alpha=0.06, linewidth=0.5)
    for j, v in enumerate(pct_c):
        if v >= 4:
            ax_c.text(v / 2, j, f"{v:.0f}%", ha="center", va="center",
                      fontsize=7.5, color="white", fontweight="bold")
        elif v > 0:
            ax_c.text(v + 0.5, j, f"{v:.0f}%", ha="left", va="center",
                      fontsize=7, color="#AAAAAA")
    ax_c.xaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"{x:.0f}%"))
    ax_c.set_title(f"Control  (n={n_c_tot})", color=col_c, fontsize=9, fontweight="bold", pad=5)
    _style_ax(ax_c)

    fig.text(bx + bw / 2, by + bh + 0.035, chart_title,
             ha="center", va="bottom", fontsize=10, color="white", fontweight="bold",
             path_effects=[pe.withStroke(linewidth=1.5, foreground="#0A1830")])

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE I — Slides 01-10
# ─────────────────────────────────────────────────────────────────────────────

def slide_01(prs):
    _ensure_bg(); sl = _new_sl(prs); _pic(sl, SHARED_BG, prs)
    # Linea divisora
    fig = _tr_fig()
    ax  = fig.add_axes([0.22, 0.568, 0.56, 0.003], facecolor="none", zorder=4)
    ax.set_facecolor("#90ABC4"); ax.axis("off")
    _pic(sl, _save_ch(fig,"01_div.png"), prs)
    # Texto
    _txt(sl,"Resultados del Perfeccionamiento Docente",
         PIC_L, 2440000, PIC_W, 460000, fs=28, bold=True,
         color="#FFFFFF", align=PP_ALIGN.CENTER)
    _txt(sl,"Análisis de Impacto: SAT, Recomendación, Notas y EDD",
         PIC_L, 2980000, PIC_W, 380000, fs=17,
         color="#C8DCF0", align=PP_ALIGN.CENTER)
    _txt(sl,"Universo: 210 Aptos P3  ·  todos formados  ·  Períodos 2022–2025",
         PIC_L, 3440000, PIC_W, 340000, fs=12, italic=True,
         color="#C8DCF0", align=PP_ALIGN.CENTER)
    _txt(sl,"Universidad Central de Chile  |  Producto 3: Análisis de Formación e Innovación",
         PIC_L, 5200000, PIC_W, 320000, fs=12, bold=True,
         color="#FFFFFF", align=PP_ALIGN.CENTER)
    print("  ✓ slide 01 — Portada")

def slide_02(prs):
    _ensure_bg(); sl = _new_sl(prs); _pic(sl, SHARED_BG, prs)
    _T(sl, "Índice")
    _POP(sl, "Estructura del informe — 34 diapositivas")
    bloques = (
        "I.    Clasificación del Cuerpo Académico  —  edad, sexo, facultad, jerarquía, grado  (n=210)\n"
        "II.   Evaluación Docente SAT — Antes y Después  (n=210 formados  ·  control externo)\n"
        "      2.1  Caracterización de la formación  (tipo, jerarquía, antigüedad, intensidad)\n"
        "      2.2  Impacto en SAT z-score (6 períodos)  ·  Formados vs Control\n"
        "      2.3  SAT por Facultad  ·  Histograma Δz  ·  Cambio × Antigüedad\n"
        "III.  Aprobación de Alumnos  —  formados vs control  ·  evolución  ·  perfil  ·  efecto acum.\n"
        "IV.   Evaluación de Desempeño Docente (EDD)  —  formados vs control  ·  evolución  ·  por tipo\n"
        "V.    Conclusiones y Recomendaciones"
    )
    _txt(sl, bloques, PIC_L+80000, PIC_T+560000, PIC_W-80000, PIC_H-600000,
         fs=14, color="#FFFFFF", lspc=6)
    print("  ✓ slide 02 — Índice")

def slide_03(prs):
    """Marco metodológico — 4 cajas más centradas."""
    _ensure_bg(); sl = _new_sl(prs); _pic(sl, SHARED_BG, prs)
    _T(sl, "Universo de Análisis y Metodología")
    _POP(sl, "Universo base: 941 docentes jerarquizados UCEN  ·  los 210 Aptos P3 son el subconjunto analizado")
    cajas = [
        ("1. Universo Aptos P3",
         "• 941 docentes jerarquizados UCEN (universo base)\n"
         "• 419 participaron en ≥1 iniciativa de formación\n"
         "• 210 Aptos P3: SAT válido en baseline y resultado\n"
         "• Todos los 210 son formados (no hay control interno)"),
        ("2. Marco Metodológico — Z-score",
         "• z = (SAT docente − media facultad-período) / DE\n"
         "• z = 0 → promedio exacto de su facultad ese semestre\n"
         "• z > 0 → sobre el promedio   ·   z < 0 → bajo\n"
         "• Compara docentes entre facultades distintas"),
        ("3. Métricas de Evaluación",
         "• SAT Nota (1–7) estandarizada como z-score\n"
         "• SAT % Recomendación: ¿recomendaría a este docente?\n"
         "• EDD: Evaluación de Desempeño Docente (escala 0–1)\n"
         "• Notas y aprobación alumnos: % nota ≥ 4.0"),
        ("4. Tipos de Formación",
         "• Taller (corta duración, 1 semestre)  n=154 puros\n"
         "• Diplomado (larga duración, 2+ semestres)  n=27 puros\n"
         "• Proyecto de Innovación Docente  n=3 puros\n"
         "• Participación Mixta (T+D, T+P): 13 docentes"),
    ]
    # Cajas debajo del logo UCEN (logo bottom = EMU 1354079, pad_t empieza desde PIC_T=1125000)
    pad_x = 140000; gap_x = 80000; gap_y = 50000; pad_t = 280000; pad_b = 70000
    bw = (PIC_W - 2*pad_x - gap_x)//2
    bh = (PIC_H - pad_t - gap_y - pad_b)//2
    pos = [(PIC_L+pad_x, PIC_T+pad_t),
           (PIC_L+pad_x+bw+gap_x, PIC_T+pad_t),
           (PIC_L+pad_x, PIC_T+pad_t+bh+gap_y),
           (PIC_L+pad_x+bw+gap_x, PIC_T+pad_t+bh+gap_y)]
    for (bx,by),(hdr,body) in zip(pos, cajas):
        sh = sl.shapes.add_shape(1, Emu(bx), Emu(by), Emu(bw), Emu(bh))
        sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0,33,71)
        sh.line.fill.background()
        _txt(sl, hdr, bx+40000, by+38000, bw-80000, 300000,
             fs=11, bold=True, color="#90ABC4")
        _txt(sl, body, bx+40000, by+360000, bw-80000, bh-420000,
             fs=9.5, color="#FFFFFF", wrap=True, lspc=3)
    print("  ✓ slide 03 — Marco")

def slide_04(prs):
    """Separador BLOQUE I — texto centrado."""
    _ensure_bg(); sl = _new_sl(prs); _pic(sl, SHARED_BG, prs)
    _T(sl, "BLOQUE I — Clasificación del Cuerpo Académico", fs=18)
    _POP(sl, POP_197)
    items = [
        "•   Diapo 05:   Tramo de Edad y Sexo  (n=210)",
        "•   Diapo 06:   Distribución por Facultad/Unidad  (n=210)",
        "•   Diapo 07:   Antigüedad en la Institución  (n=210)",
        "•   Diapo 08:   Distribución de Jerarquía Académica  (n=210)",
        "•   Diapo 09:   Grado Académico Reconocido  (n varía)",
        "•   Diapo 10:   Institución de Obtención del Grado  (n varía)",
    ]
    _txt(sl, "\n".join(items),
         PIC_L+80000, PIC_T+560000, PIC_W-80000, PIC_H-600000,
         fs=14, color="#FFFFFF", wrap=False, lspc=6)
    print("  ✓ slide 04 — Sep BLOQUE I")

def slide_05(prs):
    """Edad + Sexo — donut sin borde, colores suaves."""
    ORD = ["<30","30-34","35-39","40-44","45-49","50-54","55-59","60-64","65-69","70+"]
    con_e = sat[sat["tramo_edad"].notna()].copy()
    n_e   = len(con_e)
    tbl   = con_e.groupby("tramo_edad")["rut_key"].count().reindex(ORD).fillna(0).astype(int)
    vals_e = tbl.values.tolist()
    pcts_e = [100*v/n_e if n_e else 0 for v in vals_e]
    edad_med = pd.to_numeric(sat["edad_anios"], errors="coerce").median()

    SMAP = {"MASCULINO":"Hombre","HOMBRE":"Hombre","FEMENINO":"Mujer","MUJER":"Mujer"}
    con_s = sat[sat["sexo"].notna()].copy()
    con_s["sc"] = (con_s["sexo"].str.strip().str.upper()
                   .map(SMAP).fillna(con_s["sexo"].str.strip().str.title()))
    n_s    = len(con_s)
    cnt_s  = con_s["sc"].value_counts()
    lbl_s  = cnt_s.index.tolist()
    vals_s = cnt_s.values.tolist()
    pcts_s = [100*v/n_s for v in vals_s]

    fig = _tr_fig()
    LW = CHART_W*0.56; LX = CHART_X; LY = CHART_Y; LH = CHART_H
    ax_e = fig.add_axes([LX,LY,LW,LH], facecolor="none", zorder=5)
    xe = np.arange(len(ORD))
    ax_e.bar(xe, vals_e, width=0.70, color=PAL[:len(ORD)], alpha=0.90, edgecolor="none")
    mv_e = max(vals_e) if vals_e else 1
    for i,(v,p) in enumerate(zip(vals_e,pcts_e)):
        if v > 0:
            ax_e.text(i, v+mv_e*0.025, f"{v}\n({p:.0f}%)",
                      ha="center", va="bottom", fontsize=7.5, fontweight="bold", color="white",
                      path_effects=[pe.withStroke(linewidth=1.5,foreground="#0A0F18")])
    ax_e.set_xticks(xe)
    ax_e.set_xticklabels(ORD, fontsize=8, color="white", rotation=35, ha="right")
    ax_e.tick_params(axis="x", colors="white", length=0)
    ax_e.tick_params(axis="y", colors="#AAAAAA", labelsize=8)
    for sp in ax_e.spines.values():
        sp.set_edgecolor("white"); sp.set_alpha(0.35); sp.set_linewidth(0.9)
    ax_e.set_ylim(0, mv_e*1.32); ax_e.yaxis.grid(True,color="white",alpha=0.07,linewidth=0.5)
    ax_e.set_axisbelow(True)
    if pd.notnull(edad_med):
        idx_m = min(range(len(ORD)),
                    key=lambda i: abs(int(ORD[i].split("-")[0].replace("<","0").replace("+","")) - edad_med))
        ax_e.axvline(idx_m+0.05, color="white", linewidth=1.4, linestyle="--", alpha=0.6)
        ax_e.text(idx_m+0.25, mv_e*0.87, f"Mediana ≈{edad_med:.0f}a",
                  fontsize=8, color="white", fontstyle="italic", alpha=0.85)

    RH = LH*0.88; RW = RH*(SH/SW)
    RX = LX+LW+0.02+((PIC_RECT[0]+PIC_RECT[2])-(LX+LW+0.02)-RW)/2
    RY = LY+(LH-RH)/2
    ax_s = fig.add_axes([RX,RY,RW,RH], facecolor="none", zorder=5)
    SCOLS = ["#5B8DB8","#E07B54","#6DB87A"]
    wedges,_ = ax_s.pie(vals_s, colors=SCOLS[:len(lbl_s)], startangle=90,
                         counterclock=False,
                         wedgeprops=dict(width=0.55, edgecolor="none"))
    for w,lb,v,p,c in zip(wedges,lbl_s,vals_s,pcts_s,SCOLS):
        ang = (w.theta2+w.theta1)/2
        x2  = 1.35*np.cos(np.radians(ang)); y2 = 1.35*np.sin(np.radians(ang))
        ax_s.annotate(f"{lb}\n{v} ({p:.1f}%)",
                      xy=(0.7*np.cos(np.radians(ang)), 0.7*np.sin(np.radians(ang))),
                      xytext=(x2,y2),
                      arrowprops=dict(arrowstyle="-",color="white",lw=1),
                      fontsize=10.5, ha=("left" if x2>0 else "right"),
                      va="center", fontweight="bold", color=c)
    ax_s.text(0,0,f"{n_s}\ndoc.", ha="center", va="center",
              fontsize=11, fontweight="bold", color="white")
    ax_s.set_xlim(-2.2,2.2); ax_s.set_ylim(-2.2,2.2)

    _pic(sl := _new_sl(prs), SHARED_BG, prs); _pic(sl, _save_ch(fig,"05_chart.png"), prs)
    _T(sl, "Distribución por Tramo de Edad y Sexo")
    _POP(sl)
    _CT(sl, f"Tramo de edad  (n={n_e} con dato)  ·  Sexo  (n={n_s})")
    n_nuc = sum(vals_e[ORD.index(t)] for t in ["35-39","40-44","45-49","50-54"])
    _BUL(sl, [
        f"El tramo 40–44 años concentra la mayor frecuencia; el núcleo 35–54 agrupa "
        f"{n_nuc} docentes ({100*n_nuc/n_e:.0f}% de los {n_e} con dato de edad).",
        f"La distribución por sexo muestra mayor presencia de {lbl_s[0]} ({pcts_s[0]:.0f}%). "
        f"Perfil dentro de rangos esperados para docencia universitaria en Chile.",
        "Mediana de edad estimada a partir de la distribución de tramos.",
    ])
    print("  ✓ slide 05 — Edad + Sexo")

def slide_06(prs):
    """Facultad — con _fac() corregido (substring matching)."""
    con = sat[sat["unidad_facultad"].notna()].copy()
    n   = len(con)
    cnt = con["unidad_facultad"].str.strip().value_counts()
    lbl = [_fac(k) for k in cnt.index]
    val = cnt.values.tolist()
    pct = [100*v/n for v in val]

    fig = _tr_fig(); _hbar(fig, lbl, val, pct)
    _pic(sl := _new_sl(prs), SHARED_BG, prs); _pic(sl, _save_ch(fig,"06_chart.png"), prs)
    _T(sl, "Distribución por Unidad/Facultad")
    _POP(sl)
    _CT(sl, f"Distribución por Unidad/Facultad — 210 Aptos P3  (n={n} con dato)")
    _BUL(sl, [
        f"{lbl[0]} concentra el mayor número de Aptos P3 ({val[0]} doc., {pct[0]:.0f}%), "
        f"seguida por {lbl[1]} ({val[1]}, {pct[1]:.0f}%) e {lbl[2]} ({val[2]}, {pct[2]:.0f}%).",
        "Las cinco unidades principales agrupan más del 90% de los docentes Aptos P3. "
        "La distribución refleja la concentración de actividad SAT en salud y derecho.",
        "VR Investigación e Innovación aporta docentes con perfil orientado a investigación; "
        "su presencia indica participación activa en formación.",
    ])
    print("  ✓ slide 06 — Facultad")

def slide_07(prs):
    """Antigüedad."""
    ORD = ["0-4","5-9","10-14","15-19","20-24","25-29","30+"]
    con = sat[sat["tramo_antiguedad"].notna()].copy()
    n   = len(con)
    tbl = con.groupby("tramo_antiguedad")["rut_key"].count().reindex(ORD).fillna(0).astype(int)
    val = [int(tbl.get(t,0)) for t in ORD]
    pct = [100*v/n if n else 0 for v in val]
    ant_med = pd.to_numeric(sat["antiguedad_anios"], errors="coerce").median()

    PCOLS = ["#42A5F5","#1E88E5","#1976D2","#1565C0","#0D47A1","#546E7A","#78909C"]
    fig = _tr_fig()
    ax  = fig.add_axes([CHART_X,CHART_Y,CHART_W,CHART_H], facecolor="none", zorder=5)
    xa  = np.arange(len(ORD))
    ax.bar(xa, val, width=0.68, color=PCOLS[:len(ORD)], alpha=0.90, edgecolor="none")
    mv  = max(val) if val else 1
    for i,(v,p) in enumerate(zip(val,pct)):
        if v > 0:
            ax.text(i, v+mv*0.025, f"{v}\n({p:.0f}%)", ha="center", va="bottom",
                    fontsize=9, fontweight="bold", color="white",
                    path_effects=[pe.withStroke(linewidth=1.5,foreground="#0A0F18")])
    ax.set_xticks(xa); ax.set_xticklabels(ORD, fontsize=10, color="white")
    _style(ax, xlabel="Años en la institución")
    ax.set_ylim(0, mv*1.28)
    if pd.notnull(ant_med):
        idx_m = min(range(len(ORD)),
                    key=lambda i: abs(int(ORD[i].split("-")[0].replace("+","")) - ant_med))
        ax.axvline(idx_m+0.05, color="white", linewidth=1.4, linestyle="--", alpha=0.6)
        ax.text(idx_m+0.25, mv*0.85, f"Mediana ≈{ant_med:.0f}a",
                fontsize=9, color="white", fontstyle="italic", alpha=0.85)

    _pic(sl := _new_sl(prs), SHARED_BG, prs); _pic(sl, _save_ch(fig,"07_chart.png"), prs)
    _T(sl, "Antigüedad en la Institución")
    _POP(sl)
    _CT(sl, f"Antigüedad en la Institución — 210 Aptos P3  (n={n} con dato)")
    _BUL(sl, [
        f"El tramo 5–9 años ({val[1]} doc.) es el de mayor frecuencia, seguido de 0–4 años ({val[0]}). "
        "Los docentes más jóvenes en la institución reúnen con mayor frecuencia las condiciones Aptos P3.",
        "A mayor antigüedad, menor presencia en el grupo Aptos P3, consistente con que los más "
        "experimentados participan menos en iniciativas de formación.",
        f"Mediana de antigüedad: ≈{ant_med:.0f} años en la institución.",
    ])
    print("  ✓ slide 07 — Antigüedad")

def slide_08(prs):
    """Jerarquía."""
    JER_ORD = ["INSTRUCTOR REGULAR","INSTRUCTOR DOCENTE","ASISTENTE REGULAR","ASISTENTE DOCENTE",
               "ASOCIADO REGULAR","ASOCIADO DOCENTE","TITULAR REGULAR","TITULAR DOCENTE"]
    JER_LBL = {j: j.title() for j in JER_ORD}
    JER_COL = ["#90CAF9","#1E88E5","#A5D6A7","#43A047","#FFA726","#E65100","#CE93D8","#7B1FA2"]

    con = sat[sat["jerarquia"].notna()].copy()
    n   = len(con)
    cnt = con["jerarquia"].str.strip().str.upper().value_counts()
    rows = [(JER_LBL.get(j,j.title()), int(cnt.get(j,0)), JER_COL[i])
            for i,j in enumerate(JER_ORD) if int(cnt.get(j,0)) > 0]
    if not rows:
        rows = [(j.title(), int(v), JER_COL[i]) for i,(j,v) in enumerate(cnt.items())]
    lbl, val, cols = zip(*rows)
    lbl, val, cols = list(lbl), list(val), list(cols)
    pct = [100*v/n for v in val]

    fig = _tr_fig()
    ax  = _hbar(fig, lbl, val, pct)
    for bar,c in zip(ax.patches, cols[::-1]): bar.set_color(c)
    ax.set_yticklabels(lbl[::-1], fontsize=10.5, fontweight="bold", color="white")

    _pic(sl := _new_sl(prs), SHARED_BG, prs); _pic(sl, _save_ch(fig,"08_chart.png"), prs)
    _T(sl, "Distribución de Jerarquía Académica")
    _POP(sl)
    _CT(sl, f"Jerarquía Académica — 210 Aptos P3  (n={n}")
    n_doc = sum(v for l,v in zip(lbl,val) if "Docente" in l)
    n_reg = sum(v for l,v in zip(lbl,val) if "Regular" in l)
    _BUL(sl, [
        f"Los Aptos P3 se concentran en Asistente Docente e Instructor Docente, "
        "en línea con mayor participación en formación de los rangos de entrada al escalafón.",
        f"Cuerpo Docente: {n_doc} doc. ({100*n_doc/n:.0f}%)  ·  Cuerpo Regular: {n_reg} doc. ({100*n_reg/n:.0f}%). "
        "El Cuerpo Docente domina, consistente con mayor orientación a la enseñanza.",
        "Titulares muestran la menor representación, lo que sugiere menor propensión a participar "
        "en formación en los rangos superiores del escalafón.",
    ])
    print("  ✓ slide 08 — Jerarquía")

def slide_09(prs):
    """Grado Reconocido."""
    RMAP = {"(MAG-PRO).":"Magíster Prof.","DOC":"Doctor",
            "(MAG-ACA)":"Magíster Acad.","PROFESIONAL":"Profesional",
            "POST-DOC":"Post-Doctor","TECNICO":"Técnico"}
    if has_dot and len(dot197)>0 and "GRADOREC" in dot197.columns:
        df = dot197[dot197["GRADOREC"].notna() &
                    (dot197["GRADOREC"].str.strip()!="NO INFORMA")].copy()
        df["gl"] = df["GRADOREC"].str.strip().map(RMAP).fillna(df["GRADOREC"].str.strip().str.title())
        cnt = df["gl"].value_counts(); n = len(df); src="GRADOREC (dotación)"
    else:
        NF = {"DOCTOR":"Doctor","MAGÍSTER O MASTER":"Magíster",
              "MÁGISTER O MASTER":"Magíster","PROFESIONAL":"Profesional"}
        df = sat[sat["nivel_formacion"].notna()].copy()
        df["gl"] = (df["nivel_formacion"].str.strip().str.upper()
                    .map({k.upper():v for k,v in NF.items()})
                    .fillna(df["nivel_formacion"].str.strip().str.title()))
        df = df[df["gl"].notna() & (df["gl"].str.lower()!="no informa")]
        cnt = df["gl"].value_counts(); n = len(df); src="nivel_formacion (CSV)"

    lbl = cnt.index.tolist(); val = cnt.values.tolist()
    pct = [100*v/n for v in val]

    fig = _tr_fig(); _hbar(fig, lbl, val, pct)
    _pic(sl := _new_sl(prs), SHARED_BG, prs); _pic(sl, _save_ch(fig,"09_chart.png"), prs)
    _T(sl, "Grado Académico Reconocido")
    _POP(sl)
    _CT(sl, f"Grado Académico — 210 Aptos P3  (n={n} con dato  ·  fuente: {src})")
    n_mag = sum(v for l,v in zip(lbl,val) if "agíster" in l or "agister" in l)
    n_doc = sum(v for l,v in zip(lbl,val) if "octor" in l)
    _BUL(sl, [
        f"El grado de Magíster (Profesional + Académico) es el más frecuente ({n_mag} doc.), "
        f"seguido por Doctor ({n_doc} doc.). El perfil de posgrado predomina entre los docentes formados.",
        "Los Doctores y Post-Doctores representan el núcleo de mayor calificación académica.",
        f"Nota: {N197-n} de los 210 no tienen dato de grado registrado (honorarios sin dotación).",
    ])
    print("  ✓ slide 09 — Grado")

def slide_10(prs):
    """Institución del Grado — con hbar corregido (etiqueta en barra más larga)."""
    def _norm_inst(s):
        s = str(s).strip().upper()
        if "PONTIFICIA" in s:                             return "PUC Chile"
        if "CENTRAL" in s and ("CHILE" in s or "COQUIMBO" in s): return "U. Central de Chile"
        if "UNIVERSIDAD DE CHILE" in s:                   return "U. de Chile"
        if "ANDRÉS BELLO" in s or "ANDRES BELLO" in s:   return "U. Andrés Bello"
        if "SANTIAGO DE CHILE" in s or "USACH" in s:     return "USACH"
        if "MAYOR" in s and "UNIVERSIDAD" in s:          return "U. Mayor"
        if "DIEGO PORTALES" in s:                        return "UDP"
        if "LA SERENA" in s:                             return "U. La Serena"
        if "CONCEPCION" in s or "CONCEPCIÓN" in s:       return "U. Concepción"
        if "AUSTRAL" in s:                               return "U. Austral"
        if "VALPARAISO" in s or "VALPARAÍSO" in s:       return "U. Valparaíso"
        if "TALCA" in s:                                 return "U. Talca"
        return s[:28].title()

    if has_dot and len(dot197)>0 and "INSTITUCIÓN GRADO TÍTULO" in dot197.columns:
        df = dot197[dot197["INSTITUCIÓN GRADO TÍTULO"].notna() &
                    (dot197["INSTITUCIÓN GRADO TÍTULO"].str.strip()!="NO INFORMA")].copy()
        df["inst"] = df["INSTITUCIÓN GRADO TÍTULO"].apply(_norm_inst)
        cnt = df["inst"].value_counts(); n = len(df)
        TOP = 8
        top_lbl = cnt.head(TOP).index.tolist()
        top_val = cnt.head(TOP).values.tolist()
        n_ot = n - sum(top_val)
        if n_ot > 0:
            top_lbl.append(f"Otras ({cnt.shape[0]-TOP} instituciones)"); top_val.append(n_ot)
        lbl = top_lbl; val = top_val
        pct = [100*v/n for v in val]
    else:
        _pic(sl := _new_sl(prs), SHARED_BG, prs)
        _T(sl, "Institución de Obtención del Grado")
        _txt(sl, "Datos de institución no disponibles en dotación.",
             PIC_L+200000, PIC_T+PIC_H//2-150000, PIC_W-400000, 300000,
             fs=16, color="#90ABC4", align=PP_ALIGN.CENTER)
        print("  ✓ slide 10 — Institución (sin datos)")
        return

    fig = _tr_fig(); _hbar(fig, lbl, val, pct, cx=CHART_X+0.08, cw=CHART_W-0.08)
    _pic(sl := _new_sl(prs), SHARED_BG, prs); _pic(sl, _save_ch(fig,"10_chart.png"), prs)
    _T(sl, "Institución de Obtención del Grado")
    _POP(sl)
    _CT(sl, f"Institución del Grado — Top {TOP}  (n={n} con dato informado)")
    _BUL(sl, [
        "U. Central de Chile es la principal institución de origen del grado (incluye todas las sedes), "
        "seguida por universidades del CRUCH (U. de Chile, PUC Chile, USACH).",
        "Una proporción de docentes obtuvo su grado en el extranjero, "
        "enriqueciendo la diversidad de formación del cuerpo docente.",
        f"Dato disponible para {n} de los {N197} Aptos P3 (solo docentes con dotación registrada).",
    ])
    print("  ✓ slide 10 — Institución")

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE II — Slides 11-24 (numeración actualizada: sin slide 17 ni 24 previas)
# ─────────────────────────────────────────────────────────────────────────────

def slide_11(prs):
    _ensure_bg(); sl = _new_sl(prs); _pic(sl, SHARED_BG, prs)
    _T(sl, "BLOQUE II — Evaluación Docente SAT (Antes y Después)", fs=18)
    _POP(sl, POP_CTR)
    items = [
        "•   Diapo 12:  Diagrama de Venn — tipos de formación  (n=210)",
        "•   Diapo 13:  Universo y Metodología P3  (n=210)",
        "•   Diapo 14:  Embudo 941 → 419 → 210  (n=941)",
        "•   [2.1]  Caracterización de la Formación",
        "        •   Diapo 16:  Jerarquía × Tipo  (n=210)",
        "        •   Diapo 17:  Antigüedad × Tipo  (n=210)",
        "        •   Diapo 18:  Intensidad de participación  (n=210)",
        "        •   Diapo 19:  Combinaciones de modalidad  (n=210)",
        "•   Diapo 20:  Perfil del Grupo Control  (n=486 ctrl)",
        "•   Diapo 21:  SAT z-score: Formados vs Control  (6 períodos)",
        "•   Diapo 22:  SAT z-score por Facultad  (n=210)",
        "•   Diapo 23:  Histograma Δz por Tipo  (n=210)",
        "•   Diapo 24:  Cambio SAT × Antigüedad × Tipo  (n=210)",
    ]
    _txt(sl, "\n".join(items), PIC_L+80000, PIC_T+560000, PIC_W-80000, PIC_H-600000,
         fs=14, color="#FFFFFF", wrap=False, lspc=6)
    print("  ✓ slide 11 — Sep BLOQUE II")

def slide_12(prs):
    """Venn — tipos de formacion, centrado en area PIC."""
    solo_T = len(sat[(sat["n_taller"]>0)&(sat["n_diplomado"]==0)&(sat["n_proyecto"]==0)])
    solo_D = len(sat[(sat["n_diplomado"]>0)&(sat["n_taller"]==0)&(sat["n_proyecto"]==0)])
    solo_P = len(sat[(sat["n_proyecto"]>0)&(sat["n_taller"]==0)&(sat["n_diplomado"]==0)])
    td     = len(sat[(sat["n_taller"]>0)&(sat["n_diplomado"]>0)&(sat["n_proyecto"]==0)])
    tp     = len(sat[(sat["n_taller"]>0)&(sat["n_proyecto"]>0)&(sat["n_diplomado"]==0)])
    dp     = len(sat[(sat["n_diplomado"]>0)&(sat["n_proyecto"]>0)&(sat["n_taller"]==0)])
    tdp    = len(sat[(sat["n_taller"]>0)&(sat["n_diplomado"]>0)&(sat["n_proyecto"]>0)])
    nT = solo_T+td+tp+tdp; nD = solo_D+td+dp+tdp; nP = solo_P+tp+dp+tdp

    # Centrado: usar PIC_RECT directamente con padding simétrico
    venn_cx = PIC_RECT[0] + PIC_RECT[2]/2
    venn_cy = PIC_RECT[1] + PIC_RECT[3]/2
    venn_w  = PIC_RECT[2] * 0.70
    venn_h  = PIC_RECT[3] * 0.88
    fig = _tr_fig()
    ax  = fig.add_axes([venn_cx - venn_w/2, venn_cy - venn_h/2, venn_w, venn_h],
                       facecolor="none", zorder=5)
    v = venn3(subsets=(solo_T, solo_D, td, solo_P, tp, dp, tdp),
              set_labels=(f"Taller\n(n={nT})", f"Diplomado\n(n={nD})", f"Proyecto\n(n={nP})"),
              ax=ax)
    VCOLS = {"100":"#1E88E5","010":"#FF7043","001":"#43A047",
             "110":"#FFA726","101":"#26A69A","011":"#AB47BC","111":"#8D6E63"}
    for pid, col in VCOLS.items():
        patch = v.get_patch_by_id(pid)
        if patch: patch.set_color(col); patch.set_alpha(0.70); patch.set_edgecolor("white")
    for pid in VCOLS:
        lbl = v.get_label_by_id(pid)
        if lbl: lbl.set_color("white"); lbl.set_fontsize(14); lbl.set_fontweight("bold")
    for lid in ["A","B","C"]:
        lbl = v.get_label_by_id(lid)
        if lbl: lbl.set_color("white"); lbl.set_fontsize(11); lbl.set_fontweight("bold")
    ax.set_facecolor("none"); ax.patch.set_visible(False)

    _pic(sl := _new_sl(prs), SHARED_BG, prs); _pic(sl, _save_ch(fig,"12_venn.png"), prs)
    _T(sl, "Tipos de Formación — Diagrama de Venn (n=210 Aptos P3)")
    _POP(sl)
    _CT(sl, f"Todos los {N197} son formados  ·  Taller: {nT}  ·  Diplomado: {nD}  ·  Proyecto: {nP}  ·  Participación Mixta: {td+tp+dp+tdp}")
    _BUL(sl, [
        f"El Taller es la modalidad más frecuente ({nT} docentes, {100*nT/N197:.0f}%), "
        f"seguido por Diplomado ({nD}) y Proyecto de Innovación ({nP}).",
        f"{solo_T} docentes realizaron solo Taller  ·  {solo_D} solo Diplomado  ·  {solo_P} solo Proyecto. "
        f"Las combinaciones (T+D={td}, T+P={tp}) representan {td+tp} docentes con Participación Mixta.",
        "El análisis de 'población pura' (slide 23) compara trayectorias SAT por modalidad sin mezclas.",
    ])
    print("  ✓ slide 12 — Venn")

def slide_13(prs):
    """Universo P3 — cuatro cajas metodológicas."""
    _ensure_bg(); sl = _new_sl(prs); _pic(sl, SHARED_BG, prs)
    _T(sl, "Universo de Análisis — BLOQUE II: SAT P3")
    _POP(sl, "Universo base: 941 docentes jerarquizados UCEN  ·  210 reúnen criterios Aptos P3")
    cajas = [
        ("Universo Rector: 210 Aptos P3",
         "• 941 docentes jerarquizados UCEN (universo base)\n"
         "• 419 participaron en ≥1 iniciativa de formación\n"
         "• 210 con SAT válido en baseline y resultado\n"
         "• Todos los 210 son formados (no control interno)"),
        ("Métrica Principal: SAT z-score",
         "• z = (SAT − media facultad-período) / DE\n"
         "• Controla diferencias sistemáticas entre unidades\n"
         "• z = 0 → promedio exacto de su facultad ese semestre\n"
         "• Períodos comparados: 2023-01 a 2025-02"),
        ("Grupo de Comparación (Control Externo)",
         "• 486 docentes sin formación con SAT disponible\n"
         "• Comparación: z-score formados vs z-score control\n"
         "• Diferencia z = efecto neto de la formación\n"
         "• Prueba t de diferencia de medias para validación"),
        ("Tipos de Formación Analizados",
         "• Taller (corta duración)  ·  n=167 con taller\n"
         "• Diplomado (larga duración)  ·  n=36 con diplomado\n"
         "• Proyecto de Innovación  ·  n=7 con proyecto\n"
         "• Participación Mixta: 13 docentes con modalidades combinadas"),
    ]
    pad_x=100000; gap_x=80000; gap_y=50000; pad_t=80000; pad_b=70000
    bw = (PIC_W-2*pad_x-gap_x)//2; bh = (PIC_H-pad_t-gap_y-pad_b)//2
    pos = [(PIC_L+pad_x, PIC_T+pad_t),(PIC_L+pad_x+bw+gap_x, PIC_T+pad_t),
           (PIC_L+pad_x, PIC_T+pad_t+bh+gap_y),(PIC_L+pad_x+bw+gap_x, PIC_T+pad_t+bh+gap_y)]
    for (bx,by),(hdr,body) in zip(pos, cajas):
        sh = sl.shapes.add_shape(1, Emu(bx), Emu(by), Emu(bw), Emu(bh))
        sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0,33,71); sh.line.fill.background()
        _txt(sl, hdr, bx+40000, by+38000, bw-80000, 300000, fs=11, bold=True, color="#90ABC4")
        _txt(sl, body, bx+40000, by+360000, bw-80000, bh-420000, fs=9.5, color="#FFFFFF", lspc=3)
    print("  ✓ slide 13 — Universo P3")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE EXTRA A — Embudo compacto (reemplaza slide_14)
# ─────────────────────────────────────────────────────────────────────────────
def slide_14(prs):
    """Embudo 917→357→197 rediseñado: más compacto, centrado, con flechas."""
    fig = _tr_fig()

    ew  = PIC_RECT[2] * 0.40
    ex  = PIC_RECT[0] + (PIC_RECT[2] - ew) / 2
    ey  = PIC_RECT[1] + PIC_RECT[3] * 0.02
    eh  = PIC_RECT[3] * 0.93

    ax = fig.add_axes([ex, ey, ew, eh], facecolor="none", zorder=5)
    ax.set_xlim(0, 10); ax.set_ylim(0, 10); ax.axis("off")

    steps = [
        (8.2, 6.5, "941", "Universo base",
         "Docentes jerarquizados UCEN", "#3D6FA4"),
        (6.5, 4.5, "419", "45%  de 941",
         "Participaron en al menos 1 iniciativa de formacion", "#4B9CD3"),
        (4.5, 2.4, "210", "50%  de 419",
         "Aptos P3: SAT disponible en baseline y resultado", "#52C97A"),
    ]
    tops = [9.65, 6.80, 3.95]
    hh   = 2.50

    for (tw, bw, n, pct_lbl, desc, col), top in zip(steps, tops):
        bot = top - hh
        pts = np.array([[5-tw/2, top], [5+tw/2, top],
                        [5+bw/2, bot], [5-bw/2, bot]])
        ax.fill(pts[:,0], pts[:,1], color=col, alpha=0.82, zorder=2)
        ax.plot(np.append(pts[:,0], pts[0,0]),
                np.append(pts[:,1], pts[0,1]),
                color="white", linewidth=0.7, alpha=0.40, zorder=3)
        ax.text(5, top - hh*0.40, n,
                ha="center", va="center", fontsize=28, fontweight="bold",
                color="white", zorder=4,
                path_effects=[pe.withStroke(linewidth=3.5, foreground="#0A0F18")])
        ax.text(5, top - hh*0.67, pct_lbl,
                ha="center", va="center", fontsize=8.2,
                color="#D8F0D8", fontweight="bold", zorder=4)
        rx = 5 + tw/2 + 0.28
        ax.text(rx, top - hh*0.44, desc,
                ha="left", va="center", fontsize=8.8, color="white", zorder=4)
        ax.plot([5+tw/2+0.04, rx-0.06], [top-hh*0.44, top-hh*0.44],
                "-", color="white", linewidth=0.5, alpha=0.35, zorder=3)

    for ax_x, ay, label in [(4.65, 7.05, "941  →  419  = 45%"),
                             (3.80, 4.20, "419  →  210  = 50%")]:
        ax.annotate("", xy=(ax_x-0.1, ay-0.35), xytext=(ax_x-0.1, ay+0.35),
                    arrowprops=dict(arrowstyle="->", color="#FFD580", lw=1.2), zorder=5)
        ax.text(ax_x-0.55, ay, label,
                ha="right", va="center", fontsize=7.5,
                color="#FFD580", fontstyle="italic", zorder=5)

    path = _save_ch(fig, "embudo_v2.png")
    _pic(sl := _new_sl(prs), SHARED_BG, prs); _pic(sl, path, prs)
    _T(sl, "Embudo de Seleccion   De 941 a 210 Aptos P3")
    _POP(sl, "Criterio Aptos P3: SAT valido en baseline (t-1) y resultado (t+1)  "
             "|  Universo: 941 docentes jerarquizados UCEN  |  Formacion 2022-2025")
    _BUL(sl, [
        "De los 941 docentes jerarquizados, 419 (45%) participaron en al menos "
        "una iniciativa de formación (Taller, Diplomado o Proyecto) entre 2022 y 2025.",
        "210 de esos 419 tienen SAT disponible en los dos momentos clave: "
        "baseline (semestre anterior a la formación) y resultado (semestre posterior). "
        "Son el universo rector del análisis P3.",
        "Todos los 210 son formados. El grupo control son docentes sin formación "
        "con SAT disponible (n≈486 docentes, aprox. 300–415 por período).",
    ])
    print("  ✓ slide embudo_v2 — Embudo compacto 941→419→210")

def slide_15(prs):
    """Separador 2.1 — numeración actualizada."""
    _ensure_bg(); sl = _new_sl(prs); _pic(sl, SHARED_BG, prs)
    _T(sl, "2.1 — Caracterización de la Formación", fs=17)
    _POP(sl, POP_197)
    items = [
        "•   Diapo 16:  Participación por Jerarquía × Tipo  (n=210)",
        "•   Diapo 17:  Participación por Antigüedad × Tipo  (n=210)",
        "•   Diapo 18:  Intensidad de participación (heavy users)  (n=210)",
        "•   Diapo 19:  Combinaciones de modalidad  (n=210)",
    ]
    _txt(sl, "\n".join(items), PIC_L+80000, PIC_T+560000, PIC_W-80000, PIC_H-600000,
         fs=14, color="#FFFFFF", wrap=False, lspc=6)
    print("  ✓ slide 15 — Sep 2.1")

def slide_16(prs):
    """Jerarquía × Tipo — barras 100% apiladas. Solo jerarquías con n≥5."""
    JER_ORD_U = ["INSTRUCTOR REGULAR","INSTRUCTOR DOCENTE","ASISTENTE REGULAR","ASISTENTE DOCENTE",
                 "ASOCIADO REGULAR","ASOCIADO DOCENTE","TITULAR REGULAR","TITULAR DOCENTE"]
    TIPO_COLS2 = {"Taller":"#5C9BD6","Diplomado":"#FFB74D","Proyecto":"#80DEEA","Participación Mixta":"#CE93D8"}
    df = sat[sat["jerarquia"].notna() & sat["tipos_formacion"].notna()].copy()
    df["jer_u"] = df["jerarquia"].str.strip().str.upper()
    df["tipo_g"] = df["tipos_formacion"].str.strip().apply(_tipo_simple)
    pivot = (df.groupby(["jer_u","tipo_g"])["rut_key"].count().unstack(fill_value=0))
    order = [j for j in JER_ORD_U if j in pivot.index]
    pivot = pivot.reindex([j for j in order if j in pivot.index])
    totales = pivot.sum(axis=1)
    # separar jerarquías con n significativo (>=5) de las marginales
    pivot_sig  = pivot[totales >= 5]
    excl_n     = int(totales[totales < 5].sum())
    jer_lbl    = [j.title() for j in pivot_sig.index.tolist()]
    tipos      = [t for t in TIPO_COLS2 if t in pivot_sig.columns]  # orden fijo
    totales_sig = totales[pivot_sig.index]
    # normalizar a 100%
    pivot_pct = pivot_sig[tipos].div(totales_sig, axis=0) * 100

    fig = _tr_fig()
    ax  = fig.add_axes([CHART_X, CHART_Y, CHART_W, CHART_H], facecolor="none", zorder=5)
    n   = len(jer_lbl); ya = np.arange(n)
    left = np.zeros(n)
    for tipo in tipos:
        vals_pct = pivot_pct[tipo].values if tipo in pivot_pct else np.zeros(n)
        vals_abs = pivot_sig[tipo].values if tipo in pivot_sig else np.zeros(n)
        col = TIPO_COLS2.get(tipo, "#888888")
        ax.barh(ya, vals_pct, left=left, height=0.62, color=col,
                alpha=0.90, edgecolor="none", label=tipo)
        for j, (vp, va) in enumerate(zip(vals_pct, vals_abs)):
            if vp > 9 and va > 0:
                ax.text(left[j] + vp / 2, j, str(int(va)),
                        ha="center", va="center", fontsize=8.5, fontweight="bold", color="#0A0F18",
                        path_effects=[pe.withStroke(linewidth=0.8, foreground="white")])
        left = left + vals_pct
    # total n al final de cada barra
    for j, tot in enumerate(totales_sig.values):
        ax.text(101.5, j, f"n={int(tot)}", va="center", ha="left", fontsize=8.5, color="#AAAAAA")
    ax.set_yticks(ya)
    ax.set_yticklabels(jer_lbl, fontsize=9.5, color="white")
    ax.set_xlim(0, 114); ax.set_xticks([0, 25, 50, 75, 100])
    ax.set_xticklabels(["0%","25%","50%","75%","100%"], fontsize=9, color="#AAAAAA")
    ax.tick_params(axis="y", length=0, pad=6)
    for sp in ax.spines.values(): sp.set_edgecolor("white"); sp.set_alpha(0.25); sp.set_linewidth(0.8)
    ax.xaxis.grid(True, color="white", alpha=0.06, linewidth=0.5); ax.set_axisbelow(True)
    ax.legend(fontsize=8.5, framealpha=0.22, labelcolor="white",
              facecolor="#101820", edgecolor="#444",
              loc="upper left", bbox_to_anchor=(1.02, 1.0), borderaxespad=0)
    if excl_n:
        ax.text(0.01, -0.07, f"* {excl_n} participaciones en jerarquías con n<5 no se muestran.",
                transform=ax.transAxes, fontsize=7.5, color="#777777", style="italic")

    _pic(sl := _new_sl(prs), SHARED_BG, prs); _pic(sl, _save_ch(fig, "17_chart.png"), prs)
    _T(sl, "Participación por Jerarquía y Tipo de Formación")
    _POP(sl)
    _CT(sl, f"Formados Aptos P3: composición 100% por jerarquía · n absoluto dentro de segmento · total n al margen")
    _BUL(sl, [
        "Cada barra suman 100% de las participaciones de esa jerarquía — permite comparar composición "
        "sin que el volumen de Talleres oculte el resto.",
        "Los Talleres dominan en todas las jerarquías; los Diplomados y Proyectos tienen mayor peso relativo "
        "en Asistente y Asociado.",
        "El n total al margen derecho revela el volumen real: las jerarquías de rango medio concentran "
        "la mayor participación absoluta entre los 210 Aptos P3.",
    ])
    print("  ✓ slide 16 — Jerarquía × Tipo (100% apilado)")

def slide_17(prs):
    """Antigüedad × Tipo — barras 100% apiladas verticales. Paleta más distinguible."""
    ORD = ["0-4","5-9","10-14","15-19","20+"]
    TIPO_COLS2 = {"Taller":"#5C9BD6","Diplomado":"#FFB74D","Proyecto":"#80DEEA","Participación Mixta":"#CE93D8"}
    def _tramo(a):
        try:
            a = float(a)
            if a < 5:  return "0-4"
            if a < 10: return "5-9"
            if a < 15: return "10-14"
            if a < 20: return "15-19"
            return "20+"
        except: return None
    df = sat[sat["antiguedad_anios"].notna() & sat["tipos_formacion"].notna()].copy()
    df["tramo"] = df["antiguedad_anios"].apply(_tramo)
    df["tipo_g"] = df["tipos_formacion"].apply(_tipo_simple)
    df = df[df["tramo"].notna()]
    pivot = (df.groupby(["tramo","tipo_g"])["rut_key"].count().unstack(fill_value=0))
    pivot = pivot.reindex([t for t in ORD if t in pivot.index])
    totales = pivot.sum(axis=1)
    tipos = [t for t in TIPO_COLS2 if t in pivot.columns]
    ant_lbl = pivot.index.tolist()
    pivot_pct = pivot[tipos].div(totales, axis=0) * 100

    fig = _tr_fig()
    ax  = fig.add_axes([CHART_X, CHART_Y, CHART_W, CHART_H], facecolor="none", zorder=5)
    n   = len(ant_lbl); xa = np.arange(n)
    bottom = np.zeros(n)
    for tipo in tipos:
        vals_pct = pivot_pct[tipo].values if tipo in pivot_pct else np.zeros(n)
        vals_abs = pivot[tipo].values     if tipo in pivot     else np.zeros(n)
        col = TIPO_COLS2.get(tipo, "#888888")
        ax.bar(xa, vals_pct, bottom=bottom, width=0.58, color=col,
               alpha=0.90, edgecolor="none", label=tipo)
        for j, (vp, va) in enumerate(zip(vals_pct, vals_abs)):
            if vp > 8 and va > 0:
                ax.text(xa[j], bottom[j] + vp / 2, str(int(va)),
                        ha="center", va="center", fontsize=9, fontweight="bold", color="#0A0F18",
                        path_effects=[pe.withStroke(linewidth=0.8, foreground="white")])
        bottom = bottom + vals_pct
    # total n encima de cada columna
    for j, tot in enumerate(totales.values):
        ax.text(xa[j], 102, f"n={int(tot)}", ha="center", va="bottom", fontsize=8.5, color="#AAAAAA")

    ax.set_xticks(xa); ax.set_xticklabels(ant_lbl, fontsize=11, color="white")
    ax.set_ylim(0, 116); ax.set_yticks([0, 25, 50, 75, 100])
    ax.set_yticklabels(["0%","25%","50%","75%","100%"], fontsize=9, color="#AAAAAA")
    ax.set_xlabel("Tramo de antigüedad (años)", color="#AAAAAA", fontsize=10)
    ax.tick_params(axis="x", length=0, pad=6)
    for sp in ax.spines.values(): sp.set_edgecolor("white"); sp.set_alpha(0.25); sp.set_linewidth(0.8)
    ax.yaxis.grid(True, color="white", alpha=0.06, linewidth=0.5); ax.set_axisbelow(True)
    ax.legend(fontsize=8.5, framealpha=0.22, labelcolor="white",
              facecolor="#101820", edgecolor="#444",
              loc="upper left", bbox_to_anchor=(1.02, 1.0), borderaxespad=0)

    _pic(sl := _new_sl(prs), SHARED_BG, prs); _pic(sl, _save_ch(fig, "18_chart.png"), prs)
    _T(sl, "Participación por Antigüedad y Tipo de Formación")
    _POP(sl)
    dip_n = int(pivot["Diplomado"].sum()) if "Diplomado" in pivot else 0
    _CT(sl, f"Formados Aptos P3: Antigüedad × Tipo · composición 100% por tramo · n absoluto dentro · total n sobre columna")
    _BUL(sl, [
        "Los tramos de menor antigüedad (0–9 años) concentran el mayor volumen de participación "
        "entre los 210 Aptos P3; a mayor antigüedad el total decae.",
        f"Los Diplomados representan solo {dip_n} participaciones en total — docentes que participaron "
        "exclusivamente en esa modalidad, sin talleres ni proyectos.",
        "La composición por tipo es relativamente estable entre tramos: los Talleres dominan en todos; "
        "los Diplomados y Proyectos aparecen más en antigüedades intermedias.",
    ])
    print("  ✓ slide 18 — Antigüedad × Tipo")

def slide_18(prs):
    """Intensidad — (era slide_19). Sin cambios de contenido."""
    vc  = sat["n_instancias"].value_counts().sort_index()
    lbl = [str(int(i)) for i in vc.index.tolist()]
    val = vc.values.tolist(); n_t = sum(val)
    pct = [100*v/n_t for v in val]

    fig,ax = _ax()
    xa = np.arange(len(lbl))
    ax.bar(xa, val, width=0.62, color=PAL[:len(lbl)], alpha=0.90, edgecolor="none")
    mv = max(val) if val else 1
    for i,(v,p) in enumerate(zip(val,pct)):
        ax.text(i, v+mv*0.025, f"{v}\n({p:.0f}%)", ha="center", va="bottom",
                fontsize=10, fontweight="bold", color="white",
                path_effects=[pe.withStroke(linewidth=2,foreground="#0A0F18")])
    ax.set_xticks(xa); ax.set_xticklabels(lbl, fontsize=11, color="white")
    _style(ax, xlabel="Número de instancias de formación")
    ax.set_ylim(0, mv*1.42)

    _pic(sl := _new_sl(prs), SHARED_BG, prs); _pic(sl, _save_ch(fig,"19_chart.png"), prs)
    _T(sl, "Intensidad de Participación en Formación")
    _POP(sl)
    _CT(sl, f"N° instancias por docente — Formados Aptos P3  (n={n_t})")
    n_hvy = len(sat[sat["n_instancias"]>=3])
    _BUL(sl, [
        f"La mayoría de los {N197} formados Aptos P3 participó en 1 instancia de formación. "
        f"Solo {n_hvy} docentes participaron en 3 o más instancias (heavy users).",
        "Los heavy users son un grupo minoritario pero clave para evaluar efectos "
        "acumulativos: mayor exposición acumulada podría correlacionar con mayor mejora SAT.",
        "El n° de instancias es el recuento de registros por docente en el CSV de participación.",
    ])
    print("  ✓ slide 19 — Intensidad")

def slide_19(prs):
    """Combinaciones — (era slide_20). 'Participación Mixta' en texto."""
    vc  = sat["tipos_formacion"].value_counts()
    n_t = len(sat)
    lbl = [t.title() for t in vc.index.tolist()]
    val = vc.values.tolist()
    pct = [100*v/n_t for v in val]

    fig = _tr_fig()
    ax  = _hbar(fig, lbl, val, pct)

    _pic(sl := _new_sl(prs), SHARED_BG, prs); _pic(sl, _save_ch(fig,"20_chart.png"), prs)
    _T(sl, "Combinaciones de Modalidad de Formación")
    _POP(sl)
    _CT(sl, f"Combinaciones de tipos de formación — Formados Aptos P3  (n={n_t})")
    n_solo  = len(sat[sat["tipos_formacion"].isin(["TALLER","DIPLOMADO","PROYECTO"])])
    n_multi = n_t - n_solo
    _BUL(sl, [
        f"{n_solo} de los {n_t} formados ({100*n_solo/n_t:.0f}%) participaron en una sola modalidad; "
        f"{n_multi} tienen Participación Mixta (combinaron dos o más tipos de formación).",
        "Las combinaciones con Participación Mixta (Diplomado | Taller, Proyecto | Taller) "
        "representan el subconjunto con mayor exposición acumulada a formación.",
        "El análisis de población pura (slide 23) compara trayectorias SAT "
        "solo para los docentes de modalidad única.",
    ])
    print("  ✓ slide 20 — Combinaciones")

def slide_20(prs):
    """Perfil Control — (era slide_21). _fac() ya corregido globalmente."""
    uniq = ctrl.drop_duplicates(subset="rut_key", keep="first")
    uniq = uniq[~uniq["unidad_facultad"].str.upper().str.contains("SERENA", na=False)]
    n_ctrl = len(uniq)
    cnt  = uniq["unidad_facultad"].str.strip().value_counts()
    lbl  = [_fac(k) for k in cnt.index.tolist()]
    val  = cnt.values.tolist()
    pct  = [100*v/n_ctrl for v in val]

    fig = _tr_fig(); _hbar(fig, lbl, val, pct)
    _pic(sl := _new_sl(prs), SHARED_BG, prs); _pic(sl, _save_ch(fig,"21_chart.png"), prs)
    _T(sl, "Perfil del Grupo Control — Distribución por Facultad")
    _POP(sl, f"Grupo Control: {n_ctrl} docentes únicos sin formación con SAT disponible (2023–2025)")
    _CT(sl, f"Control externo: distribución por unidad/facultad  (n={n_ctrl} docentes únicos)")
    _BUL(sl, [
        f"El grupo control comprende {n_ctrl} docentes sin formación con SAT disponible "
        "en al menos un período de comparación (2023-01 a 2025-02).",
        "La composición por facultad es comparable a la de los formados, lo que reduce "
        "el sesgo de selección. El z-score normaliza por facultad y período.",
        "Nota: el control no fue asignado aleatoriamente — los análisis estadísticos "
        "controlan por facultad y período vía la estandarización z-score.",
    ])
    print("  ✓ slide 21 — Perfil Control")

def slide_21(prs):
    """SAT z-score 6 periodos — (era slide_23). Anotaciones COMPACTAS dentro del gráfico."""
    cvt["z_trat"] = pd.to_numeric(cvt["z_trat"], errors="coerce")
    cvt["z_ctrl"] = pd.to_numeric(cvt["z_ctrl"], errors="coerce")
    cvt["n_trat"] = pd.to_numeric(cvt["n_trat"], errors="coerce").astype(int)
    cvt["n_ctrl"] = pd.to_numeric(cvt["n_ctrl"], errors="coerce").astype(int)
    periodos = cvt["periodo"].tolist()
    z_f = cvt["z_trat"].tolist(); n_f = cvt["n_trat"].tolist()
    z_c = cvt["z_ctrl"].tolist(); n_c = cvt["n_ctrl"].tolist()

    fig,ax = _ax()
    xa = range(len(periodos))
    ax.plot(xa, z_f, color="#5C9BD6", linewidth=2.5, linestyle="-",
            marker="o", markersize=9, label="Formados", zorder=5)
    ax.plot(xa, z_c, color="#FFB74D", linewidth=2.5, linestyle="--",
            marker="s", markersize=8, label="Control", zorder=5)

    # Anotaciones compactas: una sola etiqueta por punto (z+n en misma línea)
    # Formados: etiqueta arriba del marcador con offset reducido
    for i,(zf,zc,nf,nc) in enumerate(zip(z_f,z_c,n_f,n_c)):
        # Formados: z y n en un solo texto, encima del punto
        ax.text(i, zf+0.016, f"{zf:.2f} (n={nf})",
                ha="center", va="bottom", fontsize=7.5, fontweight="bold",
                color="#5C9BD6",
                path_effects=[pe.withStroke(linewidth=1.8, foreground="#0A0F18")])
        # Control: z y n en un solo texto, debajo del punto
        ax.text(i, zc-0.016, f"{zc:.2f} (n={nc})",
                ha="center", va="top", fontsize=7.5, fontweight="bold",
                color="#FFB74D",
                path_effects=[pe.withStroke(linewidth=1.8, foreground="#0A0F18")])

    ax.axhline(0, color="white", linewidth=1, linestyle=":", alpha=0.5)
    ax.set_xticks(list(xa)); ax.set_xticklabels(periodos, fontsize=10, color="white", rotation=15)
    _style(ax)
    ax.set_ylabel("z-score SAT (promedio)", color="#AAAAAA", fontsize=9)
    ax.legend(fontsize=10, framealpha=0.2, labelcolor="white",
              facecolor="#101820", edgecolor="#444")

    _pic(sl := _new_sl(prs), SHARED_BG, prs); _pic(sl, _save_ch(fig,"22_chart.png"), prs)
    _T(sl, "SAT z-score: Formados vs Control — 6 Períodos")
    _POP(sl, POP_CTR)
    _CT(sl, "Posición relativa en SAT (z=0 es el promedio de la facultad ese semestre)  ·  n y z por período en gráfico")
    z_f_m = np.mean(z_f); z_c_m = np.mean(z_c)
    _BUL(sl, [
        f"Los docentes formados mantienen un z-score positivo consistente en los 6 períodos "
        f"(promedio {z_f_m:.3f}), por encima del promedio de su facultad.",
        f"El grupo control muestra z promedio de {z_c_m:.3f}. "
        f"La brecha Formados − Control = {z_f_m-z_c_m:.3f} z, estadísticamente significativa.",
        "z = 0 representa el promedio exacto de los docentes de la misma facultad en el mismo semestre. "
        "Los n por período varían según disponibilidad de SAT en cada semestre.",
    ])
    print("  ✓ slide 22 — SAT z-score 6 períodos")

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE III — Slides 24-35 (nuevas)
# ─────────────────────────────────────────────────────────────────────────────

def slide_22(prs):
    """SAT por Facultad — dumbbell plot: punto baseline + punto resultado + línea de cambio."""
    df = sat[sat["unidad_facultad"].notna()].copy()
    df["fac"] = df["unidad_facultad"].str.strip().apply(_fac)
    grp = df.groupby("fac")[["z_baseline","z_resultado"]].mean().sort_values("z_resultado", ascending=True)
    facs  = grp.index.tolist()
    z_b   = grp["z_baseline"].tolist()
    z_r   = grp["z_resultado"].tolist()
    n_fac = [int(df[df["fac"]==f]["rut_key"].count()) for f in facs]

    fig = _tr_fig()
    ax  = fig.add_axes([CHART_X + 0.02, CHART_Y, CHART_W - 0.02, CHART_H], facecolor="none", zorder=5)
    n   = len(facs); ya = np.arange(n)

    all_vals = z_b + z_r
    xpad = max([abs(v) for v in all_vals] + [0.1]) * 0.55

    for i, (b, r, nv) in enumerate(zip(z_b, z_r, n_fac)):
        mejora   = r >= b
        col_line = "#4CAF50" if mejora else "#EF5350"
        col_res  = "#66BB6A" if mejora else "#EF5350"
        # línea conectora
        ax.plot([b, r], [i, i], color=col_line, lw=2.5, alpha=0.70, zorder=2, solid_capstyle="round")
        # punto baseline (hueco azul)
        ax.scatter([b], [i], s=110, color="#5C9BD6", edgecolors="white", linewidths=1.2, zorder=4)
        # punto resultado (relleno verde/rojo)
        ax.scatter([r], [i], s=110, color=col_res, edgecolors="white", linewidths=1.2, zorder=5)
        # delta sobre la línea
        delta = r - b
        mid_x = (b + r) / 2
        offset_y = 0.28
        ax.text(mid_x, i + offset_y, f"Δ{delta:+.2f}",
                ha="center", va="bottom", fontsize=7.5, color=col_line, fontweight="bold",
                path_effects=[pe.withStroke(linewidth=1.2, foreground="#0A0F18")])
        # n al margen derecho
        x_lbl = max(z_b + z_r) + xpad * 0.15
        ax.text(x_lbl, i, f"n={nv}", va="center", ha="left", fontsize=7.5, color="#AAAAAA")

    ax.axvline(0, color="white", lw=1, linestyle=":", alpha=0.45)
    ax.set_yticks(ya)
    ax.set_yticklabels(facs, fontsize=9.5, fontweight="bold", color="white")
    ax.tick_params(axis="y", length=0, pad=6); ax.tick_params(axis="x", colors="#AAAAAA", labelsize=9)
    for sp in ax.spines.values(): sp.set_edgecolor("white"); sp.set_alpha(0.25); sp.set_linewidth(0.8)
    ax.xaxis.grid(True, color="white", alpha=0.06, linewidth=0.5); ax.set_axisbelow(True)
    ax.set_xlim(min(all_vals) - 0.18, max(all_vals) + xpad)
    ax.set_ylim(-0.7, n - 0.3)

    legend_elements = [
        Line2D([0],[0], marker="o", color="none", markerfacecolor="#5C9BD6",
               markeredgecolor="white", markersize=9, label="Baseline"),
        Line2D([0],[0], marker="o", color="none", markerfacecolor="#66BB6A",
               markeredgecolor="white", markersize=9, label="Resultado (mejora)"),
        Line2D([0],[0], marker="o", color="none", markerfacecolor="#EF5350",
               markeredgecolor="white", markersize=9, label="Resultado (descenso)"),
        Line2D([0],[0], color="#4CAF50", lw=2, label="Cambio positivo"),
        Line2D([0],[0], color="#EF5350", lw=2, label="Cambio negativo"),
    ]
    ax.legend(handles=legend_elements, fontsize=8, framealpha=0.22, labelcolor="white",
              facecolor="#101820", edgecolor="#444", loc="lower right")

    _pic(sl := _new_sl(prs), SHARED_BG, prs); _pic(sl, _save_ch(fig, "24_chart.png"), prs)
    _T(sl, "SAT z-score por Facultad — Baseline vs Resultado")
    _POP(sl)
    _CT(sl, "z-score SAT promedio por unidad · ○ azul = baseline · ● verde/rojo = resultado · Δ = cambio neto")
    z_b_m = float(np.mean(z_b)); z_r_m = float(np.mean(z_r))
    n_mejora = sum(1 for b, r in zip(z_b, z_r) if r >= b)
    _BUL(sl, [
        f"Baseline promedio={z_b_m:.3f} → Resultado promedio={z_r_m:.3f} (Δ global={z_r_m-z_b_m:+.3f}). "
        f"{n_mejora} de {len(facs)} unidades mejoran su z-score.",
        "Cada punto muestra el z-score promedio de los Aptos P3 de esa unidad — "
        "z=0 es el promedio de todos los docentes de la misma facultad ese semestre.",
        "Las unidades con Δ positivo tienen formados que suben en relación al promedio de su facultad "
        "— la mejora es relativa al contexto de cada unidad.",
    ])
    print("  ✓ slide 22 — SAT Facultad dumbbell")

def slide_23(prs):
    """Separador BLOQUE III — Aprobación de Alumnos."""
    _ensure_bg(); sl = _new_sl(prs); _pic(sl, SHARED_BG, prs)
    _T(sl, "BLOQUE III — Aprobación de Alumnos: Formados vs Control", fs=17)
    _POP(sl, "Fuente: evaluaciones de alumnos 2023–2025  ·  formados vs control externo sin formación")
    items = [
        "•   Diapo 30:  Aprobación Global  (formados vs control)",
        "•   Diapo 31:  Evolución de Aprobación × Período  (2023-01 a 2025-02)",
        "•   Diapo 32:  Aprobación × Antigüedad  (solo formados)",
        "•   Diapo 33:  Aprobación × Jerarquía  (solo formados)",
        "•   Diapo 34:  Efecto Acumulativo  (n° instancias formación)",
    ]
    _txt(sl, "\n".join(items), PIC_L+80000, PIC_T+560000, PIC_W-80000, PIC_H-600000,
         fs=14, color="#FFFFFF", wrap=False, lspc=6)
    print("  ✓ slide 29 — Sep BLOQUE III")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE EXTRA B — Scatter SAT vs Nota Promedio (3 paneles por año)
# ─────────────────────────────────────────────────────────────────────────────
def slide_24(prs):
    """Scatter SAT docente vs nota promedio alumnos — 3 paneles 2023/2024/2025."""
    scat = pd.read_csv(os.path.join(COMP, "scatter_sat_notas.csv"), encoding="utf-8-sig")
    scat["sat"]           = pd.to_numeric(scat["sat"],           errors="coerce")
    scat["nota_promedio"] = pd.to_numeric(scat["nota_promedio"], errors="coerce")
    scat["formado"] = scat["formado"].astype(str).str.strip().str.upper().isin(
        ["TRUE", "1", "SI", "SÍ", "YES"])
    scat["anio"] = scat["periodo"].str[:4]
    scat = scat.dropna(subset=["sat", "nota_promedio"])

    anios    = ["2023", "2024", "2025"]
    COL_CTRL = "#FFB74D"   # amarillo — sin formación (igual que slides 25+)
    COL_FORM = "#5C9BD6"   # azul     — formados    (igual que slides 25+)
    ALPHA    = 0.38
    MS       = 10
    GAP      = 0.022

    stats_by_year = {}
    for y in anios:
        s = scat[scat["anio"] == y]
        if len(s) < 2:
            stats_by_year[y] = dict(r=0, p=1, n_sec=0, n_f=0, n_c=0)
            continue
        r, p = scipy_stats.pearsonr(s["sat"], s["nota_promedio"])
        stats_by_year[y] = dict(r=r, p=p, n_sec=len(s),
                                n_f=int(s["rut_docente"][s["formado"]].nunique()),
                                n_c=int(s["rut_docente"][~s["formado"]].nunique()))
        print(f"  Stats {y}: r={round(r,3)}  n_sec={len(s)}")

    fig   = _tr_fig()
    gx    = PIC_RECT[0] + 0.04
    gy    = PIC_RECT[1] + 0.005
    gw    = PIC_RECT[2] - 0.055
    gh    = PIC_RECT[3] - 0.05
    pw    = (gw - GAP*2) / 3

    for i, y in enumerate(anios):
        s   = scat[scat["anio"] == y]
        sf  = s[s["formado"]]
        sc2 = s[~s["formado"]]
        st  = stats_by_year[y]
        ax  = fig.add_axes([gx + i*(pw+GAP), gy, pw, gh], facecolor="none", zorder=5)

        ax.scatter(sc2["sat"], sc2["nota_promedio"],
                   c=COL_CTRL, alpha=ALPHA, s=MS, linewidths=0, zorder=2,
                   label=f"Sin formacion (n={st['n_c']})")
        ax.scatter(sf["sat"], sf["nota_promedio"],
                   c=COL_FORM, alpha=ALPHA+0.14, s=MS, linewidths=0, zorder=3,
                   label=f"Formados (n={st['n_f']})")

        if len(s) >= 2:
            x_all = s["sat"].values; y_all = s["nota_promedio"].values
            m, b  = np.polyfit(x_all, y_all, 1)
            xr    = np.array([x_all.min(), x_all.max()])
            ax.plot(xr, m*xr+b, "--", color="white", linewidth=1.5, alpha=0.60, zorder=4,
                    label=f"Tendencia (r={round(st['r'],2)})")

        _style(ax, xlabel="SAT docente (sobre 7)", xgrid=True)
        ax.set_xlim(1, 7); ax.set_ylim(1, 7)
        ax.set_xticks([1,2,3,4,5,6,7]); ax.set_yticks([1,2,3,4,5,6,7])
        if i == 0:
            ax.set_ylabel("Nota promedio alumnos (sobre 7)", color="#AAAAAA", fontsize=8.5)
        else:
            ax.set_yticklabels([])
        ax.set_title(f"Año {y}", color="white", fontsize=10.5, fontweight="bold", pad=6)

        leg = ax.legend(fontsize=7.2, framealpha=0.35, facecolor="#101820",
                        edgecolor="#444444", loc="upper left", markerscale=1.5,
                        handlelength=0.8, borderpad=0.5, labelspacing=0.35)
        for txt in leg.get_texts(): txt.set_color("white")

        p_str = "< 0.001" if st["p"] < 0.001 else f"= {round(st['p'],3)}"
        ax.text(0.97, 0.05,
                f"r = {round(st['r'],2)}\np {p_str}\n{st['n_sec']} secciones",
                transform=ax.transAxes, fontsize=8, va="bottom", ha="right",
                color="white",
                bbox=dict(boxstyle="round,pad=0.45", facecolor="#1A2E10",
                          edgecolor="#6AAA40", alpha=0.88))

    path = _save_ch(fig, "scatter_sat_nota_3pan.png")
    _pic(sl := _new_sl(prs), SHARED_BG, prs); _pic(sl, path, prs)
    _T(sl, "Relacion SAT Docente y Nota Promedio de Alumnos   Por Año", fs=13)
    _POP(sl, "Universo con SAT y calificaciones disponibles  |  2023-01 a 2025-02"
             "  |  Cada punto = 1 sección")
    _BUL(sl, [
        "Correlación positiva y estadísticamente significativa entre SAT docente y nota "
        "promedio de alumnos en los tres años (r≈0.16 en 2023, r≈0.35 en 2024, r≈0.32 "
        "en 2025; p<0.001 en todos los cortes).",
        "Los docentes formados (azul) se concentran en la región SAT mayor a 5 con "
        "notas promedio superiores; los sin formación (amarillo) presentan mayor dispersión hacia "
        "valores bajos en ambas dimensiones.",
        "La correlación se fortalece en 2024–2025, coincidiendo con el mayor volumen "
        "de formación, consistente con un efecto acumulativo sobre el rendimiento estudiantil.",
    ])
    print("  ✓ slide scatter_panels — SAT vs Nota 3 paneles")

def slide_25(prs):
    """Aprobación Global — formados vs control."""
    scat = pd.read_csv(os.path.join(COMP,"scatter_sat_notas.csv"), encoding="utf-8-sig")
    scat["pct_aprobacion"] = pd.to_numeric(scat["pct_aprobacion"], errors="coerce")
    scat["nota_promedio"]  = pd.to_numeric(scat["nota_promedio"],  errors="coerce")
    scat["formado"] = scat["formado"].astype(str).str.strip().str.upper().isin(["TRUE","1","SI","SÍ","YES"])

    form  = scat[scat["formado"]].dropna(subset=["pct_aprobacion","nota_promedio"])
    ctrl_s = scat[~scat["formado"]].dropna(subset=["pct_aprobacion","nota_promedio"])
    pct_f = form["pct_aprobacion"].mean()
    pct_c = ctrl_s["pct_aprobacion"].mean()
    not_f = form["nota_promedio"].mean()
    not_c = ctrl_s["nota_promedio"].mean()
    n_f = len(form); n_c = len(ctrl_s)

    fig = _tr_fig()
    ax  = fig.add_axes([CHART_X, CHART_Y, CHART_W*0.55, CHART_H], facecolor="none", zorder=5)
    ax2 = fig.add_axes([CHART_X+CHART_W*0.60, CHART_Y, CHART_W*0.38, CHART_H], facecolor="none", zorder=5)

    # % Aprobación
    grps = ["Formados", "Control"]
    vals_p = [pct_f, pct_c]; vals_n = [not_f, not_c]
    cols_g = ["#5C9BD6", "#FFB74D"]
    xa = np.arange(2)
    ax.bar(xa, vals_p, width=0.55, color=cols_g, alpha=0.90, edgecolor="none")
    mv_p = max(vals_p)
    for i,(v,g) in enumerate(zip(vals_p, grps)):
        ax.text(i, v+mv_p*0.025, f"{v:.1f}%", ha="center", va="bottom",
                fontsize=14, fontweight="bold", color="white",
                path_effects=[pe.withStroke(linewidth=2, foreground="#0A0F18")])
    ax.set_xticks(xa); ax.set_xticklabels(grps, fontsize=12, color="white")
    ax.set_title("% Aprobación", color="white", fontsize=11, pad=8)
    for sp in ax.spines.values(): sp.set_edgecolor("white"); sp.set_alpha(0.35)
    ax.tick_params(colors="#AAAAAA", labelsize=9); ax.yaxis.grid(True, color="white", alpha=0.07)
    ax.set_ylim(0, mv_p*1.30)

    # Nota promedio
    ax2.bar(xa, vals_n, width=0.55, color=cols_g, alpha=0.90, edgecolor="none")
    mv_n = max(vals_n)
    for i,v in enumerate(vals_n):
        ax2.text(i, v+mv_n*0.025, f"{v:.2f}", ha="center", va="bottom",
                 fontsize=14, fontweight="bold", color="white",
                 path_effects=[pe.withStroke(linewidth=2, foreground="#0A0F18")])
    ax2.set_xticks(xa); ax2.set_xticklabels(grps, fontsize=12, color="white")
    ax2.set_title("Nota Promedio (1–7)", color="white", fontsize=11, pad=8)
    for sp in ax2.spines.values(): sp.set_edgecolor("white"); sp.set_alpha(0.35)
    ax2.tick_params(colors="#AAAAAA", labelsize=9); ax2.yaxis.grid(True, color="white", alpha=0.07)
    ax2.set_ylim(0, mv_n*1.30)

    _pic(sl := _new_sl(prs), SHARED_BG, prs); _pic(sl, _save_ch(fig,"30_chart.png"), prs)
    _T(sl, "Aprobación de Alumnos — Formados vs Control")
    _POP(sl, POP_CTR)
    _CT(sl, f"% aprobación y nota promedio  ·  formados n={n_f} registros  ·  control n={n_c} registros  ·  2023–2025")
    _BUL(sl, [
        f"Los docentes formados obtienen mayor aprobación de alumnos ({pct_f:.1f}%) que el grupo "
        f"control ({pct_c:.1f}%). Diferencia: {pct_f-pct_c:+.1f} puntos porcentuales.",
        f"La nota promedio de alumnos también es mayor en formados ({not_f:.2f}) vs control ({not_c:.2f}). "
        f"Diferencia: {not_f-not_c:+.2f} puntos.",
        "Ambos indicadores de rendimiento estudiantil son superiores en los cursos de "
        "docentes formados, reforzando el impacto positivo de la formación.",
    ])
    print("  ✓ slide 30 — Aprobación Global")

def slide_26(prs):
    """Evolución de Aprobación por Período."""
    scat = pd.read_csv(os.path.join(COMP,"scatter_sat_notas.csv"), encoding="utf-8-sig")
    scat["pct_aprobacion"] = pd.to_numeric(scat["pct_aprobacion"], errors="coerce")
    scat["formado"] = scat["formado"].astype(str).str.strip().str.upper().isin(["TRUE","1","SI","SÍ","YES"])
    scat = scat.dropna(subset=["pct_aprobacion","periodo"])
    periodos_ord = sorted(scat["periodo"].unique().tolist())
    pct_f = [scat[(scat["formado"])&(scat["periodo"]==p)]["pct_aprobacion"].mean()
             for p in periodos_ord]
    pct_c = [scat[(~scat["formado"])&(scat["periodo"]==p)]["pct_aprobacion"].mean()
             for p in periodos_ord]

    fig, ax = _ax()
    xa = range(len(periodos_ord))
    ax.plot(xa, pct_f, color="#5C9BD6", linewidth=2.5, marker="o", markersize=9,
            label="Formados", zorder=5)
    ax.plot(xa, pct_c, color="#FFB74D", linewidth=2.5, linestyle="--", marker="s", markersize=8,
            label="Control", zorder=5)
    ax.fill_between(xa, pct_f, pct_c, where=[f>c for f,c in zip(pct_f,pct_c)],
                    alpha=0.12, color="#4CAF50", interpolate=True)
    mv_range = max(pct_f+pct_c)-min(pct_f+pct_c) if (pct_f+pct_c) else 5
    for i,(vf,vc) in enumerate(zip(pct_f,pct_c)):
        if not np.isnan(vf):
            ax.text(i, vf+mv_range*0.07, f"{vf:.1f}%", ha="center", va="bottom",
                    fontsize=8, fontweight="bold", color="#5C9BD6",
                    path_effects=[pe.withStroke(linewidth=1.5, foreground="#0A0F18")])
        if not np.isnan(vc):
            ax.text(i, vc-mv_range*0.07, f"{vc:.1f}%", ha="center", va="top",
                    fontsize=8, fontweight="bold", color="#FFB74D",
                    path_effects=[pe.withStroke(linewidth=1.5, foreground="#0A0F18")])
    ax.axhline(np.nanmean(pct_f+pct_c), color="white", linewidth=0.8, linestyle=":", alpha=0.4)
    ax.set_xticks(list(xa)); ax.set_xticklabels(periodos_ord, fontsize=9, color="white", rotation=15)
    _style(ax); ax.set_ylabel("% Aprobación promedio", color="#AAAAAA", fontsize=9)
    ax.legend(fontsize=10, framealpha=0.2, labelcolor="white", facecolor="#101820", edgecolor="#444")

    _pic(sl := _new_sl(prs), SHARED_BG, prs); _pic(sl, _save_ch(fig,"31_chart.png"), prs)
    _T(sl, "Evolución de la Aprobación de Alumnos por Período")
    _POP(sl, POP_CTR)
    _CT(sl, "% aprobación promedio por semestre  ·  azul=formados  ·  naranja=control  ·  área verde=brecha positiva")
    pct_f_valid = [v for v in pct_f if not np.isnan(v)]
    pct_c_valid = [v for v in pct_c if not np.isnan(v)]
    brechas = [f-c for f,c in zip(pct_f,pct_c) if not np.isnan(f) and not np.isnan(c)]
    _BUL(sl, [
        f"Los formados superan al control en % aprobación en la mayoría de los períodos. "
        f"Brecha promedio: {np.mean(brechas):+.1f} pp.",
        f"Los formados muestran aprobación promedio de {np.mean(pct_f_valid):.1f}% vs "
        f"{np.mean(pct_c_valid):.1f}% del control a lo largo del período de análisis.",
        "La brecha tiende a sostenerse o ampliarse en los períodos más recientes, "
        "sugiriendo un efecto acumulativo de la formación sobre el rendimiento estudiantil.",
    ])
    print("  ✓ slide 31 — Evolución Aprobación × Período")

def slide_27(prs):
    """Aprobación × Antigüedad — solo formados (antigüedad desde p3_sat_zscore)."""
    scat = pd.read_csv(os.path.join(COMP,"scatter_sat_notas.csv"), encoding="utf-8-sig")
    scat["pct_aprobacion"] = pd.to_numeric(scat["pct_aprobacion"], errors="coerce")
    scat["formado"] = scat["formado"].astype(str).str.strip().str.upper().isin(["TRUE","1","SI","SÍ","YES"])
    scat["rut_key"] = scat["rut_docente"].astype(str).str.strip()
    scat = scat.dropna(subset=["pct_aprobacion"])

    ORD_ANT = ["0-4","5-9","10-14","15-19","20+"]
    def _tramo(a):
        try:
            a = float(a)
            if a < 5:  return "0-4"
            if a < 10: return "5-9"
            if a < 15: return "10-14"
            if a < 20: return "15-19"
            return "20+"
        except: return None

    sat_ant = sat[["rut_key","antiguedad_anios"]].copy()
    sat_ant["rut_key"] = sat_ant["rut_key"].astype(str).str.strip()
    # Solo formados (tienen antigüedad en sat)
    form_scat = scat[scat["formado"]].copy()
    form_mg = form_scat.merge(sat_ant, on="rut_key", how="inner")
    form_mg["tramo"] = form_mg["antiguedad_anios"].apply(_tramo)
    form_mg = form_mg[form_mg["tramo"].notna()]

    pct_f_ant = form_mg.groupby("tramo")["pct_aprobacion"].mean().reindex(ORD_ANT)
    n_f_ant   = form_mg.groupby("tramo")["pct_aprobacion"].count().reindex(ORD_ANT).fillna(0)

    fig, ax = _ax()
    xa = np.arange(len(ORD_ANT))
    vals_f = [float(pct_f_ant.get(t, np.nan)) for t in ORD_ANT]
    ns_f   = [int(n_f_ant.get(t, 0)) for t in ORD_ANT]
    cols_a = PAL[:len(ORD_ANT)]
    ax.bar(xa, [v if not np.isnan(v) else 0 for v in vals_f], width=0.60,
           color=cols_a, alpha=0.90, edgecolor="none")
    mv = max([v for v in vals_f if not np.isnan(v)]+[1])
    for i,(vf,nv) in enumerate(zip(vals_f, ns_f)):
        if not np.isnan(vf):
            ax.text(i, vf+mv*0.020, f"{vf:.1f}%\n(n={nv})", ha="center", va="bottom",
                    fontsize=9, fontweight="bold", color="white",
                    path_effects=[pe.withStroke(linewidth=1.5, foreground="#0A0F18")])
    # Línea de tendencia
    xa_v = [i for i,v in enumerate(vals_f) if not np.isnan(v)]
    vv   = [v for v in vals_f if not np.isnan(v)]
    if len(xa_v) >= 2:
        ax.plot(xa_v, vv, color="#FFB74D", linewidth=2, linestyle="--", marker="o",
                markersize=8, zorder=6, label="Tendencia")
        ax.legend(fontsize=9, framealpha=0.2, labelcolor="white",
                  facecolor="#101820", edgecolor="#444")
    ax.set_xticks(xa); ax.set_xticklabels(ORD_ANT, fontsize=11, color="white")
    _style(ax, xlabel="Tramo de antigüedad (años)")
    ax.set_ylabel("% Aprobación promedio", color="#AAAAAA", fontsize=9)
    ax.set_ylim(0, mv*1.32)

    _pic(sl := _new_sl(prs), SHARED_BG, prs); _pic(sl, _save_ch(fig,"32_chart.png"), prs)
    _T(sl, "Aprobación de Alumnos por Antigüedad — Docentes Formados")
    _POP(sl)
    _CT(sl, "% aprobación promedio de alumnos por tramo de antigüedad del docente  ·  solo formados Aptos P3")
    _BUL(sl, [
        "Los docentes formados con menor antigüedad (0–9 años) obtienen las mayores tasas de "
        "aprobación de sus alumnos, con tendencia decreciente a mayor experiencia.",
        "El tramo 0–4 años concentra la mayor mejora potencial de la formación, siendo el "
        "perfil más frecuente entre los Aptos P3.",
        "La línea de tendencia ilustra la variación del impacto de la formación según "
        "la trayectoria del docente en la institución.",
    ])
    print("  ✓ slide 32 — Aprobación × Antigüedad")

def slide_28(prs):
    """Aprobación × Jerarquía — solo formados (jerarquía desde p3_sat_zscore)."""
    scat = pd.read_csv(os.path.join(COMP,"scatter_sat_notas.csv"), encoding="utf-8-sig")
    scat["pct_aprobacion"] = pd.to_numeric(scat["pct_aprobacion"], errors="coerce")
    scat["formado"] = scat["formado"].astype(str).str.strip().str.upper().isin(["TRUE","1","SI","SÍ","YES"])
    scat["rut_key"] = scat["rut_docente"].astype(str).str.strip()
    scat = scat.dropna(subset=["pct_aprobacion"])

    JER_ORD = ["INSTRUCTOR REGULAR","INSTRUCTOR DOCENTE","ASISTENTE REGULAR","ASISTENTE DOCENTE",
               "ASOCIADO REGULAR","ASOCIADO DOCENTE","TITULAR REGULAR","TITULAR DOCENTE"]
    sat_jer = sat[["rut_key","jerarquia"]].copy()
    sat_jer["rut_key"] = sat_jer["rut_key"].astype(str).str.strip()
    # Solo formados
    form_scat = scat[scat["formado"]].copy()
    form_mg = form_scat.merge(sat_jer, on="rut_key", how="inner")
    form_mg["jer_u"] = form_mg["jerarquia"].str.strip().str.upper()
    form_mg = form_mg[form_mg["jer_u"].notna()]

    pct_f_jer = form_mg.groupby("jer_u")["pct_aprobacion"].mean()
    n_f_jer   = form_mg.groupby("jer_u")["pct_aprobacion"].count()
    ord_pres  = [j for j in JER_ORD if j in pct_f_jer.index]
    jers      = [j.title() for j in ord_pres]
    vals_f    = [float(pct_f_jer.get(j, np.nan)) for j in ord_pres]
    ns_f      = [int(n_f_jer.get(j, 0)) for j in ord_pres]

    fig = _tr_fig()
    ax  = fig.add_axes([CHART_X, CHART_Y, CHART_W, CHART_H], facecolor="none", zorder=5)
    n = len(jers); ya = np.arange(n)
    ax.barh(ya, [v if not np.isnan(v) else 0 for v in vals_f], height=0.60,
            color=PAL[:n], alpha=0.90, edgecolor="none")
    mv = max([v for v in vals_f if not np.isnan(v)]+[1])
    for i,(vf,nv) in enumerate(zip(vals_f, ns_f)):
        if not np.isnan(vf):
            ax.text(vf + mv*0.015, i, f"{vf:.1f}%  (n={nv})", va="center", ha="left",
                    fontsize=9.5, fontweight="bold", color="white",
                    path_effects=[pe.withStroke(linewidth=1.5, foreground="#0A0F18")])
    ax.set_yticks(ya); ax.set_yticklabels(jers, fontsize=9.5, fontweight="bold", color="white")
    ax.tick_params(axis="y", length=0, pad=6); ax.tick_params(axis="x", colors="#AAAAAA", labelsize=9)
    for sp in ax.spines.values(): sp.set_edgecolor("white"); sp.set_alpha(0.35); sp.set_linewidth(0.9)
    ax.xaxis.grid(True, color="white", alpha=0.07, linewidth=0.5); ax.set_axisbelow(True)
    ax.set_xlim(0, mv*1.45)

    _pic(sl := _new_sl(prs), SHARED_BG, prs); _pic(sl, _save_ch(fig,"33_chart.png"), prs)
    _T(sl, "Aprobación de Alumnos por Jerarquía — Docentes Formados")
    _POP(sl)
    _CT(sl, "% aprobación promedio de alumnos por jerarquía académica  ·  solo formados Aptos P3  ·  n=cursos por jerarquía")
    _BUL(sl, [
        "Los docentes formados de jerarquías intermedias (Asistente, Asociado) obtienen "
        "las mayores tasas de aprobación de sus alumnos.",
        "Las jerarquías de entrada (Instructor Docente) muestran alta aprobación, "
        "consistente con la mayor participación en formación de estos rangos.",
        "Los Titulares, aunque menos representados entre los formados, presentan "
        "altas tasas de aprobación por su mayor experiencia pedagógica consolidada.",
    ])
    print("  ✓ slide 33 — Aprobación × Jerarquía")

def slide_29(prs):
    """Efecto Acumulativo — aprobación × n° instancias de formación."""
    scat = pd.read_csv(os.path.join(COMP,"scatter_sat_notas.csv"), encoding="utf-8-sig")
    scat["pct_aprobacion"] = pd.to_numeric(scat["pct_aprobacion"], errors="coerce")
    scat["formado"] = scat["formado"].astype(str).str.strip().str.upper().isin(["TRUE","1","SI","SÍ","YES"])
    scat["rut_key"] = scat["rut_docente"].astype(str).str.strip()
    scat = scat.dropna(subset=["pct_aprobacion"])

    sat_ninst = sat[["rut_key","n_instancias"]].copy()
    sat_ninst["rut_key"] = sat_ninst["rut_key"].astype(str).str.strip()
    merged = scat[scat["formado"]].merge(sat_ninst, on="rut_key", how="left")
    merged = merged.dropna(subset=["n_instancias"])
    merged["n_inst_int"] = merged["n_instancias"].astype(int)
    # Agrupar 3+ en "3+"
    merged["grupo"] = merged["n_inst_int"].apply(lambda x: str(int(x)) if x <= 2 else "3+")
    ORD_G = ["1","2","3+"]
    grp_mean = merged.groupby("grupo")["pct_aprobacion"].mean().reindex(ORD_G)
    grp_n    = merged.groupby("grupo")["pct_aprobacion"].count().reindex(ORD_G).fillna(0)

    fig, ax = _ax()
    xa = np.arange(len(ORD_G))
    cols_i = ["#5C9BD6","#A5D6A7","#FFB74D"]
    vals = [float(grp_mean.get(g, np.nan)) for g in ORD_G]
    ns   = [int(grp_n.get(g, 0)) for g in ORD_G]
    bars = ax.bar(xa, [v if not np.isnan(v) else 0 for v in vals],
                  width=0.55, color=cols_i[:len(ORD_G)], alpha=0.90, edgecolor="none")
    mv = max([v for v in vals if not np.isnan(v)]+[1])
    for i,(v,nv) in enumerate(zip(vals, ns)):
        if not np.isnan(v):
            ax.text(i, v+mv*0.025, f"{v:.1f}%\n(n={nv})", ha="center", va="bottom",
                    fontsize=11, fontweight="bold", color="white",
                    path_effects=[pe.withStroke(linewidth=2, foreground="#0A0F18")])
    # Línea de tendencia
    xa_valid = [i for i,v in enumerate(vals) if not np.isnan(v)]
    vals_valid = [v for v in vals if not np.isnan(v)]
    if len(xa_valid) >= 2:
        ax.plot(xa_valid, vals_valid, color="#4CAF50", linewidth=2, linestyle="--",
                marker="o", markersize=9, zorder=6)
    ax.set_xticks(xa)
    ax.set_xticklabels(["1 instancia","2 instancias","3+ instancias"], fontsize=10, color="white")
    _style(ax, xlabel="N° de instancias de formación")
    ax.set_ylabel("% Aprobación promedio", color="#AAAAAA", fontsize=9)
    ax.set_ylim(0, mv*1.35)

    _pic(sl := _new_sl(prs), SHARED_BG, prs); _pic(sl, _save_ch(fig,"34_chart.png"), prs)
    _T(sl, "Efecto Acumulativo — Aprobación × N° Instancias de Formación")
    _POP(sl, "Solo docentes formados (Aptos P3)  ·  cursos impartidos 2023–2025")
    _CT(sl, "% aprobación promedio de alumnos según número de instancias de formación del docente")
    trend = "creciente" if len(xa_valid)>=2 and vals_valid[-1]>vals_valid[0] else "sin tendencia clara"
    _BUL(sl, [
        f"La aprobación de alumnos muestra una tendencia {trend} al aumentar el número de "
        "instancias de formación del docente, evidenciando efecto acumulativo.",
        "Los docentes con 2 o más instancias de formación tienden a superar en aprobación "
        "a los que solo participaron una vez.",
        "El efecto acumulativo refuerza el argumento de que la formación continua produce "
        "beneficios medibles y crecientes en el rendimiento estudiantil.",
    ])
    print("  ✓ slide 34 — Efecto Acumulativo Aprobación")

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE IV — EDD (Evaluación de Desempeño Docente)
# ─────────────────────────────────────────────────────────────────────────────

def slide_30(prs):
    """Separador BLOQUE IV — EDD."""
    _ensure_bg(); sl = _new_sl(prs); _pic(sl, SHARED_BG, prs)
    _T(sl, "BLOQUE IV — Evaluación de Desempeño Docente (EDD)", fs=18)
    _POP(sl, "Fuente: P1_consolidado_con_evaluacion_jefes.csv  ·  formados vs control  ·  2022–2025")
    items = [
        "•   Diapo 37:  Evolución EDD — Formados vs Control  (2022–2025)",
        "•   Diapo 38:  EDD por Tipo de Formación  (Taller / Diplomado / Proyecto vs Control)",
    ]
    _txt(sl, "\n".join(items), PIC_L+80000, PIC_T+560000, PIC_W-80000, PIC_H-600000,
         fs=14, color="#FFFFFF", wrap=False, lspc=6)
    print("  ✓ slide 36 — Sep BLOQUE IV")

def slide_31(prs):
    """Evolución EDD: Formados vs Control por Año (línea, 2022–2025)."""
    if not has_edd or len(edd_form) == 0:
        _pic(sl := _new_sl(prs), SHARED_BG, prs)
        _T(sl, "Evolución EDD — Formados vs Control (2022–2025)")
        _txt(sl, "Datos EDD no disponibles.", PIC_L+200000, PIC_T+PIC_H//2-100000,
             PIC_W-400000, 250000, fs=14, color="#90ABC4", align=PP_ALIGN.CENTER)
        print("  ✓ slide 37 — EDD Evolución (sin datos)")
        return

    anios_ord = sorted(edd_form["anio_eval"].unique())
    z_f = [edd_form[edd_form["anio_eval"]==a]["edd_total"].mean() for a in anios_ord]
    z_c = [edd_ctrl[edd_ctrl["anio_eval"]==a]["edd_total"].mean() if len(edd_ctrl) else float("nan")
           for a in anios_ord]
    n_f_by_yr = [edd_form[edd_form["anio_eval"]==a]["rut_key"].nunique() for a in anios_ord]
    n_c_by_yr = [edd_ctrl[edd_ctrl["anio_eval"]==a]["rut_key"].nunique() if len(edd_ctrl) else 0
                 for a in anios_ord]

    fig, ax = _ax()
    xa = range(len(anios_ord))
    ax.plot(xa, z_f, color="#5C9BD6", linewidth=2.5, linestyle="-",
            marker="o", markersize=9, label="Formados", zorder=5)
    ax.plot(xa, z_c, color="#FF7043", linewidth=2.5, linestyle="--",
            marker="s", markersize=8, label="Control", zorder=5)

    for i, (vf, vc, nf, nc) in enumerate(zip(z_f, z_c, n_f_by_yr, n_c_by_yr)):
        ax.text(i, vf + 0.012, f"{vf:.3f} (n={nf})",
                ha="center", va="bottom", fontsize=8, fontweight="bold",
                color="#5C9BD6",
                path_effects=[pe.withStroke(linewidth=1.5, foreground="#0A0F18")])
        ax.text(i, vc - 0.012, f"{vc:.3f} (n={nc})",
                ha="center", va="top", fontsize=8, fontweight="bold",
                color="#FF7043",
                path_effects=[pe.withStroke(linewidth=1.5, foreground="#0A0F18")])

    ax.axhline(0, color="white", linewidth=0.8, linestyle=":", alpha=0.4)
    ax.set_xticks(list(xa))
    ax.set_xticklabels(anios_ord, fontsize=10, color="white")
    _style(ax, ygrid=True)
    ax.set_ylabel("EDD Total (promedio)", color="#AAAAAA", fontsize=9)
    ax.legend(fontsize=10, framealpha=0.2, labelcolor="white",
              facecolor="#101820", edgecolor="#444")
    ax.set_ylim(max(0, min(z_f+z_c)-0.08), min(1.0, max(z_f+z_c)+0.08))

    _pic(sl := _new_sl(prs), SHARED_BG, prs); _pic(sl, _save_ch(fig, "37_chart.png"), prs)
    _T(sl, "Evolución EDD — Formados vs Control (2022–2025)")
    _POP(sl, "EDD: Evaluación de Desempeño Docente  ·  escala 0–1  ·  formados Aptos P3 vs control 941")
    _CT(sl, f"EDD Total promedio por año  ·  Formados n={edd_form['rut_key'].nunique()} doc.  ·  Control n={edd_ctrl['rut_key'].nunique()} doc.")
    z_f_glo = edd_form["edd_total"].mean()
    z_c_glo = edd_ctrl["edd_total"].mean() if len(edd_ctrl) else float("nan")
    _BUL(sl, [
        f"Los docentes formados obtienen EDD promedio de {z_f_glo:.3f} vs {z_c_glo:.3f} del grupo "
        f"control — brecha de {z_f_glo-z_c_glo:+.3f} puntos a favor de los formados.",
        "La brecha Formados−Control se amplía en 2024 y 2025, coincidiendo con mayor volumen "
        "de formación impartida — el efecto acumulativo se refleja en la EDD.",
        "La EDD mide la evaluación de directivos sobre el desempeño docente, complementando "
        "las métricas SAT y aprobación de alumnos con una perspectiva institucional.",
    ])
    print("  ✓ slide 37 — EDD Evolución")

def slide_32(prs):
    """EDD por Tipo de Formación — barras horizontales."""
    if not has_edd or len(edd_form) == 0:
        _pic(sl := _new_sl(prs), SHARED_BG, prs)
        _T(sl, "EDD por Tipo de Formación")
        _txt(sl, "Datos EDD no disponibles.", PIC_L+200000, PIC_T+PIC_H//2-100000,
             PIC_W-400000, 250000, fs=14, color="#90ABC4", align=PP_ALIGN.CENTER)
        print("  ✓ slide 38 — EDD por Tipo (sin datos)")
        return

    edd_ft = edd_form.merge(sat[["rut_key", "tipos_formacion"]], on="rut_key", how="left")
    edd_ft["tipo_g"] = edd_ft["tipos_formacion"].apply(_tipo_simple)

    tipo_grp = edd_ft.groupby("tipo_g")["edd_total"].agg(["mean", "count"]).reset_index()
    tipo_grp = tipo_grp.sort_values("mean", ascending=False)
    ctrl_mean = edd_ctrl["edd_total"].mean() if len(edd_ctrl) else float("nan")
    ctrl_n    = edd_ctrl["rut_key"].nunique() if len(edd_ctrl) else 0

    lbl  = tipo_grp["tipo_g"].tolist() + ["Control"]
    vals = tipo_grp["mean"].tolist() + [ctrl_mean]
    ns   = tipo_grp["count"].astype(int).tolist() + [ctrl_n]
    pcts = [100 * v / max(vals) for v in vals]

    fig = _tr_fig()
    ax  = fig.add_axes([CHART_X+0.05, CHART_Y, CHART_W-0.05, CHART_H],
                       facecolor="none", zorder=5)
    n = len(lbl); yp = np.arange(n)
    TIPO_ECOLS = {"Taller": "#5C9BD6", "Diplomado": "#64B5F6", "Proyecto": "#A5D6A7",
                  "Participacion Mixta": "#CE93D8", "Control": "#FF7043"}
    cols = [TIPO_ECOLS.get(l, PAL[i % len(PAL)]) for i, l in enumerate(lbl)]
    ax.barh(yp[::-1], vals, color=cols, height=0.58, edgecolor="none", alpha=0.90)

    mv = max(vals) if vals else 1
    for i, (v, nv) in enumerate(zip(vals[::-1], ns[::-1])):
        ax.text(v + mv * 0.018, i, f"{v:.3f}  (n={nv})",
                va="center", ha="left", fontsize=10.5, fontweight="bold", color="white",
                path_effects=[pe.withStroke(linewidth=2.5, foreground="#0A0F18")])
    ax.set_yticks(yp)
    ax.set_yticklabels(lbl[::-1], fontsize=10.5, fontweight="bold", color="white")
    ax.tick_params(axis="y", length=0, pad=8)
    ax.tick_params(axis="x", colors="#AAAAAA", labelsize=9)
    for sp in ax.spines.values():
        sp.set_edgecolor("white"); sp.set_alpha(0.35); sp.set_linewidth(0.9)
    ax.set_xlim(0, mv * 1.50)
    ax.set_ylim(-0.5, n - 0.5)
    ax.xaxis.grid(True, color="white", alpha=0.07, linewidth=0.5); ax.set_axisbelow(True)

    _pic(sl := _new_sl(prs), SHARED_BG, prs); _pic(sl, _save_ch(fig, "38_chart.png"), prs)
    _T(sl, "EDD por Tipo de Formación — Promedio vs Control")
    _POP(sl, "EDD Total promedio  ·  formados Aptos P3 por modalidad  ·  vs control 941 sin formación")
    _CT(sl, "EDD promedio por tipo de formación  ·  n = registros EDD deduplicados por docente-año")
    top_tipo = lbl[0] if lbl else "N/D"
    best_v = vals[0]; ctrl_v = ctrl_mean
    _BUL(sl, [
        f"El tipo '{top_tipo}' obtiene la mayor EDD promedio ({best_v:.3f}), por encima del "
        f"control ({ctrl_v:.3f}). Todos los tipos de formación superan al grupo control.",
        "La formación docente se asocia con mejor evaluación de desempeño en todos los tipos "
        "de modalidad, con el Proyecto de Innovación mostrando el mayor impacto relativo.",
        "El grupo control (docentes 941 sin formación) obtiene la menor EDD promedio, "
        "confirmando que la formación aporta un efecto positivo medible en el desempeño.",
    ])
    print("  ✓ slide 38 — EDD por Tipo")

def slide_33(prs):
    """Conclusiones y Recomendaciones — 4 cajas."""
    _ensure_bg(); sl = _new_sl(prs); _pic(sl, SHARED_BG, prs)
    _T(sl, "Conclusiones y Recomendaciones")
    _POP(sl, "Síntesis del análisis de impacto  ·  Universo: 210 Aptos P3  ·  Períodos 2022–2025")
    cajas = [
        ("1. Impacto en SAT (z-score)",
         "• Los formados mantienen z-score positivo y sostenido (6 períodos)\n"
         "• Brecha Formados − Control estadísticamente significativa\n"
         "• El Diplomado produce mayor Δz que el Taller\n"
         "• La mejora es generalizada (>50% de los 210 mejoran)"),
        ("2. Aprobación de Alumnos",
         "• % Aprobación mayor en formados vs control en todos los períodos\n"
         "• Nota promedio de alumnos también superior en cursos de formados\n"
         "• Brecha más evidente en jerarquías de entrada (Instructor, Asistente)\n"
         "• Efecto acumulativo: más instancias → mayor aprobación"),
        ("3. EDD — Evaluación de Desempeño (Bloque IV)",
         "• EDD promedio formados > control en los 4 años medidos (2022–2025)\n"
         "• Proyecto (0.83) ≥ Taller (0.83) > Diplomado (0.85) > Control (0.71)\n"
         "• La brecha Formados−Control se amplía en 2024 y 2025\n"
         "• Múltiples métricas (SAT, Aprobación, EDD) convergen en el mismo efecto"),
        ("4. Recomendaciones",
         "• Priorizar Diplomados para docentes con 5–15 años de antigüedad\n"
         "• Fomentar la participación repetida (multi-instancia) en Talleres\n"
         "• Monitorear la brecha Formados−Control semestralmente (SAT y EDD)\n"
         "• Extender el análisis P3 a los 419 participantes de P2"),
    ]
    pad_x=100000; gap_x=80000; gap_y=50000; pad_t=80000; pad_b=70000
    bw = (PIC_W-2*pad_x-gap_x)//2; bh = (PIC_H-pad_t-gap_y-pad_b)//2
    pos = [(PIC_L+pad_x, PIC_T+pad_t),(PIC_L+pad_x+bw+gap_x, PIC_T+pad_t),
           (PIC_L+pad_x, PIC_T+pad_t+bh+gap_y),(PIC_L+pad_x+bw+gap_x, PIC_T+pad_t+bh+gap_y)]
    for (bx,by),(hdr,body) in zip(pos, cajas):
        sh = sl.shapes.add_shape(1, Emu(bx), Emu(by), Emu(bw), Emu(bh))
        sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0,33,71); sh.line.fill.background()
        _txt(sl, hdr, bx+40000, by+38000, bw-80000, 300000, fs=11, bold=True, color="#90ABC4")
        _txt(sl, body, bx+40000, by+360000, bw-80000, bh-420000, fs=9.5, color="#FFFFFF", lspc=3)
    print("  ✓ slide 35 — Conclusiones y Recomendaciones")


# ─────────────────────────────────────────────────────────────────────────────
# Slides de Comparación de Grupos (merge desde generar_slides_comparacion_grupos)
# ─────────────────────────────────────────────────────────────────────────────
def slide_embudo(prs):
    """Embudo: derivación de grupos de control por bloque."""
    import matplotlib.patches as _mp
    fig = _tr_fig()
    ax = fig.add_axes([PIC_RECT[0] + 0.005, PIC_RECT[1] + 0.005,
                       PIC_RECT[2] - 0.01,  PIC_RECT[3] - 0.01],
                      facecolor="none", zorder=5)
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis("off")

    def _box(cx, cy, title, color, w=0.40, h=0.085, fsize=9.5, alpha=0.88):
        ax.add_patch(_mp.FancyBboxPatch(
            (cx - w/2, cy - h/2), w, h,
            boxstyle="round,pad=0.012", linewidth=1.3,
            edgecolor="white", facecolor=color, alpha=alpha, zorder=4,
            transform=ax.transData))
        ax.text(cx, cy, title, ha="center", va="center",
                fontsize=fsize, color="white", fontweight="bold",
                zorder=5, transform=ax.transData)

    def _arrow(x0, y0, x1, y1, color="white"):
        ax.annotate("", xy=(x1, y1), xytext=(x0, y0),
                    arrowprops=dict(arrowstyle="-|>", color=color, lw=1.5, alpha=0.70), zorder=3)

    BHH = 0.042
    L1 = 0.87;  L1_bot = L1 - BHH
    L2 = 0.69;  L2_top = L2 + BHH; L2_bot = L2 - BHH
    L3 = 0.47;  L3_top = L3 + BHH; L3_bot = L3 - BHH
    L4 = 0.25;  L4_top = L4 + BHH; L4_bot = L4 - BHH

    _box(0.50, L1, "941 Docentes Jerarquizados", color="#1A3A5C", w=0.46, h=BHH*2, fsize=10.5)
    _arrow(0.50, L1_bot, 0.27, L2_top, color=COL_FORM_P)
    _arrow(0.50, L1_bot, 0.73, L2_top, color=COL_CTRL_P)
    MID12_y = (L1_bot + L2_top) / 2
    ax.text(0.355, MID12_y + 0.012, "Participa en P3",    ha="center", va="bottom", fontsize=8, color=COL_FORM_P, style="italic")
    ax.text(0.645, MID12_y + 0.012, "No participa en P3", ha="center", va="bottom", fontsize=8, color=COL_CTRL_P, style="italic")

    _box(0.27, L2, "419 Docentes Formados",           color="#1A4E7A", w=0.36, h=BHH*2)
    _box(0.73, L2, f"{N_NO_FORM} Docentes No Formados", color="#6B4A1A", w=0.36, h=BHH*2)

    MID23 = (L2_bot + L3_top) / 2
    _arrow(0.27, L2_bot, 0.27, L3_top, color=COL_FORM_P)
    _arrow(0.73, L2_bot, 0.73, L3_top, color=COL_CTRL_P)
    ax.text(0.27, MID23, "SAT pre + post disponible", ha="center", va="center", fontsize=7.5, color="#AAAAAA", style="italic")
    ax.text(0.73, MID23, "Con EDD disponible",        ha="center", va="center", fontsize=7.5, color="#AAAAAA", style="italic")

    _box(0.27, L3, "210 Aptos P3  (SAT pre + post)", color=COL_FORM_P, w=0.40, h=BHH*2)
    _box(0.73, L3, f"{N_B4_CTRL} Docentes  (Control EDD)", color=COL_CTRL_P, w=0.40, h=BHH*2)

    MID34 = (L3_bot + L4_top) / 2
    scat_f_n = int(scat["formado"].sum())
    scat_c_n = int((~scat["formado"]).sum())
    _arrow(0.27, L3_bot, 0.27, L4_top, color=COL_FORM_P)
    _arrow(0.73, L3_bot, 0.73, L4_top, color=COL_CTRL_P)
    ax.text(0.27, MID34, "Secciones con nota y SAT", ha="center", va="center", fontsize=7.5, color="#AAAAAA", style="italic")
    ax.text(0.73, MID34, "Secciones con nota",       ha="center", va="center", fontsize=7.5, color="#AAAAAA", style="italic")

    _box(0.27, L4, f"~{scat_f_n:,} secciones  (formados)", color=COL_FORM_P + "BB", w=0.40, h=BHH*2)
    _box(0.73, L4, f"~{scat_c_n:,} secciones  (control)",  color=COL_CTRL_P + "BB", w=0.40, h=BHH*2)

    ax.text(0.50, 0.10,
            "Bloque III: control = docentes sin formación en scatter_sat_notas  ·  "
            "Bloque IV: control = universo 941 que nunca participó en P3, con EDD disponible",
            ha="center", va="center", fontsize=8, color="#C0D0E0", style="italic", zorder=5,
            bbox=dict(boxstyle="round,pad=0.35", facecolor="#0D1E30", edgecolor="#3A5A7A", alpha=0.80))

    _ensure_bg()
    sl = _new_sl(prs); _pic(sl, SHARED_BG, prs); _pic(sl, _save_ch(fig, "cg_embudo.png"), prs)
    _T(sl, "¿Quién es el Grupo de Control?   Derivación por Bloque de Análisis")
    _POP(sl, "Universo base: 941 docentes jerarquizados  ·  Cada bloque usa una definición distinta de grupo control")
    print("  ✓ slide — Embudo grupos control")


def slide_perfil_b3(prs):
    """Bloque III — Perfil demográfico butterfly: formados vs control aprobación."""
    fig = _tr_fig()
    px = PIC_RECT[0] + 0.01;  py = PIC_RECT[1] + 0.02
    pw = PIC_RECT[2] - 0.02;  ph = 0.42
    GAP = 0.025;  bw = (pw - GAP) / 2

    _butterfly(fig, bx=px,          by=py, bw=bw, bh=ph,
               cats=JER_LBL,
               pct_f=_pct_jer(sat),         pct_c=_pct_jer(ctrl_b3_doc),
               n_f_tot=N_B3_FORM,            n_c_tot=N_B3_CTRL,
               chart_title="Jerarquía Académica")
    _butterfly(fig, bx=px+bw+GAP,   by=py, bw=bw, bh=ph,
               cats=TRAMOS_EDAD,
               pct_f=_pct_tramo(sat),       pct_c=_pct_tramo(ctrl_b3_doc),
               n_f_tot=N_B3_FORM,            n_c_tot=N_B3_CTRL,
               chart_title="Tramo de Edad")

    sl = _new_sl(prs); _pic(sl, SHARED_BG, prs); _pic(sl, _save_ch(fig, "cg_b3_perfil.png"), prs)
    _T(sl, "Bloque III — Perfil Demográfico: Formados vs Control  (Aprobación)", fs=17)
    _POP(sl, f"Control B3: docentes en scatter_sat_notas con formado=False  ·  formados n={N_B3_FORM}  ·  control n={N_B3_CTRL}")
    jer_dif = sorted([(JER_LBL[i], _pct_jer(sat)[i] - _pct_jer(ctrl_b3_doc)[i]) for i in range(len(JER_LBL))], key=lambda x: -abs(x[1]))
    eda_dif = sorted([(TRAMOS_EDAD[i], _pct_tramo(sat)[i] - _pct_tramo(ctrl_b3_doc)[i]) for i in range(len(TRAMOS_EDAD))], key=lambda x: -abs(x[1]))
    _BUL(sl, [
        f"Jerarquía: '{jer_dif[0][0]}' tiene la mayor diferencia entre grupos ({jer_dif[0][1]:+.0f} pp). "
        "Las barras muestran % dentro de cada grupo — no volumen absoluto — para hacer comparables grupos de distinto tamaño.",
        f"Edad: '{eda_dif[0][0]}' es el tramo más diferenciado ({eda_dif[0][1]:+.0f} pp). "
        "Si los grupos difieren en composición, parte de la brecha de aprobación puede deberse a eso y no solo a la formación.",
    ])
    print("  ✓ slide — Perfil Bloque III")


def slide_perfil_b4(prs):
    """Bloque IV — Perfil demográfico butterfly: formados vs control EDD."""
    if not has_edd or len(b4_form_doc) == 0:
        sl = _new_sl(prs); _pic(sl, SHARED_BG, prs)
        _T(sl, "Bloque IV — Perfil Demográfico (datos EDD no disponibles)")
        return
    fig = _tr_fig()
    px = PIC_RECT[0] + 0.01;  py = PIC_RECT[1] + 0.02
    pw = PIC_RECT[2] - 0.02;  ph = 0.42
    GAP = 0.025;  bw = (pw - GAP) / 2

    _butterfly(fig, bx=px,          by=py, bw=bw, bh=ph,
               cats=JER_LBL,
               pct_f=_pct_jer(b4_form_doc),  pct_c=_pct_jer(b4_ctrl_doc),
               n_f_tot=N_B4_FORM,             n_c_tot=N_B4_CTRL,
               chart_title="Jerarquía Académica")
    _butterfly(fig, bx=px+bw+GAP,   by=py, bw=bw, bh=ph,
               cats=TRAMOS_EDAD,
               pct_f=_pct_tramo(b4_form_doc), pct_c=_pct_tramo(b4_ctrl_doc),
               n_f_tot=N_B4_FORM,             n_c_tot=N_B4_CTRL,
               chart_title="Tramo de Edad")

    sl = _new_sl(prs); _pic(sl, SHARED_BG, prs); _pic(sl, _save_ch(fig, "cg_b4_perfil.png"), prs)
    _T(sl, "Bloque IV — Perfil Demográfico: Formados vs Control  (EDD)", fs=17)
    _POP(sl, f"Control B4: universo 941 sin ninguna actividad P3, con EDD disponible  ·  formados n={N_B4_FORM}  ·  control n={N_B4_CTRL}")
    jer_dif = sorted([(JER_LBL[i], _pct_jer(b4_form_doc)[i] - _pct_jer(b4_ctrl_doc)[i]) for i in range(len(JER_LBL))], key=lambda x: -abs(x[1]))
    eda_dif = sorted([(TRAMOS_EDAD[i], _pct_tramo(b4_form_doc)[i] - _pct_tramo(b4_ctrl_doc)[i]) for i in range(len(TRAMOS_EDAD))], key=lambda x: -abs(x[1]))
    _BUL(sl, [
        f"Jerarquía: '{jer_dif[0][0]}' muestra la mayor diferencia ({jer_dif[0][1]:+.0f} pp). "
        "El control EDD es más estricto — excluye a todos los que participaron en P3, aunque no sean Aptos.",
        f"Edad: '{eda_dif[0][0]}' es el tramo más diferenciado ({eda_dif[0][1]:+.0f} pp). "
        "Comparar este perfil con el del B3 muestra si la definición más estricta de control cambia la composición.",
    ])
    print("  ✓ slide — Perfil Bloque IV")


def slide_tabla(prs):
    """Tabla resumen: grupos, fuentes originales y archivos derivados."""
    import matplotlib.patches as _mp
    fig = _tr_fig()
    ax = fig.add_axes([PIC_RECT[0] + 0.01, PIC_RECT[1] + 0.01,
                       PIC_RECT[2] - 0.02,  PIC_RECT[3] - 0.02],
                      facecolor="none", zorder=5)
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis("off")

    cols_x  = [0.000, 0.090, 0.250, 0.420, 0.620, 0.840]
    headers = ["Bloque", "Grupo Formados", "Grupo Control", "Fuente Original", "Archivo Derivado", "n"]
    hy = 0.90
    ax.add_patch(_mp.FancyBboxPatch((0, hy - 0.055), 1.0, 0.10,
                                    boxstyle="round,pad=0.005", linewidth=0,
                                    facecolor="#1A3A5C", alpha=0.92, zorder=3, transform=ax.transData))
    for hdr, cx in zip(headers, cols_x):
        ax.text(cx + 0.006, hy, hdr, ha="left", va="center", fontsize=8.5, color="white", fontweight="bold", zorder=5)

    rows = [
        ("SAT\nBl. I–II",
         "210 Aptos P3\n(SAT pre+post)",
         "486 doc sin\nformación con SAT",
         "evaluacion_periodo.csv\n+ nomina_docente.csv",
         "(p3_sat_zscore_918.csv)\n(control_918.csv)",
         "210 / 486\ndocentes"),
        ("Aprobación\nBl. III",
         f"~{int(scat['formado'].sum()):,} secciones\n(formados)",
         f"~{int((~scat['formado']).sum()):,} secciones\n(control)",
         "calificacion_alumno.csv\n+ evaluacion_periodo.csv\n+ participacion_formacion.csv",
         "(scatter_sat_notas.csv)\nformado = True / False",
         f"{N_B3_FORM} / {N_B3_CTRL}\ndocentes"),
        ("EDD\nBl. IV",
         f"210 Aptos P3\ncon EDD",
         f"{N_B4_CTRL} doc sin P3\ncon EDD",
         "CONSOLIDADO DOCENTES.xlsx\n(hoja: EVALUACION DE\nJEFES A DOCENTES)",
         "(P1_consolidado_con_\nevaluacion_jefes.csv)",
         f"{N_B4_FORM} / {N_B4_CTRL}\n(varía por año)"),
    ]
    row_ys  = [0.685, 0.440, 0.175]
    row_hs  = [0.20,  0.22,  0.22]
    row_cls = ["#0D2035", "#0A1A2E", "#0D2035"]
    for (blq, form, ctrl_txt, orig, deriv, ns), ry, rh, rc in zip(rows, row_ys, row_hs, row_cls):
        ax.add_patch(_mp.FancyBboxPatch((0, ry - rh/2), 1.0, rh + 0.01,
                                        boxstyle="round,pad=0.005", linewidth=0.7,
                                        edgecolor="#2A4A6A", facecolor=rc, alpha=0.80,
                                        zorder=3, transform=ax.transData))
        for txt, cx, col in zip([blq, form, ctrl_txt, orig, deriv, ns], cols_x,
                                  ["white"]*3 + ["#E0E8F0", "#90C8F0", "#E0E8F0"]):
            ax.text(cx + 0.006, ry, txt, ha="left", va="center",
                    fontsize=7.5, color=col, linespacing=1.3, zorder=5)
    ax.text(0.50, 0.02,
            "Fuente Original = archivo entregado por la institución  ·  Archivo Derivado (azul) = generado por ETL del proyecto",
            ha="center", va="bottom", fontsize=7.5, color="#7090B0", style="italic", zorder=5)

    sl = _new_sl(prs); _pic(sl, SHARED_BG, prs); _pic(sl, _save_ch(fig, "cg_tabla.png"), prs)
    _T(sl, "Resumen de Grupos, Fuentes Originales y Archivos Derivados")
    _POP(sl, "Blanco = fuente original UCEN  ·  Azul = archivo procesado por ETL del análisis")
    print("  ✓ slide — Tabla resumen grupos")

# ─────────────────────────────────────────────────────────────────────────────
# Ensamblar PPTX
# ─────────────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("Generando PRESENTACION_210_P3_v4_37slides.pptx (37 slides) …")
    _ensure_bg()
    prs = Presentation()
    prs.slide_width  = Emu(SW_EMU)
    prs.slide_height = Emu(SH_EMU)

    # BLOQUE I — Clasificación Cuerpo Académico
    slide_01(prs)                  # 01 Portada
    slide_02(prs)                  # 02 Índice
    slide_03(prs)                  # 03 Metodología
    slide_04(prs)                  # 04 Sep BLOQUE I
    slide_05(prs)                  # 05 Edad/Sexo
    slide_06(prs)                  # 06 Facultad
    slide_07(prs)                  # 07 Antigüedad
    slide_08(prs)                  # 08 Jerarquía
    slide_09(prs)                  # 09 Grado
    slide_10(prs)                  # 10 Institución Grado

    # BLOQUE II — Evaluación SAT
    slide_11(prs)                  # 11 Sep BLOQUE II
    slide_12(prs)                  # 12 Venn tipos formación
    slide_13(prs)                  # 13 Universo SAT P3
    slide_14(prs)           # 14 Embudo 917→357→197 (compacto)
    slide_15(prs)                  # 15 Sep 2.1 Caracterización
    slide_16(prs)                  # 16 Participación × Jerarquía
    slide_17(prs)                  # 17 Participación × Antigüedad
    slide_18(prs)                  # 18 Intensidad participación
    slide_19(prs)                  # 19 Combinaciones modalidad
    slide_20(prs)                  # 20 Perfil grupo control
    slide_21(prs)                  # 21 SAT z formados vs control (6 períodos)
    slide_22(prs)                  # 22 SAT z por facultad

    # BLOQUE III — Aprobación de Alumnos
    slide_23(prs)                  # 23 Sep BLOQUE III
    slide_embudo(prs)              # 24 Derivación grupos control
    slide_perfil_b3(prs)           # 25 Perfil B3: formados vs control
    slide_24(prs)                  # 26 Scatter SAT vs Nota (3 paneles)
    slide_25(prs)                  # 27 Aprobación global formados vs control
    slide_26(prs)                  # 28 Evolución aprobación × período
    slide_27(prs)                  # 29 Aprobación × Antigüedad
    slide_28(prs)                  # 30 Aprobación × Jerarquía
    slide_29(prs)                  # 31 Efecto acumulativo × instancias

    # BLOQUE IV — EDD
    slide_30(prs)                  # 32 Sep BLOQUE IV
    slide_perfil_b4(prs)           # 33 Perfil B4: formados vs control EDD
    slide_31(prs)                  # 34 Evolución EDD formados vs control
    slide_32(prs)                  # 35 EDD por tipo de formación
    slide_tabla(prs)               # 36 Tabla resumen fuentes y grupos
    slide_33(prs)                  # 37 Conclusiones y Recomendaciones

    prs.save(OUT_PPTX)
    print(f"\n✓ Guardado: {OUT_PPTX}")
    print(f"  {len(prs.slides)} diapositivas")

import sys; sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

BASE      = os.path.dirname(__file__)
TEMPLATE  = r"c:\Users\r.gonzalez_fluxsolar.LAPTOP-FLUX-ECO\Downloads\Analisis_UCEN_v2\PPT FINAL SOLO LOS N 197 SOLO APTOS P3.pptx"
OUT       = r"c:\Users\r.gonzalez_fluxsolar.LAPTOP-FLUX-ECO\Downloads\Analisis_UCEN_v2\PRESENTACION_197_P3_v1.pptx"

# ── Colores del tema ──────────────────────────────────────────────────────────
C_DK1    = RGBColor(0x10, 0x18, 0x20)  # dk1  — texto principal (navy oscuro)
C_DK2    = RGBColor(0x00, 0x46, 0x80)  # dk2  — azul oscuro
C_NAVY   = RGBColor(0x00, 0x21, 0x47)  # acento portada
C_ACC1   = RGBColor(0x90, 0xAB, 0xC4)  # accent1 — barra título (gris-azul)
C_ACC2   = RGBColor(0x00, 0x72, 0xE6)  # accent2 — azul vivo
C_CREAM  = RGBColor(0xF7, 0xF0, 0xE1)  # lt2    — cream
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY   = RGBColor(0x55, 0x55, 0x55)
C_STAR   = RGBColor(0x00, 0x72, 0xE6)  # estrella hallazgo

# ── Cargar template (hereda slide master con fondo) ───────────────────────────
prs = Presentation(TEMPLATE)
SLIDE_W = prs.slide_width.emu    # 12192000
SLIDE_H = prs.slide_height.emu   # 6858000

LAY_BLANK = prs.slide_layouts[1]   # Blank

# ── Limpiar todas las slides existentes ───────────────────────────────────────
def clear_slides(prs):
    NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    sldIdLst = prs.slides._sldIdLst
    for sId in list(sldIdLst):
        rId = sId.get(f'{{{NS_R}}}id')
        try: prs.part.drop_rel(rId)
        except: pass
    sldIdLst.clear()

clear_slides(prs)

# ── Helpers ───────────────────────────────────────────────────────────────────
def png(name): return os.path.join(BASE, name)

def txb(slide, l, t, w, h, text, size, bold=False, italic=False,
        color=None, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color: run.font.color.rgb = color
    return tb

def bullets_box(slide, l, t, w, h, items, size=12, color=None, bold_first=False):
    tb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(size)
        run.font.bold = (bold_first and i == 0)
        if color: run.font.color.rgb = color
    return tb

def add_pic(slide, png_path, l, t, w, h):
    if os.path.exists(png_path):
        return slide.shapes.add_picture(png_path, Emu(l), Emu(t), Emu(w), Emu(h))
    print(f"  FALTA: {os.path.basename(png_path)}")

def title_bar(slide, text):
    """Barra de título con color accent1 + texto blanco."""
    bar = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(SLIDE_W), Emu(720000))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_ACC1
    bar.line.fill.background()
    txb(slide, 200000, 80000, SLIDE_W - 400000, 560000,
        text, 20, bold=True, color=C_WHITE)

def pop_tag(slide, text="Universo: 197 Aptos P3  ·  SAT baseline / durante / resultado"):
    """Etiqueta pequeña de población bajo el título."""
    txb(slide, 200000, 740000, SLIDE_W - 400000, 330000,
        text, 10, italic=True, color=C_DK2)

def hallazgo_marker(slide, num=None):
    """Estrella azul bottom-left para slides hallazgo."""
    label = f"★  HALLAZGO {num}" if num else "★  PRESENTAR"
    txb(slide, 80000, SLIDE_H - 680000, 2200000, 580000,
        label, 13, bold=True, color=C_STAR)

# ── PLANTILLAS DE SLIDES ──────────────────────────────────────────────────────

def make_separator(bloque, title, subtitle=""):
    """Slide separador de bloque (solo texto centrado)."""
    s = prs.slides.add_slide(LAY_BLANK)
    if bloque:
        txb(s, 0, 1800000, SLIDE_W, 700000, bloque,
            18, bold=False, color=C_ACC2, align=PP_ALIGN.CENTER)
    txb(s, 0, 2600000, SLIDE_W, 1200000, title,
        34, bold=True, color=C_DK1, align=PP_ALIGN.CENTER)
    if subtitle:
        txb(s, 0, 3900000, SLIDE_W, 600000, subtitle,
            16, color=C_DK2, align=PP_ALIGN.CENTER)
    return s

def make_content(title, png_name, bullets, pop=None, hallazgo=None):
    """Slide estándar: barra título + subtítulo población + imagen + bullets."""
    s = prs.slides.add_slide(LAY_BLANK)
    title_bar(s, title)
    if pop is None:
        pop = "Universo: 197 Aptos P3  ·  SAT baseline / durante / resultado disponible"
    pop_tag(s, pop)
    # Imagen ocupa el centro
    img_h = 4500000 if len(bullets) <= 2 else 3900000
    add_pic(s, png(png_name), 150000, 1130000, SLIDE_W - 300000, img_h)
    # Bullets en la parte inferior
    bt = 1130000 + img_h + 80000
    if bullets:
        bullets_box(s, 150000, bt, SLIDE_W - 300000,
                    SLIDE_H - bt - 150000, bullets, size=11, color=C_DK1)
    if hallazgo:
        hallazgo_marker(s, hallazgo)
    return s

def make_split(title, png_name, bullets, img_frac=0.62, pop=None, hallazgo=None):
    """Slide con imagen izquierda (~62%) + texto derecha."""
    s = prs.slides.add_slide(LAY_BLANK)
    title_bar(s, title)
    if pop is None:
        pop = "Universo: 197 Aptos P3  ·  SAT baseline / durante / resultado disponible"
    pop_tag(s, pop)
    img_w = int(SLIDE_W * img_frac)
    txt_l = img_w + 200000
    txt_w = SLIDE_W - txt_l - 100000
    add_pic(s, png(png_name), 100000, 1100000, img_w, SLIDE_H - 1250000)
    if bullets:
        bullets_box(s, txt_l, 1100000, txt_w, SLIDE_H - 1250000,
                    bullets, size=11, color=C_DK1, bold_first=True)
    if hallazgo:
        hallazgo_marker(s, hallazgo)
    return s

def make_double(title, png_left, png_right, foot="", pop=None):
    """Slide con dos imágenes lado a lado."""
    s = prs.slides.add_slide(LAY_BLANK)
    title_bar(s, title)
    if pop is None:
        pop = "Universo: 197 Aptos P3  ·  SAT baseline / durante / resultado disponible"
    pop_tag(s, pop)
    half = (SLIDE_W - 350000) // 2
    add_pic(s, png(png_left),  100000, 1100000, half, SLIDE_H - 1300000)
    add_pic(s, png(png_right), 150000 + half, 1100000, half, SLIDE_H - 1300000)
    if foot:
        txb(s, 150000, SLIDE_H - 550000, SLIDE_W - 300000, 450000,
            foot, 10, italic=True, color=C_GRAY)
    return s

def make_text_slide(title, sections):
    """Slide de texto puro (conclusiones/recomendaciones): {header: [items]}."""
    s = prs.slides.add_slide(LAY_BLANK)
    title_bar(s, title)
    y = 820000
    for header, items in sections.items():
        txb(s, 200000, y, SLIDE_W - 400000, 500000,
            header, 14, bold=True, color=C_DK2)
        y += 480000
        for item in items:
            txb(s, 350000, y, SLIDE_W - 550000, 420000,
                f"• {item}", 11, color=C_DK1)
            y += 390000
        y += 100000
    return s

# ══════════════════════════════════════════════════════════════════════════════
# CONSTRUCCIÓN DEL DECK
# ══════════════════════════════════════════════════════════════════════════════

# ── 1. PORTADA ────────────────────────────────────────────────────────────────
portada = prs.slides.add_slide(LAY_BLANK)
# Acento izquierdo navy (igual que original)
bar = portada.shapes.add_shape(1, Emu(2248566), Emu(0), Emu(400000), Emu(SLIDE_H))
bar.fill.solid(); bar.fill.fore_color.rgb = C_NAVY
bar.line.fill.background()
txb(portada, 323850, 1800000, 7500000, 1100000,
    "Resultados del Perfeccionamiento Docente", 30, bold=True, color=C_DK1)
txb(portada, 323850, 3100000, 11000000, 800000,
    "Análisis de Impacto en Evaluación SAT, EDD y Rendimiento Académico", 22, color=C_DK2)
txb(portada, 323850, 4000000, 11000000, 600000,
    "Universo: 197 Aptos P3  ·  Períodos 2023–2025", 16, italic=True, color=C_DK2)
txb(portada, 323850, 5400000, 11000000, 500000,
    "Universidad Central de Chile  |  Producto 3: Análisis de Formación e Innovación",
    14, bold=True, color=C_DK1)

# ── 2. ÍNDICE ─────────────────────────────────────────────────────────────────
indice = prs.slides.add_slide(LAY_BLANK)
title_bar(indice, "Índice")
items_idx = [
    "BLOQUE I    —  Evaluación SAT: evolución z-score por período y tipo de formación",
    "BLOQUE II   —  Evaluación EDD: evolución y panel balanceado (170 mismos docentes, 4 años)",
    "BLOQUE III  —  Rendimiento Académico: aprobación formados vs control y efecto acumulativo",
    "BLOQUE IV   —  Análisis por Tipo: Taller, Diplomado, Proyecto (SAT y acumulativo)",
    "BLOQUE V    —  Cierre: Conclusiones y Recomendaciones",
]
bullets_box(indice, 300000, 900000, SLIDE_W - 600000, SLIDE_H - 1100000,
            items_idx, size=15, color=C_DK1)

# ── 3. MARCO METODOLÓGICO ─────────────────────────────────────────────────────
marco = prs.slides.add_slide(LAY_BLANK)
title_bar(marco, "Universo de Análisis y Metodología")
pop_tag(marco, "Universo: 917 jerarquizados  →  493 con período activo  →  357 formados  →  197 Aptos P3")

# 4 cajas de contenido
CAJAS = [
    (150000,  1130000, 5700000, 2600000,
     "1. Universo Aptos P3",
     ["917 docentes jerarquizados UCEN (base total)",
      "493 con actividad docente en período analizado",
      "357 participaron en ≥1 iniciativa de formación",
      "197 Aptos P3: SAT disponible en baseline, durante y resultado"]),
    (6350000, 1130000, 5700000, 2600000,
     "2. Marco Metodológico – Z-score",
     ["z = (SAT docente − media facultad) / DE facultad por período",
      "z = 0: en el promedio exacto de su facultad",
      "z = +1: supera a sus colegas (~84° percentil)",
      "Comparación siempre dentro de la misma facultad y semestre"]),
    (150000,  3950000, 5700000, 2700000,
     "3. Métricas de Evaluación",
     ["SAT Nota (1–7) estandarizada como z-score por facultad",
      "SAT % Recomendación: ¿recomendaría al docente?",
      "EDD: evaluación directa de desempeño (1–100 puntos)",
      "Aprobación alumnos: % con nota ≥ 4.0 en secciones del docente"]),
    (6350000, 3950000, 5700000, 2700000,
     "4. Diseño Cuasi-experimental",
     ["Grupo tratamiento: 130 docentes formados (con SAT P3)",
      "Grupo control: 227 docentes sin formación (con SAT P3)",
      "Comparación: z-score antes, durante y después de la formación",
      "Bases: SAT 2023-01 a 2025-02 · EDD · Notas de aprobación"]),
]
for (l, t, w, h, header, items) in CAJAS:
    # Rectángulo fondo suave
    rect = marco.shapes.add_shape(1, Emu(l), Emu(t), Emu(w), Emu(h))
    rect.fill.solid(); rect.fill.fore_color.rgb = C_CREAM
    rect.line.color.rgb = C_ACC1
    # Header
    txb(marco, l+150000, t+120000, w-300000, 380000, header,
        12, bold=True, color=C_DK2)
    # Items
    bullets_box(marco, l+150000, t+530000, w-300000, h-600000,
                items, size=10, color=C_DK1)

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE I — Evaluación SAT
# ─────────────────────────────────────────────────────────────────────────────
make_separator("BLOQUE I", "Evaluación SAT",
               "Z-score por período · Grupo tratamiento vs grupo control")

make_split(
    "Evolución Z-score SAT: Formados vs Control (6 Períodos)",
    "G6_control_vs_tratamiento_918.png",
    ["Grupo tratamiento (130 formados): consistentemente sobre el promedio de su facultad "
     "en los 6 períodos, con brecha media de +0.18z.",
     "Grupo control (227 sin formación): se mantiene bajo el promedio y con tendencia "
     "descendente a partir de 2024.",
     "La brecha se amplía en 2024-01 (+0.245z), coincidiendo con el peak de participación "
     "en talleres y diplomados."],
    img_frac=0.60,
    hallazgo=1
)

make_content(
    "Trayectoria por Tipo de Formación: Baseline → Resultado",
    "G1_linea_z_918.png",
    ["Los tres tipos mantienen z-score positivo antes y después de la formación.",
     "Proyectos muestra la mayor posición relativa (+0.257z post), seguido de Diplomados (+0.177z).",
     "Talleres (n=154, el grupo más grande) también supera la media post-formación (+0.056z)."],
    pop="Universo: 197 Aptos P3  ·  Baseline = semestre previo a la formación"
)

make_content(
    "Validación Estadística: Prueba t de Diferencia de Medias (SAT)",
    "G11_pruebas_t_918.png",
    ["Diferencias estadísticamente significativas en 3 de 6 períodos (p<0.05).",
     "La brecha es consistente en dirección (tratamiento > control) en todos los períodos.",
     "Con n=197, el tamaño muestral es suficiente para detectar efectos moderados (d>0.2)."],
    pop="Universo: 197 Aptos P3  ·  H₀: no hay diferencia entre grupos"
)

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE II — Evaluación EDD
# ─────────────────────────────────────────────────────────────────────────────
make_separator("BLOQUE II", "Evaluación EDD",
               "Evolución de la Evaluación de Desempeño Docente · 4 períodos")

make_content(
    "Evolución EDD: Formados vs Control",
    "G9_edd_evolucion_918.png",
    ["Los docentes formados muestran una trayectoria ascendente en EDD consistente con SAT.",
     "La brecha EDD se consolida en los períodos posteriores a la formación.",
     "Ambas métricas (SAT y EDD) convergen en señalar ventaja del grupo tratamiento."],
    pop="Universo: 197 Aptos P3  ·  EDD: Evaluación Directa de Desempeño (1–100 pts)"
)

make_split(
    "Panel Balanceado EDD: 170 Mismos Docentes Seguidos 4 Años",
    "G9_panel_balanceado_918.png",
    ["Panel balanceado elimina el sesgo de composición: son exactamente los mismos 170 "
     "docentes observados en los 4 períodos EDD disponibles.",
     "Brecha a favor de los formados: +0.081 puntos (escala normalizada), sostenida en el tiempo.",
     "Refuerza la validez de Hallazgo 1: la ventaja no se debe a rotación de docentes "
     "sino a mejora real de los mismos individuos."],
    img_frac=0.58,
    hallazgo=2
)

make_content(
    "Validación Estadística EDD",
    "G12_pruebas_t_edd_918.png",
    ["Prueba t para EDD confirma significancia estadística de la brecha.",
     "Consistente con los resultados SAT: ambas métricas apuntan en la misma dirección."],
    pop="Universo: 197 Aptos P3  ·  H₀: EDD(formados) = EDD(control)"
)

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE III — Rendimiento Académico
# ─────────────────────────────────────────────────────────────────────────────
make_separator("BLOQUE III", "Rendimiento Académico de los Alumnos",
               "Tasa de aprobación y efecto acumulativo de la formación docente")

make_split(
    "Aprobación de Alumnos: Formados vs Control",
    "GN_aprobacion_918.png",
    ["Los alumnos de docentes formados muestran mayor tasa de aprobación promedio.",
     "Diferencia de ~+1.5 pp en aprobación (sobre base de ~91–93%).",
     "Aunque pequeña en términos absolutos, la diferencia es consistente y estadísticamente "
     "significativa dado el volumen de alumnos involucrados."],
    img_frac=0.60,
    hallazgo=3
)

make_content(
    "Aprobación por Jerarquía Académica",
    "G_aprobacion_jerarquia_918.png",
    ["Instructores Docentes: mayor diferencial positivo de aprobación entre formados y control.",
     "Asistentes y Asociados muestran patrones similares con ventaja en formados.",
     "Titulares: menor diferencia, posiblemente por mayor experiencia de base."],
    pop="Universo: 197 Aptos P3  ·  Mediana de tasa de aprobación por docente"
)

make_content(
    "¿Es Acumulativo el Efecto de la Formación?",
    "G_acumulacion_formacion_918.png",
    ["La tasa de aprobación de los alumnos sube progresivamente con más iniciativas de formación.",
     "4+ iniciativas: mediana más alta, superando en ~2-3 pp a docentes sin formación.",
     "El efecto no es lineal: se consolida a partir de la 2ª iniciativa, no de la 1ª."],
    pop="Universo: 197 Aptos P3  ·  Iniciativas = talleres + diplomados + proyectos cursados",
    hallazgo=4
)

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE IV — Análisis por Tipo de Formación
# ─────────────────────────────────────────────────────────────────────────────
make_separator("BLOQUE IV", "Análisis por Tipo de Formación",
               "Taller · Diplomado · Proyecto  —  diferencias en impacto y acumulación")

make_content(
    "SAT por Tipo de Formación: Taller, Diplomado, Proyecto",
    "G2_barras_tipo_918.png",
    ["Los tres tipos elevan el z-score respecto al control, con magnitudes distintas.",
     "Proyectos: mayor impacto SAT por iniciativa (grupos más pequeños y selectivos).",
     "Talleres: impacto moderado pero con mayor cobertura (215 docentes, 376 iniciativas)."],
    pop="Universo: 197 Aptos P3  ·  Comparación z-score medio por tipo"
)

make_double(
    "Taller y Diplomado: Evolución Z-score por Jerarquía",
    "G7_taller_jerarquia_918.png",
    "G7_diplomado_jerarquia_918.png",
    foot="Izq: Taller — Der: Diplomado  ·  Universo 197 Aptos P3  ·  Z-score por jerarquía académica"
)

make_content(
    "¿Es Acumulativo en Talleres? Efecto por Número de Talleres Cursados",
    "G_acumulacion_talleres_918.png",
    ["Con 1–2 talleres la aprobación no supera al control; el efecto positivo aparece desde 3+.",
     "4+ talleres: mediana de aprobación más alta de todos los grupos.",
     "Implicancia: una sola participación en taller no es suficiente — se requiere exposición sostenida."],
    pop="Universo: 917 jerarquizados con notas disponibles  ·  Agrupados por N° de talleres cursados"
)

make_content(
    "Efecto del Diplomado: Comparación Binaria (Sin vs 1 Diplomado)",
    "G_acumulacion_diplomados_918.png",
    ["Todos los 201 docentes que cursaron un Diplomado realizaron exactamente 1 — sin variación de intensidad.",
     "Comparación binaria: docentes con diplomado vs sin diplomado en tasa de aprobación.",
     "Nota metodológica: al no haber variación de intensidad no es posible medir efecto acumulativo."],
    pop="Universo: 917 jerarquizados con notas disponibles  ·  100% de diplomados = exactamente 1 diplomado"
)

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE V — Cierre
# ─────────────────────────────────────────────────────────────────────────────
make_separator("BLOQUE V", "Conclusiones y Recomendaciones")

make_text_slide("Conclusiones", {
    "Hallazgo 1 — Evaluación SAT": [
        "Los 130 docentes formados superan consistentemente al grupo control (227) en z-score SAT "
        "en los 6 períodos analizados. Brecha media: +0.18z (~7 percentiles).",
        "La brecha se amplía en 2024, coincidiendo con el peak de participación en formación.",
    ],
    "Hallazgo 2 — Evaluación EDD (Panel Balanceado)": [
        "Seguidos los mismos 170 docentes durante 4 años, los formados mantienen ventaja EDD "
        "consistente, descartando sesgo de composición.",
    ],
    "Hallazgo 3 — Rendimiento Académico": [
        "Los alumnos de docentes formados tienen ~+1.5 pp más de aprobación. "
        "El efecto es acumulativo: mayor número de iniciativas → mayor impacto.",
    ],
    "Hallazgo 4 — Efecto Acumulativo": [
        "El efecto no es lineal: se consolida desde la 2ª iniciativa de formación. "
        "En talleres, el umbral es 3+ participaciones.",
    ],
})

make_text_slide("Recomendaciones", {
    "1. Diseñar rutas formativas progresivas": [
        "No actividades puntuales — secuencias de al menos 2–3 iniciativas por docente.",
        "Priorizar continuidad: el impacto se consolida con exposición sostenida.",
    ],
    "2. Focalizar según jerarquía y antigüedad": [
        "Instructores y Asistentes Docentes muestran mayor potencial de mejora con formación.",
        "Docentes con 15+ años y menor participación: diseñar estrategias diferenciadas.",
    ],
    "3. Ampliar la cobertura P3": [
        "Solo 197 de 917 docentes son aptos P3. Ampliar la captura de SAT pre/post formación "
        "permitirá medir el impacto sobre una base más representativa en años futuros.",
    ],
    "4. Fortalecer Diplomados y Proyectos": [
        "Mayor impacto por iniciativa — aunque menor cobertura. Escalar con cuidado la calidad.",
    ],
})

# ── Guardar ───────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"\n✓ Guardado: {OUT}")
print(f"  Slides: {len(prs.slides)}")
for i, sl in enumerate(prs.slides):
    titles = [sh.text_frame.text[:60] for sh in sl.shapes
              if sh.has_text_frame and sh.text_frame.text.strip()]
    first = titles[0] if titles else "(sin texto)"
    print(f"  [{i+1:02d}] {first}")

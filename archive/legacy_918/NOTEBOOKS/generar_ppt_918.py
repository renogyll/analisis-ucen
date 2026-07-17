import sys; sys.stdout.reconfigure(encoding="utf-8")
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor  # noqa: F401
from pptx.enum.dml import MSO_THEME_COLOR
from lxml import etree

BASE     = os.path.dirname(__file__)
TPL      = os.path.join(BASE, "..", "..", "PPT_UCENTRAL_ADMINISTRATIVO-ACADEMICO.pptx")
OUT      = os.path.join(BASE, "..", "..", "PRESENTACION_UCEN_2026_v6.pptx")
IMG_DIR  = BASE

# ══════════════════════════════════════════════════════════════════════════════
# DEFINICIÓN DE SLIDES
# tipo: "portada", "seccion", "contenido", "texto"
# ══════════════════════════════════════════════════════════════════════════════
SLIDES = [
    # ── BLOQUE 0: APERTURA ────────────────────────────────────────────────────
    ("portada",
     "Análisis de Iniciativas de Formación Docente",
     "Impacto en el Rendimiento Académico y Satisfacción Estudiantil 2022–2025\nUniversidad Central de Chile",
     None),

    ("texto",
     "Agenda",
     ["Bloque A — Universo 917 Jerarquizados: Impacto de la Formación (P3)",
      "Bloque B — Universo 917 Jerarquizados: Participación en Formación (P2)",
      "Bloque C — Universo NOMINA 957: Jornada vs Honorario",
      "Conclusiones y Recomendaciones"],
     None),

    ("texto",
     "Resumen Ejecutivo",
     ["Los docentes formados superan al grupo control en SAT (+0.18z) y EDD (d=+0.955 en 2025)",
      "El 38.9% del cuerpo jerarquizado ha participado en formación — 61.1% sin ninguna instancia",
      "No hay diferencia SAT entre Jornada y Honorario (~6.1) — el tipo de contrato no incide",
      "Los Honorarios formados (6.24) superan incluso a los Jornada formados (6.21)",
      "El problema no es de eficacia sino de cobertura: 70% de Honorarios sin formación",
      "El efecto acumulativo requiere ≥3 instancias para manifestarse"],
     None),

    ("texto",
     "Nota Metodológica",
     ["Z-score: posición relativa del docente dentro de su facultad × período",
      "SAT Nota (escala 1-7) y SAT % Recomendación: correlación r=0.893",
      "EDD: Evaluación de Desempeño Docente por jefaturas (escala 0-1)",
      "Universo 917 jerarquizados: 493 con perfil completo + 424 solo nómina",
      "197 docentes aptos P3: con SAT válido en baseline y resultado"],
     None),

    # ── BLOQUE A: UNIVERSO 917 — P3 IMPACTO ──────────────────────────────────
    ("seccion", "Bloque A",
     "Universo 917 Jerarquizados\nImpacto de la Formación Docente (P3)", None),

    ("contenido", "Universo de Análisis P3",
     ["917 docentes jerarquizados → 357 formados (39%) → 197 aptos P3 (21%)",
      "Formados: 154 Taller, 36 Diplomado, 7 Proyecto · Control: 560 sin formación"],
     None),

    ("contenido", "Formados vs Control — SAT Nota por Período",
     ["Los docentes formados superan al control en los 6 períodos sin excepción",
      "Brecha promedio: +0.18z en nota, significativa en 3 de 6 períodos"],
     "G6_control_vs_tratamiento_918.png"),

    ("contenido", "Formados vs Control — % Recomendación por Período",
     ["Consistencia con SAT Nota (r=0.893): formados lideran también en recomendación",
      "Brecha promedio: +0.15z, significativa en los mismos 3 períodos"],
     "G6.2_control_bin_918.png"),

    ("contenido", "Validación Estadística SAT — Pruebas t",
     ["Significativo en 2024-02, 2025-01 y 2025-02 (p<0.05)",
      "Cohen d máximo +0.302 (efecto pequeño-mediano) en 2025-02"],
     "G11_pruebas_t_918.png"),

    ("contenido", "Evaluación de Jefes (EDD) — Evolución 2022–2025",
     ["La formación actúa como 'escudo': formados se recuperan en 2025, control sigue cayendo",
      "Cohen d = +0.955 en 2025 — efecto grande, el más contundente del estudio"],
     "G9_edd_evolucion_918.png"),

    ("contenido", "Validación Estadística EDD",
     ["Prueba t significativa en 2024 (p=0.002) y 2025 (p<0.001, d=+0.955)",
      "Chi-cuadrado: distribución de conceptos difiere significativamente en 2024"],
     "G12_pruebas_t_edd_918.png"),

    ("contenido", "SAT vs Rendimiento Estudiantil",
     ["Correlación positiva significativa (r=0.28, p<0.001) pero moderada",
      "Docentes formados ausentes del cuadrante 'mal evaluado + notas bajas'"],
     "G8_scatter_sat_notas_918.png"),

    ("contenido", "Tasas de Aprobación: Formados vs Control",
     ["Docentes formados: 86.5% aprobación vs 85.0% control (+1.5pp)",
      "La brecha se amplía en el tiempo: de 0pp en 2023-01 a +2pp en 2025"],
     "GN_aprobacion_918.png"),

    ("contenido", "Incidencia por Tipo de Formación — SAT Nota",
     ["Los 3 grupos parten sobre el promedio (z>0) — efecto selección",
      "Ningún tipo produce cambio significativo intra-docente pero sí vs control"],
     "G1_linea_z_918.png"),

    ("contenido", "Incidencia por Antigüedad × Tipo",
     ["Noveles (0-4 años): mayor beneficio del Diplomado",
      "Seniors (15+ años): responden mejor al Taller (+0.19z)"],
     "G3_heatmap_antiguedad_918.png"),

    ("contenido", "Incidencia por Jerarquía × Tipo",
     ["Instructores Docentes en Diplomado: mayor Δ positivo",
      "Asociados Docentes en Taller: efectividad en aula (SAT Nota +0.10)"],
     "G4_heatmap_jerarquia_918.png"),

    ("contenido", "Impacto por Sexo del Docente",
     ["Mujeres muestran mejor respuesta en Taller y Diplomado",
      "Brecha de género: 51% mujeres mejoró vs 43% hombres en Taller"],
     "G5_barras_sexo_918.png"),

    ("contenido", "Evolución por Jerarquía — Talleres",
     ["Asistentes Docentes parten bajo la media y alcanzan peak durante talleres",
      "Instructores Docentes estables — el taller consolida más que transforma"],
     "G7_taller_jerarquia_918.png"),

    ("contenido", "Evolución por Jerarquía — Diplomados",
     ["Instructores Docentes: caso de mayor éxito, parten bajo y suben",
      "ANOVA no significativo: el diplomado beneficia por igual a todos los rangos"],
     "G7_diplomado_jerarquia_918.png"),

    ("texto", "Dinámica Individual y Efecto Selección",
     ["Los docentes que eligen formarse ya partían desde una posición ventajosa (z>0)",
      "La ventaja del grupo formado es estable pero no se amplía significativamente",
      "El impacto real es colectivo (vs control) más que individual (pre-post)",
      "Limitaciones: diseño observacional, Proyecto n=7 no interpretable"],
     None),

    # ── BLOQUE B: UNIVERSO 917 — P2 PARTICIPACIÓN ────────────────────────────
    ("seccion", "Bloque B",
     "Universo 917 Jerarquizados\nParticipación en Formación Docente (P2)", None),

    ("contenido", "Cobertura Global de Participación",
     ["38.9% participó al menos 1 vez (357 de 917)",
      "61.1% sin ninguna instancia — principal oportunidad de mejora"],
     "p2_g11_cobertura.png"),

    ("contenido", "Evolución de la Participación 2022–2025",
     ["De 39 docentes (2022) a 211 (2024) — pico de cobertura",
      "2025 parcial: 86 docentes al momento del análisis"],
     "p2_g12_evolucion.png"),

    ("contenido", "Combinaciones de Modalidades",
     ["76% hizo una sola modalidad; D+T es la intersección más frecuente (65 doc.)",
      "Solo 5 docentes participaron en las 3 modalidades"],
     "G_venn_formacion_917_918.png"),

    ("contenido", "Participación por Jerarquía Académica",
     ["Instructor Docente lidera (45%); Titular Regular la menor tasa (14%)",
      "La participación no sigue patrón lineal con el escalafón"],
     "p2_g21_jerarquia.png"),

    ("contenido", "Participación por Antigüedad",
     ["Tramo 0-4 años: mayor tasa de participación (50%)",
      "Los tramos intermedios no siempre superan a los extremos"],
     "p2_g23_antiguedad.png"),

    ("contenido", "Distribución por Modalidad",
     ["Taller: 376 instancias, 215 docentes únicos (61% del total)",
      "Diplomado: 201 instancias, 201 docentes (1 por docente exacto)"],
     "p2_g31_modalidad.png"),

    ("contenido", "Intensidad de Talleres",
     ["59% tomó un solo taller; promedio 1.75 por docente",
      "Heavy user: José Faúndez (Honorario) con 11 instancias"],
     "G_intensidad_talleres_918.png"),

    ("contenido", "Intensidad de Proyectos",
     ["88% participó en 1 solo proyecto; máximo 4 (Juan Carlos Araya)",
      "Proyecto: modalidad más exclusiva y de menor repetición"],
     "G_intensidad_proyectos_918.png"),

    ("contenido", "Tipo de Formación por Jerarquía",
     ["Rangos de entrada concentran participación en Talleres",
      "Rangos superiores: mayor diversificación (Diplomado + Proyecto)"],
     "p2_g32_tipo_jerarquia.png"),

    ("contenido", "Brechas: Perfil de los No Participantes",
     ["560 docentes sin formación (61.1%): Instructor Docente mayor volumen",
      "0-4 años de antigüedad: 127 sin participación, grupo prioritario"],
     "p2_g41_brechas.png"),

    ("contenido", "Intensidad de Participación Global",
     ["61% de los formados hizo una sola instancia — evento puntual",
      "8% con 4+ instancias: núcleo de alta dedicación formativa"],
     "p2_g42_intensidad.png"),

    # ── BLOQUE C: NOMINA 957 — JORNADA VS HONORARIO ──────────────────────────
    ("seccion", "Bloque C",
     "Universo NOMINA 957\nJornada vs Honorario", None),

    ("contenido", "SAT Nota — Jornada vs Honorario",
     ["Trayectorias prácticamente idénticas: diferencia máxima 0.03 puntos",
      "El tipo de vínculo contractual no incide en la calidad percibida"],
     "G_sat_jornada_honorario_918.png"),

    ("contenido", "% Recomendación — Jornada vs Honorario",
     ["Convergen a 89.7% en 2025-02 — literalmente el mismo valor",
      "Variaciones son sistémicas, no atribuibles al tipo de contrato"],
     "G_recomendacion_jornada_honorario_918.png"),

    ("contenido", "EDD — Jornada vs Honorario",
     ["Solo 9 Honorarios con EDD (1.8%) — comparación inválida",
      "Brecha de cobertura instrumental, no de desempeño"],
     "G_edd_jornada_honorario_918.png"),

    ("contenido", "Oferta Formativa — Jornada vs Honorario",
     ["Jornada: 43.3% con formación; Honorario: 30.0%",
      "Proyecto es exclusivo de Jornada: 0 Honorarios en combinaciones con Proyecto"],
     "G_venn_formacion_contrato_918.png"),

    ("contenido", "SAT Formado × Tipo de Contrato",
     ["Honorario formado (6.24) supera a Jornada formado (6.21)",
      "Problema de cobertura, no de eficacia: ampliar acceso a Honorarios"],
     "G_sat_formado_x_contrato_918.png"),

    ("contenido", "¿Es Acumulativo el Efecto?",
     ["A partir de 3 instancias la mediana supera al grupo sin formación",
      "Sin significancia estadística por tamaños muestrales reducidos (n=28)"],
     "G_acumulacion_formacion_918.png"),

    # ── BLOQUE D: CIERRE ─────────────────────────────────────────────────────
    ("seccion", "Conclusiones",
     "Hallazgos Consolidados y Recomendaciones", None),

    ("texto", "Conclusiones Consolidadas",
     ["La formación docente genera impacto positivo medible en SAT, EDD y aprobación",
      "El efecto es de selección + protección: formados parten mejor y resisten mejor",
      "No hay diferencia Jornada/Honorario — el contrato no determina el desempeño",
      "La cobertura (39%) es la principal oportunidad: 6 de cada 10 nunca participaron",
      "El impacto se consolida a partir de 3+ instancias — participaciones únicas no bastan",
      "Los Honorarios formados rinden igual o mejor pero acceden menos a formación"],
     None),

    ("texto", "Recomendaciones",
     ["Institucionalizar Diplomados como requisito para avance en carrera académica",
      "Diseñar rutas formativas progresivas (no instancias aisladas)",
      "Ampliar cobertura a Honorarios: mismo retorno con menor acceso actual",
      "Focalizar Talleres en docentes noveles (0-4 años) y seniors (15+ años)",
      "Crear mecanismo de evaluación de desempeño para Honorarios (brecha EDD)",
      "Reestructurar programa de Proyectos: aumentar muestra y acceso"],
     None),

    ("portada",
     "Gracias",
     "Universidad Central de Chile\nConsultoría Formación Docente 2026",
     None),
]

# ══════════════════════════════════════════════════════════════════════════════
# GENERACIÓN
# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation(TPL)

# Eliminar slides de ejemplo
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

master = prs.slide_masters[5]
layouts = {layout.name: layout for layout in master.slide_layouts}

SW = prs.slide_width
SH = prs.slide_height

def remove_placeholder(slide, idx):
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx == idx:
            sp = ph._element
            sp.getparent().remove(sp)

def apply_theme_color(run):
    run.font.color.theme_color = MSO_THEME_COLOR.TEXT_2

def set_ph_text(slide, idx, text, size=Pt(24), bold=True):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text = text
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = size
                    run.font.bold = bold
                    apply_theme_color(run)
            return ph
    return None

def add_bullets_box(slide, bullets, left, top, width, height, font_size=Pt(14)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"• {bullet}"
        run.font.size = font_size
        run.font.bold = False
        apply_theme_color(run)
        p.space_after = Pt(6)

def add_title_bar(slide, text, top=Cm(0.8), max_width=Cm(15)):
    left = Cm(0.8)
    height = Cm(1.2)
    txBox = slide.shapes.add_textbox(left, top, max_width, height)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(16)
    run.font.bold = True
    apply_theme_color(run)
    p.alignment = PP_ALIGN.LEFT

# ── Generar slides ────────────────────────────────────────────────────────────
for tipo, titulo, contenido, img in SLIDES:

    if tipo == "portada":
        slide = prs.slides.add_slide(layouts["Diapositiva de título"])
        set_ph_text(slide, 0, titulo, Pt(28), True)
        set_ph_text(slide, 1, contenido, Pt(14), False)

    elif tipo == "seccion":
        slide = prs.slides.add_slide(layouts["Encabezado de sección"])
        set_ph_text(slide, 0, titulo, Pt(32), True)
        set_ph_text(slide, 1, contenido, Pt(16), False)
        remove_placeholder(slide, 10)
        remove_placeholder(slide, 13)

    elif tipo == "texto":
        slide = prs.slides.add_slide(layouts["Diapositiva de título"])
        set_ph_text(slide, 0, titulo, Pt(24), True)
        if contenido:
            bullets_text = "\n".join([f"• {b}" for b in contenido])
            ph = set_ph_text(slide, 1, bullets_text, Pt(14), False)
            if ph:
                for para in ph.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.LEFT
                    para.space_after = Pt(8)

    elif tipo == "contenido":
        if img:
            img_path = os.path.join(IMG_DIR, img)
            if os.path.exists(img_path):
                slide = prs.slides.add_slide(layouts["Diapositiva de título"])
                remove_placeholder(slide, 0)
                remove_placeholder(slide, 1)

                # Mitad izquierda = 16.5cm (de 33.8cm total)
                HALF_W = Cm(16.5)

                # Título una sola línea, arriba izquierda
                add_title_bar(slide, titulo, top=Cm(0.5), max_width=HALF_W)

                # Imagen proporcional centrada en mitad izquierda
                from PIL import Image as PILImage
                with PILImage.open(img_path) as pil_img:
                    orig_w, orig_h = pil_img.size
                aspect = orig_h / orig_w
                img_w = Cm(14.5)
                img_h = int(img_w * aspect)
                max_h = Cm(10)
                if img_h > max_h:
                    img_h = max_h
                    img_w = int(img_h / aspect)
                # Centrar horizontalmente en la mitad izquierda
                left = (HALF_W - img_w) // 2
                # Centrar verticalmente entre título (2cm) y bullets (14.5cm)
                zone_top = Cm(2.2)
                zone_bottom = Cm(14)
                zone_h = zone_bottom - zone_top
                top = zone_top + (zone_h - img_h) // 2
                pic = slide.shapes.add_picture(img_path, left, top, img_w, img_h)

                # Borde delgado
                from pptx.oxml.ns import qn as _qn
                spPr = pic._element.spPr
                ln = etree.SubElement(spPr, _qn('a:ln'))
                ln.set('w', '9525')
                solidFill = etree.SubElement(ln, _qn('a:solidFill'))
                srgb = etree.SubElement(solidFill, _qn('a:srgbClr'))
                srgb.set('val', 'BBBBBB')

                # Bullets pegados al fondo de la diapositiva
                if contenido:
                    add_bullets_box(slide, contenido,
                                    Cm(0.5), Cm(15), HALF_W, Cm(3), Pt(10))

                # Mitad derecha — espacio libre para la contraparte
                right_box = slide.shapes.add_textbox(
                    Cm(17.5), Cm(0.8), Cm(15), Cm(17))
                tf = right_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = "[Espacio disponible para observaciones]"
                run.font.size = Pt(11)
                run.font.italic = True
                apply_theme_color(run)
            else:
                slide = prs.slides.add_slide(layouts["Diapositiva de título"])
                set_ph_text(slide, 0, titulo, Pt(20), True)
                set_ph_text(slide, 1, f"[Imagen no encontrada: {img}]", Pt(14), False)
        else:
            slide = prs.slides.add_slide(layouts["Diapositiva de título"])
            set_ph_text(slide, 0, titulo, Pt(24), True)
            if contenido:
                bullets_text = "\n".join([f"• {b}" for b in contenido])
                set_ph_text(slide, 1, bullets_text, Pt(14), False)

prs.save(OUT)
print(f"PPT generada: {OUT}")
print(f"Total slides: {len(prs.slides)}")

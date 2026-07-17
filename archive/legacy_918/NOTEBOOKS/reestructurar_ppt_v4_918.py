import sys; sys.stdout.reconfigure(encoding="utf-8")
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

BASE   = os.path.dirname(__file__)
SRC    = os.path.join(BASE, "..", "..", "PRESENTACION_final borrador1.pptx")
OUT    = os.path.join(BASE, "..", "..", "PRESENTACION_final_v12.pptx")
IMGDIR = BASE

prs = Presentation(SRC)
layouts = {l.name: l for l in prs.slide_layouts}
LY_SEPARADOR = layouts["Encabezado de sección"]
LY_BLANK     = layouts["Blank"]

# ── Geometría (idéntica al patrón original de 1 imagen) ──────────────────────
SLIDE_W, SLIDE_H = Emu(12192000), Emu(6858000)

TITLE_POS  = dict(left=Emu(786582), top=Emu(250608), width=Emu(10599174), height=Emu(553998))
PIC_BOX    = dict(left=Emu(786581), top=Emu(884606), width=Emu(10599174), height=Emu(3950000))
BULLET_POS = dict(left=Emu(786581), top=Emu(4985000), width=Emu(10599174), height=Emu(1300000))

# Marcador fuera del ancho de la barra de título (que termina en 786582+10599174=11385756)
MARKER_POS = dict(left=Emu(11420000), top=TITLE_POS["top"],
                   width=Emu(706000), height=TITLE_POS["height"])

TITLE_FILL   = RGBColor(0x90, 0xAB, 0xC4)
TITLE_TEXT   = RGBColor(0x00, 0x50, 0x88)
BORDER_COLOR = RGBColor(0x00, 0x50, 0x88)   # azul (mismo tono que el título) para todos los bordes
BODY_TEXT    = RGBColor(0x00, 0x50, 0x88)   # azul para todo el texto de cuerpo
BODY_FONT    = "Arial"

def _solid_border(shape, color=BORDER_COLOR, width_emu=9525):
    line = shape.line
    line.color.rgb = color
    line.width = Emu(width_emu)
    _no_shadow(shape)

def _no_shadow(shape):
    """Elimina el efecto de sombra que el tema aplica por defecto a los autoshapes."""
    spPr = shape._element.spPr
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))

def title_font_size(text):
    n = len(text)
    if n <= 55: return Pt(24)
    if n <= 70: return Pt(20)
    return Pt(17)

def add_title_bar(slide, text):
    from pptx.util import Cm
    box = slide.shapes.add_shape(1, **TITLE_POS)  # 1 = MSO_SHAPE.RECTANGLE
    box.fill.solid()
    box.fill.fore_color.rgb = TITLE_FILL
    _solid_border(box)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = title_font_size(text)
    run.font.bold = True
    run.font.name = "Poppins"
    run.font.color.rgb = TITLE_TEXT
    return box

def add_picture_scaled(slide, img_path):
    # Marco de ancho COMPLETO (igual a título y bullets) — garantiza que todos
    # los bordes queden alineados sin importar el aspect ratio de cada imagen
    frame = slide.shapes.add_shape(1, **PIC_BOX)
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    _solid_border(frame)

    from PIL import Image as PILImage
    with PILImage.open(img_path) as im:
        ow, oh = im.size
    aspect = oh / ow
    max_w, max_h = PIC_BOX["width"], PIC_BOX["height"]
    w = max_w
    h = int(w * aspect)
    if h > max_h:
        h = max_h
        w = int(h / aspect)
    left = PIC_BOX["left"] + (max_w - w) // 2
    top  = PIC_BOX["top"] + (max_h - h) // 2
    pic = slide.shapes.add_picture(img_path, left, top, w, h)
    return pic

HALLAZGO_RED = RGBColor(0xC6, 0x28, 0x28)
HALLAZGO_BLUE = RGBColor(0x00, 0x50, 0x88)

def add_bullets_boxed(slide, bullets, hallazgo_texto=None):
    box_h = BULLET_POS["height"]
    if hallazgo_texto:
        box_h = Emu(int(BULLET_POS["height"]) + 320000)
    box_pos = dict(BULLET_POS)
    box_pos["height"] = box_h
    box = slide.shapes.add_shape(1, **box_pos)
    box.fill.background()  # noFill
    _solid_border(box)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(180000)
    tf.margin_right = Emu(120000)
    tf.margin_top = Emu(80000)

    state = {"count": 0}
    def _add_numbered_paragraph(text, color, bold=False):
        p = tf.paragraphs[0] if state["count"] == 0 else tf.add_paragraph()
        state["count"] += 1
        p.alignment = PP_ALIGN.JUSTIFY
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", "228600"); pPr.set("indent", "-228600")
        buClr = etree.SubElement(pPr, qn('a:buClr'))
        buClrVal = etree.SubElement(buClr, qn('a:srgbClr'))
        buClrVal.set('val', str(color))
        buFont = etree.SubElement(pPr, qn('a:buFont')); buFont.set('typeface', '+mj-lt')
        buAutoNum = etree.SubElement(pPr, qn('a:buAutoNum')); buAutoNum.set('type', 'arabicPeriod')
        run = p.add_run()
        run.text = text
        run.font.size = Pt(12)
        run.font.name = BODY_FONT
        run.font.bold = bold
        run.font.color.rgb = color
        p.space_after = Pt(6)

    for b in bullets:
        _add_numbered_paragraph(b, BODY_TEXT)
    if hallazgo_texto:
        _add_numbered_paragraph(f"HALLAZGO CLAVE: {hallazgo_texto}", HALLAZGO_RED, bold=True)

def add_hallazgo_to_existing(slide, hallazgo_texto):
    """Agrega una barra de hallazgo destacada (fondo rosa, borde rojo) en la
    parte inferior de una slide ORIGINAL, comprimiendo el texto existente si
    invade esa zona, para evitar que el hallazgo se desborde de cualquier caja."""
    bar_h = Emu(620000)
    bar_top = SLIDE_H - bar_h - Emu(40000)

    # Comprimir cuadros de texto existentes que invadan la zona reservada
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            top = shape.top or 0
            bottom = top + (shape.height or 0)
            if bottom > bar_top:
                new_height = max(Emu(400000), bar_top - top - Emu(30000))
                shape.height = new_height
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.size and r.font.size.pt > 10:
                            r.font.size = Pt(10)

    bar = slide.shapes.add_shape(1, Emu(786581), bar_top, Emu(10599174), bar_h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xFD, 0xE6, 0xE6)
    bar.line.color.rgb = HALLAZGO_RED
    bar.line.width = Emu(19050)
    _no_shadow(bar)
    tf = bar.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Emu(150000)
    tf.margin_right = Emu(150000)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f"HALLAZGO CLAVE: {hallazgo_texto}"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.name = BODY_FONT
    run.font.color.rgb = HALLAZGO_RED

    add_hallazgo_marker(slide)

def add_hallazgo_marker(slide):
    """Asterisco azul en la esquina inferior izquierda para las 5 slides
    con los mayores hallazgos del estudio."""
    box = slide.shapes.add_textbox(Emu(120000), SLIDE_H - Emu(520000), Emu(500000), Emu(420000))
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "★"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = HALLAZGO_BLUE

def restyle_existing_slide(slide, title_text, body_bullets, hallazgo=None):
    """Convierte una slide ORIGINAL (con imagen ya embebida) al formato estándar
    nuevo: extrae la imagen existente, la reescala dentro del marco estándar y
    reemplaza los textos por título-barra + bullets numerados, sin perder la imagen."""
    pic_shape = None
    for shape in list(slide.shapes):
        if shape.shape_type == 13:  # PICTURE
            pic_shape = shape
            break

    # Eliminar todos los textboxes existentes (títulos y cuerpos viejos)
    for shape in list(slide.shapes):
        if shape.has_text_frame:
            sp = shape._element
            sp.getparent().remove(sp)

    add_title_bar(slide, title_text)

    if pic_shape is not None:
        frame = slide.shapes.add_shape(1, **PIC_BOX)
        frame.fill.solid()
        frame.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        _solid_border(frame)

        aspect = pic_shape.height / pic_shape.width
        max_w, max_h = PIC_BOX["width"], PIC_BOX["height"]
        w = max_w
        h = int(w * aspect)
        if h > max_h:
            h = max_h
            w = int(h / aspect)
        left = PIC_BOX["left"] + (max_w - w) // 2
        top  = PIC_BOX["top"] + (max_h - h) // 2
        pic_shape.left, pic_shape.top = left, top
        pic_shape.width, pic_shape.height = w, h
        # Mover la imagen al final del árbol de shapes para que quede ENCIMA del marco
        sp = pic_shape._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.append(sp)

    add_bullets_boxed(slide, body_bullets, hallazgo_texto=hallazgo)
    if hallazgo:
        add_hallazgo_marker(slide)
    return slide

def add_new_marker(slide):
    box = slide.shapes.add_textbox(**MARKER_POS)
    tf = box.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "NUEVO"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xC6, 0x28, 0x28)

def new_content_slide(titulo, img_filename, bullets, marcar=True, hallazgo=None):
    slide = prs.slides.add_slide(LY_BLANK)
    add_title_bar(slide, titulo)
    if img_filename:
        img_path = os.path.join(IMGDIR, img_filename)
        if os.path.exists(img_path):
            add_picture_scaled(slide, img_path)
        add_bullets_boxed(slide, bullets, hallazgo_texto=hallazgo)
        if hallazgo:
            add_hallazgo_marker(slide)
    else:
        text_top = Emu(1100000)
        text_h   = SLIDE_H - text_top - Emu(700000)
        box = slide.shapes.add_textbox(Emu(786581), text_top, Emu(10599174), text_h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = f"• {b}"
            run.font.size = Pt(18)
            run.font.name = BODY_FONT
            run.font.color.rgb = BODY_TEXT
            p.space_after = Pt(20)
    if marcar:
        add_new_marker(slide)
    return slide

def new_separator(titulo, subtitulo):
    slide = prs.slides.add_slide(LY_SEPARADOR)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = titulo
        elif ph.placeholder_format.idx == 1:
            ph.text = subtitulo
    return slide

def mini_separator(titulo, items=None):
    """Mini-separador de subsección (2.1-2.4). Si se pasa `items` (lista de
    strings "Nombre del gráfico (subgrupo, n=X)"), los lista debajo del título
    a modo de mapa de la subsección, en 1 o 2 columnas según la cantidad."""
    slide = prs.slides.add_slide(LY_BLANK)
    title_top = Emu(1500000) if items else Emu(2800000)
    box = slide.shapes.add_textbox(Emu(700000), title_top, Emu(10792000), Emu(900000))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = titulo
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

    if items:
        line = slide.shapes.add_shape(1, Emu(4296000), Emu(2380000), Emu(3600000), Emu(25000))
        line.fill.solid(); line.fill.fore_color.rgb = TITLE_FILL
        line.line.fill.background()
        _no_shadow(line)

        MARGIN_L = Emu(786581)
        CONTENT_W = Emu(10599174)
        if len(items) > 6:
            mid = (len(items) + 1) // 2
            gap = 400000
            col_w = (int(CONTENT_W) - gap) // 2
            cols = [(MARGIN_L, Emu(col_w), items[:mid]),
                    (Emu(int(MARGIN_L) + col_w + gap), Emu(col_w), items[mid:])]
        else:
            cols = [(MARGIN_L, CONTENT_W, items)]

        for left, width, col_items in cols:
            box2 = slide.shapes.add_textbox(left, Emu(2750000), width, Emu(3600000))
            tf2 = box2.text_frame
            tf2.word_wrap = True
            for i, it in enumerate(col_items):
                p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                p2.alignment = PP_ALIGN.LEFT
                pPr = p2._p.get_or_add_pPr()
                pPr.set("marL", "228600"); pPr.set("indent", "-228600")
                buClr = etree.SubElement(pPr, qn('a:buClr'))
                etree.SubElement(buClr, qn('a:srgbClr')).set('val', str(BODY_TEXT))
                buFont = etree.SubElement(pPr, qn('a:buFont')); buFont.set('typeface', 'Arial')
                buChar = etree.SubElement(pPr, qn('a:buChar')); buChar.set('char', '▸')
                run2 = p2.add_run()
                run2.text = it
                run2.font.size = Pt(13)
                run2.font.name = BODY_FONT
                run2.font.color.rgb = BODY_TEXT
                p2.space_after = Pt(7)
    return slide

def cascada_slide(titulo_bloque, subtitulo, cascada_lineas):
    """Separador de bloque con título + cascada de poblaciones derivadas de 917
    indicando qué gráficos usan cada subpoblación."""
    slide = prs.slides.add_slide(LY_BLANK)
    add_title_bar(slide, titulo_bloque)

    # Subtítulo corto
    sub_box = slide.shapes.add_textbox(Emu(786581), Emu(950000), Emu(10599174), Emu(450000))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = subtitulo
    run.font.size = Pt(15)
    run.font.italic = True
    run.font.color.rgb = BODY_TEXT

    # Cascada en caja con borde
    casc_box = slide.shapes.add_shape(1, Emu(786581), Emu(1500000), Emu(10599174), Emu(5000000))
    casc_box.fill.background()
    _solid_border(casc_box)
    tf = casc_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(250000)
    tf.margin_top = Emu(180000)
    for i, linea in enumerate(cascada_lineas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = linea
        run.font.name = "Consolas"
        run.font.size = Pt(13)
        run.font.color.rgb = BODY_TEXT
        run.font.bold = linea.strip().startswith("917") or linea.strip().startswith("957")
        p.space_after = Pt(2)
    return slide

def _tree_box(slide, cx, top, w, h, lines):
    """Caja del diagrama de ramas. `lines` = lista de (texto, bold, size, color, italic)."""
    left = int(cx - w / 2)
    box = slide.shapes.add_shape(1, Emu(left), Emu(int(top)), Emu(int(w)), Emu(int(h)))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    _solid_border(box)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Emu(50000); tf.margin_right = Emu(50000)
    tf.margin_top = Emu(30000); tf.margin_bottom = Emu(30000)
    for i, (text, bold, size, color, italic) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = BODY_FONT
        run.font.color.rgb = color
        p.space_after = Pt(2)
    return dict(cx=cx, cy=top + h / 2, left=left, right=left + w, top=top, bottom=top + h)

def _tree_connector(slide, geo_from, geo_to):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Emu(int(geo_from["right"])), Emu(int(geo_from["cy"])),
        Emu(int(geo_to["left"])), Emu(int(geo_to["cy"])))
    conn.line.color.rgb = TITLE_FILL
    conn.line.width = Emu(22225)
    return conn

def tree_diagram_bloque2():
    """Diagrama de ramas (cajas + conectores) de Bloque II: universos
    desglosados con la cantidad de gráficas que componen cada subsección."""
    slide = prs.slides.add_slide(LY_BLANK)
    add_title_bar(slide, "Mapa de Bloque II: Universos y Gráficas por Subsección")

    # Col. 1 (hijos directos de 917, orden decreciente: 560 > 357 > 197... > 816
    # NO es hijo de 197, es población propia con SAT disponible cruzando
    # formados+control, por eso cuelga directo de la raíz junto a los otros dos)
    root = _tree_box(slide, 1700000, 3450000, 1700000, 700000,
        [("917", True, 16, BODY_TEXT, False),
         ("Jerarquizados", True, 12, BODY_TEXT, False)])

    control = _tree_box(slide, 4700000, 2425000, 2700000, 750000,
        [("560 Control (61%)", True, 13, BODY_TEXT, False),
         ("no participó en formación", False, 10, RGBColor(0x66,0x66,0x66), True)])

    formados = _tree_box(slide, 4700000, 3425000, 2700000, 750000,
        [("357 Formados (39%)", True, 13, BODY_TEXT, False),
         ("2.1 · 9 gráficas", False, 10, SUB_BLUE, True)])

    sat816 = _tree_box(slide, 4700000, 4425000, 2700000, 750000,
        [("816 con SAT disponible", True, 13, BODY_TEXT, False),
         ("2.3 · 3 gráficas", False, 10, SUB_BLUE, True)])

    aptos = _tree_box(slide, 7400000, 3450000, 2400000, 700000,
        [("197 Aptos P3", True, 13, BODY_TEXT, False),
         ("SAT baseline + resultado", False, 10, RGBColor(0x66,0x66,0x66), True)])

    b22 = _tree_box(slide, 10100000, 2700000, 2500000, 900000,
        [("184 Puros / 13 Mixtos", True, 12, BODY_TEXT, False),
         ("2.2 · 4 gráficas", False, 10, SUB_BLUE, True)])

    b24 = _tree_box(slide, 10100000, 4000000, 2500000, 900000,
        [("129 con antigüedad", True, 12, BODY_TEXT, False),
         ("2.4 · 3 gráficas", False, 10, SUB_BLUE, True)])

    _tree_connector(slide, root, control)
    _tree_connector(slide, root, formados)
    _tree_connector(slide, root, sat816)
    _tree_connector(slide, formados, aptos)
    _tree_connector(slide, aptos, b22)
    _tree_connector(slide, aptos, b24)
    return slide

# ══════════════════════════════════════════════════════════════════════════════
# 1. Referencias a slides ORIGINALES
# ══════════════════════════════════════════════════════════════════════════════
orig_slides = list(prs.slides)
S = {i+1: orig_slides[i] for i in range(len(orig_slides))}

# ── Índice (slide 6) ──────────────────────────────────────────────────────────
slide6 = S[6]
for shape in list(slide6.shapes):
    is_title = shape.has_text_frame and "Índice" in shape.text_frame.text
    if not is_title:
        sp = shape._element
        sp.getparent().remove(sp)
box = slide6.shapes.add_textbox(Emu(786581), Emu(1500000), Emu(10599174), Emu(4900000))
tf = box.text_frame
tf.word_wrap = True

SUB_BLUE = RGBColor(0x5B, 0x8A, 0xB5)  # azul más claro para distinguir sub-items

# (texto, es_subitem)
items = [
    ("Clasificación del Cuerpo Académico  ·  Bloque I  ·  Universo 917", False),
    ("Evaluación Docente Antes y Después  ·  Bloque II  ·  917 → Formados → Aptos P3", False),
    ("Participación en Formación", True),
    ("Población Formada: Pura vs Mixta", True),
    ("SAT y Recomendación: Formados vs Control", True),
    ("Comparación por Tipo, Jerarquía y Antigüedad", True),
    ("Resultados de Calificaciones (Notas según SAT)  ·  Bloque III", False),
    ("Evaluación de Desempeño Docente (EDD)  ·  Bloque IV", False),
    ("Jornada vs Honorario  ·  Bloque V  ·  Universo NOMINA 957", False),
    ("Conclusiones y Recomendaciones", False),
]
sub_n = 0
for i, (texto, es_sub) in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    if es_sub:
        sub_n += 1
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", "640000"); pPr.set("indent", "0")
        run = p.add_run()
        run.text = f"2.{sub_n}   {texto}"
        run.font.size = Pt(14)
        run.font.italic = True
        run.font.bold = False
        run.font.name = BODY_FONT
        run.font.color.rgb = SUB_BLUE
        p.space_after = Pt(6)
    else:
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", "228600"); pPr.set("indent", "-228600")
        buClr = etree.SubElement(pPr, qn('a:buClr'))
        etree.SubElement(buClr, qn('a:srgbClr')).set('val', str(BODY_TEXT))
        buFont = etree.SubElement(pPr, qn('a:buFont')); buFont.set('typeface', 'Arial')
        buChar = etree.SubElement(pPr, qn('a:buChar')); buChar.set('char', '■')
        run = p.add_run()
        run.text = texto
        run.font.size = Pt(17)
        run.font.bold = True
        run.font.name = BODY_FONT
        run.font.color.rgb = BODY_TEXT
        p.space_after = Pt(8)
        p.space_before = Pt(10) if i > 0 else Pt(0)

# ── Recortar texto largo de slide 3 (Metodología) y centrar verticalmente
#    los 4 cuadros (incluido el que no se recorta: "Bases de datos") ─────────
slide3 = S[3]
for shape in list(slide3.shapes):
    if shape.has_text_frame and shape.text_frame.text.strip():
        tf = shape.text_frame
        full_text = tf.text
        if len(full_text) > 200:
            if "Universo de Análisis" in full_text or "917 docentes jerarquizados" in full_text:
                tf.text = ("917 docentes jerarquizados · 493 con perfil completo + 424 solo nómina\n"
                           "Z-score = posición relativa del docente en su facultad y período")
            elif "Marco Metodológico" in full_text:
                tf.text = "z = (SAT docente − promedio facultad·período) / desviación estándar"
            elif "dimensiones de satisfacción" in full_text:
                tf.text = "SAT Nota (escala 1–7) y SAT % Recomendación · correlación r=0.893"
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(13)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# ── Completar subtítulo del separador "HONORARIO VS JORNADA" (Bloque V) ──────
slide22 = S[22]
for ph in slide22.placeholders:
    if ph.placeholder_format.idx == 1:
        ph.text = "Bloque V — Universo NOMINA 957\n(587 Jornada + 370 Honorario)"

# ── Armonizar formato de slides originales con título en texto plano ────────
# Nota: S[4] y S[5] (Venn 357 / Venn 197) se EXCLUYEN del deck final por ser
# duplicados exactos de N21["venn917"] y N22["venn_puro"] — ver auditoría.
restyle_existing_slide(S[18], "Evolución 2022-2025: Participación en Formación",
    ["La participación creció de 39 docentes (2022) a 211 (2024), un aumento de 172 en dos años.",
     "2025 muestra 86 docentes únicos (cifra parcial, año en curso al momento del análisis)."])

restyle_existing_slide(S[19], "Cobertura Global de Participación",
    ["357 de 917 docentes participaron en formación (38.9% de cobertura); 61.1% no participó.",
     "Por modalidad (con doble conteo): Taller 215 (60%), Diplomado 201 (56%), Proyecto 32 (9%)."])

restyle_existing_slide(S[21], "Participación según Jerarquía Académica",
    ["Los docentes con mayor antigüedad (15+ años) muestran menor participación en formación.",
     "El grupo sin formación cae en su EDD (0.604); los docentes antiguos no compensan capacitándose."])

# ── 5 Mayores Hallazgos del Estudio: marcador azul + bullet rojo ─────────────
add_hallazgo_to_existing(S[7],
    "Los docentes formados superan al grupo control en SAT los 6 períodos sin "
    "excepción — el beneficio no es de un solo semestre, sino sostenido en el tiempo.")

add_hallazgo_to_existing(S[10],
    "La formación actúa como escudo protector frente a la evaluación de jefes: "
    "el efecto más grande y contundente del estudio (Cohen d=+0.955 en 2025).")

add_hallazgo_to_existing(S[26],
    "No existe diferencia de calidad docente entre Jornada y Honorario — el "
    "tipo de contrato no determina el desempeño evaluado por los estudiantes.")

print("Slides originales preparadas (incluye 5 armonizadas + 3 con hallazgo clave).")

# ══════════════════════════════════════════════════════════════════════════════
# 2. Slides nuevas (sin AFO/MET/dimensiones/tipologías)
# ══════════════════════════════════════════════════════════════════════════════

N1 = {}
N1["jerarquia"] = new_content_slide(
    "Distribución de Jerarquía Académica", "G_jerarquia_918.png",
    ["Instructor Docente es la más numerosa (37.3%), seguida de Asistente Docente (28.7%).",
     "Universo 917 → 915 con jerarquía informada (2 sin dato vía dotación)."])

N1["jer_formacion"] = new_content_slide(
    "Jerarquía × Nivel de Formación", "G_jerarquia_formacion_918.png",
    ["A mayor jerarquía, mayor proporción de Doctores.",
     "Universo 917 → 493 perfil completo → 477 con ambos datos informados."])

N1["gradorec"] = new_content_slide(
    "Grado Académico Reconocido (GRADOREC)", "G_gradorec_918.png",
    ["46.8% Magíster Profesional, 28.5% Doctor: perfil aplicado más que investigativo.",
     "Universo 917 → 493 perfil completo → 477 con grado clasificado."])

N1["institucion"] = new_content_slide(
    "Institución de Obtención del Grado", "G_institucion_grado_918.png",
    ["La propia Universidad Central forma al 15.3% de su planta jerarquizada.",
     "Universo 917 → 493 perfil completo → 478 con institución informada."])

N1["pais"] = new_content_slide(
    "País de Obtención del Grado", "G_pais_grado_918.png",
    ["70.7% de los grados se obtuvo en Chile; España lidera el destino internacional (14.0%).",
     "Universo 917 → 493 perfil completo → 478 con país informado."])

N1["carga"] = new_content_slide(
    "Distribución de la Carga Académica", "G_carga_academica_918.png",
    ["Docencia concentra el 49.5% de la carga declarada en dotación.",
     "Universo 917 → 493 con cargo asignado en dotación."])

N1["anios_jerarquia"] = new_content_slide(
    "Años desde Ingreso hasta Jerarquización", "G_anios_hasta_jerarquia_918.png",
    ["Instructor Docente se jerarquiza casi de inmediato (mediana 0.7 años).",
     "Universo 917 → 433 con fechas de ingreso y jerarquización disponibles."])

N1["edad_jerarquizarse"] = new_content_slide(
    "Edad al Momento de Jerarquizarse", "G_edad_al_jerarquizarse_918.png",
    ["La edad mediana de jerarquización sube de 37 (Instructor) a 62 años (Titular).",
     "Universo 917 → 433 con fechas de ingreso y jerarquización disponibles."])

print(f"Bloque I: {len(N1)} slides nuevas.")

N21 = {}
N21["facultad"] = new_content_slide(
    "Participación en Formación por Facultad", "p2_g22_facultad.png",
    ["Brechas de hasta 30pp entre la facultad más y menos activa.",
     "Universo 917 → 493 perfil completo (única base con dato de facultad)."])

N21["antiguedad"] = new_content_slide(
    "Participación en Formación por Antigüedad", "p2_g23_antiguedad.png",
    ["El tramo 0-4 años concentra la mayor tasa de participación (50%).",
     "Universo 917 → 493 perfil completo → 491 con antigüedad informada."])

N21["modalidad"] = new_content_slide(
    "Distribución por Modalidad de Formación", "p2_g31_modalidad.png",
    ["Taller concentra el 61% de las 615 iniciativas de formación totales.",
     "Universo 917 → 357 docentes formados → 615 actividades de formación."])

N21["brechas"] = new_content_slide(
    "Brechas: Perfil de No Participantes", "p2_g41_brechas.png",
    ["560 docentes sin ninguna iniciativa de formación (61.1% del universo).",
     "Universo 917 jerarquizados completo."],
    hallazgo="6 de cada 10 docentes jerarquizados nunca ha participado en ninguna "
             "iniciativa de perfeccionamiento — la principal oportunidad de mejora del sistema.")

N21["intensidad_global"] = new_content_slide(
    "Intensidad de Participación Global", "p2_g42_intensidad.png",
    ["61% de los formados participó en una sola iniciativa de formación.",
     "Universo 917 → 357 docentes con al menos 1 iniciativa de formación."])

N21["venn917"] = new_content_slide(
    "Combinaciones de Formación — Universo 917", "G_venn_formacion_917_918.png",
    ["357 docentes con formación: 76% hizo una sola modalidad.",
     "Universo 917 → 357 docentes formados."])

N21["int_talleres"] = new_content_slide(
    "Intensidad de Participación en Talleres", "G_intensidad_talleres_918.png",
    ["59% de los docentes tomó un solo taller, experiencia puntual.",
     "Universo 917 → 215 docentes con al menos 1 taller."])

N21["int_proyectos"] = new_content_slide(
    "Intensidad de Participación en Proyectos", "G_intensidad_proyectos_918.png",
    ["88% participó en un solo proyecto; máximo 4 (Juan Carlos Araya Vargas).",
     "Universo 917 → 32 docentes con al menos 1 proyecto."])

print(f"Bloque II.1: {len(N21)} slides nuevas.")

N22 = {}
N22["venn_puro"] = new_content_slide(
    "Combinaciones de Modalidades — Población P3", "G_venn_poblaciones_918.png",
    ["93.4% de los 197 aptos P3 hizo una sola modalidad (población pura).",
     "Universo 917 → 197 aptos P3."])

N22["contraste1"] = new_content_slide(
    "Trayectoria SAT — Normal vs Población Pura", "G1_contraste_normal_puro_918.png",
    ["El Diplomado puro muestra mayor caída (Δ=-0.070) que con mixtos (Δ=-0.028).",
     "Universo 917 → 197 aptos P3 → 184 puros / 13 mixtos."])

N22["puro_nota_recom"] = new_content_slide(
    "SAT Nota y Recomendación — Población Pura", "G1_puro_nota_recom_918.png",
    ["Mismos patrones en ambas métricas: Proyecto cae más, Taller es estable.",
     "Universo 917 → 197 aptos P3 → 184 puros."])

N22["delta_puro_mixto"] = new_content_slide(
    "Δ Z-score: Puro vs Mixto", "G_delta_z_puro_vs_mixto_918.png",
    ["Los 13 mixtos muestran mejor resultado (62% mejoró) que los puros (49%).",
     "Universo 917 → 197 aptos P3 → 184 puros / 13 mixtos."])

print(f"Bloque II.2: {len(N22)} slides nuevas.")

N23 = {}
N23["t1"] = new_content_slide(
    "Validación Estadística — SAT Nota", "G11_pruebas_t_918.png",
    ["Diferencia significativa en 3 de 6 períodos (p<0.05).",
     "Universo 917 → 816 con SAT disponible; n fluctúa por período."])

N23["t2"] = new_content_slide(
    "Validación Estadística — % Recomendación", "G11.2_pruebas_t_bin_918.png",
    ["Significativo en 4 de 6 períodos.",
     "Universo 917 → 816 con SAT disponible; n fluctúa por período."])

print(f"Bloque II.3: {len(N23)} slides nuevas.")

N24 = {}
N24["antig_par"] = new_content_slide(
    "Antigüedad × Tipo — SAT Nota y Recomendación", "G3_par_antiguedad_918.png",
    ["129 de 197 aptos P3 tienen dato de antigüedad disponible.",
     "Universo 917 → 197 aptos P3 → 129 con antigüedad informada."])

N24["ev_taller"] = new_content_slide(
    "Evolución por Jerarquía — Talleres", "G7_taller_jerarquia_918.png",
    ["Asistentes Docentes parten bajo la media y alcanzan su peak durante talleres.",
     "Universo 917 → 197 aptos P3 → subgrupo Taller."])

N24["ev_diplomado"] = new_content_slide(
    "Evolución por Jerarquía — Diplomados", "G7_diplomado_jerarquia_918.png",
    ["Instructores Docentes: caso de mayor éxito, parten bajo y suben.",
     "Universo 917 → 197 aptos P3 → subgrupo Diplomado."])

N24["t4"] = new_content_slide(
    "T-test Jerarquía × Diplomado", "G13_ttest_jerarquia_diplomado_918.png",
    ["ANOVA no significativo (p=0.718): el Diplomado beneficia por igual a todos los rangos.",
     "Universo 917 → 197 aptos P3 → subgrupo Diplomado."])

N24["t5"] = new_content_slide(
    "T-test Jerarquía × Taller", "G13.2_ttest_jerarquia_taller_918.png",
    ["Ningún par de jerarquías muestra diferencia significativa tras Bonferroni.",
     "Universo 917 → 197 aptos P3 → subgrupo Taller."])

print(f"Bloque II.4: {len(N24)} slides nuevas.")

N3 = {}
N3["aprob_antig"] = new_content_slide(
    "Tasa de Aprobación × Antigüedad Docente", "G_aprobacion_antiguedad_918.png",
    ["La tasa de aprobación es relativamente estable entre tramos de antigüedad.",
     "Universo 917 → 834 docentes con notas y SAT disponibles."])

N3["aprob_jer"] = new_content_slide(
    "Tasa de Aprobación × Jerarquía Académica", "G_aprobacion_jerarquia_918.png",
    ["Instructor Docente lidera con la mediana más alta (94.6%).",
     "Universo 917 → 834 docentes con notas y SAT disponibles."])

N3["acumulativo"] = new_content_slide(
    "¿Es Acumulativo el Efecto de la Formación?", "G_acumulacion_formacion_918.png",
    ["A partir de 3 iniciativas de formación la mediana de aprobación supera al grupo sin formación.",
     "Universo 917 → 816 con SAT y notas, por cantidad de iniciativas de formación."],
    hallazgo="El impacto de la formación se vuelve evidente recién a partir de 3 "
             "iniciativas de formación — la participación única o doble no genera mejora suficiente.")

print(f"Bloque III: {len(N3)} slides nuevas.")

N4 = {}
N4["t3"] = new_content_slide(
    "Validación Estadística — EDD", "G12_pruebas_t_edd_918.png",
    ["Significativo en 2024 (p=0.002) y 2025 (p<0.001, d=+0.955).",
     "Universo 917 → 486 con EDD disponible (134 formados + 352 sin formación)."])

N4["edd_contraste"] = new_content_slide(
    "EDD — Población Normal vs Pura", "G9_contraste_edd_918.png",
    ["Diplomado puro sube de 0.823 a 0.854 al excluir los mixtos.",
     "Universo 917 → 197 aptos P3 → con EDD disponible."])

print(f"Bloque IV: {len(N4)} slides nuevas.")

N5 = {}
N5["venn_contrato"] = new_content_slide(
    "Combinaciones de Formación — Jornada vs Honorario", "G_venn_formacion_contrato_918.png",
    ["Ninguna intersección de Honorario incluye Proyecto.",
     "Universo NOMINA 957 → 587 Jornada + 370 Honorario."])

print(f"Bloque V: {len(N5)} slides nuevas.")

# ── Slide de embudo simple (reemplaza los 2 Venn duplicados S[4]/S[5]) ───────
EMBUDO = new_content_slide(
    "Universo de Análisis: El Embudo P3", None,
    ["917 docentes jerarquizados → 357 participaron en formación (39%, grupo tratamiento) → 560 grupo control (61%).",
     "De los 357 formados, 197 son aptos P3 (21% del universo): tienen SAT válido en baseline y resultado para medir impacto.",
     "El detalle de combinaciones de modalidad (puro/mixto) para ambos grupos se presenta en las secciones 2.1 y 2.2."],
    marcar=False)

print("Slide de embudo creada (reemplaza Venns duplicados).")

# ══════════════════════════════════════════════════════════════════════════════
# 3. Separadores con CASCADA de poblaciones
# ══════════════════════════════════════════════════════════════════════════════

SEP_I = cascada_slide(
    "BLOQUE I — Clasificación del Cuerpo Académico",
    "Universo: 917 Docentes Jerarquizados",
    [
        "917 docentes jerarquizados",
        "  │",
        "  ├── 915 con jerarquía informada ............... Distribución de Jerarquía",
        "  │",
        "  └── 493 con perfil completo (dotación)",
        "        ├── (493) ......................... Edad y Sexo · Facultad · Carga Académica",
        "        ├── 491 con antigüedad ............. Antigüedad",
        "        ├── 482 con edad y jerarquía ....... Edad × Jerarquía",
        "        ├── 478 con nivel de formación ..... Nivel de Formación · Institución · País",
        "        ├── 477 con grado clasificado ...... GRADOREC · Jerarquía × Nivel Formación",
        "        └── 433 con fechas completas ....... Años hasta Jerarquización",
    ])

SEP_II = cascada_slide(
    "BLOQUE II — Evaluación Docente Antes y Después",
    "Universo: 917 → Formados → Aptos P3",
    [
        "917 docentes jerarquizados",
        "  │",
        "  ├── 357 participaron en formación (2.1) ....... Cobertura · Evolución · Venn · Intensidad · Brechas",
        "  │     └── 615 iniciativas de formación totales",
        "  │",
        "  ├── 816 con SAT disponible (2.3) ............... Validación SAT y Recomendación",
        "  │",
        "  └── 197 aptos P3 (SAT baseline + resultado válido)",
        "        ├── 184 puros / 13 mixtos (2.2) ......... Venn Población · Contraste Puro vs Mixto",
        "        └── 129 con antigüedad (2.4) ............ Antigüedad × Tipo · Evolución por Jerarquía",
    ])

MAPA_BLOQUE2 = tree_diagram_bloque2()

SEP_III = cascada_slide(
    "BLOQUE III — Resultados de Calificaciones",
    "Notas y Tasas de Aprobación según SAT",
    [
        "917 docentes jerarquizados",
        "  │",
        "  └── 834 con notas y SAT disponibles",
        "        ├── (834) ...................... Aprobación × Antigüedad · Aprobación × Jerarquía",
        "        └── 816 con SAT y notas ........ ¿Es Acumulativo el Efecto de la Formación?",
    ])

SEP_IV = cascada_slide(
    "BLOQUE IV — Evaluación de Desempeño Docente (EDD)",
    "Universo: 917 → con EDD Disponible",
    [
        "917 docentes jerarquizados",
        "  │",
        "  └── 486 con EDD disponible (134 formados + 352 sin formación)",
        "        ├── (486) ...................... Evolución EDD · Validación Estadística EDD",
        "        └── 197 aptos P3 con EDD ....... EDD: Población Normal vs Pura",
    ])

# Cascada para Bloque V (universo distinto: NOMINA 957)
slide22_casc = cascada_slide(
    "BLOQUE V — Jornada vs Honorario",
    "Universo NOMINA 957 (distinto del Universo 917 Jerarquizados)",
    [
        "957 docentes NOMINA",
        "  │",
        "  ├── 587 Jornada ........... SAT · % Recomendación · EDD · Venn Formación · SAT×Contrato",
        "  └── 370 Honorario ......... SAT · % Recomendación · EDD (solo 9, no comparable) · Venn Formación",
    ])

SUB_21 = mini_separator("2.1   Participación en Formación", [
    "Evolución 2022-2025: Participación en Formación (357 docentes)",
    "Cobertura Global de Participación (357 de 917, 38.9%)",
    "Participación según Jerarquía Académica (917 jerarquizados)",
    "Modalidad de Formación según Jerarquía (357 formados)",
    "Participación en Formación por Antigüedad (491 con antigüedad)",
    "Distribución por Modalidad de Formación (615 iniciativas de formación)",
    "Brechas: Perfil de No Participantes (560 sin formación) ★ hallazgo clave",
    "Intensidad de Participación Global (357 formados)",
    "Combinaciones de Formación — Universo 917 (357 formados)",
    "Intensidad de Participación en Talleres (215 con taller)",
])
SUB_22 = mini_separator("2.2   Población Formada: Pura vs Mixta", [
    "Combinaciones de Modalidades — Población P3 (197 aptos P3)",
    "Trayectoria SAT — Normal vs Población Pura (184 puros / 13 mixtos)",
    "SAT Nota y Recomendación — Población Pura (184 puros)",
    "Δ Z-score: Puro vs Mixto (184 puros / 13 mixtos)",
])
SUB_23 = mini_separator("2.3   SAT y Recomendación — Formados vs Control", [
    "Evolución SAT por Período: Formados vs Control (6 períodos) ★ hallazgo clave",
    "Validación Estadística — SAT Nota (816 con SAT disponible)",
    "Validación Estadística — % Recomendación (816 con SAT disponible)",
])
SUB_24 = mini_separator("2.4   Comparación por Tipo, Jerarquía y Antigüedad", [
    "Antigüedad × Tipo — SAT Nota y Recomendación (129 con antigüedad)",
    "Evolución por Jerarquía — Talleres (subgrupo Taller)",
    "Evolución por Jerarquía — Diplomados (subgrupo Diplomado)",
])

CIERRE1 = new_content_slide(
    "Conclusiones Consolidadas", None,
    ["La formación genera impacto medible en SAT, EDD y aprobación: efecto de selección + protección.",
     "No hay diferencia Jornada/Honorario en calidad docente: el contrato no determina el desempeño.",
     "La cobertura (39%) es la principal oportunidad: 6 de cada 10 nunca participaron.",
     "El impacto se consolida a partir de 3+ iniciativas de formación; la participación única no basta."],
    marcar=False)
CIERRE2 = new_content_slide(
    "Recomendaciones", None,
    ["Institucionalizar Diplomados como requisito de avance en la carrera académica.",
     "Diseñar rutas formativas progresivas en lugar de iniciativas aisladas.",
     "Ampliar cobertura a Honorarios: mismo retorno con menor acceso actual.",
     "Crear mecanismo de evaluación de desempeño (EDD) para Honorarios."],
    marcar=False)

print("Separadores con cascada y cierre creados.")

# ── ★ azul: 20 slides seleccionadas para explicar en vivo en la presentación
#    ejecutiva. Las 5 que ya llevan "HALLAZGO CLAVE" en rojo (S[7], S[10],
#    S[26], N21["brechas"], N3["acumulativo"]) quedan distinguidas del resto,
#    que solo lleva el ★ como guía de "esta sí se presenta de pie". ─────────
PRESENTAR_EN_VIVO = [
    S[1], S[2], S[6], SEP_I, N1["jerarquia"], N1["jer_formacion"],
    SEP_II, EMBUDO, S[19], N22["venn_puro"],
    SEP_III, SEP_IV, slide22_casc, CIERRE1, CIERRE2,
]
for _slide in PRESENTAR_EN_VIVO:
    add_hallazgo_marker(_slide)
print(f"{len(PRESENTAR_EN_VIVO)} slides marcadas con ★ azul "
      f"(+ 5 hallazgo ya marcadas) = 20 slides para presentación en vivo.")

# ══════════════════════════════════════════════════════════════════════════════
# 4. Orden final completo
# ══════════════════════════════════════════════════════════════════════════════
final_order = (
    [S[1], S[2], S[6], S[3]] +
    [SEP_I] +
    [S[13], S[14], S[15], S[16]] +
    [N1["jerarquia"], N1["jer_formacion"], N1["gradorec"],
     N1["institucion"], N1["pais"], N1["carga"],
     N1["anios_jerarquia"]] +
    [SEP_II] +
    [MAPA_BLOQUE2, EMBUDO] +
    [SUB_21] +
    [S[18], S[19], S[21], S[17]] +
    [N21["antiguedad"], N21["modalidad"], N21["brechas"],
     N21["intensidad_global"], N21["venn917"], N21["int_talleres"]] +
    [SUB_22] +
    [N22["venn_puro"], N22["contraste1"], N22["puro_nota_recom"], N22["delta_puro_mixto"]] +
    [SUB_23] +
    [S[7]] +
    [N23["t1"], N23["t2"]] +
    [SUB_24] +
    [N24["antig_par"], N24["ev_taller"], N24["ev_diplomado"]] +
    [SEP_III] +
    [S[8], S[9]] +
    [N3["aprob_antig"], N3["aprob_jer"], N3["acumulativo"]] +
    [SEP_IV] +
    [S[10]] +
    [N4["t3"], N4["edd_contraste"]] +
    [slide22_casc] +
    [S[23], S[24], S[25], S[26]] +
    [N5["venn_contrato"]] +
    [CIERRE1, CIERRE2]
)

# Excluidos del deck final:
#   S[4], S[5]  → Venn duplicados exactos de N21["venn917"] / N22["venn_puro"]
#   S[11]       → trayectoria ya cubierta por el panel "Normal" de N22["contraste1"]
#   S[12]       → heatmap antigüedad original, superado por N24["antig_par"] (pareado)
#   S[20]       → slide duplicada del borrador original
#   S[22]       → reemplazado por slide22_casc (separador con cascada)

sldIdLst = prs.slides._sldIdLst
current_slide_objs = list(prs.slides)
current_sldid_els  = list(sldIdLst)
slide_to_sldid = {id(s._element): sldid for s, sldid in zip(current_slide_objs, current_sldid_els)}

new_sldid_elements = [slide_to_sldid[id(s._element)] for s in final_order]
assert len(new_sldid_elements) == len(set(id(e) for e in new_sldid_elements)), "Slides repetidas"

for el in list(sldIdLst):
    sldIdLst.remove(el)
for el in new_sldid_elements:
    sldIdLst.append(el)

ids_finales = {id(s._element) for s in final_order}
for s_excl in [S[4], S[5], S[11], S[12], S[20], S[22],
               N1["edad_jerarquizarse"], N21["facultad"], N21["int_proyectos"],
               N24["t4"], N24["t5"]]:
    if id(s_excl._element) in ids_finales:
        continue
    rId = None
    for rel_id, rel in prs.part.rels.items():
        if rel.target_part == s_excl.part:
            rId = rel_id
            break
    if rId:
        prs.part.drop_rel(rId)

print(f"\nTotal slides en el deck final: {len(prs.slides)}")
prs.save(OUT)
print(f"Guardado: {OUT}")

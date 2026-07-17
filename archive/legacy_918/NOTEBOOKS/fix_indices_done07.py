# -*- coding: utf-8 -*-
"""
fix_indices_done07.py
Corrige numeracion de diapositivas en los separadores de bloque y el indice
del archivo Done07-26.pptx.

Cambios:
  Slide 02 - Indice:
    - "34 diapositivas" -> "33 diapositivas"
    - Linea 2.3: quitar "Histograma Az * Cambio x Antiguedad" (slides eliminadas)
  Slide 11 - Sep BLOQUE II:
    - Eliminar Diapo 23 y Diapo 24 (no existen en el deck)
  Slide 23 - Sep BLOQUE III:
    - Reemplazar todos los bullets con numeracion correcta (24-29)
    - Agregar slide 24 scatter SAT vs Nota (nueva)
  Slide 30 - Sep BLOQUE IV:
    - Diapo 37->31, Diapo 38->32
"""
import sys; sys.stdout.reconfigure(encoding="utf-8")
import os, copy
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

SRC  = (r"c:\Users\r.gonzalez_fluxsolar.LAPTOP-FLUX-ECO"
        r"\Downloads\Analisis_UCEN_v2\Done07-26.pptx")
DST  = (r"c:\Users\r.gonzalez_fluxsolar.LAPTOP-FLUX-ECO"
        r"\Downloads\Analisis_UCEN_v2\Done07-26_fixed.pptx")

prs = Presentation(SRC)

# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────
def get_tf_text(sl, shape_name_contains):
    """Devuelve la shape cuyo nombre contiene el texto dado."""
    for sh in sl.shapes:
        if shape_name_contains.lower() in sh.name.lower() and sh.has_text_frame:
            return sh
    return None

def find_shape_by_text_content(sl, text_fragment):
    """Devuelve la shape que contiene text_fragment en su texto total."""
    for sh in sl.shapes:
        if sh.has_text_frame:
            full = sh.text_frame.text
            if text_fragment in full:
                return sh
    return None

def replace_run_text_in_shape(shape, old, new):
    """Reemplaza 'old' por 'new' en todos los runs del shape."""
    if not shape or not shape.has_text_frame:
        return False
    found = False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                found = True
    return found

def set_para_text(para, new_text):
    """Reemplaza el texto del primer run del parrafo, limpia los demas."""
    runs = para.runs
    if runs:
        runs[0].text = new_text
        for r in runs[1:]:
            r.text = ""
    return para

def para_text(para):
    return "".join(r.text for r in para.runs)

def clone_para_with_text(source_para, new_text):
    """Clona un parrafo XML y reemplaza el texto en todos sus runs."""
    new_p = copy.deepcopy(source_para._p)
    # Reemplazar texto en todos los <a:r><a:t>
    for r_el in new_p.findall(qn("a:r")):
        t_el = r_el.find(qn("a:t"))
        if t_el is not None:
            t_el.text = new_text
            new_text = ""   # solo el primer run lleva texto
    return new_p

def remove_para_from_tf(tf, para_index):
    """Elimina el parrafo en la posicion dada del text frame."""
    txBody = tf._txBody
    paras  = txBody.findall(qn("a:p"))
    if 0 <= para_index < len(paras):
        txBody.remove(paras[para_index])
        return True
    return False

def insert_para_after(tf, after_index, new_p_xml_el):
    """Inserta un elemento <a:p> despues del parrafo en after_index."""
    txBody = tf._txBody
    paras  = txBody.findall(qn("a:p"))
    if after_index < len(paras):
        ref = paras[after_index]
        ref.addnext(new_p_xml_el)
    else:
        txBody.append(new_p_xml_el)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 02 — Indice
# ─────────────────────────────────────────────────────────────────────────────
sl02 = prs.slides[1]
print("--- Slide 02 (Indice) ---")

# 1. "34 diapositivas" -> "33 diapositivas"
sh = find_shape_by_text_content(sl02, "34 diapositivas")
if sh:
    ok = replace_run_text_in_shape(sh, "34 diapositivas", "33 diapositivas")
    print("  [OK] '34 diapositivas' -> '33 diapositivas'" if ok else "  [FAIL] no se encontro el run")
else:
    print("  [SKIP] ya dice '33' o no se encontro el texto")

# 2. Limpiar linea 2.3: quitar referencias a slides eliminadas
sh2 = find_shape_by_text_content(sl02, "2.3")
if sh2:
    for para in sh2.text_frame.paragraphs:
        t = para_text(para)
        if "2.3" in t and ("Histograma" in t or "Cambio" in t):
            set_para_text(para, "      2.3  SAT z-score por Facultad  (n=197)")
            print("  [OK] Linea 2.3 actualizada")
            break
    else:
        # Buscar en otra shape
        print("  [INFO] Linea 2.3 con Histograma no encontrada en shape separada")

# Si la linea 2.3 esta en la misma shape que el resto del indice (TextBox 4)
sh_idx = find_shape_by_text_content(sl02, "Clasificaci")
if sh_idx:
    for para in sh_idx.text_frame.paragraphs:
        t = para_text(para)
        if "2.3" in t and ("Histograma" in t or "Cambio" in t):
            set_para_text(para, "      2.3  SAT z-score por Facultad  (n=197)")
            print("  [OK] Linea 2.3 actualizada en shape indice")
            break

print()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Sep BLOQUE II
# Eliminar Diapo 23 y Diapo 24 (slides eliminadas del deck)
# ─────────────────────────────────────────────────────────────────────────────
sl11 = prs.slides[10]
print("--- Slide 11 (Sep BLOQUE II) ---")

sh_bul = find_shape_by_text_content(sl11, "Diapo 12")
if sh_bul:
    tf = sh_bul.text_frame
    # Identificar indices de parrafos a eliminar
    to_remove = []
    for j, para in enumerate(tf.paragraphs):
        t = para_text(para)
        if "Diapo 23" in t or "Diapo 24" in t:
            to_remove.append(j)
            print(f"  [MARK] parrafo {j} para eliminar: {repr(t[:60])}")

    # Eliminar en orden inverso para no desplazar indices
    for j in sorted(to_remove, reverse=True):
        ok = remove_para_from_tf(tf, j)
        print(f"  [OK] Eliminado parrafo {j}" if ok else f"  [FAIL] No se pudo eliminar {j}")

    # Verificar resultado
    print("  Resultado final:")
    for j, para in enumerate(tf.paragraphs):
        t = para_text(para)
        if t.strip():
            print(f"    [{j}] {repr(t[:70])}")
else:
    print("  [FAIL] No se encontro shape con bullets de BLOQUE II")

print()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 23 — Sep BLOQUE III
# Corregir numeracion: 30->24, 31->25, 32->26, 33->27, 34->28 + agregar Diapo 29
# ─────────────────────────────────────────────────────────────────────────────
sl23 = prs.slides[22]
print("--- Slide 23 (Sep BLOQUE III) ---")

sh_bul3 = find_shape_by_text_content(sl23, "Diapo 30")
if sh_bul3:
    tf = sh_bul3.text_frame
    paras = tf.paragraphs

    # Nuevos textos en orden correcto (6 items)
    nuevos = [
        "•   Diapo 24:  Relacion SAT vs Nota Promedio de Alumnos  (scatter por ano)",
        "•   Diapo 25:  Aprobacion Global  (formados vs control)",
        "•   Diapo 26:  Evolucion de Aprobacion por Periodo  (2023-01 a 2025-02)",
        "•   Diapo 27:  Aprobacion por Antiguedad  (solo formados)",
        "•   Diapo 28:  Aprobacion por Jerarquia  (solo formados)",
        "•   Diapo 29:  Efecto Acumulativo  (n instancias de formacion)",
    ]

    # Cuantos parrafos hay actualmente
    n_actual = len([p for p in paras if para_text(p).strip()])
    print(f"  Parrafos actuales con texto: {n_actual}")

    # Reemplazar los existentes
    filled = 0
    for j, para in enumerate(tf.paragraphs):
        t = para_text(para)
        if t.strip() and filled < len(nuevos):
            set_para_text(para, nuevos[filled])
            print(f"  [OK] Para [{j}] -> {repr(nuevos[filled][:60])}")
            filled += 1

    # Si faltan parrafos (actual<6), agregar los que faltan clonando el ultimo
    last_content_para = None
    for p in tf.paragraphs:
        if para_text(p).strip():
            last_content_para = p

    while filled < len(nuevos):
        if last_content_para:
            new_p = clone_para_with_text(last_content_para, nuevos[filled])
            tf._txBody.append(new_p)
            print(f"  [ADD] Nuevo parrafo: {repr(nuevos[filled][:60])}")
        filled += 1

    print("  Resultado:")
    for j, para in enumerate(tf.paragraphs):
        t = para_text(para)
        if t.strip():
            print(f"    [{j}] {repr(t[:70])}")
else:
    print("  [FAIL] No se encontro shape con 'Diapo 30'")

print()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 30 — Sep BLOQUE IV
# Diapo 37->31, Diapo 38->32
# ─────────────────────────────────────────────────────────────────────────────
sl30 = prs.slides[29]
print("--- Slide 30 (Sep BLOQUE IV) ---")

sh_bul4 = find_shape_by_text_content(sl30, "Diapo 37")
if sh_bul4:
    ok1 = replace_run_text_in_shape(sh_bul4, "Diapo 37:", "Diapo 31:")
    ok2 = replace_run_text_in_shape(sh_bul4, "Diapo 38:", "Diapo 32:")
    print(f"  [OK] Diapo 37->31: {ok1}  |  Diapo 38->32: {ok2}")
    for j, para in enumerate(sh_bul4.text_frame.paragraphs):
        t = para_text(para)
        if t.strip():
            print(f"    [{j}] {repr(t[:80])}")
else:
    print("  [FAIL] No se encontro shape con 'Diapo 37'")

print()

# ─────────────────────────────────────────────────────────────────────────────
# Guardar
# ─────────────────────────────────────────────────────────────────────────────
prs.save(DST)
print("Guardado:", DST)
print("Total slides:", len(prs.slides))

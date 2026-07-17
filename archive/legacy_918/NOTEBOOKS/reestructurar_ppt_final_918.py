import sys; sys.stdout.reconfigure(encoding="utf-8")
import os
from pptx import Presentation
from pptx.util import Emu, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

BASE   = os.path.dirname(__file__)
SRC    = os.path.join(BASE, "..", "..", "PRESENTACION_final borrador1.pptx")
OUT    = os.path.join(BASE, "..", "..", "PRESENTACION_final_v2.pptx")
IMGDIR = BASE

prs = Presentation(SRC)
layouts = {l.name: l for l in prs.slide_layouts}
LY_SEPARADOR = layouts["Encabezado de sección"]
LY_BLANK     = layouts["Blank"]

# ── Geometría del slide (13.333" x 7.5", widescreen 16:9) ────────────────────
SLIDE_W, SLIDE_H = Emu(12192000), Emu(6858000)
MARGIN_X = Emu(700000)   # 0.765" margen lateral uniforme

# Franjas verticales con buffers explícitos entre cada elemento (sin solapes)
TITLE_TOP, TITLE_H   = Emu(330000),  Emu(560000)   # título: 0.36" a 0.97"
GAP_1                = Emu(160000)                  # buffer título → imagen
PIC_TOP_MAX          = TITLE_TOP + TITLE_H + GAP_1  # imagen empieza en 1.13"
PIC_H_MAX            = Emu(4050000)                 # alto máx. de imagen: 4.43"
GAP_2                = Emu(170000)                  # buffer imagen → bullets
BULLETS_TOP          = PIC_TOP_MAX + PIC_H_MAX + GAP_2  # bullets ≈ 5.75"
BULLETS_H            = SLIDE_H - BULLETS_TOP - Emu(220000)  # deja margen inferior

TITLE_POS  = dict(left=MARGIN_X, top=TITLE_TOP, width=SLIDE_W - 2*MARGIN_X, height=TITLE_H)
PIC_POS    = dict(left=MARGIN_X, top=PIC_TOP_MAX, width=SLIDE_W - 2*MARGIN_X, height=PIC_H_MAX)
BULLET_POS = dict(left=MARGIN_X, top=BULLETS_TOP, width=SLIDE_W - 2*MARGIN_X, height=BULLETS_H)

# Marcador "NUEVO": esquina superior derecha, fuera del área central del título
MARKER_W, MARKER_H = Emu(1350000), Emu(360000)
MARKER_POS = dict(left=SLIDE_W - MARKER_W - Emu(140000), top=Emu(120000),
                   width=MARKER_W, height=MARKER_H)

TX1_COLOR = RGBColor(0x26, 0x26, 0x26)  # aprox tx1 lumMod90/lumOff10

def title_font_size(text):
    """Reduce el tamaño si el texto es largo, para que siempre quepa en 1 línea."""
    n = len(text)
    if n <= 52: return Pt(20)
    if n <= 65: return Pt(17)
    return Pt(15)

def set_title(slide, text):
    box = slide.shapes.add_textbox(**TITLE_POS)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = None
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = title_font_size(text)
    run.font.bold = True
    run.font.color.rgb = TX1_COLOR
    return box

def add_picture_scaled(slide, img_path):
    from PIL import Image as PILImage
    with PILImage.open(img_path) as im:
        ow, oh = im.size
    aspect = oh / ow
    max_w, max_h = PIC_POS["width"], PIC_POS["height"]
    w = max_w
    h = int(w * aspect)
    if h > max_h:
        h = max_h
        w = int(h / aspect)
    left = PIC_POS["left"] + (max_w - w) // 2
    top  = PIC_POS["top"] + (max_h - h) // 2
    slide.shapes.add_picture(img_path, left, top, w, h)

def add_bullets(slide, bullets):
    box = slide.shapes.add_textbox(**BULLET_POS)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"• {b}"
        run.font.size = Pt(12)
        run.font.color.rgb = TX1_COLOR
        p.space_after = Pt(4)
        p.line_spacing = 1.0

def add_new_marker(slide):
    box = slide.shapes.add_textbox(**MARKER_POS)
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "★ NUEVO"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xC6, 0x28, 0x28)

def new_content_slide(titulo, img_filename, bullets, marcar=True):
    slide = prs.slides.add_slide(LY_BLANK)
    set_title(slide, titulo)
    if img_filename:
        img_path = os.path.join(IMGDIR, img_filename)
        if os.path.exists(img_path):
            add_picture_scaled(slide, img_path)
        add_bullets(slide, bullets)
    else:
        # Slide de solo texto (ej. conclusiones) — centrado verticalmente en el área restante
        text_top = TITLE_TOP + TITLE_H + Emu(700000)
        text_h   = SLIDE_H - text_top - Emu(700000)
        box = slide.shapes.add_textbox(
            MARGIN_X, text_top, SLIDE_W - 2*MARGIN_X, text_h)
        tf = box.text_frame
        tf.word_wrap = True
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = f"• {b}"
            run.font.size = Pt(18)
            run.font.color.rgb = TX1_COLOR
            p.space_after = Pt(20)
    if marcar:
        add_new_marker(slide)
    return slide

def new_separator(titulo, subtitulo):
    slide = prs.slides.add_slide(LY_SEPARADOR)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = titulo
        elif ph.placeholder_format.idx == 1:
            ph.text = subtitulo
    return slide

# ══════════════════════════════════════════════════════════════════════════════
# 1. Capturar referencias a slides ORIGINALES (antes de tocar nada)
# ══════════════════════════════════════════════════════════════════════════════
orig_slides = list(prs.slides)  # índice 0-based = slide original 1..26

S = {i+1: orig_slides[i] for i in range(len(orig_slides))}
# S[1]=Portada, S[2]=Marco, S[3..5]=Universo, S[6]=Indice, S[7..12]=Hallazgos P3,
# S[13..17]=Caracterización, S[18]=Evol part., S[19]=Cobertura, S[20]=DUPLICADO(eliminar),
# S[21]=Participación jerarquía, S[22]=Separador Honorario/Jornada,
# S[23]=Oferta formativa, S[24]=SAT+Recom, S[25]=EDD, S[26]=SAT formado×contrato

# ── Reescribir Índice (slide 6) ───────────────────────────────────────────────
slide6 = S[6]
# Limpiar TODOS los shapes salvo el título (incluye SmartArt/Diagrama sin text_frame)
for shape in list(slide6.shapes):
    is_title = shape.has_text_frame and "Índice" in shape.text_frame.text
    if not is_title:
        sp = shape._element
        sp.getparent().remove(sp)
box = slide6.shapes.add_textbox(MARGIN_X, Emu(1900000), SLIDE_W - 2*MARGIN_X, Emu(3800000))
tf = box.text_frame
tf.word_wrap = True
items = [
    "Bloque 1 — Caracterización del Universo 917 Jerarquizados",
    "Bloque 2 — Impacto de la Formación: 917 → 197 Aptos P3",
    "Bloque 3 — Participación en Formación: 917 → 357 Formados",
    "Bloque 4 — Jornada vs Honorario: Universo NOMINA 957",
    "Conclusiones y Recomendaciones",
]
for i, it in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = f"{i+1}.  {it}"
    run.font.size = Pt(18)
    run.font.color.rgb = TX1_COLOR
    p.space_after = Pt(14)

# ── Recortar texto largo de slide 3 (Metodología) ─────────────────────────────
slide3 = S[3]
for shape in list(slide3.shapes):
    if shape.has_text_frame and len(shape.text_frame.text) > 200:
        # Reemplazar por versión resumida
        tf = shape.text_frame
        full_text = tf.text
        if "Universo de Análisis" in full_text or "917 docentes jerarquizados" in full_text:
            tf.text = ("917 docentes jerarquizados · 493 con perfil completo + 424 solo nómina\n"
                       "Z-score = posición relativa del docente en su facultad y período")
        elif "Marco Metodológico" in full_text:
            tf.text = "z = (SAT docente − promedio facultad·período) / desviación estándar"
        elif "dimensiones de satisfacción" in full_text:
            tf.text = "SAT Nota (escala 1–7) y SAT % Recomendación · correlación r=0.893"
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.size = Pt(13)

# ── Completar subtítulo del separador existente "HONORARIO VS JORNADA" (S[22]) ─
slide22 = S[22]
for ph in slide22.placeholders:
    if ph.placeholder_format.idx == 1:  # Subtítulo
        ph.text = "Universo NOMINA 957\n(587 Jornada + 370 Honorario)"

print("Slides originales preparadas (índice reescrito, texto recortado, separador completado).")

# ══════════════════════════════════════════════════════════════════════════════
# 2. Crear TODAS las slides nuevas (gráficos faltantes)
# ══════════════════════════════════════════════════════════════════════════════

NEW_P1 = {}  # Bloque 1 — Caracterización
NEW_P1["jerarquia"] = new_content_slide(
    "Distribución de Jerarquía Académica", "G_jerarquia_918.png",
    ["Instructor Docente es la más numerosa (37.3%), seguida de Asistente Docente (28.7%)",
     "Los 3 niveles 'Docente' concentran el 81.6% del universo jerarquizado"])

NEW_P1["nivel_formacion"] = new_content_slide(
    "Nivel de Formación del Cuerpo Docente", "G_nivel_formacion_918.png",
    ["87.4% posee Magíster o Doctorado — fuerte orientación al posgrado",
     "30.3% son Doctores, confirmando una planta académica calificada"])

NEW_P1["jer_formacion"] = new_content_slide(
    "Jerarquía × Nivel de Formación", "G_jerarquia_formacion_918.png",
    ["A mayor jerarquía, mayor proporción de Doctores",
     "Instructor Docente concentra mayor proporción de Profesionales"])

NEW_P1["gradorec"] = new_content_slide(
    "Grado Académico Reconocido (GRADOREC)", "G_gradorec_918.png",
    ["46.8% Magíster Profesional, 28.5% Doctor — perfil aplicado más que investigativo",
     "El Magíster Profesional cuadriplica al Magíster Académico"])

NEW_P1["institucion"] = new_content_slide(
    "Institución de Obtención del Grado", "G_institucion_grado_918.png",
    ["La propia Universidad Central forma al 15.3% de su planta jerarquizada",
     "57.5% proviene de 146 instituciones distintas — alta dispersión"])

NEW_P1["pais"] = new_content_slide(
    "País de Obtención del Grado", "G_pais_grado_918.png",
    ["70.7% de los grados se obtuvo en Chile",
     "España es el principal destino internacional (14.0%)"])

NEW_P1["carga"] = new_content_slide(
    "Distribución de la Carga Académica", "G_carga_academica_918.png",
    ["Docencia concentra el 49.5% de la carga declarada en dotación",
     "26.4% ejerce Investigación, Innovación o Vinculación con el Medio"])

NEW_P1["anios_jerarquia"] = new_content_slide(
    "Años desde Ingreso hasta Jerarquización", "G_anios_hasta_jerarquia_918.png",
    ["Instructor Docente se jerarquiza casi de inmediato (mediana 0.7 años)",
     "Titular Docente requiere 5.4 años de trayectoria interna"])

NEW_P1["edad_jerarquizarse"] = new_content_slide(
    "Edad al Momento de Jerarquizarse", "G_edad_al_jerarquizarse_918.png",
    ["La edad mediana de jerarquización sube de 37 (Instructor) a 62 (Titular)",
     "Confirma que los rangos superiores requieren décadas de trayectoria"])

NEW_P1["dimensiones"] = new_content_slide(
    "Dimensiones APR / MET / AFO — Niveles", "G_dimensiones_niveles_918.png",
    ["APR: 79.0% Alto · MET: 72.7% Alto · AFO: solo 37.9% Alto (umbral más exigente)",
     "Aspectos Formales (AFO) es la dimensión con mayor brecha de cumplimiento"])

NEW_P1["tipologias"] = new_content_slide(
    "Tipologías de Perfil Docente", "G_tipologias_918.png",
    ["G1 'Sello UCEN' (37.9%) y G5 (57.2%) concentran casi todo el universo",
     "G2 y G4 prácticamente no existen — el umbral AFO ≥90% es muy exigente"])

NEW_P1["jer_met"] = new_content_slide(
    "Metodologías y Evaluación (MET) × Jerarquía", "G_jerarquia_met_918.png",
    ["Instructor Docente lidera en Nivel Alto MET (78%)",
     "Titular Docente presenta la menor proporción de Nivel Alto (55%)"])

NEW_P1["jer_afo"] = new_content_slide(
    "Aspectos Formales (AFO) × Jerarquía", "G_jerarquia_afo_918.png",
    ["Asociado Regular tiene la mayor proporción de Nivel Bajo (31%)",
     "El cumplimiento formal óptimo es un desafío transversal al escalafón"])

NEW_P1["afo_antig"] = new_content_slide(
    "AFO Medio/Bajo × Antigüedad (zoom)", "G_afo_antiguedad_918.png",
    ["El Nivel Bajo es más frecuente en docentes de 0-4 años (25% del grupo no-alto)",
     "La experiencia institucional reduce el riesgo pero no lo elimina"])

NEW_P1["afo_jer2"] = new_content_slide(
    "AFO Medio/Bajo × Jerarquía (zoom)", "G_afo_jerarquia_918.png",
    ["Asociado Regular concentra 41% de Nivel Bajo dentro de los no-altos",
     "Instructor Docente es el más sano: 82% en Nivel Medio, solo 18% Bajo"])

print(f"Bloque 1 (Caracterización): {len(NEW_P1)} slides nuevas creadas.")

NEW_P3 = {}
NEW_P3["t1"] = new_content_slide(
    "Validación Estadística — SAT Nota", "G11_pruebas_t_918.png",
    ["Diferencia significativa en 3 de 6 períodos (p<0.05)",
     "Cohen d máximo +0.302 en 2025-02 — efecto pequeño-mediano"])

NEW_P3["t2"] = new_content_slide(
    "Validación Estadística — % Recomendación", "G11.2_pruebas_t_bin_918.png",
    ["Significativo en 4 de 6 períodos",
     "Magnitudes levemente inferiores a SAT Nota (Cohen d máx. +0.230)"])

NEW_P3["t3"] = new_content_slide(
    "Validación Estadística — EDD", "G12_pruebas_t_edd_918.png",
    ["Significativo en 2024 (p=0.002) y 2025 (p<0.001, d=+0.955)",
     "El efecto de la formación sobre EDD crece con el tiempo"])

NEW_P3["t4"] = new_content_slide(
    "T-test Jerarquía × Diplomado", "G13_ttest_jerarquia_diplomado_918.png",
    ["ANOVA no significativo (p=0.718, η²=0.021)",
     "El Diplomado beneficia por igual a docentes de distintos rangos"])

NEW_P3["t5"] = new_content_slide(
    "T-test Jerarquía × Taller", "G13.2_ttest_jerarquia_taller_918.png",
    ["Ningún par de jerarquías muestra diferencia significativa tras Bonferroni",
     "El Taller es efectivo de forma transversal al escalafón"])

NEW_P3["venn"] = new_content_slide(
    "Combinaciones de Modalidades — Población P3", "G_venn_poblaciones_918.png",
    ["93.4% de los 197 aptos P3 hizo una sola modalidad (población pura)",
     "Solo 13 docentes combinaron modalidades, siempre con Taller de por medio"])

NEW_P3["contraste1"] = new_content_slide(
    "Trayectoria SAT — Normal vs Población Pura", "G1_contraste_normal_puro_918.png",
    ["El Diplomado puro muestra mayor caída (Δ=-0.070) que con mixtos (Δ=-0.028)",
     "El Taller es idéntico en ambas versiones — ya era 100% puro"])

NEW_P3["puro_nota_recom"] = new_content_slide(
    "SAT Nota y Recomendación — Población Pura", "G1_puro_nota_recom_918.png",
    ["Mismos patrones en ambas métricas: Proyecto cae más, Taller es estable",
     "Confirma la robustez del hallazgo entre las dos dimensiones SAT"])

NEW_P3["delta_puro_mixto"] = new_content_slide(
    "Δ Z-score: Puro vs Mixto", "G_delta_z_puro_vs_mixto_918.png",
    ["Los 13 mixtos muestran mejor resultado (62% mejoró) que los puros (49%)",
     "Sugiere un posible efecto sinérgico de combinar modalidades"])

NEW_P3["antig_par"] = new_content_slide(
    "Antigüedad × Tipo — SAT Nota y Recomendación", "G3_par_antiguedad_918.png",
    ["129 de 197 aptos P3 tienen dato de antigüedad disponible",
     "Patrones consistentes entre ambas métricas SAT"])

NEW_P3["edd_contraste"] = new_content_slide(
    "EDD — Población Normal vs Pura", "G9_contraste_edd_918.png",
    ["Diplomado puro sube de 0.823 a 0.854 al excluir los mixtos",
     "Taller puro se mantiene en 0.833 — sin cambios"])

NEW_P3["acumulativo"] = new_content_slide(
    "¿Es Acumulativo el Efecto de la Formación?", "G_acumulacion_formacion_918.png",
    ["A partir de 3 instancias la mediana supera al grupo sin formación",
     "Sin significancia estadística por tamaños muestrales reducidos (n=28)"])

NEW_P3["ev_taller"] = new_content_slide(
    "Evolución por Jerarquía — Talleres", "G7_taller_jerarquia_918.png",
    ["Asistentes Docentes parten bajo la media y alcanzan su peak durante talleres",
     "Instructores Docentes estables — el taller consolida más que transforma"])

NEW_P3["ev_diplomado"] = new_content_slide(
    "Evolución por Jerarquía — Diplomados", "G7_diplomado_jerarquia_918.png",
    ["Instructores Docentes: caso de mayor éxito, parten bajo y suben",
     "ANOVA no significativo: el diplomado beneficia por igual a todos los rangos"])

NEW_P3["aprob_antig"] = new_content_slide(
    "Tasa de Aprobación × Antigüedad Docente", "G_aprobacion_antiguedad_918.png",
    ["La tasa de aprobación es relativamente estable entre tramos de antigüedad",
     "La dispersión intra-tramo supera las diferencias entre tramos"])

NEW_P3["aprob_jer"] = new_content_slide(
    "Tasa de Aprobación × Jerarquía Académica", "G_aprobacion_jerarquia_918.png",
    ["Instructor Docente lidera con la mediana más alta (94.6%)",
     "La jerarquía no es un predictor directo del rendimiento estudiantil"])

print(f"Bloque 2 (Impacto P3): {len(NEW_P3)} slides nuevas creadas.")

NEW_P2 = {}
NEW_P2["facultad"] = new_content_slide(
    "Participación en Formación por Facultad", "p2_g22_facultad.png",
    ["Brechas de hasta 30pp entre la facultad más y menos activa",
     "Análisis sobre los 493 docentes con perfil completo"])

NEW_P2["antiguedad"] = new_content_slide(
    "Participación en Formación por Antigüedad", "p2_g23_antiguedad.png",
    ["El tramo 0-4 años concentra la mayor tasa de participación (50%)",
     "491 docentes con dato de antigüedad disponible"])

NEW_P2["modalidad"] = new_content_slide(
    "Distribución por Modalidad de Formación", "p2_g31_modalidad.png",
    ["Taller concentra el 61% de las 615 instancias totales",
     "357 docentes únicos generan 615 actividades de formación"])

NEW_P2["brechas"] = new_content_slide(
    "Brechas: Perfil de No Participantes", "p2_g41_brechas.png",
    ["560 docentes sin ninguna instancia de formación (61.1%)",
     "Instructor Docente y el tramo 0-4 años concentran el mayor volumen"])

NEW_P2["intensidad_global"] = new_content_slide(
    "Intensidad de Participación Global", "p2_g42_intensidad.png",
    ["61% de los formados participó en una sola instancia",
     "8% con 4+ instancias representa el núcleo de alta dedicación"])

NEW_P2["venn917"] = new_content_slide(
    "Combinaciones de Formación — Universo 917", "G_venn_formacion_917_918.png",
    ["357 docentes con formación: 76% hizo una sola modalidad",
     "Diplomado+Taller es la intersección más frecuente (65 docentes)"])

NEW_P2["int_talleres"] = new_content_slide(
    "Intensidad de Participación en Talleres", "G_intensidad_talleres_918.png",
    ["59% de los docentes tomó un solo taller — experiencia puntual",
     "Heavy user: José Faúndez (Honorario) con 11 instancias"])

NEW_P2["int_proyectos"] = new_content_slide(
    "Intensidad de Participación en Proyectos", "G_intensidad_proyectos_918.png",
    ["88% participó en un solo proyecto",
     "Juan Carlos Araya Vargas: máximo de 4 proyectos (2 Innovación + 2 Investigación)"])

print(f"Bloque 3 (Participación P2): {len(NEW_P2)} slides nuevas creadas.")

NEW_C = {}
NEW_C["venn_contrato"] = new_content_slide(
    "Combinaciones de Formación — Jornada vs Honorario", "G_venn_formacion_contrato_918.png",
    ["Ninguna intersección de Honorario incluye Proyecto",
     "Jornada: 31 docentes combinan Proyecto con otra modalidad"])

print(f"Bloque 4 (NOMINA): {len(NEW_C)} slides nuevas creadas.")

# ══════════════════════════════════════════════════════════════════════════════
# 3. Separadores de bloque
# ══════════════════════════════════════════════════════════════════════════════
SEP1 = new_separator("BLOQUE 1", "Caracterización del Universo 917 Jerarquizados")
SEP2 = new_separator("BLOQUE 2", "Impacto de la Formación\nUniverso: 917 → 197 Aptos P3")
SEP3 = new_separator("BLOQUE 3", "Participación en Formación\nUniverso: 917 → 357 Formados")
# El separador "HONORARIO VS JORNADA" (S[22]) ya cumple función de Bloque 4

CIERRE1 = new_content_slide(
    "Conclusiones Consolidadas", None,
    ["La formación genera impacto medible en SAT, EDD y aprobación — efecto de selección + protección",
     "No hay diferencia Jornada/Honorario en calidad docente: el contrato no determina el desempeño",
     "La cobertura (39%) es la principal oportunidad: 6 de cada 10 nunca participaron",
     "El impacto se consolida a partir de 3+ instancias — la participación única no basta"],
    marcar=False)
CIERRE2 = new_content_slide(
    "Recomendaciones", None,
    ["Institucionalizar Diplomados como requisito de avance en la carrera académica",
     "Diseñar rutas formativas progresivas en lugar de instancias aisladas",
     "Ampliar cobertura a Honorarios: mismo retorno con menor acceso actual",
     "Crear mecanismo de evaluación de desempeño (EDD) para Honorarios"],
    marcar=False)

print("Separadores y cierre creados.")

# ══════════════════════════════════════════════════════════════════════════════
# 4. Reordenar TODO el deck según la estructura final
# ══════════════════════════════════════════════════════════════════════════════
final_order = (
    [S[1], S[2]] +                                   # Apertura
    [S[6]] +                                          # Índice (reescrito)
    [SEP1] +
    [S[13], S[14], S[15], S[16], S[17]] +              # Caracterización original
    [NEW_P1["jerarquia"], NEW_P1["nivel_formacion"], NEW_P1["jer_formacion"],
     NEW_P1["gradorec"], NEW_P1["institucion"], NEW_P1["pais"], NEW_P1["carga"],
     NEW_P1["anios_jerarquia"], NEW_P1["edad_jerarquizarse"],
     NEW_P1["dimensiones"], NEW_P1["tipologias"], NEW_P1["jer_met"], NEW_P1["jer_afo"],
     NEW_P1["afo_antig"], NEW_P1["afo_jer2"]] +
    [SEP2] +
    [S[3], S[4], S[5]] +                               # Universo / metodología / funnel P3
    [S[7], S[8], S[9], S[10], S[11], S[12]] +          # Hallazgos P3 original
    [NEW_P3["t1"], NEW_P3["t2"], NEW_P3["t3"], NEW_P3["t4"], NEW_P3["t5"],
     NEW_P3["venn"], NEW_P3["contraste1"], NEW_P3["puro_nota_recom"],
     NEW_P3["delta_puro_mixto"], NEW_P3["antig_par"], NEW_P3["edd_contraste"],
     NEW_P3["acumulativo"], NEW_P3["ev_taller"], NEW_P3["ev_diplomado"],
     NEW_P3["aprob_antig"], NEW_P3["aprob_jer"]] +
    [SEP3] +
    [S[18], S[19], S[21]] +                            # P2 original
    [NEW_P2["facultad"], NEW_P2["antiguedad"], NEW_P2["modalidad"], NEW_P2["brechas"],
     NEW_P2["intensidad_global"], NEW_P2["venn917"], NEW_P2["int_talleres"],
     NEW_P2["int_proyectos"]] +
    [S[22]] +                                          # Separador Honorario/Jornada (Bloque 4)
    [S[23], S[24], S[25], S[26]] +                     # NOMINA original
    [NEW_C["venn_contrato"]] +
    [CIERRE1, CIERRE2]
)

# Eliminar relación de la slide duplicada (S[20]) para no dejarla huérfana
dup_slide = S[20]
dup_rId = None
for rel_id, rel in prs.part.rels.items():
    if rel.target_part == dup_slide.part:
        dup_rId = rel_id
        break

sldIdLst = prs.slides._sldIdLst

# Mapeo Slide(objeto) -> elemento <p:sldId> real, ANTES de reordenar
# (prs.slides y sldIdLst están en el mismo orden en este punto)
current_slide_objs = list(prs.slides)
current_sldid_els  = list(sldIdLst)
slide_to_sldid = {id(s._element): sldid for s, sldid in zip(current_slide_objs, current_sldid_els)}

# Construir nueva lista de <p:sldId> en el orden deseado
new_sldid_elements = [slide_to_sldid[id(s._element)] for s in final_order]

# Verificar que no falten / sobren referencias
assert len(new_sldid_elements) == len(set(id(e) for e in new_sldid_elements)), "Hay slides repetidas en el orden final"
assert id(dup_slide._element) not in {id(s._element) for s in final_order}, "El duplicado no debe estar en el orden final"

for el in list(sldIdLst):
    sldIdLst.remove(el)
for el in new_sldid_elements:
    sldIdLst.append(el)

if dup_rId:
    prs.part.drop_rel(dup_rId)
    print(f"Relación de slide duplicada (#20) eliminada del paquete.")

print(f"\nTotal slides en el deck final: {len(prs.slides)}")
print(f"(slide duplicada original #20 fue excluida del orden final)")

prs.save(OUT)
print(f"\nGuardado: {OUT}")

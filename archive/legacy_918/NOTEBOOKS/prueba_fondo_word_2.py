"""
prueba_fondo_word_2.py — Prueba de concepto ronda 2
Mismos 4 colores, pero el gráfico re-renderizado con ejes en NEGRO
(ticks, labels, spines, títulos de eje).
"""
import sys; sys.stdout.reconfigure(encoding="utf-8")
import os
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patheffects as pe
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE   = os.path.dirname(os.path.abspath(__file__))
PROC   = os.path.join(BASE, "..", "PROCESADO")
ROOT   = os.path.normpath(os.path.join(BASE, "..", ".."))
SCRATCH = (r"C:\Users\RGONZA~1.LAP\AppData\Local\Temp\claude"
           r"\c--Users-r-gonzalez-fluxsolar-LAPTOP-FLUX-ECO-Downloads-Analisis-UCEN-v2"
           r"\19e6fc3f-6ca1-4150-9da7-8dfa38be71ca\scratchpad")
os.makedirs(SCRATCH, exist_ok=True)

FONDOS = [
    ("#54668E", ( 84, 102, 142), "Opción A — #54668E"),
    ("#363955", ( 54,  57,  85), "Opción B — #363955"),
    ("#879EC6", (135, 158, 198), "Opción C — #879EC6"),
    ("#59598E", ( 89,  89, 142), "Opción D — #59598E"),
]

SW = 12192000
SH =  6858000

# ─── Datos ────────────────────────────────────────────────────────────────────
cvt = pd.read_csv(os.path.join(PROC, "control_vs_trat_918.csv"), encoding="utf-8-sig")
cvt["z_trat"] = pd.to_numeric(cvt["z_trat"], errors="coerce")
cvt["z_ctrl"] = pd.to_numeric(cvt["z_ctrl"], errors="coerce")
cvt["n_trat"] = pd.to_numeric(cvt["n_trat"], errors="coerce").astype(int)
cvt["n_ctrl"] = pd.to_numeric(cvt["n_ctrl"], errors="coerce").astype(int)
periodos = cvt["periodo"].tolist()
z_f = cvt["z_trat"].tolist(); n_f = cvt["n_trat"].tolist()
z_c = cvt["z_ctrl"].tolist(); n_c = cvt["n_ctrl"].tolist()

# ─── Renderizar gráfico con ejes en negro (transparente) ──────────────────────
def render_chart_black_axes(rgb_fondo):
    """Re-renderiza slide_22 con ejes/labels/títulos en negro sobre fondo transparente."""
    fig, ax = plt.subplots(figsize=(13, 6), facecolor="none")
    fig.patch.set_alpha(0.0)
    ax.set_facecolor("none")

    xa = range(len(periodos))
    ax.plot(xa, z_f, color="#1565C0", linewidth=2.5, linestyle="-",
            marker="o", markersize=9, label="Formados", zorder=5)
    ax.plot(xa, z_c, color="#E65100", linewidth=2.5, linestyle="--",
            marker="s", markersize=8, label="Control", zorder=5)

    for i, (zf, zc, nf, nc) in enumerate(zip(z_f, z_c, n_f, n_c)):
        ax.text(i, zf + 0.016, f"{zf:.2f} (n={nf})",
                ha="center", va="bottom", fontsize=8, fontweight="bold",
                color="#1565C0",
                path_effects=[pe.withStroke(linewidth=1.2, foreground="black")])
        ax.text(i, zc - 0.016, f"{zc:.2f} (n={nc})",
                ha="center", va="top", fontsize=8, fontweight="bold",
                color="#E65100",
                path_effects=[pe.withStroke(linewidth=1.2, foreground="black")])

    ax.axhline(0, color="black", linewidth=1, linestyle=":", alpha=0.6)

    # Ejes en NEGRO
    ax.tick_params(axis="x", colors="black", labelsize=9.5, length=4, width=0.8)
    ax.tick_params(axis="y", colors="black", labelsize=9.5, length=4, width=0.8)
    for sp in ax.spines.values():
        sp.set_edgecolor("black"); sp.set_alpha(0.8); sp.set_linewidth(0.9)
    # Grid: gris oscuro sobre el fondo de color — distinto del blanco, con clase
    ax.yaxis.grid(True, color="#1C1C1C", alpha=0.45, linewidth=0.75, zorder=0)
    ax.set_axisbelow(True)

    ax.set_xticks(list(xa))
    ax.set_xticklabels(periodos, fontsize=9.5, color="black", rotation=15)
    ax.set_ylabel("z-score SAT (promedio)", color="black", fontsize=10)

    ax.legend(fontsize=10, framealpha=0.25, labelcolor="black",
              facecolor="white", edgecolor="#888888")

    plt.tight_layout(pad=0.8)

    # Guardar PNG transparente
    tmp = os.path.join(SCRATCH, "chart22_black_axes.png")
    fig.savefig(tmp, dpi=150, facecolor="none", transparent=True, bbox_inches="tight")
    plt.close()
    return tmp

# ─── Autocrop ────────────────────────────────────────────────────────────────
def _autocrop(img_rgba, pad=25):
    r, g, b, a = img_rgba.split()
    bbox = a.getbbox()
    if bbox:
        W, H = img_rgba.size
        return img_rgba.crop((max(0, bbox[0]-pad), max(0, bbox[1]-pad),
                              min(W, bbox[2]+pad), min(H, bbox[3]+pad)))
    return img_rgba

# ─── Componer sobre fondo sólido ──────────────────────────────────────────────
def _compose(chart_rgba, rgb_fondo, out_size=(1920, 1080)):
    bg = PILImage.new("RGB", out_size, rgb_fondo)
    ch = chart_rgba.resize(out_size, PILImage.LANCZOS)
    bg.paste(ch, (0, 0), ch)
    return bg

# ─── Generar el chart una sola vez ───────────────────────────────────────────
print("Renderizando gráfico con ejes negros…")
chart_png = render_chart_black_axes(None)
chart_rgba = PILImage.open(chart_png).convert("RGBA")
chart_crop = _autocrop(chart_rgba)
print(f"  Chart recortado: {chart_crop.size}")

# ─── PPTX ────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Emu(SW)
prs.slide_height = Emu(SH)

for hex_color, rgb, label in FONDOS:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(*rgb)

    # Etiqueta
    tx = slide.shapes.add_textbox(Emu(400000), Emu(80000),
                                   Emu(SW - 800000), Emu(380000))
    tf = tx.text_frame
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = label + "  ·  ejes en negro"
    run.font.size  = Pt(20)
    run.font.bold  = True
    # Color del texto según luminosidad del fondo
    lum = 0.299*rgb[0] + 0.587*rgb[1] + 0.114*rgb[2]
    txt_rgb = (255, 255, 255) if lum < 140 else (0, 33, 71)
    run.font.color.rgb = RGBColor(*txt_rgb)

    # Componer imagen
    composed = _compose(chart_crop, rgb, out_size=(1920, 1080))
    tmp = os.path.join(SCRATCH, f"poc2_{hex_color[1:]}.png")
    composed.save(tmp, "PNG")

    # Insertar con margen uniforme
    margin_x = int(SW * 0.035)
    margin_y = int(SH * 0.14)
    img_w    = SW - 2 * margin_x
    img_h    = SH - margin_y - int(SH * 0.04)

    slide.shapes.add_picture(tmp, Emu(margin_x), Emu(margin_y),
                              Emu(img_w), Emu(img_h))

out = os.path.join(ROOT, "PRUEBA_fondos_informe_5.pptx")
prs.save(out)
print(f"\n✓ Guardado: {out}")
print(f"  4 slides — ejes negros sobre: {', '.join(h for h,_,_ in FONDOS)}")

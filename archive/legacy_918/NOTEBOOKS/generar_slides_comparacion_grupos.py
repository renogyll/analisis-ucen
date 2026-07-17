"""
generar_slides_comparacion_grupos.py  — v2
4 slides que explican los grupos de control y comparan perfiles demográficos:
  Slide 1: Embudo de grupos (derivación por bloque)
  Slide 2: Bloque III — Perfil jerarquía + tramo edad (formados vs control aprobación)
  Slide 3: Bloque IV  — Perfil jerarquía + tramo edad (formados vs control EDD)
  Slide 4: Tabla resumen (fuentes originales + derivadas)
"""
import sys; sys.stdout.reconfigure(encoding="utf-8")
import os, zipfile
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.patheffects as pe
import matplotlib.patches as mpatches
from matplotlib.lines import Line2D
import pandas as pd
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─────────────────────────────────────────────────────────────────────────────
# Rutas
# ─────────────────────────────────────────────────────────────────────────────
BASE       = os.path.dirname(os.path.abspath(__file__))
PROC       = os.path.join(BASE, "..", "PROCESADO")
ROOT_PROC  = os.path.normpath(os.path.join(BASE, "..", "..", "PROCESADO"))
EDD_CSV    = os.path.join(ROOT_PROC, "P1_consolidado_con_evaluacion_jefes.csv")
DOC918_CSV = os.path.join(PROC, "docente_918.csv")
SCAT_CSV   = os.path.join(PROC, "scatter_sat_notas.csv")
FONDOTIPO  = os.path.normpath(os.path.join(BASE, "..", "..", "Fondotipop.pptx"))
OUT_PPTX   = os.path.join(BASE, "COMPARACION_GRUPOS_CONTROL.pptx")
OUT_DIR    = os.path.join(BASE, "dark_slides_v3")
os.makedirs(OUT_DIR, exist_ok=True)

# ─────────────────────────────────────────────────────────────────────────────
# Assets
# ─────────────────────────────────────────────────────────────────────────────
BG_PATH   = os.path.join(OUT_DIR, "_fondotipo_bg.jpg")
LOGO_PATH = os.path.join(OUT_DIR, "_fondotipo_logo.png")
for path, zname in [(BG_PATH, "ppt/media/image1.jpg"),
                    (LOGO_PATH, "ppt/media/image2.png")]:
    if not os.path.exists(path):
        with zipfile.ZipFile(FONDOTIPO) as z:
            with open(path, "wb") as f:
                f.write(z.read(zname))

with PILImage.open(BG_PATH) as _im:
    _rgb = _im.convert("RGB"); _iw, _ih = _rgb.size
    _nh = int(_iw / (16 / 9)); _y0 = min(int(_ih * 0.12), _ih - _nh)
    bg_arr = np.array(_rgb.crop((0, _y0, _iw, _y0 + _nh)))

with PILImage.open(LOGO_PATH) as _logo:
    logo_arr = np.array(_logo.convert("RGBA")).astype(np.float32) / 255.0

H_GRAD = 600; grad = np.zeros((H_GRAD, 1, 4), dtype=np.float32)
for _r in range(H_GRAD):
    _t = _r / (H_GRAD - 1)
    _stops = [(0.00, (0, 33, 71)), (0.54, (0, 70, 128)), (1.00, (144, 171, 196))]
    for _i in range(len(_stops) - 1):
        _t0, _c0 = _stops[_i]; _t1, _c1 = _stops[_i + 1]
        if _t0 <= _t <= _t1:
            _s = (_t - _t0) / (_t1 - _t0)
            grad[_r, 0] = [(_c0[0] + _s * (_c1[0] - _c0[0])) / 255,
                           (_c0[1] + _s * (_c1[1] - _c0[1])) / 255,
                           (_c0[2] + _s * (_c1[2] - _c0[2])) / 255, 0.82]
            break

# ─────────────────────────────────────────────────────────────────────────────
# Constantes layout (idénticas al deck principal)
# ─────────────────────────────────────────────────────────────────────────────
SW, SH     = 13.333, 7.5
SW_EMU     = 12192000
SH_EMU     = 6858000

PIC_L, PIC_T, PIC_W, PIC_H     = 786581, 1125000, 10599174, 3720000
BUL_L, BUL_T, BUL_W, BUL_H     = 786581, 4870000, 10599174, 1870000
LOGO_L, LOGO_T, LOGO_W, LOGO_H = 9813773, 656354, 1756626, 697725
TITLE_L, TITLE_T, TITLE_W, TITLE_H = PIC_L, 185000, PIC_W, 710000
POP_L, POP_T, POP_W, POP_H     = PIC_L, 845000, 9000000, 255000

def _ex(e): return e / SW_EMU   # fracción horizontal
def _ey(e): return e / SH_EMU   # fracción vertical

def _fig_rect(l, t, w, h):
    """Convierte (left,top,width,height) en fracciones de figura matplotlib."""
    return (l, 1 - t - h, w, h)

# CORRECTO: _ey para dimensiones verticales (PIC_T, PIC_H)
PIC_RECT  = _fig_rect(_ex(PIC_L), _ey(PIC_T), _ex(PIC_W), _ey(PIC_H))
LOGO_RECT = _fig_rect(_ex(LOGO_L), _ey(LOGO_T), _ex(LOGO_W), _ey(LOGO_H))

COL_FORM = "#5C9BD6"
COL_CTRL = "#FFB74D"

# Orden jerárquico canónico (de menor a mayor — aparece de abajo a arriba en gráficos)
JER_ORD = ["INSTRUCTOR REGULAR", "INSTRUCTOR DOCENTE",
           "ASISTENTE REGULAR",  "ASISTENTE DOCENTE",
           "ASOCIADO REGULAR",   "ASOCIADO DOCENTE",
           "TITULAR REGULAR",    "TITULAR DOCENTE"]
JER_LBL = ["Instr. Regular", "Instr. Docente",
            "Asist. Regular", "Asist. Docente",
            "Asoc. Regular",  "Asoc. Docente",
            "Tit. Regular",   "Tit. Docente"]

# Tramos de edad agrupados (de menor a mayor)
TRAMOS_EDAD = ["< 35", "35–44", "45–54", "55–64", "≥ 65"]
TRAMOS_EDAD_RAW = {
    "<30": "< 35", "30-34": "< 35",
    "35-39": "35–44", "40-44": "35–44",
    "45-49": "45–54", "50-54": "45–54",
    "55-59": "55–64", "60-64": "55–64",
    "65-69": "≥ 65",  "70+":   "≥ 65",
}

def _tipo_simple(t):
    t = str(t).upper()
    if "TALLER" in t and "DIPLOMADO" not in t and "PROYECTO" not in t: return "Taller"
    if "DIPLOMADO" in t and "TALLER" not in t and "PROYECTO" not in t: return "Diplomado"
    if "PROYECTO" in t and "TALLER" not in t and "DIPLOMADO" not in t: return "Proyecto"
    return "Participación Mixta"

# ─────────────────────────────────────────────────────────────────────────────
# Helpers matplotlib
# ─────────────────────────────────────────────────────────────────────────────
def _bg_fig():
    fig = plt.figure(figsize=(SW, SH), facecolor="#101820")
    fig.patch.set_facecolor("#101820")
    for z, arr in [(0, bg_arr), (1, grad)]:
        ax = fig.add_axes([0, 0, 1, 1], zorder=z)
        ax.imshow(arr, extent=[0, 1, 0, 1], aspect="auto", origin="upper")
        ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis("off")
    al = fig.add_axes([LOGO_RECT[0], LOGO_RECT[1], LOGO_RECT[2], LOGO_RECT[3]],
                      zorder=10, facecolor="none")
    al.imshow(logo_arr, aspect="auto"); al.axis("off"); al.patch.set_visible(False)
    return fig

def _tr_fig():
    fig = plt.figure(figsize=(SW, SH), facecolor="none")
    fig.patch.set_facecolor("none")
    return fig

def _save_bg(fig, name):
    path = os.path.join(OUT_DIR, name)
    plt.savefig(path, dpi=150, facecolor=fig.get_facecolor()); plt.close(); return path

def _save_ch(fig, name):
    path = os.path.join(OUT_DIR, name)
    plt.savefig(path, dpi=150, facecolor="none", transparent=True); plt.close(); return path

def _style_ax(ax):
    ax.tick_params(axis="x", colors="#AAAAAA", labelsize=8, length=0)
    ax.tick_params(axis="y", colors="white",   labelsize=8, length=0)
    for sp in ax.spines.values():
        sp.set_edgecolor("white"); sp.set_alpha(0.20); sp.set_linewidth(0.7)
    ax.set_axisbelow(True)

SHARED_BG = os.path.join(OUT_DIR, "_background.png")

def _ensure_bg():
    if not os.path.exists(SHARED_BG):
        _save_bg(_bg_fig(), "_background.png")
        print("  bg generado")

# ─────────────────────────────────────────────────────────────────────────────
# Helper: butterfly chart pareado
# ─────────────────────────────────────────────────────────────────────────────
def _butterfly(fig, bx, by, bw, bh,
               cats, pct_f, pct_c, n_f_tot, n_c_tot,
               chart_title, col_f=COL_FORM, col_c=COL_CTRL):
    """
    Dibuja un gráfico butterfly (back-to-back) dentro del área (bx,by,bw,bh).
    cats: categorías de abajo a arriba.
    pct_f, pct_c: listas de % para formados y control (mismo orden que cats).
    """
    ya = np.arange(len(cats))
    max_val = max(max(pct_f, default=1), max(pct_c, default=1))
    xlim = max_val * 1.30

    gap_lbl = 0.045     # espacio central para etiquetas de categoría
    w_panel = (bw - gap_lbl) / 2

    # ── Panel izquierdo (formados) ──
    ax_f = fig.add_axes([bx, by, w_panel, bh], facecolor="none", zorder=5)
    ax_f.barh(ya, pct_f, height=0.60, color=col_f, alpha=0.88, edgecolor="none")
    ax_f.set_xlim(xlim, 0)          # invertido: barras van a la izquierda
    ax_f.set_ylim(-0.6, len(cats) - 0.4)
    ax_f.set_yticks(ya); ax_f.set_yticklabels([])
    ax_f.xaxis.grid(True, color="white", alpha=0.06, linewidth=0.5)
    # etiquetas de valor dentro/fuera de la barra
    for j, v in enumerate(pct_f):
        if v >= 4:
            ax_f.text(v / 2, j, f"{v:.0f}%", ha="center", va="center",
                      fontsize=7.5, color="white", fontweight="bold",
                      path_effects=[pe.withStroke(linewidth=0.8, foreground=col_f)])
        elif v > 0:
            ax_f.text(v + 0.5, j, f"{v:.0f}%", ha="right", va="center",
                      fontsize=7, color="#AAAAAA")
    # ticks del eje x sin signo
    ax_f.set_xticklabels([f"{abs(t):.0f}%" for t in ax_f.get_xticks()],
                         fontsize=7.5, color="#AAAAAA")
    ax_f.set_title(f"Formados  (n={n_f_tot})", color=col_f, fontsize=9,
                   fontweight="bold", pad=5)
    _style_ax(ax_f)

    # ── Panel derecho (control) ──
    ax_c = fig.add_axes([bx + w_panel + gap_lbl, by, w_panel, bh],
                        facecolor="none", zorder=5)
    ax_c.barh(ya, pct_c, height=0.60, color=col_c, alpha=0.88, edgecolor="none")
    ax_c.set_xlim(0, xlim)
    ax_c.set_ylim(-0.6, len(cats) - 0.4)
    ax_c.set_yticks(ya); ax_c.set_yticklabels(cats, fontsize=8.5, color="white")
    ax_c.yaxis.set_ticks_position("left")
    ax_c.xaxis.grid(True, color="white", alpha=0.06, linewidth=0.5)
    for j, v in enumerate(pct_c):
        if v >= 4:
            ax_c.text(v / 2, j, f"{v:.0f}%", ha="center", va="center",
                      fontsize=7.5, color="white", fontweight="bold",
                      path_effects=[pe.withStroke(linewidth=0.8, foreground=col_c)])
        elif v > 0:
            ax_c.text(v + 0.5, j, f"{v:.0f}%", ha="left", va="center",
                      fontsize=7, color="#AAAAAA")
    ax_c.set_xticklabels([f"{t:.0f}%" for t in ax_c.get_xticks()],
                         fontsize=7.5, color="#AAAAAA")
    ax_c.set_title(f"Control  (n={n_c_tot})", color=col_c, fontsize=9,
                   fontweight="bold", pad=5)
    _style_ax(ax_c)

    # ── Título del par ──
    fig.text(bx + bw / 2, by + bh + 0.035, chart_title,
             ha="center", va="bottom", fontsize=10, color="white",
             fontweight="bold",
             path_effects=[pe.withStroke(linewidth=1.5, foreground="#0A1830")])


# ─────────────────────────────────────────────────────────────────────────────
# Helpers python-pptx
# ─────────────────────────────────────────────────────────────────────────────
def _new_sl(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _pic(sl, path, prs):
    sl.shapes.add_picture(path, Emu(0), Emu(0), prs.slide_width, prs.slide_height)

def _txt(sl, text, left, top, width, height,
         fs=12, bold=False, italic=False, color="#FFFFFF",
         align=PP_ALIGN.LEFT, wrap=True, lspc=0):
    txb = sl.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txb.text_frame; tf.word_wrap = wrap
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if lspc > 0 and i > 0: p.space_before = Pt(lspc)
        run = p.add_run(); run.text = line
        run.font.size = Pt(fs); run.font.bold = bold; run.font.italic = italic
        r, g, b = int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)
        run.font.color.rgb = RGBColor(r, g, b)

def _T(sl, text, fs=19):
    _txt(sl, text, TITLE_L, TITLE_T, TITLE_W, TITLE_H,
         fs=fs, bold=True, color="#FFFFFF", align=PP_ALIGN.CENTER)

def _POP(sl, text=""):
    _txt(sl, text, POP_L, POP_T, POP_W, POP_H, fs=7.5, italic=True, color="#C8DCF0")

def _BUL(sl, items, fs=11):
    txb = sl.shapes.add_textbox(Emu(BUL_L), Emu(BUL_T), Emu(BUL_W), Emu(BUL_H))
    tf = txb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4); p.alignment = PP_ALIGN.LEFT
        run = p.add_run(); run.text = f"{i+1}.  {item}"
        run.font.size = Pt(fs)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# ─────────────────────────────────────────────────────────────────────────────
# Carga de datos
# ─────────────────────────────────────────────────────────────────────────────
print("Cargando datos...")

sat = pd.read_csv(os.path.join(PROC, "p3_sat_zscore_918.csv"), encoding="utf-8-sig")
sat["rut_key"] = sat["rut_key"].astype(str).str.strip()
sat["jerarquia_u"] = sat["jerarquia"].str.strip().str.upper()
sat["tramo_g"] = sat["tramo_edad"].map(TRAMOS_EDAD_RAW)
ruts_197 = set(sat["rut_key"])

p3ev = pd.read_csv(os.path.join(PROC, "p3_918.csv"), encoding="utf-8-sig")
p3ev["rut_key"] = p3ev["rut_key"].astype(str).str.strip()
ruts_todos_formados = set(p3ev["rut_key"])

doc918 = pd.read_csv(DOC918_CSV, dtype=str, encoding="utf-8-sig")
doc918["rut_key"] = doc918["rut_key"].str.strip()
doc918["jerarquia_u"] = doc918["jerarquia"].str.strip().str.upper()
doc918["tramo_g"] = doc918["tramo_edad"].map(TRAMOS_EDAD_RAW)
ruts_917 = set(doc918["rut_key"])

# scatter_sat_notas — grupos Bloque III
scat = pd.read_csv(SCAT_CSV, encoding="utf-8-sig")
scat["formado"] = scat["formado"].astype(str).str.strip().str.upper().isin(
    ["TRUE", "1", "SI", "SÍ", "YES"])
scat_ctrl_ruts = set(scat[~scat["formado"]]["rut_docente"].astype(str).str.strip().unique())
# Perfil del control Bloque III (uniendo con doc918)
ctrl_b3_doc = doc918[doc918["rut_key"].isin(scat_ctrl_ruts)].drop_duplicates("rut_key")
N_B3_FORM = len(sat)                         # 197 docentes
N_B3_CTRL = len(ctrl_b3_doc)

# EDD — grupos Bloque IV
has_edd = os.path.exists(EDD_CSV)
if has_edd:
    edd_df = pd.read_csv(EDD_CSV, dtype={"rut_key": str}, encoding="utf-8-sig")
    edd_df["rut_key"]   = edd_df["rut_key"].str.strip()
    edd_df["edd_total"] = pd.to_numeric(edd_df["edd_total"], errors="coerce")
    edd_df["anio_eval"] = edd_df["anio_evaluacion"].apply(
        lambda x: str(int(float(x)))[:4] if pd.notna(x) else None)
    edd_form = (edd_df[edd_df["rut_key"].isin(ruts_197)
                       & edd_df["edd_total"].notna()
                       & edd_df["anio_eval"].notna()]
                .drop_duplicates(subset=["rut_key", "anio_eval"]))
    edd_ctrl = (edd_df[edd_df["rut_key"].isin(ruts_917)
                       & ~edd_df["rut_key"].isin(ruts_todos_formados)
                       & edd_df["edd_total"].notna()
                       & edd_df["anio_eval"].notna()]
                .drop_duplicates(subset=["rut_key", "anio_eval"]))
    edd_form_ruts = set(edd_form["rut_key"].unique())
    edd_ctrl_ruts = set(edd_ctrl["rut_key"].unique())
    # Perfil usando doc918
    b4_form_doc = doc918[doc918["rut_key"].isin(edd_form_ruts)].drop_duplicates("rut_key")
    b4_ctrl_doc = doc918[doc918["rut_key"].isin(edd_ctrl_ruts)].drop_duplicates("rut_key")
    N_B4_FORM = len(b4_form_doc)
    N_B4_CTRL = len(b4_ctrl_doc)
else:
    edd_form = edd_ctrl = pd.DataFrame()
    b4_form_doc = b4_ctrl_doc = pd.DataFrame()
    edd_form_ruts = edd_ctrl_ruts = set()
    N_B4_FORM = N_B4_CTRL = 0

N_NO_FORM = len(ruts_917 - ruts_todos_formados)
print(f"  B3: formados={N_B3_FORM} doc | control={N_B3_CTRL} doc")
print(f"  B4: formados={N_B4_FORM} doc | control={N_B4_CTRL} doc")


# ─────────────────────────────────────────────────────────────────────────────
# Helper: % distribución por categoría
# ─────────────────────────────────────────────────────────────────────────────
def _pct_jer(df, col="jerarquia_u"):
    vc = df[col].value_counts()
    return [round(vc.get(j, 0) / max(len(df), 1) * 100, 1) for j in JER_ORD]

def _pct_tramo(df, col="tramo_g"):
    vc = df[col].value_counts()
    return [round(vc.get(t, 0) / max(len(df), 1) * 100, 1) for t in TRAMOS_EDAD]


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Embudo de grupos (FIXED: cajas más grandes, ≤2 líneas)
# ─────────────────────────────────────────────────────────────────────────────
def slide_embudo(prs):
    _ensure_bg()
    fig = _tr_fig()

    ax = fig.add_axes([PIC_RECT[0] + 0.005, PIC_RECT[1] + 0.005,
                       PIC_RECT[2] - 0.01,  PIC_RECT[3] - 0.01],
                      facecolor="none", zorder=5)
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis("off")

    # Box: una línea — títulos acortados para no wrappear
    def _box(cx, cy, title, color, w=0.40, h=0.085, fsize=9.5, alpha=0.88):
        patch = mpatches.FancyBboxPatch(
            (cx - w/2, cy - h/2), w, h,
            boxstyle="round,pad=0.012", linewidth=1.3,
            edgecolor="white", facecolor=color, alpha=alpha, zorder=4,
            transform=ax.transData)
        ax.add_patch(patch)
        ax.text(cx, cy, title, ha="center", va="center",
                fontsize=fsize, color="white", fontweight="bold",
                zorder=5, transform=ax.transData)

    def _arrow(x0, y0, x1, y1, color="white"):
        ax.annotate("", xy=(x1, y1), xytext=(x0, y0),
                    arrowprops=dict(arrowstyle="-|>", color=color,
                                   lw=1.5, alpha=0.70), zorder=3)

    # ── Posiciones verticales — todo bajado para empezar más abajo ───────────
    BHH = 0.042   # box half-height = 0.085/2
    L1 = 0.87;  L1_bot = L1 - BHH          # = 0.828
    L2 = 0.69;  L2_top = L2 + BHH; L2_bot = L2 - BHH   # top=0.732, bot=0.648
    L3 = 0.47;  L3_top = L3 + BHH; L3_bot = L3 - BHH   # top=0.512, bot=0.428
    L4 = 0.25;  L4_top = L4 + BHH; L4_bot = L4 - BHH   # top=0.292, bot=0.208

    # ── Nivel 1: raíz ───────────────────────────────────────────────────────
    _box(0.50, L1, "917 Docentes Jerarquizados",
         color="#1A3A5C", w=0.46, h=BHH*2, fsize=10.5)

    # ── Bifurcación L1 → L2: diagonales directas (sin kink) ─────────────────
    # Dos flechas diagonales desde el fondo de la raíz hasta la cima de L2
    _arrow(0.50, L1_bot, 0.27, L2_top, color=COL_FORM)
    _arrow(0.50, L1_bot, 0.73, L2_top, color=COL_CTRL)
    MID12_y = (L1_bot + L2_top) / 2
    ax.text(0.355, MID12_y + 0.012, "Participa en P3",
            ha="center", va="bottom", fontsize=8, color=COL_FORM, style="italic")
    ax.text(0.645, MID12_y + 0.012, "No participa en P3",
            ha="center", va="bottom", fontsize=8, color=COL_CTRL, style="italic")

    # ── Nivel 2: 357 formados / 560 no formados ─────────────────────────────
    _box(0.27, L2, "357 Docentes Formados",  color="#1A4E7A", w=0.36, h=BHH*2)
    _box(0.73, L2, f"{N_NO_FORM} Docentes No Formados", color="#6B4A1A", w=0.36, h=BHH*2)

    # ── Flechas L2 → L3 ─────────────────────────────────────────────────────
    MID23 = (L2_bot + L3_top) / 2
    _arrow(0.27, L2_bot, 0.27, L3_top, color=COL_FORM)
    _arrow(0.73, L2_bot, 0.73, L3_top, color=COL_CTRL)
    ax.text(0.27, MID23, "SAT pre + post disponible",
            ha="center", va="center", fontsize=7.5, color="#AAAAAA", style="italic")
    ax.text(0.73, MID23, "Con EDD disponible",
            ha="center", va="center", fontsize=7.5, color="#AAAAAA", style="italic")

    # ── Nivel 3: 197 Aptos P3 / ctrl EDD ───────────────────────────────────
    _box(0.27, L3, "197 Aptos P3  (SAT pre + post)",
         color=COL_FORM, w=0.40, h=BHH*2)
    _box(0.73, L3, f"{N_B4_CTRL} Docentes  (Control EDD)",
         color=COL_CTRL, w=0.40, h=BHH*2)

    # ── Flechas L3 → L4 ─────────────────────────────────────────────────────
    MID34 = (L3_bot + L4_top) / 2
    scat_f_n = int(scat["formado"].sum())
    scat_c_n = int((~scat["formado"]).sum())
    _arrow(0.27, L3_bot, 0.27, L4_top, color=COL_FORM)
    _arrow(0.73, L3_bot, 0.73, L4_top, color=COL_CTRL)
    ax.text(0.27, MID34, "Secciones con nota y SAT",
            ha="center", va="center", fontsize=7.5, color="#AAAAAA", style="italic")
    ax.text(0.73, MID34, "Secciones con nota",
            ha="center", va="center", fontsize=7.5, color="#AAAAAA", style="italic")

    # ── Nivel 4: secciones ──────────────────────────────────────────────────
    _box(0.27, L4, f"~{scat_f_n:,} secciones  (formados)",
         color=COL_FORM + "BB", w=0.40, h=BHH*2)
    _box(0.73, L4, f"~{scat_c_n:,} secciones  (control)",
         color=COL_CTRL + "BB", w=0.40, h=BHH*2)

    # ── Nota inferior ─────────────────────────────────────────────────────────
    ax.text(0.50, 0.10,
            "Bloque III: control = docentes sin formación en scatter_sat_notas  ·  "
            "Bloque IV: control = universo 918 que nunca participó en P3, con EDD disponible",
            ha="center", va="center", fontsize=8, color="#C0D0E0",
            style="italic", zorder=5,
            bbox=dict(boxstyle="round,pad=0.35", facecolor="#0D1E30",
                      edgecolor="#3A5A7A", alpha=0.80))

    sl = _new_sl(prs); _pic(sl, SHARED_BG, prs); _pic(sl, _save_ch(fig, "cg_embudo.png"), prs)
    _T(sl, "¿Quién es el Grupo de Control?   Derivación por Bloque de Análisis")
    _POP(sl, "Universo base: 917 docentes jerarquizados (docente_918.csv)  ·  Cada bloque usa una definición distinta de control")
    print("  ✓ Slide 1 — Embudo")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Bloque III: Perfil jerarquía + tramo edad (aprobación)
# ─────────────────────────────────────────────────────────────────────────────
def slide_perfil_b3(prs):
    _ensure_bg()
    fig = _tr_fig()

    # Área PIC disponible — py más bajo y ph más corto para no tapar logo UCEN (y≈0.80)
    px = PIC_RECT[0] + 0.01
    py = PIC_RECT[1] + 0.02   # base: 0.31 desde abajo de figura
    pw = PIC_RECT[2] - 0.02
    ph = 0.42                  # tope del gráfico: 0.31+0.42=0.73 → bajo el logo

    GAP = 0.025  # gap entre los dos pares
    bw = (pw - GAP) / 2  # ancho de cada butterfly

    # Datos formados (197 docentes — sat)
    pct_jer_f = _pct_jer(sat)
    pct_eda_f = _pct_tramo(sat)

    # Datos control Bloque III (docentes en scatter formado=False → doc918)
    pct_jer_c = _pct_jer(ctrl_b3_doc)
    pct_eda_c = _pct_tramo(ctrl_b3_doc)

    # ── Butterfly 1: Jerarquía ──
    _butterfly(fig,
               bx=px, by=py, bw=bw, bh=ph,
               cats=JER_LBL, pct_f=pct_jer_f, pct_c=pct_jer_c,
               n_f_tot=N_B3_FORM, n_c_tot=N_B3_CTRL,
               chart_title="Jerarquía Académica")

    # ── Butterfly 2: Tramo Edad ──
    _butterfly(fig,
               bx=px + bw + GAP, by=py, bw=bw, bh=ph,
               cats=TRAMOS_EDAD, pct_f=pct_eda_f, pct_c=pct_eda_c,
               n_f_tot=N_B3_FORM, n_c_tot=N_B3_CTRL,
               chart_title="Tramo de Edad")

    sl = _new_sl(prs); _pic(sl, SHARED_BG, prs); _pic(sl, _save_ch(fig, "cg_b3_perfil.png"), prs)
    _T(sl, "Bloque III — Perfil Demográfico: Formados vs Control  (Aprobación de Alumnos)", fs=16)
    _POP(sl, f"Control B3: docentes en scatter_sat_notas con formado=False, unidos con doc918  ·  "
             f"formados n={N_B3_FORM}  ·  control n={N_B3_CTRL}  ·  % sobre total de cada grupo")

    # Calcular diferencias notables para bullets
    jer_dif = [(JER_LBL[i], pct_jer_f[i] - pct_jer_c[i]) for i in range(len(JER_LBL))]
    jer_dif.sort(key=lambda x: abs(x[1]), reverse=True)
    top_jer = jer_dif[0]
    eda_dif = [(TRAMOS_EDAD[i], pct_eda_f[i] - pct_eda_c[i]) for i in range(len(TRAMOS_EDAD))]
    eda_dif.sort(key=lambda x: abs(x[1]), reverse=True)
    top_eda = eda_dif[0]

    _BUL(sl, [
        f"Jerarquía: '{top_jer[0]}' muestra la mayor diferencia entre grupos "
        f"({top_jer[1]:+.0f} pp {'más en formados' if top_jer[1]>0 else 'más en control'}). "
        "Los formados tienden a concentrarse en jerarquías de rango medio.",
        f"Tramo de edad: '{top_eda[0]}' es el tramo con mayor diferencia "
        f"({top_eda[1]:+.0f} pp {'más en formados' if top_eda[1]>0 else 'más en control'}). "
        "Comprender el perfil permite separar efectos de la formación de diferencias de composición.",
        "Las barras muestran % dentro de cada grupo — no volúmenes absolutos — para hacer "
        "comparables grupos de diferente tamaño.",
    ])
    print("  ✓ Slide 2 — Perfil Bloque III")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Bloque IV: Perfil jerarquía + tramo edad (EDD)
# ─────────────────────────────────────────────────────────────────────────────
def slide_perfil_b4(prs):
    _ensure_bg()
    if not has_edd or len(b4_form_doc) == 0:
        sl = _new_sl(prs); _pic(sl, SHARED_BG, prs)
        _T(sl, "Bloque IV — Perfil Demográfico (datos EDD no disponibles)")
        print("  ✓ Slide 3 — Perfil B4 (sin datos)")
        return

    fig = _tr_fig()

    px = PIC_RECT[0] + 0.01
    py = PIC_RECT[1] + 0.02   # mismo ajuste que slide 2: no tapar logo UCEN
    pw = PIC_RECT[2] - 0.02
    ph = 0.42

    GAP = 0.025
    bw = (pw - GAP) / 2

    pct_jer_f = _pct_jer(b4_form_doc)
    pct_eda_f = _pct_tramo(b4_form_doc)
    pct_jer_c = _pct_jer(b4_ctrl_doc)
    pct_eda_c = _pct_tramo(b4_ctrl_doc)

    _butterfly(fig,
               bx=px, by=py, bw=bw, bh=ph,
               cats=JER_LBL, pct_f=pct_jer_f, pct_c=pct_jer_c,
               n_f_tot=N_B4_FORM, n_c_tot=N_B4_CTRL,
               chart_title="Jerarquía Académica")

    _butterfly(fig,
               bx=px + bw + GAP, by=py, bw=bw, bh=ph,
               cats=TRAMOS_EDAD, pct_f=pct_eda_f, pct_c=pct_eda_c,
               n_f_tot=N_B4_FORM, n_c_tot=N_B4_CTRL,
               chart_title="Tramo de Edad")

    sl = _new_sl(prs); _pic(sl, SHARED_BG, prs); _pic(sl, _save_ch(fig, "cg_b4_perfil.png"), prs)
    _T(sl, "Bloque IV — Perfil Demográfico: Formados vs Control  (EDD Evaluación de Jefes)", fs=16)
    _POP(sl, f"Control B4: docentes en universo 918 que nunca participaron en P3, con EDD disponible  ·  "
             f"formados n={N_B4_FORM}  ·  control n={N_B4_CTRL}  ·  % sobre total de cada grupo")

    jer_dif = [(JER_LBL[i], pct_jer_f[i] - pct_jer_c[i]) for i in range(len(JER_LBL))]
    jer_dif.sort(key=lambda x: abs(x[1]), reverse=True)
    top_jer = jer_dif[0]
    eda_dif = [(TRAMOS_EDAD[i], pct_eda_f[i] - pct_eda_c[i]) for i in range(len(TRAMOS_EDAD))]
    eda_dif.sort(key=lambda x: abs(x[1]), reverse=True)
    top_eda = eda_dif[0]

    _BUL(sl, [
        f"Jerarquía: '{top_jer[0]}' muestra la mayor diferencia entre grupos EDD "
        f"({top_jer[1]:+.0f} pp {'más en formados' if top_jer[1]>0 else 'más en control'}). "
        "Si los grupos tienen composiciones muy distintas, la brecha de EDD puede reflejar "
        "diferencias de jerarquía, no solo de formación.",
        f"Tramo de edad: '{top_eda[0]}' es el tramo más diferenciado "
        f"({top_eda[1]:+.0f} pp {'más en formados' if top_eda[1]>0 else 'más en control'}). "
        "El control EDD es más estricto — excluye cualquier participante en P3 aunque no sea Apto.",
        "Comparar este perfil con el de Bloque III permite ver si la definición más estricta "
        "del control EDD cambia la composición demográfica del grupo de comparación.",
    ])
    print("  ✓ Slide 3 — Perfil Bloque IV")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Tabla resumen (fuente original + derivada)
# ─────────────────────────────────────────────────────────────────────────────
def slide_tabla(prs):
    _ensure_bg()
    fig = _tr_fig()

    ax = fig.add_axes([PIC_RECT[0] + 0.01, PIC_RECT[1] + 0.01,
                       PIC_RECT[2] - 0.02,  PIC_RECT[3] - 0.02],
                      facecolor="none", zorder=5)
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis("off")

    # Columnas: Bloque | Formados | Control | Fuente ORIGINAL | Derivada | n
    cols_x  = [0.000, 0.090, 0.250, 0.420, 0.620, 0.840]
    cols_w  = [0.088, 0.158, 0.168, 0.198, 0.218, 0.158]
    headers = ["Bloque", "Grupo Formados", "Grupo Control",
               "Fuente Original", "Archivo Derivado", "n Form. / Control"]

    # Header
    hy = 0.90
    hp = mpatches.FancyBboxPatch((0, hy - 0.055), 1.0, 0.10,
                                  boxstyle="round,pad=0.005", linewidth=0,
                                  facecolor="#1A3A5C", alpha=0.92, zorder=3,
                                  transform=ax.transData)
    ax.add_patch(hp)
    for hdr, cx in zip(headers, cols_x):
        ax.text(cx + 0.006, hy, hdr, ha="left", va="center",
                fontsize=8.5, color="white", fontweight="bold", zorder=5)

    rows = [
        ("SAT\nBloques\nI – II",
         "197 Aptos P3\n(SAT pre+post\ndisponibles)",
         "486 doc sin\nformación\ncon SAT",
         "evaluacion_periodo.csv\n+ nomina_docente.csv",
         "(p3_sat_zscore_918.csv)\n(control_918.csv)",
         "197 / 486\ndocentes"),
        ("Aprobación\nalumnos\nBloque III",
         "Secciones de\nlos 197 Aptos P3\ncon nota y SAT",
         "Secciones de\ndoc sin formación\ncon nota y SAT",
         "calificacion_alumno.csv\n+ evaluacion_periodo.csv\n+ participacion_formacion.csv",
         "(scatter_sat_notas.csv)\nformado = True / False",
         f"~{len(scat):,} filas\n({N_B3_FORM}/{N_B3_CTRL} doc)"),
        ("EDD\nBloque IV",
         f"197 Aptos P3\ncon EDD\ndisponible",
         f"{N_B4_CTRL} doc del\nuniverso 918 que\nnunca participaron en P3",
         "CONSOLIDADO DOCENTES.xlsx\n(hoja: EVALUACION DE\nJEFES A DOCENTES)",
         "(P1_consolidado_con_\nevaluacion_jefes.csv)",
         f"{N_B4_FORM} / {N_B4_CTRL}\ndocentes\n(varía por año)"),
    ]

    row_ys  = [0.685, 0.440, 0.175]
    row_hs  = [0.20,  0.22,  0.22]
    row_cls = ["#0D2035", "#0A1A2E", "#0D2035"]

    for (blq, form, ctrl, orig, deriv, ns), ry, rh, rc in zip(rows, row_ys, row_hs, row_cls):
        rp = mpatches.FancyBboxPatch((0, ry - rh/2), 1.0, rh + 0.01,
                                      boxstyle="round,pad=0.005", linewidth=0.7,
                                      edgecolor="#2A4A6A", facecolor=rc, alpha=0.80,
                                      zorder=3, transform=ax.transData)
        ax.add_patch(rp)
        # Fuente original en blanco, derivada en color más tenue
        for txt, cx, col in zip([blq, form, ctrl, orig, deriv, ns],
                                  cols_x,
                                  ["white"]*3 + ["#E0E8F0", "#90C8F0", "#E0E8F0"]):
            ax.text(cx + 0.006, ry, txt, ha="left", va="center",
                    fontsize=7.5, color=col, linespacing=1.3, zorder=5)
        # Separador vertical
        ax.axvline(cols_x[4] - 0.005, ymin=ry - rh/2, ymax=ry + rh/2,
                   color="#3A5A7A", lw=0.6, alpha=0.5, zorder=3)

    # Nota al pie sobre fuente derivada
    ax.text(0.50, 0.02,
            "Fuente Original = archivo entregado por la institución  ·  "
            "Archivo Derivado (en azul) = generado por ETL del proyecto a partir de la fuente original",
            ha="center", va="bottom", fontsize=7.5, color="#7090B0", style="italic", zorder=5)

    sl = _new_sl(prs); _pic(sl, SHARED_BG, prs); _pic(sl, _save_ch(fig, "cg_tabla.png"), prs)
    _T(sl, "Resumen de Grupos, Fuentes Originales y Archivos Derivados")
    _POP(sl, "Blanco = fuente original entregada por UCEN  ·  Azul = archivo procesado/derivado por ETL del análisis")
    _BUL(sl, [
        "Bloque I–II (SAT): fuente madre es evaluacion_periodo.csv; los z-scores se calculan "
        "por facultad×período y se almacenan en p3_sat_zscore_918.csv (formados) y control_918.csv.",
        "Bloque III (Aprobación): scatter_sat_notas.csv es un archivo ad-hoc que cruza notas, "
        "SAT y el flag de formación — sin ETL documentado; reconstruible desde las tres fuentes originales.",
        "Bloque IV (EDD): la fuente original es la hoja 'EVALUACION DE JEFES A DOCENTES' del Excel "
        "institucional; el ETL (etl_consolidado_con_jefes.py) la renombra y consolida con la nómina.",
    ])
    print("  ✓ Slide 4 — Tabla resumen")


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    _ensure_bg()
    prs = Presentation()
    prs.slide_width  = Emu(SW_EMU)
    prs.slide_height = Emu(SH_EMU)

    print("Generando slides...")
    slide_embudo(prs)
    slide_perfil_b3(prs)
    slide_perfil_b4(prs)
    slide_tabla(prs)

    prs.save(OUT_PPTX)
    print(f"\n✅  {OUT_PPTX}")
    print(f"   {len(prs.slides)} slides")

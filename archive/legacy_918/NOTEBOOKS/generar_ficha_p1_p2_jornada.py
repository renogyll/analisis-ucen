# -*- coding: utf-8 -*-
"""
generar_ficha_p1_p2_jornada.py  v2
Genera FICHA_Ampliar_P1_P2_Jornada.pptx en formato dark v3.
UNIVERSO: solo docentes Jornada / Planta (excluye Honorarios).
"""
import sys; sys.stdout.reconfigure(encoding="utf-8")
import os, zipfile
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.patheffects as pe
import pandas as pd
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─────────────────────────────────────────────────────────────────────────────
# Rutas
# ─────────────────────────────────────────────────────────────────────────────
BASE      = os.path.dirname(os.path.abspath(__file__))
PROC918   = os.path.join(BASE, "..", "PROCESADO")
ROOT      = os.path.normpath(os.path.join(BASE, "..", ".."))
FONDOTIPO = os.path.join(ROOT, "Fondotipop.pptx")
OUT_DIR   = os.path.join(BASE, "dark_slides_v3")
OUT_PPTX  = os.path.join(ROOT, "FICHA_Ampliar_P1_P2_Jornada_v4.pptx")
os.makedirs(OUT_DIR, exist_ok=True)

# ─────────────────────────────────────────────────────────────────────────────
# Datos
# ─────────────────────────────────────────────────────────────────────────────
doc_all = pd.read_csv(os.path.join(PROC918, "docente_918.csv"), encoding="utf-8-sig")
p2_all  = pd.read_csv(os.path.join(PROC918, "participacion_p2_918.csv"), encoding="utf-8-sig")

jor_mask = doc_all["tipo_contrato"].str.strip().str.lower().str.startswith("jornada")
doc      = doc_all[jor_mask].copy().reset_index(drop=True)
jor_ruts = set(doc["rut_key"])
p2       = p2_all[p2_all["rut_key"].isin(jor_ruts)].copy()

N_TOTAL    = len(doc)                                          # 545
N_CON_JER  = int(doc["jerarquia"].notna().sum())              # 545
dot_mask   = doc["fuente"].isin(["NOMINA_DOTACION", "SOLO_DOTACION"])
jor_dot    = doc[dot_mask].copy()
N_DOTACION = len(jor_dot)                                     # 485
N_SIN_DOT  = int((~dot_mask).sum())                          # 60

def _ni(ser):
    return ser.notna() & ~ser.isin({"NO INFORMA", "NO INFORMA "})

m_ant  = jor_dot["fecha_ingreso"].notna()
m_jer  = m_ant  & jor_dot["fecha_nacimiento"].notna() & _ni(jor_dot["jerarquia_dot"])
m_niv  = m_jer  & _ni(jor_dot["nivel_formacion"])
m_grd  = m_niv  & _ni(jor_dot["nombre_grado"])

c0    = N_DOTACION
c_jer = int(m_jer.sum())    # 474
c_niv = int(m_niv.sum())    # 469
c_grd = int(m_grd.sum())    # 469
c_fec = c_grd - int(round(45 * c_grd / 967))  # ~447

D_JER = c0 - c_jer
D_NIV = c_jer - c_niv
D_GRD = c_niv - c_grd
D_FEC = c_grd - c_fec

SUBCASC = [
    (c0,   "Total dotacion",     "Edad · Sexo · Facultad · Antiguedad · Carga"),
    (c_jer,"Jerarquia clasif.",  "Edad x Jerarquia"),
    (c_niv,"Nivel formacion",    "Nivel Formacion · Institucion · Pais"),
    (c_grd,"Grado clasificado",  "GRADOREC · Jerarquia x Nivel Formacion"),
    (c_fec,"Fechas completas*",  "Anos hasta Jerarquizacion  (*estimado)"),
]

N_FORMADOS    = int(p2["rut_key"].nunique())
N_INICIATIVAS = len(p2)
N_TALL        = int((p2["tipo_formacion"] == "TALLER").sum())
N_DIP         = int((p2["tipo_formacion"] == "DIPLOMADO").sum())
N_PROY        = int((p2["tipo_formacion"] == "PROYECTO").sum())
N_SIN_FORM    = N_TOTAL - N_FORMADOS

F_DOC = "CONSOLIDADO DOCENTES 3-05-2026.xlsx"

print(f"Jornada: {N_TOTAL} | {N_DOTACION} dotacion | {N_SIN_DOT} sin dotacion")
print(f"Sub-cascadas: {c0}→{c_jer}→{c_niv}→{c_grd}→{c_fec}")
print(f"P2: {N_FORMADOS} formados | {N_INICIATIVAS} init | T={N_TALL} D={N_DIP} P={N_PROY}")

# ─────────────────────────────────────────────────────────────────────────────
# Assets
# ─────────────────────────────────────────────────────────────────────────────
BG_PATH   = os.path.join(OUT_DIR, "_fondotipo_bg.jpg")
LOGO_PATH = os.path.join(OUT_DIR, "_fondotipo_logo.png")
for path, zname in [(BG_PATH, "ppt/media/image1.jpg"),
                    (LOGO_PATH, "ppt/media/image2.png")]:
    if not os.path.exists(path):
        with zipfile.ZipFile(FONDOTIPO) as z:
            with open(path, "wb") as f:
                f.write(z.read(zname))

with PILImage.open(BG_PATH) as _im:
    _rgb = _im.convert("RGB"); _iw, _ih = _rgb.size
    _nh  = int(_iw / (16/9)); _y0 = min(int(_ih * 0.12), _ih - _nh)
    bg_arr = np.array(_rgb.crop((0, _y0, _iw, _y0 + _nh)))

with PILImage.open(LOGO_PATH) as _logo:
    logo_arr = np.array(_logo.convert("RGBA")).astype(np.float32) / 255.0

H_GRAD = 600; grad = np.zeros((H_GRAD, 1, 4), dtype=np.float32)
for _r in range(H_GRAD):
    _t = _r / (H_GRAD - 1)
    _stops = [(0.00, (0, 33, 71)), (0.54, (0, 70, 128)), (1.00, (144, 171, 196))]
    for _i in range(len(_stops) - 1):
        _t0, _c0 = _stops[_i]; _t1, _c1 = _stops[_i + 1]
        if _t0 <= _t <= _t1:
            _s = (_t - _t0) / (_t1 - _t0)
            grad[_r, 0] = [(_c0[0]+_s*(_c1[0]-_c0[0]))/255,
                           (_c0[1]+_s*(_c1[1]-_c0[1]))/255,
                           (_c0[2]+_s*(_c1[2]-_c0[2]))/255, 0.82]; break

# ─────────────────────────────────────────────────────────────────────────────
# Layout constants
# ─────────────────────────────────────────────────────────────────────────────
SW, SH         = 13.333, 7.5
SW_EMU, SH_EMU = 12192000, 6858000
PIC_L, PIC_T, PIC_W, PIC_H         = 786581, 1125000, 10599174, 3720000
BUL_L, BUL_T, BUL_W, BUL_H         = 786581, 4870000, 10599174, 1870000
LOGO_L, LOGO_T, LOGO_W, LOGO_H     = 9813773, 656354, 1756626, 697725
TITLE_L, TITLE_T, TITLE_W, TITLE_H = PIC_L, 185000, PIC_W, 710000
POP_L,   POP_T,   POP_W,   POP_H   = PIC_L, 845000, 9000000, 255000

def _ex(e): return e / SW_EMU
def _ey(e): return e / SH_EMU
def _fig_rect(l, t, w, h): return (l, 1-t-h, w, h)
PIC_RECT  = _fig_rect(_ex(PIC_L), _ey(PIC_T), _ex(PIC_W), _ey(PIC_H))
LOGO_RECT = _fig_rect(_ex(LOGO_L), _ey(LOGO_T), _ex(LOGO_W), _ey(LOGO_H))

# Logo threshold in axes y-coords (content above this covers the logo)
# Logo fig_y_bottom = 1 - _ey(LOGO_T) - _ey(LOGO_H) = ~0.8025
# With standard PIC_RECT (gy~0.294, gh~0.532): threshold = (0.8025-0.294)/0.532*10 = 9.56
LOGO_Y_THRESH = 9.30  # keep all axes content strictly below this

# ─────────────────────────────────────────────────────────────────────────────
# Matplotlib helpers
# ─────────────────────────────────────────────────────────────────────────────
def _bg_fig():
    fig = plt.figure(figsize=(SW, SH), facecolor="#101820")
    fig.patch.set_facecolor("#101820")
    for z, arr in [(0, bg_arr), (1, grad)]:
        ax = fig.add_axes([0, 0, 1, 1], zorder=z)
        ax.imshow(arr, extent=[0, 1, 0, 1], aspect="auto", origin="upper")
        ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis("off")
    al = fig.add_axes([LOGO_RECT[0], LOGO_RECT[1], LOGO_RECT[2], LOGO_RECT[3]],
                      zorder=10, facecolor="none")
    al.imshow(logo_arr, aspect="auto"); al.axis("off"); al.patch.set_visible(False)
    return fig

def _tr_fig():
    fig = plt.figure(figsize=(SW, SH), facecolor="none")
    fig.patch.set_facecolor("none"); return fig

SHARED_BG = os.path.join(OUT_DIR, "_background_jor.png")

def _ensure_bg():
    if not os.path.exists(SHARED_BG):
        fig = _bg_fig()
        plt.savefig(SHARED_BG, dpi=150, facecolor=fig.get_facecolor()); plt.close()

def _save_ch(fig, name):
    path = os.path.join(OUT_DIR, name)
    plt.savefig(path, dpi=150, facecolor="none", transparent=True)
    plt.close(); return path

# ─────────────────────────────────────────────────────────────────────────────
# pptx helpers
# ─────────────────────────────────────────────────────────────────────────────
def _new_sl(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _pic(sl, path, prs):
    sl.shapes.add_picture(path, Emu(0), Emu(0), prs.slide_width, prs.slide_height)

def _txt(sl, text, left, top, width, height,
         fs=12, bold=False, italic=False, color="#FFFFFF",
         align=PP_ALIGN.LEFT, wrap=True):
    txb = sl.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf  = txb.text_frame; tf.word_wrap = wrap
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run(); run.text = line
        run.font.size = Pt(fs); run.font.bold = bold; run.font.italic = italic
        r, g, b = int(color[1:3],16), int(color[3:5],16), int(color[5:7],16)
        run.font.color.rgb = RGBColor(r, g, b)

def _T(sl, text, fs=20):
    _txt(sl, text, TITLE_L, TITLE_T, TITLE_W, TITLE_H,
         fs=fs, bold=True, color="#FFFFFF", align=PP_ALIGN.CENTER)

def _POP(sl, text):
    _txt(sl, text, POP_L, POP_T, POP_W, POP_H,
         fs=7.5, italic=True, color="#C8DCF0")

def _BUL(sl, items, fs=12):
    txb = sl.shapes.add_textbox(Emu(BUL_L), Emu(BUL_T), Emu(BUL_W), Emu(BUL_H))
    tf  = txb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after  = Pt(12)
        p.space_before = Pt(2)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run(); run.text = "[>]  " + item
        run.font.size = Pt(fs)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Portada
# ─────────────────────────────────────────────────────────────────────────────
def slide_portada(prs):
    sl = _new_sl(prs)
    _pic(sl, SHARED_BG, prs)
    _txt(sl, "Ficha Tecnica  |  Jornada / Planta",
         TITLE_L, 1_700_000, TITLE_W, 600_000,
         fs=28, bold=True, color="#FFFFFF", align=PP_ALIGN.CENTER)
    _txt(sl, "Ampliar los Datos de P1 y P2",
         TITLE_L, 2_350_000, TITLE_W, 700_000,
         fs=26, bold=True, color="#7EC8E3", align=PP_ALIGN.CENTER)
    _txt(sl,
         "Cascadas de datos, brechas de cobertura y solicitudes — universo Jornada/Planta",
         TITLE_L, 3_150_000, TITLE_W, 500_000,
         fs=12, italic=True, color="#C8DCF0", align=PP_ALIGN.CENTER)
    _txt(sl,
         f"P1 — Caracterizacion: {N_TOTAL} Jornada  |  {N_DOTACION} con perfil completo  |  {N_SIN_DOT} sin dotacion",
         TITLE_L, 3_850_000, TITLE_W, 380_000,
         fs=11.5, color="#FFD580", align=PP_ALIGN.CENTER)
    _txt(sl,
         f"P2 — Formacion: {N_FORMADOS} Jornada formados de {N_TOTAL}  "
         f"({N_INICIATIVAS} iniciativas: {N_TALL} Talleres · {N_DIP} Diplomados · {N_PROY} Proyectos)",
         TITLE_L, 4_300_000, TITLE_W, 380_000,
         fs=11.5, color="#A8E8A0", align=PP_ALIGN.CENTER)
    _txt(sl,
         f"Archivo fuente: {F_DOC}  (hojas: NOMINA, DOTACION, DIPLOMADO, TALLERES, PROYECTOS)",
         TITLE_L, 4_800_000, TITLE_W, 350_000,
         fs=9.5, italic=True, color="#A0B8D0", align=PP_ALIGN.CENTER)
    print("  slide 1 - Portada OK")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Cascada P1
# ─────────────────────────────────────────────────────────────────────────────
def slide_cascada_p1(prs):
    fig = _tr_fig()
    gx = PIC_RECT[0] + 0.01
    gy = PIC_RECT[1] + 0.00
    gw = PIC_RECT[2] - 0.02
    gh = PIC_RECT[3] - 0.01
    ax = fig.add_axes([gx, gy, gw, gh], facecolor="none", zorder=5)
    ax.set_xlim(0, 10); ax.set_ylim(0, 10); ax.axis("off")

    BW, BH = 3.90, 1.50
    BX = 0.15

    def _box(ax, x0, y0, w, h, n, label, line1, line2, col, fs_n=22):
        ax.add_patch(mpatches.FancyBboxPatch(
            (x0, y0), w, h, boxstyle="round,pad=0.10",
            facecolor=col, edgecolor="white", linewidth=0.8, alpha=0.90, zorder=3))
        ax.text(x0+0.60, y0+h*0.60, n,
                ha="center", va="center", fontsize=fs_n, fontweight="bold",
                color="white", zorder=5,
                path_effects=[pe.withStroke(linewidth=3, foreground="#050D1A")])
        ax.text(x0+1.35, y0+h*0.75, label,
                ha="left", va="center", fontsize=9.5, fontweight="bold",
                color="white", zorder=5)
        ax.text(x0+1.35, y0+h*0.48, line1,
                ha="left", va="center", fontsize=7.0, color="#B8D0E8", zorder=5)
        if line2:
            ax.text(x0+1.35, y0+h*0.24, line2,
                    ha="left", va="center", fontsize=7.0, color="#B8D0E8", zorder=5)

    def _arrow(ax, x, y_from, y_to, lbl, col="#FFD580"):
        mid = (y_from + y_to) / 2
        ax.annotate("", xy=(x, y_to+0.05), xytext=(x, y_from-0.05),
                    arrowprops=dict(arrowstyle="->", color=col, lw=1.4,
                                   mutation_scale=16), zorder=6)
        ax.text(x-0.10, mid, lbl, ha="right", va="center",
                fontsize=7.5, color=col, fontstyle="italic", zorder=5)

    # Node 1: 545 (y0=7.50, top at 9.00 — below logo threshold 9.30)
    _box(ax, BX, 7.50, BW, BH, str(N_TOTAL),
         "docentes Jornada jerarquizados",
         "Hoja NOMINA",
         f"({F_DOC})", "#1B4E8F")

    # "100% con jerarquia" — pequeña nota verde en la parte baja de Node 1
    ax.text(BX+0.28, 7.50+0.22,
            "[OK]  100% con jerarquia informada",
            ha="left", va="center", fontsize=6.8, color="#80FF80",
            fontstyle="italic", zorder=6)

    # Node 2: 485 (y0=5.20)
    _box(ax, BX, 5.20, BW, BH, str(N_DOTACION),
         "Jornada con perfil completo",
         "Hoja DOTACION",
         f"({F_DOC})", "#1B7A4A", fs_n=22)

    _arrow(ax, BX+BW/2-0.2, 7.50, 5.20+BH,
           f"-{N_SIN_DOT} sin dotacion", col="#E05A3A")

    ax.text(BX+BW/2-0.2-0.15, 6.65,
            f"  {N_SIN_DOT} solo en NOMINA  ",
            ha="right", va="center", fontsize=8.0, fontweight="bold",
            color="#E05A3A", zorder=5,
            bbox=dict(boxstyle="round,pad=0.3", facecolor="#2A0808",
                      edgecolor="#E05A3A", alpha=0.85))

    # Sub-cascade — calcular posiciones primero para saber hasta donde llega la spine
    midbox  = 5.20 + BH/2   # punto medio Node 2
    y_start = 8.70           # centro caja 1 (tope=9.225 < 9.30 ✓)
    spine_x = BX+BW+0.35    # columna vertical del arbol
    sx0     = BX+BW+0.80    # inicio de las cajas verdes
    sh      = 1.05
    y_gap   = 1.38
    sub_ys  = [y_start - i*y_gap for i in range(len(SUBCASC))]

    # Linea horizontal desde Node 2 hasta la columna vertical
    ax.plot([BX+BW, spine_x], [midbox, midbox],
            "-", color="white", linewidth=2.2, alpha=0.72, zorder=4)
    # Columna vertical desde la caja mas baja hasta la mas alta
    ax.plot([spine_x, spine_x], [sub_ys[-1], y_start],
            "-", color="white", linewidth=2.2, alpha=0.72, zorder=4)
    sub_colors = ["#1B7A4A","#27824A","#348A52","#42925A","#529A62"]

    for (n, lbl, analisis), y, col in zip(SUBCASC, sub_ys, sub_colors):
        ax.add_patch(mpatches.FancyBboxPatch(
            (sx0, y-sh/2), 5.00, sh,
            boxstyle="round,pad=0.07",
            facecolor=col, edgecolor="white",
            linewidth=0.5, alpha=0.82, zorder=3))
        ax.text(sx0+0.50, y+0.10, str(n),
                ha="center", va="center", fontsize=16, fontweight="bold",
                color="white", zorder=5,
                path_effects=[pe.withStroke(linewidth=2.5, foreground="#050D1A")])
        ax.text(sx0+1.00, y+0.24, lbl,
                ha="left", va="center", fontsize=8.5, fontweight="bold",
                color="white", zorder=5)
        ax.text(sx0+1.00, y-0.18, f"-> {analisis}",
                ha="left", va="center", fontsize=7.5,
                color="#C8E8C8", zorder=5)
        # Rama horizontal + punto en la union
        ax.plot([spine_x, sx0-0.02], [y, y],
                "-", color="white", linewidth=1.8, alpha=0.65, zorder=4)
        ax.plot([spine_x], [y], "o", color="white",
                markersize=5, alpha=0.85, zorder=5)

    path = _save_ch(fig, "jor_p1_cascada.png")
    sl = _new_sl(prs)
    _pic(sl, SHARED_BG, prs)
    _pic(sl, path, prs)
    _T(sl, "P1 — Cascada de Datos Jornada/Planta: 545 docentes a 485 con perfil completo", fs=12)
    _POP(sl, f"Fuente: {F_DOC}  |  "
             f"Hoja NOMINA ({N_TOTAL} Jornada) cruzada con hoja DOTACION ({N_DOTACION} con perfil)")
    _BUL(sl, [
        f"Los {N_TOTAL} docentes Jornada/planta tienen 100% con jerarquia informada. "
        f"El salto critico es la dotacion: {N_SIN_DOT} Jornada estan en NOMINA pero NO en DOTACION, "
        f"quedando sin edad, antiguedad, nivel de formacion y carga academica.",
        f"Dentro de los {N_DOTACION} con dotacion las sub-cascadas son brechas menores: "
        f"{D_JER} con jerarquia 'NO INFORMA', {D_NIV} con nivel formacion 'NO INFORMA', "
        f"y ~{D_FEC} sin fecha de jerarquizacion (estimado desde hoja NOMINA).",
    ])
    print("  slide 2 - Cascada P1 OK")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Brecha: 60 Jornada sin dotacion
# ─────────────────────────────────────────────────────────────────────────────
def slide_brecha_dotacion(prs):
    fig = _tr_fig()
    gx = PIC_RECT[0] + 0.01
    gy = PIC_RECT[1] + 0.00
    gw = PIC_RECT[2] - 0.02
    gh = PIC_RECT[3] - 0.01
    ax = fig.add_axes([gx, gy, gw, gh], facecolor="none", zorder=5)
    ax.set_xlim(0, 10); ax.set_ylim(0, 10); ax.axis("off")

    PH = 8.70  # panel height — y0=0.15 → top at 8.85 (≪ 9.30) ✓

    def _panel(ax, x0, y0, w, h, n, title, col, items_ok, items_no=None):
        ax.add_patch(mpatches.FancyBboxPatch(
            (x0, y0), w, h, boxstyle="round,pad=0.12",
            facecolor=col, edgecolor="white", linewidth=0.8, alpha=0.28, zorder=2))
        ax.add_patch(mpatches.FancyBboxPatch(
            (x0, y0+h-0.62), w, 0.62,
            boxstyle="round,pad=0.05",
            facecolor=col, edgecolor="none", alpha=0.88, zorder=3))
        ax.text(x0+w/2, y0+h-0.31, f"{n}  {title}",
                ha="center", va="center", fontsize=11, fontweight="bold",
                color="white", zorder=5)
        y_item = y0+h-0.88
        for it in items_ok:
            ax.text(x0+0.22, y_item, it,
                    ha="left", va="top", fontsize=8.0, color="white",
                    zorder=5, linespacing=1.3)
            y_item -= 0.56
        if items_no:
            ax.text(x0+0.22, y_item-0.08, "NO disponible:",
                    ha="left", va="top", fontsize=8.0, color="#E08080",
                    fontweight="bold", zorder=5)
            y_item -= 0.50
            for it in items_no:
                ax.text(x0+0.22, y_item, it,
                        ha="left", va="top", fontsize=8.0, color="#E08080", zorder=5)
                y_item -= 0.52

    ok_vars_dot = [
        "[OK]  Sexo",
        "[OK]  Jerarquia academica",
        "[OK]  Tipo contrato (Jornada)",
        "[OK]  Edad / fecha nacimiento",
        "[OK]  Antiguedad / fecha ingreso",
        "[OK]  Unidad / Facultad",
        "[OK]  Nivel formacion (grado)",
        "[OK]  Institucion y pais grado",
        "[OK]  Carga horaria (jornada hrs)",
        "[OK]  Cargo especifico",
    ]
    ok_vars_nom = [
        "[OK]  Sexo",
        "[OK]  Jerarquia academica",
        "[OK]  Tipo contrato (Jornada)",
        "[OK]  Funcion principal academica",
        "[OK]  Fecha jerarquizacion (~95%)",
    ]
    no_vars_nom = [
        "[---]  Edad / fecha nacimiento",
        "[---]  Antiguedad / fecha ingreso",
        "[---]  Unidad / Facultad detallada",
        "[---]  Nivel de formacion",
        "[---]  Nombre grado / titulo",
        "[---]  Institucion y pais grado",
        "[---]  Carga horaria semanal",
        "[---]  Cargo especifico",
    ]

    _panel(ax, 0.08, 0.15, 5.40, PH, N_DOTACION,
           "Jornada con DOTACION", "#1B7A4A", ok_vars_dot)

    ax.axvline(5.80, ymin=0.02, ymax=0.95,
               color="white", linewidth=0.6, alpha=0.30, zorder=4)

    _panel(ax, 5.95, 0.15, 3.80, PH, N_SIN_DOT,
           "Jornada SOLO en NOMINA", "#8B2020", ok_vars_nom, no_vars_nom)

    ax.text(5.98, 0.05,
            "Posible causa: nuevos ingresos 2026 o discrepancia de RUT",
            ha="left", va="center", fontsize=8.0, color="#E08080",
            fontstyle="italic", zorder=5)

    path = _save_ch(fig, "jor_p1_brecha.png")
    sl = _new_sl(prs)
    _pic(sl, SHARED_BG, prs)
    _pic(sl, path, prs)
    _T(sl, f"Brecha de Datos — {N_SIN_DOT} Jornada Sin Dotacion (de {N_TOTAL} en total)", fs=14)
    _POP(sl, f"NOMINA ({F_DOC}): {N_TOTAL} Jornada  |  DOTACION: {N_DOTACION}  |  "
             f"Sin dotacion: {N_SIN_DOT} ({100*N_SIN_DOT/N_TOTAL:.1f}% del universo)")
    _BUL(sl, [
        f"A diferencia del universo total (brecha de 424, mayoritariamente Honorarios), "
        f"en Jornada/planta la brecha es solo de {N_SIN_DOT} docentes sin dotacion: "
        f"posiblemente nuevos ingresos 2026 o errores de cruce por RUT entre fuentes.",
        f"Solicitar la dotacion de estos {N_SIN_DOT} Jornada es factible y recuperaria "
        f"el perfil completo (edad, antiguedad, nivel de formacion) elevando la "
        f"cobertura de {N_DOTACION} a {N_TOTAL} Jornada.",
    ])
    print("  slide 3 - Brecha dotacion OK")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Sub-cascadas dentro de los 485
# ─────────────────────────────────────────────────────────────────────────────
def slide_subcascadas(prs):
    fig = _tr_fig()
    gx = PIC_RECT[0] + 0.12
    gy = PIC_RECT[1] + 0.02
    gw = PIC_RECT[2] * 0.50   # compact so ax2 gets ~0.34 figure-units of width
    gh = PIC_RECT[3] - 0.04
    ax = fig.add_axes([gx, gy, gw, gh], facecolor="none", zorder=5)

    ns     = [sc[0] for sc in SUBCASC]
    labels = [sc[1] for sc in SUBCASC]
    anals  = [sc[2] for sc in SUBCASC]
    cols   = ["#1B7A4A","#27924A","#35A452","#44B660","#4A9040"]

    bars = ax.barh(range(len(ns)), ns, color=cols, alpha=0.88,
                   edgecolor="white", linewidth=0.6, height=0.58, zorder=3)

    for bar, n, lbl, col in zip(bars, ns, labels, cols):
        ax.text(n + 1, bar.get_y() + bar.get_height()/2,
                str(n), va="center", ha="left", fontsize=16, fontweight="bold",
                color="white",
                path_effects=[pe.withStroke(linewidth=2, foreground="#050D1A")])
        diff = c0 - n
        if diff > 0:
            ax.text(n - 1, bar.get_y() + bar.get_height()/2,
                    f"-{diff}", va="center", ha="right", fontsize=9,
                    color="#FFD580", fontweight="bold",
                    path_effects=[pe.withStroke(linewidth=1.5, foreground="#050D1A")])

    x_min = min(ns) - 12
    x_max = c0 + 18
    ax.set_xlim(x_min, x_max)
    ax.set_xlabel("N docentes Jornada disponibles para cada analisis",
                  color="#AAAAAA", fontsize=9)
    ax.set_yticks(range(len(ns)))
    ax.set_yticklabels(labels, fontsize=10, color="white")
    ax.tick_params(axis="x", colors="#AAAAAA", labelsize=8.5)
    for sp in ax.spines.values():
        sp.set_edgecolor("white"); sp.set_alpha(0.25); sp.set_linewidth(0.7)
    ax.xaxis.grid(True, color="white", alpha=0.08, linewidth=0.5)
    ax.set_axisbelow(True)

    # Panel derecho: analisis habilitados
    ax2 = fig.add_axes([gx+gw+0.012, gy+gh*0.04, 0.34, gh*0.93],
                       facecolor="none", zorder=5)
    ax2.axis("off")
    ax2.text(0.0, 1.0, "ANALISIS HABILITADO",
             transform=ax2.transAxes, va="top", ha="left",
             fontsize=9.5, fontweight="bold", color="#FFD580", zorder=5)

    y_pos = 0.85
    for n, lbl, ana, col in zip(ns, labels, anals, cols):
        circ = mpatches.Circle((0.05, y_pos), 0.025,
                                facecolor=col, edgecolor="white",
                                linewidth=0.5, alpha=0.85,
                                transform=ax2.transAxes, zorder=3)
        ax2.add_patch(circ)
        ax2.text(0.12, y_pos, ana,
                 transform=ax2.transAxes, va="center", ha="left",
                 fontsize=7.5, color="white", zorder=5)
        y_pos -= 0.175

    path = _save_ch(fig, "jor_p1_subcascadas.png")
    sl = _new_sl(prs)
    _pic(sl, SHARED_BG, prs)
    _pic(sl, path, prs)
    _T(sl, f"P1 — Calidad de Datos Dentro de los {N_DOTACION} Jornada con Dotacion", fs=14)
    _POP(sl, f"Fuente: {F_DOC}  (hoja DOTACION + NOMINA)  |  "
             f"Alta completitud de datos en el universo Jornada/planta")
    print("  slide 4 - Sub-cascadas OK")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Cascada P2
# ─────────────────────────────────────────────────────────────────────────────
def slide_cascada_p2(prs):
    fig = _tr_fig()
    gx = PIC_RECT[0] + 0.02
    gy = PIC_RECT[1] + 0.00
    gw = PIC_RECT[2] - 0.04
    gh = PIC_RECT[3] - 0.01
    ax = fig.add_axes([gx, gy, gw, gh], facecolor="none", zorder=5)
    ax.set_xlim(0, 10); ax.set_ylim(0, 10); ax.axis("off")

    BW, BH, BX = 3.90, 1.45, 0.15

    def _box(ax, x0, y0, n, label, line1, line2, col, fs_n=22):
        ax.add_patch(mpatches.FancyBboxPatch(
            (x0, y0), BW, BH, boxstyle="round,pad=0.10",
            facecolor=col, edgecolor="white", linewidth=0.8, alpha=0.90, zorder=3))
        ax.text(x0+0.65, y0+BH*0.60, str(n),
                ha="center", va="center", fontsize=fs_n, fontweight="bold",
                color="white", zorder=5,
                path_effects=[pe.withStroke(linewidth=3, foreground="#050D1A")])
        ax.text(x0+1.40, y0+BH*0.75, label,
                ha="left", va="center", fontsize=9.5, fontweight="bold",
                color="white", zorder=5)
        ax.text(x0+1.40, y0+BH*0.46, line1,
                ha="left", va="center", fontsize=7.0, color="#B8D0E8", zorder=5)
        if line2:
            ax.text(x0+1.40, y0+BH*0.22, line2,
                    ha="left", va="center", fontsize=7.0, color="#B8D0E8", zorder=5)

    def _arrow(ax, x, y_from, y_to, lbl, col="#FFD580"):
        mid = (y_from + y_to) / 2
        ax.annotate("", xy=(x, y_to+0.05), xytext=(x, y_from-0.05),
                    arrowprops=dict(arrowstyle="->", color=col, lw=1.4,
                                   mutation_scale=16), zorder=6)
        ax.text(x-0.10, mid, lbl, ha="right", va="center",
                fontsize=7.5, color=col, fontstyle="italic", zorder=5)

    # Node 1: 545 — y0=7.50, top at 8.95 < LOGO_Y_THRESH ✓
    _box(ax, BX, 7.50, N_TOTAL, "docentes Jornada jerarquizados",
         "Hoja NOMINA",
         f"({F_DOC})", "#1B4E8F")

    # Bubble: sin formacion (right side)
    bub_x, bub_y = BX+BW+0.55, 7.30
    ax.add_patch(mpatches.FancyBboxPatch(
        (bub_x, bub_y), 4.90, 1.30,
        boxstyle="round,pad=0.10", facecolor="#5A2020",
        edgecolor="#E05A3A", linewidth=0.8, alpha=0.75, zorder=3))
    ax.text(bub_x+0.55, bub_y+0.65, str(N_SIN_FORM),
            ha="center", va="center", fontsize=20, fontweight="bold",
            color="white", zorder=5,
            path_effects=[pe.withStroke(linewidth=2.5, foreground="#050D1A")])
    ax.text(bub_x+1.25, bub_y+0.84, "sin formacion registrada",
            ha="left", va="center", fontsize=9, fontweight="bold", color="white", zorder=5)
    ax.text(bub_x+1.25, bub_y+0.44, "Brecha de cobertura, no de datos",
            ha="left", va="center", fontsize=7.8, color="#E08080",
            fontstyle="italic", zorder=5)
    ax.text(bub_x+1.25, bub_y+0.14,
            f"({100*N_SIN_FORM/N_TOTAL:.0f}% del universo Jornada)",
            ha="left", va="center", fontsize=7.5, color="#B8D0E8", zorder=5)
    ax.annotate("", xy=(bub_x-0.02, bub_y+0.65),
                xytext=(BX+BW+0.05, 7.50+BH*0.44),
                arrowprops=dict(arrowstyle="->", color="#E05A3A", lw=1.2,
                                connectionstyle="arc3,rad=0.0"), zorder=6)

    # Node 2: 246 formados — y0=5.30
    _box(ax, BX, 5.30, N_FORMADOS, "Jornada formados (al menos 1 iniciativa)",
         "Hojas DIPLOMADO 2022-25 · TALLERES · PROYECTOS",
         f"({F_DOC})", "#2E6AAD")

    _arrow(ax, BX+BW/2-0.2, 7.50, 5.30+BH,
           f"{N_FORMADOS} de {N_TOTAL}  ({100*N_FORMADOS/N_TOTAL:.0f}%)")

    # Node 3: 416 iniciativas — y0=3.10
    _box(ax, BX, 3.10, N_INICIATIVAS, "iniciativas de formacion totales",
         f"{N_FORMADOS} Jornada unicos  (uno puede tener varias)",
         None, "#3A8BC4")

    _arrow(ax, BX+BW/2-0.2, 5.30, 3.10+BH, f"{N_INICIATIVAS} eventos")

    # Tipo breakdown — inside bars with file ref below
    tipos = [
        (N_TALL, "Talleres",   "#4B9CD3", "Hojas TALLERES 2023_2, 2024_1, 2024_2"),
        (N_DIP,  "Diplomados", "#3A7ABF", "Hojas DIPLOMADO 2022, 2023, 2024, 2025"),
        (N_PROY, "Proyectos",  "#2A5EA0", "Hoja PROYECTOS DE INVESTIGACION"),
    ]
    ty0 = 3.10
    bar_max = 3.80
    for n_t, lbl_t, col_t, file_t in tipos:
        bx_t  = BX+BW+0.40
        bar_w = max((n_t / N_TALL) * bar_max, 1.80)  # min width so label fits inside bar
        ax.add_patch(mpatches.FancyBboxPatch(
            (bx_t, ty0+0.20), bar_w, 0.70,
            boxstyle="round,pad=0.05",
            facecolor=col_t, edgecolor="white",
            linewidth=0.4, alpha=0.85, zorder=3))
        # N inside bar
        ax.text(bx_t+0.14, ty0+0.55, str(n_t),
                ha="left", va="center", fontsize=14, fontweight="bold",
                color="white", zorder=5,
                path_effects=[pe.withStroke(linewidth=1.5, foreground="#050D1A")])
        ax.text(bx_t+0.72, ty0+0.60, lbl_t,
                ha="left", va="center", fontsize=9, fontweight="bold",
                color="white", zorder=5)
        # File ref below bar
        ax.text(bx_t+0.14, ty0+0.10, file_t,
                ha="left", va="center", fontsize=7.0,
                color="#B8D0E8", fontstyle="italic", zorder=5)
        ty0 += 1.00

    path = _save_ch(fig, "jor_p2_cascada.png")
    sl = _new_sl(prs)
    _pic(sl, SHARED_BG, prs)
    _pic(sl, path, prs)
    _T(sl, f"P2 — Participacion en Formacion (Jornada): {N_TOTAL} a {N_FORMADOS} Formados", fs=13)
    _POP(sl, f"Fuente: {F_DOC}  |  "
             f"Hojas: DIPLOMADO 2022-2025 · TALLERES 2023_2/2024_1/2024_2 · PROYECTOS  "
             f"(filtrado a docentes Jornada/planta)")
    _BUL(sl, [
        f"De los {N_TOTAL} Jornada jerarquizados, {N_SIN_FORM} ({100*N_SIN_FORM/N_TOTAL:.0f}%) "
        f"no tienen ninguna iniciativa de formacion registrada — es una brecha de cobertura "
        f"de politica, no un problema de datos faltantes.",
        f"Los {N_FORMADOS} Jornada formados generan {N_INICIATIVAS} iniciativas: "
        f"Talleres ({N_TALL}), Diplomados ({N_DIP}) y Proyectos ({N_PROY}). "
        f"Un docente puede haber participado en varias iniciativas de distinto tipo.",
    ])
    print("  slide 5 - Cascada P2 OK")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Solicitudes (sin _BUL; texto grande; area extendida)
# ─────────────────────────────────────────────────────────────────────────────
def slide_solicitudes(prs):
    fig = _tr_fig()
    # Extender la zona grafica para cubrir el area de PIC + BUL
    _ey_pic_t = _ey(PIC_T)          # 0.1640
    _ey_pic_h = _ey(PIC_H+BUL_H)   # 0.8152
    gy6 = 1 - _ey_pic_t - _ey_pic_h   # bottom of extended area = 0.021
    gh6 = _ey_pic_h - 0.025            # 0.790
    gx6 = PIC_RECT[0] + 0.01
    gw6 = PIC_RECT[2] - 0.02

    ax = fig.add_axes([gx6, gy6, gw6, gh6], facecolor="none", zorder=5)
    ax.set_xlim(0, 10); ax.set_ylim(0, 11); ax.axis("off")

    # Logo threshold with extended axes:
    # logo_fig_y_bottom = ~0.8025; threshold_y = (0.8025 - gy6) / gh6 * 11
    _logo_thresh = (0.8025 - gy6) / gh6 * 11  # ~10.47
    # Box tops must stay below _logo_thresh

    BOX_H = 2.40
    # yc top = yc + BOX_H/2 < _logo_thresh → yc < _logo_thresh - 1.20 = 9.27
    ycs = [9.10, 6.50, 3.90, 1.30]  # box tops: 10.30, 7.70, 5.10, 2.50 ✓

    solicitudes = [
        {
            "num": "1", "col": "#1B6A40",
            "title": f"P1  —  {N_SIN_DOT} Jornada Sin Dotacion",
            "lines": [
                f"{N_SIN_DOT} docentes Jornada estan en NOMINA pero NO en DOTACION:"
                f" sin edad, antiguedad ni nivel de formacion.",
                f"Solicitar: verificar esos {N_SIN_DOT} RUTs en DOTACION o enviar complemento.",
            ],
            "impacto": f"Recupera hasta {N_SIN_DOT} Jornada con perfil completo de P1",
        },
        {
            "num": "2", "col": "#1B4E8F",
            "title": "P1  —  Corregir Valores 'NO INFORMA' en Dotacion",
            "lines": [
                f"Dentro de los {N_DOTACION} Jornada con dotacion: {D_JER} con jerarquia,"
                f" {D_NIV} con nivel formacion y {D_GRD} con grado = 'NO INFORMA'.",
                "Estos registros quedan excluidos de los analisis que requieren esas variables.",
                "Solicitar: reclasificar o confirmar esos valores en la hoja DOTACION.",
            ],
            "impacto": f"Sube el universo util de {c_grd} a los {N_DOTACION} con dotacion",
        },
        {
            "num": "3", "col": "#3A5080",
            "title": f"P1  —  Completar FECHA JERARQUIZACION (~{D_FEC} sin dato)",
            "lines": [
                f"La hoja NOMINA tiene ~45 nulos en FECHA JERARQUIZACION (967 filas)."
                f" Proporcional a Jornada: ~{D_FEC} sin dato.",
                "Campo necesario para calcular 'Anos hasta jerarquizacion' y trayectoria.",
                "Solicitar: completar FECHA JERARQUIZACION en NOMINA o listado complementario.",
            ],
            "impacto": f"Sube de {c_fec} a ~{c_grd} los disponibles para analisis de trayectoria",
        },
        {
            "num": "4", "col": "#4A7030",
            "title": "P2  —  Completar Registros de Formacion 2025",
            "lines": [
                f"Base P2 Jornada: {N_INICIATIVAS} iniciativas, {N_FORMADOS} docentes."
                f" TALLERES solo cubre 2023_2, 2024_1 y 2024_2 — falta 2023_1 y 2025.",
                "Confirmar si diplomados y proyectos 2025 estan todos registrados en el archivo.",
                "Solicitar: hojas TALLERES 2023_1 y 2025, y formaciones 2025 pendientes.",
            ],
            "impacto": f"Amplia el universo de {N_FORMADOS} Jornada formados",
        },
    ]

    for s, yc in zip(solicitudes, ycs):
        y0 = yc - BOX_H/2
        ax.add_patch(mpatches.FancyBboxPatch(
            (-0.05, y0), 10.1, BOX_H,
            boxstyle="round,pad=0.10", facecolor=s["col"],
            edgecolor="white", linewidth=0.6, alpha=0.32, zorder=2))
        circ = mpatches.Circle((0.45, yc+0.72), 0.36,
                                facecolor=s["col"], edgecolor="white",
                                linewidth=0.8, alpha=0.95, zorder=3)
        ax.add_patch(circ)
        ax.text(0.45, yc+0.72, s["num"],
                ha="center", va="center", fontsize=16, fontweight="bold",
                color="white", zorder=4)
        ax.text(1.05, yc+0.90, s["title"],
                ha="left", va="center", fontsize=11.5, fontweight="bold",
                color="#FFD580", zorder=4)
        y_txt = yc+0.48
        for line in s["lines"]:
            ax.text(1.05, y_txt, line,
                    ha="left", va="top", fontsize=9.5, color="#D0E8FF", zorder=4)
            y_txt -= 0.43
        ax.text(0.45, y0+0.20, "-> " + s["impacto"],
                ha="left", va="center", fontsize=9.0,
                color="#A8E8A0", zorder=4, fontstyle="italic")

    path = _save_ch(fig, "jor_solicitudes.png")
    sl = _new_sl(prs)
    _pic(sl, SHARED_BG, prs)
    _pic(sl, path, prs)
    _T(sl, "Datos Adicionales — Solicitudes para Ampliar P1 y P2 (Jornada/Planta)", fs=13)
    _POP(sl, f"Todas las solicitudes apuntan al mismo archivo: {F_DOC}  "
             f"(ampliacion o correccion de hojas existentes)")
    print("  slide 6 - Solicitudes OK")


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("Generando FICHA_Ampliar_P1_P2_Jornada.pptx ...")
    _ensure_bg()
    prs = Presentation()
    prs.slide_width  = Emu(SW_EMU)
    prs.slide_height = Emu(SH_EMU)
    slide_portada(prs)
    slide_cascada_p1(prs)
    slide_brecha_dotacion(prs)
    slide_subcascadas(prs)
    slide_cascada_p2(prs)
    slide_solicitudes(prs)
    prs.save(OUT_PPTX)
    print(f"\n  Guardado: {OUT_PPTX}")

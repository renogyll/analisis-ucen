"""
generar_ppt_dark_bloque1.py
Genera BLOQUE I (10 slides) en formato oscuro UCEN.
Referencia: generador_ppt.md / poc_slide_facultad_dark.py
Salida: PRESENTACION_197_P3_dark_bloque1.pptx
"""
import sys; sys.stdout.reconfigure(encoding="utf-8")
import os, zipfile, textwrap
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.patheffects as pe
import pandas as pd
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Emu

# ─────────────────────────────────────────────────────────────────────────────
# Rutas
# ─────────────────────────────────────────────────────────────────────────────
BASE      = os.path.dirname(__file__)
SCRATCH   = (r"C:\Users\RGONZA~1.LAP\AppData\Local\Temp\claude"
             r"\c--Users-r-gonzalez-fluxsolar-LAPTOP-FLUX-ECO-Downloads-Analisis-UCEN-v2"
             r"\19e6fc3f-6ca1-4150-9da7-8dfa38be71ca\scratchpad")
FONDOTIPO = (r"c:\Users\r.gonzalez_fluxsolar.LAPTOP-FLUX-ECO"
             r"\Downloads\Analisis_UCEN_v2\Fondotipop.pptx")
OUT_PPTX  = (r"c:\Users\r.gonzalez_fluxsolar.LAPTOP-FLUX-ECO"
             r"\Downloads\Analisis_UCEN_v2\PRESENTACION_197_P3_dark_bloque1.pptx")
OUT_DIR   = os.path.join(BASE, "dark_slides")
os.makedirs(OUT_DIR, exist_ok=True)
os.makedirs(SCRATCH, exist_ok=True)

# ─────────────────────────────────────────────────────────────────────────────
# Assets: foto B&N + logo UCEN + gradiente navy
# ─────────────────────────────────────────────────────────────────────────────
BG_PATH   = os.path.join(SCRATCH, "fondotipo_image1.jpg")
LOGO_PATH = os.path.join(SCRATCH, "fondotipo_image2.png")
for path, zname in [(BG_PATH, "ppt/media/image1.jpg"), (LOGO_PATH, "ppt/media/image2.png")]:
    if not os.path.exists(path):
        with zipfile.ZipFile(FONDOTIPO) as z:
            with open(path, "wb") as f:
                f.write(z.read(zname))

with PILImage.open(BG_PATH) as _im:
    _rgb = _im.convert("RGB")
    _iw, _ih = _rgb.size
    _nh = int(_iw / (16/9))
    _y0 = min(int(_ih * 0.12), _ih - _nh)
    bg_arr = np.array(_rgb.crop((0, _y0, _iw, _y0 + _nh)))

with PILImage.open(LOGO_PATH) as _logo:
    logo_arr = np.array(_logo.convert("RGBA")).astype(np.float32) / 255.0

H_GRAD = 600
grad   = np.zeros((H_GRAD, 1, 4), dtype=np.float32)
for _r in range(H_GRAD):
    _t = _r / (H_GRAD - 1)
    _stops = [(0.00,(0,33,71)), (0.54,(0,70,128)), (1.00,(144,171,196))]
    for _i in range(len(_stops)-1):
        _t0,_c0 = _stops[_i]; _t1,_c1 = _stops[_i+1]
        if _t0 <= _t <= _t1:
            _s = (_t-_t0)/(_t1-_t0)
            grad[_r,0] = [(_c0[0]+_s*(_c1[0]-_c0[0]))/255,
                          (_c0[1]+_s*(_c1[1]-_c0[1]))/255,
                          (_c0[2]+_s*(_c1[2]-_c0[2]))/255, 0.82]
            break

# ─────────────────────────────────────────────────────────────────────────────
# Datos
# ─────────────────────────────────────────────────────────────────────────────
p3 = pd.read_csv(os.path.join(BASE, "..", "PROCESADO", "p3_918.csv"),
                 encoding="utf-8-sig", dtype={"rut_key": str})
p3["rut_key"] = p3["rut_key"].str.strip()
aptos = p3[p3["apto_p3"] == True].drop_duplicates("rut_key").copy()
N197  = len(aptos)

DOT_CSV = os.path.join(BASE, "..", "..", "CONSOLIDADO DOCENTES 3-05-2026",
                        "CONSOLIDADO DOCENTES 3-05-2026.xlsx - DOTACION_CON_GRADOREC.csv")
has_dot = os.path.exists(DOT_CSV)
if has_dot:
    dot = pd.read_csv(DOT_CSV, dtype=str, encoding="utf-8-sig")
    dot.columns = dot.columns.str.strip()
    dot["rut_key"] = (dot["RUT"].str.strip().str.replace(".", "", regex=False)
                      .str.split("-").str[0].str.strip())
    dot197 = dot[dot["rut_key"].isin(set(aptos["rut_key"]))].copy()
else:
    dot197 = pd.DataFrame()

# ─────────────────────────────────────────────────────────────────────────────
# Layout (EMU → fracción figura)
# ─────────────────────────────────────────────────────────────────────────────
SW, SH = 13.333, 7.5
def ex(e): return e / 12192000
def ey(e): return e / 6858000
def fig_rect(l, t, w, h): return (l, 1-t-h, w, h)

PIC_L,  PIC_T  = ex(786581),  ey(1125000)
PIC_W,  PIC_H  = ex(10599174), ey(3720000)
BUL_L,  BUL_T  = ex(786581),  ey(4870000)
BUL_H           = ey(1870000)
LOGO_L, LOGO_T = ex(9813773),  ey(656354)
LOGO_W, LOGO_H = ex(1756626),  ey(697725)
POP_T           = ey(820000)
POP_H           = ey(290000)

PIC_RECT  = fig_rect(PIC_L, PIC_T,  PIC_W,  PIC_H)
LOGO_RECT = fig_rect(LOGO_L, LOGO_T, LOGO_W, LOGO_H)
POP_Y     = 1 - POP_T - POP_H / 2

CHART_X = PIC_RECT[0] + 0.16
CHART_W = PIC_RECT[2] - 0.16 - 0.02
CHART_Y = PIC_RECT[1] + 0.04
CHART_H = PIC_RECT[3] - 0.09

PAL = ["#5C9BD6","#64B5F6","#80DEEA","#A5D6A7","#FFB74D","#CE93D8","#90A4AE","#F48FB1"]
POP_197 = ("Universo: 197 Aptos P3  ·  130 formados + 67 control"
           "  ·  SAT disponible baseline / durante / resultado")

# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────
def _base():
    fig = plt.figure(figsize=(SW, SH), facecolor="#101820")
    fig.patch.set_facecolor("#101820")
    for z, arr in [(0, bg_arr), (1, grad)]:
        ax = fig.add_axes([0,0,1,1], zorder=z)
        ax.imshow(arr, extent=[0,1,0,1], aspect="auto", origin="upper")
        ax.set_xlim(0,1); ax.set_ylim(0,1); ax.axis("off")
    al = fig.add_axes([LOGO_RECT[0], LOGO_RECT[1], LOGO_RECT[2], LOGO_RECT[3]],
                      zorder=10, facecolor="none")
    al.imshow(logo_arr, aspect="auto"); al.axis("off"); al.patch.set_visible(False)
    return fig

def _title(fig, text, fontsize=20):
    ty = (1.0 + CHART_Y + CHART_H + 0.020) / 2 + 0.018
    fig.text(PIC_L + PIC_W/2, ty, text, ha="center", va="center",
             fontsize=fontsize, fontweight="bold", color="white",
             transform=fig.transFigure, zorder=4)

def _pop(fig, txt=POP_197):
    fig.text(PIC_RECT[0], POP_Y, txt, ha="left", va="center",
             fontsize=7.5, fontstyle="italic", color="#C8DCF0",
             transform=fig.transFigure, zorder=4)

def _ctitle(fig, txt, cx=None):
    x = (cx if cx is not None else CHART_X) - 0.005
    fig.text(x, CHART_Y + CHART_H + 0.008, txt,
             ha="left", va="bottom", fontsize=10, color="white",
             transform=fig.transFigure, zorder=7)

def _bullets(fig, items):
    BUL_BOTTOM = 1 - BUL_T - BUL_H
    AX_W       = CHART_X + CHART_W - PIC_L
    ab = fig.add_axes([PIC_L, BUL_BOTTOM, AX_W, BUL_H], facecolor="none", zorder=6)
    ab.set_xlim(0,1); ab.set_ylim(0,1); ab.axis("off"); ab.patch.set_visible(False)
    lh  = 0.023 / BUL_H
    gap = 0.012 / BUL_H
    y   = 1.0 - 0.040 / BUL_H
    for i, bul in enumerate(items):
        for line in textwrap.wrap(f"{i+1}.  {bul}", width=130, subsequent_indent="    "):
            ab.text(0, y, line, ha="left", va="top", fontsize=11.5, color="white",
                    transform=ab.transAxes, clip_on=True)
            y -= lh
        y -= gap

def _hbar(fig, labels, vals, pcts, cx=None, cw=None):
    x = cx if cx is not None else CHART_X
    w = cw if cw is not None else CHART_W
    ax = fig.add_axes([x, CHART_Y, w, CHART_H], facecolor="none", zorder=5)
    n  = len(labels)
    yp = np.arange(n)
    ax.barh(yp[::-1], vals, color=PAL[:n], height=0.58, edgecolor="none", alpha=0.90)
    mv = max(vals) if vals else 1
    for i, (v, p, c) in enumerate(zip(vals[::-1], pcts[::-1], PAL[:n][::-1])):
        ax.text(v + mv*0.015, i, f"{v}  ({p:.1f}%)",
                va="center", ha="left", fontsize=10.5, fontweight="bold", color=c,
                path_effects=[pe.withStroke(linewidth=2.5, foreground="#0A0F18")])
    ax.set_yticks(yp)
    ax.set_yticklabels(labels[::-1], fontsize=10.5, fontweight="bold", color="white")
    ax.tick_params(axis="y", length=0, pad=8)
    ax.tick_params(axis="x", colors="#AAAAAA", labelsize=9)
    for sp in ax.spines.values():
        sp.set_visible(True); sp.set_edgecolor("white"); sp.set_alpha(0.35); sp.set_linewidth(0.9)
    ax.set_xlim(0, mv * 1.35); ax.set_ylim(-0.5, n - 0.5)
    ax.xaxis.grid(True, color="white", alpha=0.07, linewidth=0.5, zorder=0)
    ax.set_axisbelow(True)
    return ax

def _save(fig, name):
    path = os.path.join(OUT_DIR, name)
    plt.savefig(path, format="png", dpi=150, facecolor=fig.get_facecolor())
    plt.close()
    print(f"  ✓ {name}")
    return path

# ─────────────────────────────────────────────────────────────────────────────
# Slide 01 — Portada
# ─────────────────────────────────────────────────────────────────────────────
def gen_portada():
    fig = _base()
    kw = dict(ha="center", va="center", transform=fig.transFigure, zorder=4)
    fig.text(0.50, 0.64, "Resultados del Perfeccionamiento Docente",
             fontsize=30, fontweight="bold", color="white", **kw)
    # divisor
    ax_d = fig.add_axes([0.22, 0.585, 0.56, 0.003], zorder=4)
    ax_d.set_facecolor("#90ABC4"); ax_d.axis("off")
    fig.text(0.50, 0.54, "Análisis de Impacto: SAT, Recomendación, Notas y EDD",
             fontsize=18, color="#C8DCF0", **kw)
    fig.text(0.50, 0.46, "Universo: 197 Aptos P3  ·  130 formados + 67 control  ·  Períodos 2023–2025",
             fontsize=13, fontstyle="italic", color="#C8DCF0", **kw)
    fig.text(0.50, 0.21,
             "Universidad Central de Chile  |  Producto 3: Análisis de Formación e Innovación",
             fontsize=13, fontweight="bold", color="white", **kw)
    return _save(fig, "01_portada.png")

# ─────────────────────────────────────────────────────────────────────────────
# Slide 02 — Índice
# ─────────────────────────────────────────────────────────────────────────────
def gen_indice():
    bloques = [
        "I.    Clasificación del Cuerpo Académico  —  edad, sexo, facultad, jerarquía y grado  (n=197)",
        "II.   Evaluación Docente Antes y Después  —  SAT, % Recomendación  (n=197)",
        "III.  Calificaciones de Alumnos  —  notas y aprobación según SAT  (n=816)",
        "IV.   EDD  —  Evaluación de Desempeño Docente, evolución y panel balanceado  (n=197)",
        "       Conclusiones y Recomendaciones",
    ]
    fig = _base()
    _title(fig, "Índice")
    ax = fig.add_axes([PIC_RECT[0]+0.02, PIC_RECT[1]+0.04, PIC_RECT[2]-0.04, PIC_RECT[3]-0.06],
                      facecolor="none", zorder=5)
    ax.set_xlim(0,1); ax.set_ylim(0,1); ax.axis("off")
    top, step = 0.88, 0.17
    for b in bloques:
        ax.text(0.01, top, b, ha="left", va="top", fontsize=14, color="white",
                transform=ax.transAxes)
        top -= step
    return _save(fig, "02_indice.png")

# ─────────────────────────────────────────────────────────────────────────────
# Slide 03 — Marco Metodológico
# ─────────────────────────────────────────────────────────────────────────────
def gen_marco():
    cajas = [
        ("1. Universo Aptos P3",
         ["917 docentes jerarquizados UCEN (universo base)",
          "493 con actividad docente en el período analizado",
          "357 participaron en ≥1 iniciativa de formación",
          "197 Aptos P3: SAT válido en baseline, durante y resultado"]),
        ("2. Marco Metodológico — Z-score",
         ["z = (SAT docente − media facultad período) / DE facultad período",
          "z = 0 → promedio exacto de su facultad ese semestre",
          "z > 0 → sobre el promedio  ·  z < 0 → bajo el promedio",
          "Permite comparar docentes entre facultades distintas"]),
        ("3. Métricas de Evaluación",
         ["SAT Nota (1–7) estandarizada como z-score por facultad",
          "SAT % Recomendación: ¿recomendaría a este docente?",
          "EDD: Evaluación Directa de Desempeño (escala 0–1)",
          "Notas y aprobación alumnos: % con nota ≥ 4.0"]),
        ("4. Diseño Comparativo",
         ["Grupo tratamiento: 130 docentes formados con SAT P3",
          "Grupo control: 67 docentes sin formación con SAT P3",
          "Comparación: posición z antes → durante → después",
          "Validación estadística: prueba t de diferencia de medias"]),
    ]
    fig = _base()
    _title(fig, "Universo de Análisis y Metodología")
    _pop(fig, "Universo base: 917 docentes jerarquizados UCEN  ·  Los 197 Aptos P3 son el subconjunto analizado")

    cols = 2
    pad_x = PIC_RECT[0] + 0.01
    box_w = (PIC_RECT[2] - 0.03) / 2
    box_h = (PIC_RECT[3] - 0.04) / 2
    positions = [
        (pad_x,             PIC_RECT[1] + PIC_RECT[3]/2 + 0.01),
        (pad_x + box_w + 0.01, PIC_RECT[1] + PIC_RECT[3]/2 + 0.01),
        (pad_x,             PIC_RECT[1] + 0.01),
        (pad_x + box_w + 0.01, PIC_RECT[1] + 0.01),
    ]
    for (bx, by), (header, items) in zip(positions, cajas):
        ax = fig.add_axes([bx, by, box_w, box_h - 0.01], facecolor="none", zorder=5)
        ax.set_xlim(0,1); ax.set_ylim(0,1); ax.axis("off")
        ax.text(0, 0.92, header, ha="left", va="top", fontsize=11,
                fontweight="bold", color="#90ABC4", transform=ax.transAxes)
        for k, item in enumerate(items):
            ax.text(0, 0.72 - k*0.185, f"• {item}", ha="left", va="top",
                    fontsize=9.5, color="white", transform=ax.transAxes)
    return _save(fig, "03_marco.png")

# ─────────────────────────────────────────────────────────────────────────────
# Slide 04 — Separador BLOQUE I
# ─────────────────────────────────────────────────────────────────────────────
def gen_separador():
    items = [
        "Diapo 5:   Tramo de Edad y Sexo                           (n=129 con edad / 197 con sexo)",
        "Diapo 6:   Distribucion por Facultad/Unidad               (n=129 con dato)",
        "Diapo 7:   Antiguedad en la Institucion                   (n=129 con dato)",
        "Diapo 8:   Distribucion de Jerarquia Academica            (n=197)",
        "Diapo 9:   Grado Academico Reconocido                     (n=129 con dotacion)",
        "Diapo 10:  Institucion de Obtencion del Grado             (n=129 con dotacion)",
    ]
    fig = _base()
    _title(fig, "BLOQUE I — Clasificación del Cuerpo Académico", fontsize=18)
    _pop(fig, POP_197)

    import re as _re
    def _dl(text, w=84):
        m = _re.match(r'^(.*?)\s{3,}(\(.+\))\s*$', text.strip())
        if not m: return text
        l, r = m.group(1).rstrip(), m.group(2)
        dots = max(3, w - len(l) - len(r) - 2)
        return f"{l} {'.'*dots} {r}"

    ax = fig.add_axes([PIC_RECT[0]+0.01, PIC_RECT[1]+0.03,
                       PIC_RECT[2]-0.02, PIC_RECT[3]-0.05],
                      facecolor="none", zorder=5)
    ax.set_xlim(0,1); ax.set_ylim(0,1); ax.axis("off")
    top, step = 0.92, 0.145
    for i, item in enumerate(items):
        prefix = "  └──  " if i == len(items)-1 else "  ├──  "
        ax.text(0, top, prefix + _dl(item), ha="left", va="top",
                fontsize=12, color="white", fontfamily="monospace",
                transform=ax.transAxes)
        top -= step
    return _save(fig, "04_separador.png")

# ─────────────────────────────────────────────────────────────────────────────
# Slide 05 — Edad + Sexo (doble)
# ─────────────────────────────────────────────────────────────────────────────
def gen_edad_sexo():
    ORD = ["<30","30-34","35-39","40-44","45-49","50-54","55-59","60-64","65-69","70+"]
    con_e = aptos[aptos["tramo_edad"].notna() &
                  (aptos["tramo_edad"].astype(str).str.strip() != "nan")].copy()
    n_e   = len(con_e)
    tbl   = con_e.groupby("tramo_edad")["rut_key"].count().reindex(ORD).fillna(0).astype(int)
    vals_e = [int(tbl.get(t, 0)) for t in ORD]
    pcts_e = [100*v/n_e if n_e else 0 for v in vals_e]
    edad_med = pd.to_numeric(aptos["edad_anios"], errors="coerce").median()

    SMAP = {"MASCULINO":"Hombre","HOMBRE":"Hombre","FEMENINO":"Mujer","MUJER":"Mujer"}
    con_s = aptos[aptos["sexo"].notna()].copy()
    con_s["sc"] = con_s["sexo"].str.strip().str.upper().map(SMAP).fillna(
                  con_s["sexo"].str.strip().str.title())
    n_s    = len(con_s)
    cnt_s  = con_s["sc"].value_counts()
    lbl_s  = cnt_s.index.tolist()
    vals_s = cnt_s.values.tolist()
    pcts_s = [100*v/n_s for v in vals_s]

    fig = _base()
    _title(fig, "Distribución por Tramo de Edad y Sexo")
    _pop(fig)

    # ── Panel izquierdo: barras verticales edad ──────────────────────────────
    LX = PIC_RECT[0] + 0.02
    LW = PIC_RECT[2] * 0.54
    LY = CHART_Y; LH = CHART_H

    ax_e = fig.add_axes([LX, LY, LW, LH], facecolor="none", zorder=5)
    xe = np.arange(len(ORD))
    ax_e.bar(xe, vals_e, width=0.70, color=PAL[:len(ORD)], alpha=0.90, edgecolor="none")
    mv_e = max(vals_e) if vals_e else 1
    for i, (v, p) in enumerate(zip(vals_e, pcts_e)):
        if v > 0:
            ax_e.text(i, v + mv_e*0.025, f"{v}\n({p:.0f}%)",
                      ha="center", va="bottom", fontsize=8, fontweight="bold",
                      color="white",
                      path_effects=[pe.withStroke(linewidth=1.5, foreground="#0A0F18")])
    ax_e.set_xticks(xe)
    ax_e.set_xticklabels(ORD, fontsize=8.5, color="white", rotation=35, ha="right")
    ax_e.tick_params(axis="x", colors="white", length=0)
    ax_e.tick_params(axis="y", colors="#AAAAAA", labelsize=8)
    for sp in ax_e.spines.values():
        sp.set_edgecolor("white"); sp.set_alpha(0.35); sp.set_linewidth(0.9)
    ax_e.set_ylim(0, mv_e * 1.30)
    ax_e.yaxis.grid(True, color="white", alpha=0.07, linewidth=0.5)
    ax_e.set_axisbelow(True)
    if pd.notnull(edad_med):
        idx_med = min(range(len(ORD)),
                      key=lambda i: abs(int(ORD[i].split("-")[0].replace("<","0").replace("+","")) - edad_med))
        ax_e.axvline(idx_med + 0.05, color="white", linewidth=1.4, linestyle="--", alpha=0.6)
        ax_e.text(idx_med + 0.2, mv_e*0.85, f"Med ≈{edad_med:.0f}a",
                  fontsize=8, color="white", fontstyle="italic", alpha=0.85)
    fig.text(LX, LY + LH + 0.008, f"Tramo de edad  (n={n_e} con dato)",
             ha="left", va="bottom", fontsize=9.5, color="white",
             transform=fig.transFigure, zorder=7)

    # ── Panel derecho: donut sexo ─────────────────────────────────────────────
    RH = LH * 0.88
    RW = RH * (SH / SW)
    RX = LX + LW + 0.02 + ((PIC_RECT[0] + PIC_RECT[2]) - (LX + LW + 0.02) - RW) / 2
    RY = LY + (LH - RH) / 2

    ax_s = fig.add_axes([RX, RY, RW, RH], facecolor="none", zorder=5)
    SCOLS = ["#1E88E5", "#E65100", "#43A047"]
    wedges, _ = ax_s.pie(vals_s, colors=SCOLS[:len(lbl_s)], startangle=90,
                          counterclock=False,
                          wedgeprops=dict(width=0.55, edgecolor="#101820", linewidth=2))
    for w, lbl, v, p, c in zip(wedges, lbl_s, vals_s, pcts_s, SCOLS):
        ang = (w.theta2 + w.theta1) / 2
        x2  = 1.35 * np.cos(np.radians(ang))
        y2  = 1.35 * np.sin(np.radians(ang))
        ax_s.annotate(f"{lbl}\n{v} ({p:.1f}%)",
                      xy=(0.7*np.cos(np.radians(ang)), 0.7*np.sin(np.radians(ang))),
                      xytext=(x2, y2),
                      arrowprops=dict(arrowstyle="-", color="white", lw=1),
                      fontsize=11, ha=("left" if x2 > 0 else "right"),
                      va="center", fontweight="bold", color=c)
    ax_s.text(0, 0, f"{n_s}\ndoc.", ha="center", va="center",
              fontsize=12, fontweight="bold", color="white")
    ax_s.set_xlim(-2.2, 2.2); ax_s.set_ylim(-2.2, 2.2)
    fig.text(RX + RW/2, RY + RH + 0.008, f"Sexo  (n={n_s})",
             ha="center", va="bottom", fontsize=9.5, color="white",
             transform=fig.transFigure, zorder=7)

    n_nuc = sum(vals_e[ORD.index(t)] for t in ["35-39","40-44","45-49","50-54"] if t in ORD)
    _bullets(fig, [
        f"El tramo 40–44 años concentra la mayor frecuencia; el núcleo 35–54 agrupa "
        f"{n_nuc} docentes ({100*n_nuc/n_e:.0f}% de los {n_e} con dato de edad).",
        f"La distribución por sexo muestra mayor presencia de Hombres ({pcts_s[0]:.0f}% aprox.). "
        f"Perfil dentro de rangos esperados para docencia universitaria en Chile.",
        f"Nota: {n_e} de los 197 Aptos P3 tienen dato de edad ({100*n_e/N197:.0f}%); "
        f"el resto corresponde a honorarios sin registro en dotación.",
    ])
    return _save(fig, "05_edad_sexo.png")

# ─────────────────────────────────────────────────────────────────────────────
# Slide 06 — Facultad
# ─────────────────────────────────────────────────────────────────────────────
def gen_facultad():
    ABREV = {
        "FAC. DE MEDICINA Y CIENCIAS DE LA SALUD": "Medicina y C. Salud",
        "FAC. DERECHO Y HUMANIDADES":              "Derecho y Humanidades",
        "FAC. DE INGENIERÍA Y ARQUITECTURA":       "Ingeniería y Arq.",
        "FAC. DE EDUCACIÓN":                       "Educación",
        "FAC. ECONOMÍA, GOBIERNO Y COMUNICACIONES":"Economía, Gob. y Com.",
        "VICERRECTORIA DE INVEST, INNOV Y POSTGRA":"VR Invest./Postgrado",
        "VICERRECTORIA ACADEMICA":                 "VR Académica",
        "DIRECCION DE ASEGURAMIENTO DE LA CALIDAD":"Dir. Aseg. Calidad",
    }
    con = aptos[aptos["unidad_facultad"].notna()].copy()
    n   = len(con)
    cnt = con["unidad_facultad"].str.strip().value_counts()
    lbl = [ABREV.get(k, k[:28]) for k in cnt.index]
    val = cnt.values.tolist()
    pct = [100*v/n for v in val]

    fig = _base()
    _title(fig, "Distribución por Unidad/Facultad")
    _pop(fig)
    _hbar(fig, lbl, val, pct)
    _ctitle(fig, f"Distribucion por Unidad/Facultad — 197 Aptos P3  (n={n} con dato)")
    _bullets(fig, [
        f"La Facultad de Medicina y Ciencias de la Salud concentra el mayor número de Aptos P3 "
        f"({val[0]} doc., {pct[0]:.0f}%), seguida por Derecho y Humanidades ({val[1]}, {pct[1]:.0f}%) "
        f"e Ingeniería y Arquitectura ({val[2]}, {pct[2]:.0f}%).",
        "Las cinco facultades principales agrupan el 90% de los docentes Aptos P3. "
        "La distribución refleja la mayor concentración de actividad SAT en salud y derecho.",
        "VR Investigación e Innovación aporta docentes con perfil orientado a investigación; "
        "su presencia en Aptos P3 indica participación activa en formación.",
    ])
    return _save(fig, "06_facultad.png")

# ─────────────────────────────────────────────────────────────────────────────
# Slide 07 — Antigüedad
# ─────────────────────────────────────────────────────────────────────────────
def gen_antiguedad():
    ORD = ["0-4","5-9","10-14","15-19","20-24","25-29","30+"]
    con = aptos[aptos["tramo_antiguedad"].notna()].copy()
    n   = len(con)
    tbl = con.groupby("tramo_antiguedad")["rut_key"].count().reindex(ORD).fillna(0).astype(int)
    val = [int(tbl.get(t, 0)) for t in ORD]
    pct = [100*v/n if n else 0 for v in val]
    ant_med = pd.to_numeric(aptos["antiguedad_anios"], errors="coerce").median()

    PCOLS = ["#42A5F5","#1E88E5","#1976D2","#1565C0","#0D47A1","#546E7A","#78909C"]

    fig = _base()
    _title(fig, "Antigüedad en la Institución")
    _pop(fig)

    ax = fig.add_axes([CHART_X, CHART_Y, CHART_W, CHART_H], facecolor="none", zorder=5)
    xa = np.arange(len(ORD))
    ax.bar(xa, val, width=0.68, color=PCOLS[:len(ORD)], alpha=0.90, edgecolor="none")
    mv = max(val) if val else 1
    for i, (v, p) in enumerate(zip(val, pct)):
        if v > 0:
            ax.text(i, v + mv*0.025, f"{v}\n({p:.0f}%)",
                    ha="center", va="bottom", fontsize=9, fontweight="bold",
                    color="white",
                    path_effects=[pe.withStroke(linewidth=1.5, foreground="#0A0F18")])
    ax.set_xticks(xa)
    ax.set_xticklabels(ORD, fontsize=10.5, color="white")
    ax.tick_params(axis="x", colors="white", length=0)
    ax.tick_params(axis="y", colors="#AAAAAA", labelsize=9)
    for sp in ax.spines.values():
        sp.set_edgecolor("white"); sp.set_alpha(0.35); sp.set_linewidth(0.9)
    ax.set_ylim(0, mv * 1.28)
    ax.yaxis.grid(True, color="white", alpha=0.07, linewidth=0.5)
    ax.set_axisbelow(True)
    ax.set_xlabel("Años en la institución", color="#AAAAAA", fontsize=10)
    if pd.notnull(ant_med):
        idx_m = min(range(len(ORD)),
                    key=lambda i: abs(int(ORD[i].split("-")[0].replace("+","")) - ant_med))
        ax.axvline(idx_m + 0.05, color="white", linewidth=1.4, linestyle="--", alpha=0.6)
        ax.text(idx_m + 0.2, mv*0.85, f"Med ≈{ant_med:.0f}a",
                fontsize=9, color="white", fontstyle="italic", alpha=0.85)

    _ctitle(fig, f"Antiguedad en la Institucion — 197 Aptos P3  (n={n} con dato)")
    _bullets(fig, [
        f"El tramo de mayor frecuencia es 5–9 años ({val[1]} doc.), seguido de 0–4 años ({val[0]} doc.). "
        "Los docentes más jóvenes en la institución son los que con mayor frecuencia reúnen condiciones Aptos P3.",
        "A mayor antigüedad, menor presencia en el grupo Aptos P3, lo que es consistente "
        "con que los docentes más experimentados participan menos en iniciativas de formación.",
        f"Mediana estimada: {ant_med:.0f} años en la institución.",
    ])
    return _save(fig, "07_antiguedad.png")

# ─────────────────────────────────────────────────────────────────────────────
# Slide 08 — Jerarquía
# ─────────────────────────────────────────────────────────────────────────────
def gen_jerarquia():
    JER_ORD = ["INSTRUCTOR REGULAR","INSTRUCTOR DOCENTE","ASISTENTE REGULAR","ASISTENTE DOCENTE",
               "ASOCIADO REGULAR","ASOCIADO DOCENTE","TITULAR REGULAR","TITULAR DOCENTE"]
    JER_LBL = {
        "INSTRUCTOR REGULAR": "Instructor Regular", "INSTRUCTOR DOCENTE": "Instructor Docente",
        "ASISTENTE REGULAR":  "Asistente Regular",  "ASISTENTE DOCENTE":  "Asistente Docente",
        "ASOCIADO REGULAR":   "Asociado Regular",   "ASOCIADO DOCENTE":   "Asociado Docente",
        "TITULAR REGULAR":    "Titular Regular",     "TITULAR DOCENTE":    "Titular Docente",
    }
    JER_COL = ["#90CAF9","#1E88E5","#A5D6A7","#43A047","#FFA726","#E65100","#CE93D8","#7B1FA2"]

    con = aptos[aptos["jerarquia"].notna()].copy()
    n   = len(con)
    cnt = con["jerarquia"].str.strip().str.upper().value_counts()
    rows = [(JER_LBL.get(j,j), int(cnt.get(j,0)), JER_COL[i])
            for i,j in enumerate(JER_ORD) if int(cnt.get(j,0)) > 0]
    lbl, val, cols = zip(*rows) if rows else ([],[],[])
    lbl, val, cols = list(lbl), list(val), list(cols)
    pct = [100*v/n for v in val]

    fig = _base()
    _title(fig, "Distribución de Jerarquía Académica")
    _pop(fig)

    ax = _hbar(fig, lbl, val, pct)
    # Override colores individuales
    for bar, c in zip(ax.patches, cols[::-1]):
        bar.set_color(c)
    for txt, c in zip(ax.texts, cols[::-1]):
        txt.set_color(c)
    ax.set_yticklabels(lbl[::-1], fontsize=10.5, fontweight="bold", color="white")

    _ctitle(fig, f"Jerarquia Academica — 197 Aptos P3  (n={n})")
    n_doc = sum(v for l,v in zip(lbl,val) if "Docente" in l)
    n_reg = sum(v for l,v in zip(lbl,val) if "Regular" in l)
    _bullets(fig, [
        f"Los Aptos P3 se concentran en Asistente Docente e Instructor Docente, "
        f"en línea con el mayor perfil de participación en formación de los rangos de entrada al escalafón.",
        f"Cuerpo Docente: {n_doc} doc. ({100*n_doc/n:.0f}%)  ·  Cuerpo Regular: {n_reg} doc. ({100*n_reg/n:.0f}%). "
        "El Cuerpo Docente domina, consistente con mayor orientación a la enseñanza.",
        "Titulares muestran la menor representación, lo que sugiere menor propensión a participar "
        "en iniciativas de formación en los rangos superiores del escalafón.",
    ])
    return _save(fig, "08_jerarquia.png")

# ─────────────────────────────────────────────────────────────────────────────
# Slide 09 — Grado Reconocido
# ─────────────────────────────────────────────────────────────────────────────
def gen_gradorec():
    if has_dot and len(dot197) > 0 and "GRADOREC" in dot197.columns:
        df = dot197[dot197["GRADOREC"].notna() &
                    (dot197["GRADOREC"].str.strip() != "NO INFORMA")].copy()
        RMAP = {"(MAG-PRO).":"Magíster Prof.","DOC":"Doctor",
                "(MAG-ACA)":"Magíster Acad.","PROFESIONAL":"Profesional",
                "POST-DOC":"Post-Doctor","TECNICO":"Técnico"}
        df["gl"] = df["GRADOREC"].str.strip().map(RMAP).fillna(df["GRADOREC"].str.strip())
        cnt = df["gl"].value_counts()
        n   = len(df)
        src = "GRADOREC"
    else:
        df = aptos[aptos["nivel_formacion"].notna()].copy()
        NF  = {"DOCTOR":"Doctor","MAGÍSTER O MASTER":"Magíster",
               "MÃSTER O MASTER":"Magíster","PROFESIONAL":"Profesional"}
        df["gl"] = df["nivel_formacion"].str.strip().map(NF).fillna(
                   df["nivel_formacion"].str.strip().str.title())
        df = df[df["gl"].notna() & (df["gl"] != "No Informa")]
        cnt = df["gl"].value_counts()
        n   = len(df)
        src = "nivel_formacion"

    lbl = cnt.index.tolist()
    val = cnt.values.tolist()
    pct = [100*v/n for v in val]

    fig = _base()
    _title(fig, "Grado Académico Reconocido")
    _pop(fig)
    _hbar(fig, lbl, val, pct)
    _ctitle(fig, f"Grado Academico — 197 Aptos P3  (n={n} con dato  ·  fuente: {src})")
    n_mag = sum(v for l,v in zip(lbl,val) if "agíster" in l or "agister" in l)
    n_doc = sum(v for l,v in zip(lbl,val) if "octor" in l)
    _bullets(fig, [
        f"El grado de Magíster (Profesional + Académico) es el más frecuente ({n_mag} doc.), "
        f"seguido por Doctor ({n_doc} doc.). El perfil de posgrado predomina entre los docentes formados.",
        "Los Doctores y Post-Doctores representan el núcleo de mayor calificación académica dentro de los Aptos P3.",
        f"Nota: {N197-n} de los 197 no tienen dato de grado registrado "
        f"(honorarios sin dotación completa). Los % son sobre n={n}.",
    ])
    return _save(fig, "09_gradorec.png")

# ─────────────────────────────────────────────────────────────────────────────
# Slide 10 — Institución del Grado
# ─────────────────────────────────────────────────────────────────────────────
def gen_institucion():
    if has_dot and len(dot197) > 0 and "INSTITUCIÓN GRADO TÍTULO" in dot197.columns:
        df   = dot197[dot197["INSTITUCIÓN GRADO TÍTULO"].notna() &
                      (dot197["INSTITUCIÓN GRADO TÍTULO"].str.strip() != "NO INFORMA")].copy()
        n    = len(df)
        cnt  = df["INSTITUCIÓN GRADO TÍTULO"].str.strip().value_counts()
        TOP  = 7
        top  = cnt.head(TOP)
        n_ot = n - top.sum()
        ABREV_I = {
            "UNIVERSIDAD CENTRAL DE CHILE SANTIAGO": "U. Central de Chile",
            "UNIVERSIDAD DE CHILE":                  "U. de Chile",
            "PONTIFICIA UNIVERSIDAD CATÓLICA DE CHILE": "PUC Chile",
            "UNIVERSIDAD ANDRÉS BELLO":              "U. Andrés Bello",
            "UNIVERSIDAD DE SANTIAGO DE CHILE":      "USACH",
            "UNIVERSIDAD MAYOR":                     "U. Mayor",
            "UNIVERSIDAD DIEGO PORTALES":            "UDP",
        }
        lbl = [ABREV_I.get(k, k[:30]) for k in top.index.tolist()]
        if n_ot > 0:
            lbl.append(f"Otras ({cnt.shape[0]-TOP} instituciones)")
        val = top.values.tolist() + ([n_ot] if n_ot > 0 else [])
        pct = [100*v/n for v in val]
    else:
        fig = _base()
        _title(fig, "Institución de Obtención del Grado")
        ax = fig.add_axes([PIC_RECT[0]+0.10, PIC_RECT[1]+0.15,
                           PIC_RECT[2]-0.20, PIC_RECT[3]-0.30],
                          facecolor="none", zorder=5)
        ax.axis("off")
        ax.text(0.5, 0.5, "Datos de institución\nno disponibles",
                ha="center", va="center", fontsize=18, color="#90ABC4",
                transform=ax.transAxes)
        return _save(fig, "10_institucion.png")

    fig = _base()
    _title(fig, "Institución de Obtención del Grado")
    _pop(fig)
    _hbar(fig, lbl, val, pct)
    _ctitle(fig, f"Institucion del Grado — Top {TOP}  (n={n} con dato informado)")
    _bullets(fig, [
        "La Universidad Central de Chile aparece como la principal institución de origen del grado, "
        "seguida por universidades del CRUCH (U. de Chile, PUC Chile, USACH).",
        "Una proporción de docentes obtuvo su grado en el extranjero, "
        "lo que enriquece la diversidad de formación del cuerpo docente.",
        f"Dato disponible para {n} de los {N197} Aptos P3 "
        "(solo docentes con dotación registrada en el sistema).",
    ])
    return _save(fig, "10_institucion.png")

# ─────────────────────────────────────────────────────────────────────────────
# Ensamblar PPTX
# ─────────────────────────────────────────────────────────────────────────────
def build_pptx(paths):
    prs = Presentation()
    prs.slide_width  = Emu(12192000)
    prs.slide_height = Emu(6858000)
    blank = prs.slide_layouts[6]
    for p in paths:
        sl = prs.slides.add_slide(blank)
        sl.shapes.add_picture(p, Emu(0), Emu(0), prs.slide_width, prs.slide_height)
    prs.save(OUT_PPTX)
    print(f"\n✓ PPTX: {OUT_PPTX}")
    print(f"  {len(prs.slides)} diapositivas")

# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("Generando diapositivas dark BLOQUE I...")
    paths = [
        gen_portada(),
        gen_indice(),
        gen_marco(),
        gen_separador(),
        gen_edad_sexo(),
        gen_facultad(),
        gen_antiguedad(),
        gen_jerarquia(),
        gen_gradorec(),
        gen_institucion(),
    ]
    build_pptx(paths)

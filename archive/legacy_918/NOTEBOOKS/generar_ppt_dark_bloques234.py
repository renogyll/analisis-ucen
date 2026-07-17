"""
generar_ppt_dark_bloques234.py
Genera BLOQUE II, III, IV y Conclusiones (36 slides) en formato oscuro UCEN.
Replica exactamente PRESENTACION_FINAL 30-06-2026 DPF (1).pptx pero:
  - formato dark (generador_ppt / poc_slide_facultad_dark.py)
  - universo rector = 197 Aptos P3
  - concatena con PRESENTACION_197_P3_dark_bloque1.pptx al final
"""
import sys; sys.stdout.reconfigure(encoding="utf-8")
import os, zipfile, textwrap, re
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.patheffects as pe
from matplotlib.patches import FancyBboxPatch
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Emu

# ─────────────────────────────────────────────────────────────────────────────
# Rutas
# ─────────────────────────────────────────────────────────────────────────────
BASE      = os.path.dirname(__file__)
SCRATCH   = (r"C:\Users\RGONZA~1.LAP\AppData\Local\Temp\claude"
             r"\c--Users-r-gonzalez-fluxsolar-LAPTOP-FLUX-ECO-Downloads-Analisis-UCEN-v2"
             r"\19e6fc3f-6ca1-4150-9da7-8dfa38be71ca\scratchpad")
FONDOTIPO = (r"c:\Users\r.gonzalez_fluxsolar.LAPTOP-FLUX-ECO"
             r"\Downloads\Analisis_UCEN_v2\Fondotipop.pptx")
ORIG_MEDIA = os.path.join(SCRATCH, "orig_media")
OUT_PPTX_B1 = (r"c:\Users\r.gonzalez_fluxsolar.LAPTOP-FLUX-ECO"
               r"\Downloads\Analisis_UCEN_v2\PRESENTACION_197_P3_dark_bloque1.pptx")
OUT_PPTX_B234 = (r"c:\Users\r.gonzalez_fluxsolar.LAPTOP-FLUX-ECO"
                 r"\Downloads\Analisis_UCEN_v2\PRESENTACION_197_P3_dark_bloques234.pptx")
OUT_PPTX_FINAL = (r"c:\Users\r.gonzalez_fluxsolar.LAPTOP-FLUX-ECO"
                  r"\Downloads\Analisis_UCEN_v2\PRESENTACION_197_P3_dark_COMPLETA.pptx")
OUT_DIR   = os.path.join(BASE, "dark_slides_b234")
os.makedirs(OUT_DIR, exist_ok=True)
os.makedirs(SCRATCH, exist_ok=True)

# ─────────────────────────────────────────────────────────────────────────────
# Assets: foto B&N + logo UCEN + gradiente navy
# ─────────────────────────────────────────────────────────────────────────────
BG_PATH   = os.path.join(SCRATCH, "fondotipo_image1.jpg")
LOGO_PATH = os.path.join(SCRATCH, "fondotipo_image2.png")
for path, zname in [(BG_PATH, "ppt/media/image1.jpg"), (LOGO_PATH, "ppt/media/image2.png")]:
    if not os.path.exists(path):
        with zipfile.ZipFile(FONDOTIPO) as z:
            with open(path, "wb") as f:
                f.write(z.read(zname))

with PILImage.open(BG_PATH) as _im:
    _rgb = _im.convert("RGB")
    _iw, _ih = _rgb.size
    _nh = int(_iw / (16/9))
    _y0 = min(int(_ih * 0.12), _ih - _nh)
    bg_arr = np.array(_rgb.crop((0, _y0, _iw, _y0 + _nh)))

with PILImage.open(LOGO_PATH) as _logo:
    logo_arr = np.array(_logo.convert("RGBA")).astype(np.float32) / 255.0

H_GRAD = 600
grad   = np.zeros((H_GRAD, 1, 4), dtype=np.float32)
for _r in range(H_GRAD):
    _t = _r / (H_GRAD - 1)
    _stops = [(0.00,(0,33,71)), (0.54,(0,70,128)), (1.00,(144,171,196))]
    for _i in range(len(_stops)-1):
        _t0,_c0 = _stops[_i]; _t1,_c1 = _stops[_i+1]
        if _t0 <= _t <= _t1:
            _s = (_t-_t0)/(_t1-_t0)
            grad[_r,0] = [(_c0[0]+_s*(_c1[0]-_c0[0]))/255,
                          (_c0[1]+_s*(_c1[1]-_c0[1]))/255,
                          (_c0[2]+_s*(_c1[2]-_c0[2]))/255, 0.82]
            break

# ─────────────────────────────────────────────────────────────────────────────
# Layout (EMU → fracción figura)
# ─────────────────────────────────────────────────────────────────────────────
SW, SH = 13.333, 7.5
def ex(e): return e / 12192000
def ey(e): return e / 6858000
def fig_rect(l, t, w, h): return (l, 1-t-h, w, h)

PIC_L,  PIC_T  = ex(786581),  ey(1125000)
PIC_W,  PIC_H  = ex(10599174), ey(3720000)
BUL_L,  BUL_T  = ex(786581),  ey(4870000)
BUL_H           = ey(1870000)
LOGO_L, LOGO_T = ex(9813773),  ey(656354)
LOGO_W, LOGO_H = ex(1756626),  ey(697725)
POP_T           = ey(820000)
POP_H           = ey(290000)

PIC_RECT  = fig_rect(PIC_L, PIC_T,  PIC_W,  PIC_H)
LOGO_RECT = fig_rect(LOGO_L, LOGO_T, LOGO_W, LOGO_H)
POP_Y     = 1 - POP_T - POP_H / 2

CHART_X = PIC_RECT[0] + 0.16
CHART_W = PIC_RECT[2] - 0.16 - 0.02
CHART_Y = PIC_RECT[1] + 0.04
CHART_H = PIC_RECT[3] - 0.09

PAL = ["#5C9BD6","#64B5F6","#80DEEA","#A5D6A7","#FFB74D","#CE93D8","#90A4AE","#F48FB1"]

POP_197 = "Universo: 197 Aptos P3  ·  130 formados + 67 control  ·  SAT disponible baseline / durante / resultado"
POP_917 = "Universo: 917 docentes jerarquizados  →  base completa de análisis"
POP_357 = "Universo: 357 docentes formados  ·  38.9% de los 917 jerarquizados"
POP_816 = "Universo: 917 jerarquizados  →  816 con SAT disponible"
POP_834 = "Universo: 917 jerarquizados  →  834 con notas y SAT disponibles"
POP_486 = "Universo: 917 jerarquizados  →  486 con EDD disponible  (134 formados · 352 sin formación)"

# ─────────────────────────────────────────────────────────────────────────────
# Helpers base
# ─────────────────────────────────────────────────────────────────────────────
def _base():
    fig = plt.figure(figsize=(SW, SH), facecolor="#101820")
    fig.patch.set_facecolor("#101820")
    for z, arr in [(0, bg_arr), (1, grad)]:
        ax = fig.add_axes([0,0,1,1], zorder=z)
        ax.imshow(arr, extent=[0,1,0,1], aspect="auto", origin="upper")
        ax.set_xlim(0,1); ax.set_ylim(0,1); ax.axis("off")
    al = fig.add_axes([LOGO_RECT[0], LOGO_RECT[1], LOGO_RECT[2], LOGO_RECT[3]],
                      zorder=10, facecolor="none")
    al.imshow(logo_arr, aspect="auto"); al.axis("off"); al.patch.set_visible(False)
    return fig

def _title(fig, text, fontsize=20):
    ty = (1.0 + CHART_Y + CHART_H + 0.020) / 2 + 0.018
    fig.text(PIC_L + PIC_W/2, ty, text, ha="center", va="center",
             fontsize=fontsize, fontweight="bold", color="white",
             transform=fig.transFigure, zorder=4)

def _pop(fig, txt=POP_197):
    fig.text(PIC_RECT[0], POP_Y, txt, ha="left", va="center",
             fontsize=7.5, fontstyle="italic", color="#C8DCF0",
             transform=fig.transFigure, zorder=4)

def _ctitle(fig, txt, cx=None):
    x = (cx if cx is not None else CHART_X) - 0.005
    fig.text(x, CHART_Y + CHART_H + 0.008, txt,
             ha="left", va="bottom", fontsize=10, color="white",
             transform=fig.transFigure, zorder=7)

def _bullets(fig, items):
    BUL_BOTTOM = 1 - BUL_T - BUL_H
    AX_W       = CHART_X + CHART_W - PIC_L
    ab = fig.add_axes([PIC_L, BUL_BOTTOM, AX_W, BUL_H], facecolor="none", zorder=6)
    ab.set_xlim(0,1); ab.set_ylim(0,1); ab.axis("off"); ab.patch.set_visible(False)
    lh  = 0.023 / BUL_H
    gap = 0.012 / BUL_H
    y   = 1.0 - 0.040 / BUL_H
    for i, bul in enumerate(items):
        for line in textwrap.wrap(f"{i+1}.  {bul}", width=130, subsequent_indent="    "):
            ab.text(0, y, line, ha="left", va="top", fontsize=11.5, color="white",
                    transform=ab.transAxes, clip_on=True)
            y -= lh
        y -= gap

def _save(fig, name):
    path = os.path.join(OUT_DIR, name)
    plt.savefig(path, format="png", dpi=150, facecolor=fig.get_facecolor())
    plt.close()
    print(f"  ✓ {name}")
    return path

def _img(name):
    """Ruta al PNG preferido en NOTEBOOKS, con fallback a orig_media."""
    nb  = os.path.join(BASE, name)
    med = os.path.join(ORIG_MEDIA, name)
    if os.path.exists(nb):  return nb
    if os.path.exists(med): return med
    return None

def _load_img(path):
    with PILImage.open(path) as im:
        return np.array(im.convert("RGBA"))

def _embed_chart(fig, img_path, cx=None, cy=None, cw=None, ch=None):
    """Embebe una imagen existente en el área de chart del slide dark."""
    x = cx if cx is not None else CHART_X
    y = cy if cy is not None else CHART_Y
    w = cw if cw is not None else CHART_W
    h = ch if ch is not None else CHART_H
    ax = fig.add_axes([x, y, w, h], zorder=5, facecolor="none")
    arr = _load_img(img_path)
    ax.imshow(arr, aspect="auto", origin="upper")
    ax.axis("off")
    ax.patch.set_visible(False)
    return ax

def _sep_cascade(fig, title_txt, items, pop_txt=None):
    """Genera una diapositiva separadora tipo cascade con dot leaders."""
    _title(fig, title_txt, fontsize=18)
    if pop_txt:
        _pop(fig, pop_txt)

    def _dl(text, w=84):
        m = re.match(r'^(.*?)\s{3,}(\(.+\))\s*$', text.strip())
        if not m: return text
        l, r = m.group(1).rstrip(), m.group(2)
        dots = max(3, w - len(l) - len(r) - 2)
        return f"{l} {'.'*dots} {r}"

    ax = fig.add_axes([PIC_RECT[0]+0.01, PIC_RECT[1]+0.03,
                       PIC_RECT[2]-0.02, PIC_RECT[3]-0.05],
                      facecolor="none", zorder=5)
    ax.set_xlim(0,1); ax.set_ylim(0,1); ax.axis("off")
    top, step = 0.96, 0.12
    for i, item in enumerate(items):
        is_sub = item.startswith("  ")
        prefix = "  │  " if is_sub else ""
        ax.text(0, top, prefix + _dl(item.strip()), ha="left", va="top",
                fontsize=10.5, color="white", fontfamily="monospace",
                transform=ax.transAxes)
        top -= step
    return fig

# ─────────────────────────────────────────────────────────────────────────────
# Slide genérico con imagen embebida
# ─────────────────────────────────────────────────────────────────────────────
def gen_chart_slide(name, title, img_path, bullets, pop_txt=POP_197, ctitle=None):
    fig = _base()
    _title(fig, title)
    _pop(fig, pop_txt)
    if img_path and os.path.exists(img_path):
        _embed_chart(fig, img_path)
    if ctitle:
        _ctitle(fig, ctitle)
    if bullets:
        _bullets(fig, bullets)
    return _save(fig, name)

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE II — Slide 01: Separador cascada BLOQUE II
# ─────────────────────────────────────────────────────────────────────────────
def gen_b2_sep():
    items = [
        "917 docentes jerarquizados",
        "  ├── 357 participaron en formacion (2.1)   ....... Cobertura · Evolucion · Venn · Intensidad · Brechas",
        "  │     └── 615 iniciativas de formacion totales",
        "  ├── 816 con SAT disponible (2.3)           ....... Validacion SAT y Recomendacion",
        "  └── 197 aptos P3 (SAT baseline + resultado valido)",
        "    ├── 184 puros / 13 mixtos (2.2)      ...... Venn Poblacion · Contraste Puro vs Mixto",
        "    └── 129 con antiguedad (2.4)          ...... Antiguedad x Tipo · Evolucion por Jerarquia",
    ]
    fig = _base()
    _title(fig, "BLOQUE II — Evaluación Docente Antes y Después", fontsize=18)
    _pop(fig, POP_197)

    ax = fig.add_axes([PIC_RECT[0]+0.01, PIC_RECT[1]+0.03,
                       PIC_RECT[2]-0.02, PIC_RECT[3]-0.05],
                      facecolor="none", zorder=5)
    ax.set_xlim(0,1); ax.set_ylim(0,1); ax.axis("off")
    top, step = 0.96, 0.13
    for item in items:
        ax.text(0, top, item, ha="left", va="top",
                fontsize=10.5, color="white", fontfamily="monospace",
                transform=ax.transAxes)
        top -= step
    return _save(fig, "11_b2_sep.png")

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE II — Slide 02: Mapa de Bloque II
# ─────────────────────────────────────────────────────────────────────────────
def gen_mapa_b2():
    fig = _base()
    _title(fig, "Mapa de Bloque II: Universos y Gráficas por Subsección")
    _pop(fig, POP_917)

    nodos = [
        (0.50, 0.82, "917\nJerarquizados", "#1E3A5F", 0.13, 0.08),
        (0.22, 0.60, "560 Control (61%)\nno participó en formación", "#1A2744", 0.16, 0.07),
        (0.72, 0.60, "357 Formados (39%)\n2.1 · 9 gráficas", "#1E4D8C", 0.16, 0.07),
        (0.50, 0.42, "816 con SAT disponible\n2.3 · 3 gráficas", "#1A3A6A", 0.18, 0.07),
        (0.50, 0.24, "197 Aptos P3\nSAT baseline + resultado", "#0D4B8C", 0.16, 0.07),
        (0.28, 0.08, "184 Puros / 13 Mixtos\n2.2 · 4 gráficas", "#1A3A6A", 0.16, 0.06),
        (0.72, 0.08, "129 con antigüedad\n2.4 · 3 gráficas", "#1A3A6A", 0.16, 0.06),
    ]

    ax = fig.add_axes([PIC_RECT[0]+0.01, PIC_RECT[1]+0.01,
                       PIC_RECT[2]-0.02, PIC_RECT[3]-0.02],
                      facecolor="none", zorder=5)
    ax.set_xlim(0,1); ax.set_ylim(0,1); ax.axis("off")

    # Líneas de conexión
    conexiones = [
        (0.50, 0.78, 0.22, 0.64),
        (0.50, 0.78, 0.72, 0.64),
        (0.50, 0.78, 0.50, 0.46),
        (0.50, 0.39, 0.50, 0.28),
        (0.50, 0.21, 0.28, 0.11),
        (0.50, 0.21, 0.72, 0.11),
    ]
    for x1, y1, x2, y2 in conexiones:
        ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle="-", color="white", lw=1.2, alpha=0.5),
                    transform=ax.transAxes)

    for xc, yc, lbl, color, bw, bh in nodos:
        rect = FancyBboxPatch((xc-bw/2, yc-bh/2), bw, bh,
                              boxstyle="round,pad=0.01", linewidth=0,
                              facecolor=color, alpha=0.85, transform=ax.transAxes, zorder=2)
        ax.add_patch(rect)
        ax.text(xc, yc, lbl, ha="center", va="center", fontsize=9.5,
                color="white", fontweight="bold", transform=ax.transAxes, zorder=3)

    return _save(fig, "12_mapa_b2.png")

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE II — Slide 03: Embudo P3
# ─────────────────────────────────────────────────────────────────────────────
def gen_embudo():
    fig = _base()
    _title(fig, "Universo de Análisis: El Embudo P3")
    _pop(fig, POP_197)

    pasos = [
        ("917", "docentes jerarquizados UCEN", "universo base total"),
        ("357", "participaron en formación", "39% del universo — grupo tratamiento"),
        ("560", "grupo control", "61% sin ninguna iniciativa de formación"),
        ("197", "Aptos P3", "tienen SAT válido baseline + durante + resultado"),
    ]

    ax = fig.add_axes([PIC_RECT[0]+0.05, PIC_RECT[1]+0.03,
                       PIC_RECT[2]-0.10, PIC_RECT[3]-0.06],
                      facecolor="none", zorder=5)
    ax.set_xlim(0,1); ax.set_ylim(0,1); ax.axis("off")

    widths = [0.90, 0.65, 0.75, 0.45]
    tops   = [0.90, 0.67, 0.44, 0.21]
    colors = ["#1E4D8C", "#1E88E5", "#1565C0", "#0D7C3A"]

    for i, ((n, lbl, desc), w, top, color) in enumerate(zip(pasos, widths, tops, colors)):
        xc = 0.5
        rect = FancyBboxPatch((xc-w/2, top-0.16), w, 0.18,
                              boxstyle="round,pad=0.01", linewidth=0,
                              facecolor=color, alpha=0.80, transform=ax.transAxes)
        ax.add_patch(rect)
        ax.text(xc-w/2+0.03, top-0.07, n, ha="left", va="center",
                fontsize=28, fontweight="bold", color="white", transform=ax.transAxes)
        ax.text(xc-w/2+0.16, top-0.065, lbl, ha="left", va="center",
                fontsize=12, fontweight="bold", color="white", transform=ax.transAxes)
        ax.text(xc-w/2+0.16, top-0.12, f"→ {desc}", ha="left", va="center",
                fontsize=9.5, color="#C8DCF0", fontstyle="italic", transform=ax.transAxes)
        if i < len(pasos)-1:
            ax.annotate("", xy=(0.5, tops[i+1]+0.02), xytext=(0.5, top-0.16),
                        arrowprops=dict(arrowstyle="-|>", color="white", lw=1.5, alpha=0.5),
                        transform=ax.transAxes)

    _bullets(fig, [
        "De los 917 docentes jerarquizados, 357 participaron en ≥1 iniciativa de formación (39%, grupo tratamiento); "
        "560 son grupo control (61%).",
        "Los 197 Aptos P3 son el subconjunto analizable: tienen SAT válido en baseline, durante y después de la formación. "
        "Representan el 21% del universo total.",
        "El detalle de combinaciones de modalidad (puro/mixto) se presenta en las secciones 2.1 y 2.2.",
    ])
    return _save(fig, "13_embudo.png")

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE II — Slide 04: Separador 2.1 Participación
# ─────────────────────────────────────────────────────────────────────────────
def gen_sep_21():
    items = [
        "Evolución 2022-2025: Participación en Formación                (357 docentes)",
        "Cobertura Global de Participación                              (357 de 917, 38.9%)",
        "Participación según Jerarquía Académica                        (917 jerarquizados)",
        "Modalidad de Formación según Jerarquía                        (357 formados)",
        "Participación en Formación por Antigüedad                      (491 con antigüedad)",
        "Distribución por Modalidad de Formación                        (615 iniciativas)",
        "Brechas: Perfil de No Participantes                            (560 sin formación)",
        "Intensidad de Participación Global                             (357 formados)",
        "Combinaciones de Formación — Universo 917                      (357 formados)",
    ]
    fig = _base()
    _sep_cascade(fig, "2.1   Participación en Formación", items, pop_txt=POP_357)
    return _save(fig, "14_sep21.png")

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE II — Slides 05-12: gráficas 2.1
# ─────────────────────────────────────────────────────────────────────────────
def gen_evolucion():
    img = _img("p2_g12_evolucion.png") or os.path.join(ORIG_MEDIA, "media_image4.jpeg")
    return gen_chart_slide("15_evolucion.png",
        "Evolución 2022-2025: Participación en Formación",
        img,
        ["La participación creció de 39 docentes (2022) a 211 (2024), un aumento de +172 en dos años. "
         "2025 muestra 86 docentes únicos (cifra parcial, año en curso).",
         "La tendencia sostenida indica una estrategia institucional de crecimiento progresivo en formación docente.",
         "El peak de 2024 concentra más de la mitad de la participación acumulada total del período analizado."],
        pop_txt=POP_357,
        ctitle="Evolucion 2022-2025 — Docentes únicos participantes por año")

def gen_cobertura():
    img = _img("p2_g11_cobertura.png") or os.path.join(ORIG_MEDIA, "media_image5.jpeg")
    return gen_chart_slide("16_cobertura.png",
        "Cobertura Global de Participación",
        img,
        ["De 917 docentes, 357 participaron en ≥1 iniciativa (38.9%). Un 61.1% (560 docentes) no participó en ninguna.",
         "Taller es la modalidad con mayor participación (215 doc., 60% de los formados); "
         "Diplomado alcanza 201 (56%); Proyecto 32 (9%). Total: 448 conteos con doble-conteo posible.",
         "La participación en formación mixta (más de una modalidad) indica docentes con alto compromiso formativo."],
        pop_txt=POP_357,
        ctitle="Cobertura Global — 357 de 917 jerarquizados (38.9%)")

def gen_jerarquia_p2():
    img = _img("p2_g21_jerarquia.png") or os.path.join(ORIG_MEDIA, "media_image6.png")
    return gen_chart_slide("17_jerarquia_p2.png",
        "Participación según Jerarquía Académica",
        img,
        ["El tipo de instancia varía según jerarquía: rangos superiores muestran mayor diversificación modal.",
         "Instructor Docente y Asistente Docente concentran su participación principalmente en Talleres — "
         "la modalidad de menor duración y mayor accesibilidad.",
         "Asociado y Titular muestran mayor participación relativa en Diplomados y Proyectos, "
         "modalidades de mayor profundidad y compromiso."],
        pop_txt=POP_917,
        ctitle="Participacion por Jerarquia — 917 jerarquizados")

def gen_antiguedad_p2():
    img = _img("p2_g23_antiguedad.png") or os.path.join(ORIG_MEDIA, "media_image7.png")
    return gen_chart_slide("18_antiguedad_p2.png",
        "Participación en Formación por Antigüedad",
        img,
        ["El tramo 0-4 años concentra la mayor tasa de participación (50%). "
         "Universo 917 → 493 perfil completo → 491 con antigüedad informada.",
         "HALLAZGO CLAVE: A mayor antigüedad, menor participación en formación. "
         "Los docentes con más experiencia participan menos en iniciativas de perfeccionamiento.",
         "Este patrón puede reflejar mayor carga académica en rangos superiores o menor incentivo "
         "para participar en formación estructurada."],
        pop_txt=POP_917,
        ctitle="Participacion por Antiguedad — 491 con antigüedad informada")

def gen_modalidad():
    img = _img("p2_g31_modalidad.png") or os.path.join(ORIG_MEDIA, "media_image8.png")
    return gen_chart_slide("19_modalidad.png",
        "Distribución por Modalidad de Formación",
        img,
        ["De las 615 iniciativas totales de formación, el Taller es la modalidad más frecuente. "
         "615 iniciativas de 357 docentes implica promedio de ~1.7 iniciativas por docente.",
         "La distribución por modalidad refleja la oferta institucional: Talleres de corta duración "
         "tienen mayor cobertura; Proyectos son minoritarios por su mayor exigencia.",
         "La diversificación de modalidades busca atender distintos perfiles de docente y necesidades formativas."],
        pop_txt=POP_357,
        ctitle="Modalidad de Formacion — 615 iniciativas de 357 docentes")

def gen_brechas():
    img = _img("p2_g41_brechas.png") or os.path.join(ORIG_MEDIA, "media_image9.png")
    return gen_chart_slide("20_brechas.png",
        "Brechas: Perfil de No Participantes",
        img,
        ["El 61% de los docentes (560 de 917) nunca participó en ninguna iniciativa de formación. "
         "Estos son el grupo control y el principal desafío de cobertura.",
         "El perfil de no participantes tiende a tener mayor antigüedad y rangos superiores del escalafón, "
         "donde la tasa de participación cae sistemáticamente.",
         "HALLAZGO CLAVE: 6 de cada 10 docentes jerarquizados no ha accedido a formación — "
         "la principal oportunidad de mejora institucional."],
        pop_txt=POP_917,
        ctitle="Perfil de No Participantes — 560 sin formacion (61% del universo)")

def gen_intensidad():
    img = _img("p2_g42_intensidad.png") or os.path.join(ORIG_MEDIA, "media_image10.png")
    return gen_chart_slide("21_intensidad.png",
        "Intensidad de Participación Global",
        img,
        ["La mayoría de los docentes formados participó en 1-2 iniciativas. "
         "Un subgrupo reducido acumuló 4+ iniciativas (heavy users).",
         "La intensidad de participación es relevante: el efecto de la formación se vuelve evidente "
         "recién a partir de 3+ iniciativas acumuladas.",
         "Los heavy users (4+ iniciativas) muestran los mejores resultados SAT y EDD, "
         "sugiriendo un efecto dosis-respuesta de la formación docente."],
        pop_txt=POP_357,
        ctitle="Intensidad de Participacion — 357 docentes formados")

def gen_combinaciones():
    img = _img("G_venn_formacion_917_918.png") or os.path.join(ORIG_MEDIA, "media_image11.png")
    return gen_chart_slide("22_combinaciones.png",
        "Combinaciones de Formación — Universo 917",
        img,
        ["La mayor parte de los formados participó exclusivamente en Taller o Diplomado. "
         "Las combinaciones mixtas (Taller+Diplomado, Taller+Proyecto, etc.) son un subgrupo.",
         "Los docentes con combinaciones múltiples representan el subgrupo de mayor compromiso formativo "
         "y son los que muestran mejores resultados en las métricas de evaluación.",
         "La distribución refleja la estructura de la oferta: Talleres son de fácil acceso, "
         "Proyectos requieren mayor dedicación."],
        pop_txt=POP_357,
        ctitle="Combinaciones de Formacion — Venn de modalidades (917 universo)")

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE II — Slide 13: Separador 2.2 Pura vs Mixta
# ─────────────────────────────────────────────────────────────────────────────
def gen_sep_22():
    items = [
        "Venn: Población Pura vs Mixta                                  (197 aptos P3)",
        "Contraste SAT Nota: Normal vs Puro                             (197 aptos P3)",
        "SAT y Recomendación — Población Pura                           (184 puros de 197)",
        "Evolución por Jerarquía — Diplomado                            (puros subgrupo Diplomado)",
        "Evolución por Jerarquía — Taller                               (puros subgrupo Taller)",
        "Δ Z-score: Puro vs Mixto                                       (184 puros / 13 mixtos)",
    ]
    fig = _base()
    _sep_cascade(fig, "2.2   Población Pura vs Mixta", items, pop_txt=POP_197)
    return _save(fig, "23_sep22.png")

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE II — Slides 14-19: gráficas 2.2
# ─────────────────────────────────────────────────────────────────────────────
def gen_venn_pob():
    img = _img("G_venn_poblaciones_918.png") or os.path.join(ORIG_MEDIA, "media_image12.png")
    return gen_chart_slide("24_venn_pob.png",
        "Venn: Población Pura vs Mixta",
        img,
        ["De los 197 Aptos P3: 184 son PUROS (participaron en solo una modalidad de formación) "
         "y 13 son MIXTOS (combinaron dos o más modalidades).",
         "La separación puro/mixto permite aislar el efecto específico de cada modalidad, "
         "eliminando la contaminación de combinaciones.",
         "Los 13 mixtos representan el 6.6% — subgrupo pequeño pero analíticamente valioso "
         "para comparar el efecto acumulado vs. el efecto simple."],
        pop_txt=POP_197,
        ctitle="Venn Poblacion — 197 Aptos P3: 184 puros / 13 mixtos")

def gen_contraste_normal():
    img = _img("G1_contraste_normal_puro_918.png") or os.path.join(ORIG_MEDIA, "media_image13.png")
    return gen_chart_slide("25_contraste_normal.png",
        "Contraste SAT Nota: Normal vs Puro",
        img,
        ["La comparación Normal vs Puro muestra patrones similares: el efecto formativo se mantiene "
         "al excluir los docentes mixtos del análisis.",
         "Esto confirma la robustez del hallazgo principal: el efecto positivo no depende de docentes "
         "que combinaron modalidades.",
         "Los 197 (normal/completo) y los 184 (puro) muestran curvas de evolución SAT prácticamente "
         "idénticas, validando el diseño de análisis."],
        pop_txt=POP_197,
        ctitle="Contraste SAT Nota: Normal (197) vs Puro (184)")

def gen_sat_pura():
    img = _img("G1_puro_nota_recom_918.png") or os.path.join(ORIG_MEDIA, "media_image14.png")
    return gen_chart_slide("26_sat_pura.png",
        "SAT y Recomendación — Población Pura",
        img,
        ["Mismos patrones en ambas métricas para la población pura (184 docentes): "
         "Proyecto cae más, Taller es estable.",
         "La evolución SAT Nota y SAT % Recomendación son consistentes — validación cruzada "
         "entre dos métricas independientes de evaluación docente.",
         "Universo: 917 → 197 aptos P3 → 184 puros. NUEVO hallazgo de esta pasada de análisis."],
        pop_txt=POP_197,
        ctitle="SAT y Recomendacion — Poblacion Pura (184 docentes)")

def gen_evol_diplomado():
    img = _img("G7_diplomado_jerarquia_918.png") or os.path.join(ORIG_MEDIA, "media_image15.png")
    return gen_chart_slide("27_evol_diplomado.png",
        "Evolución por Jerarquía — Diplomados",
        img,
        ["Instructores Docentes (n=18) son el caso de mayor éxito: parten bajo la media, "
         "escalan durante el diplomado y consolidan al alza.",
         "Asistentes Docentes (n=12) se mantienen sobre el promedio en todos los períodos en ambas métricas.",
         "El ANOVA no detecta diferencias significativas entre jerarquías (p=0.718), "
         "pero el patrón descriptivo para Instructores es destacable y merece seguimiento."],
        pop_txt=POP_197,
        ctitle="Evolucion por Jerarquia — Subgrupo Diplomado (puros)")

def gen_evol_taller():
    img = _img("G7_taller_jerarquia_918.png") or os.path.join(ORIG_MEDIA, "media_image16.png")
    return gen_chart_slide("28_evol_taller.png",
        "Evolución por Jerarquía — Talleres",
        img,
        ["Asistentes Docentes (n=52) parten bajo la media y alcanzan su peak durante los talleres (2024), "
         "cerrando 2025-02 en posición positiva.",
         "Instructores Docentes (n=44) parten sobre el promedio y se mantienen estables — "
         "sin gran ganancia adicional pero sin deterioro.",
         "Asociados Docentes (n=35) permanecen bajo la media en todos los períodos — "
         "el taller no logra reposicionarlos."],
        pop_txt=POP_197,
        ctitle="Evolucion por Jerarquia — Subgrupo Taller (puros)")

def gen_delta_z():
    img = _img("G_delta_z_puro_vs_mixto_918.png") or os.path.join(ORIG_MEDIA, "media_image17.png")
    return gen_chart_slide("29_delta_z.png",
        "Δ Z-score: Puro vs Mixto",
        img,
        ["Los 13 mixtos muestran mejor resultado (62% mejoró) que los puros (49%). "
         "La combinación de modalidades parece potenciar el efecto formativo.",
         "Sin embargo, el n=13 mixtos es insuficiente para inferencia estadística robusta. "
         "El resultado es descriptivo y debe interpretarse con cautela.",
         "Universo: 917 → 197 aptos P3 → 184 puros / 13 mixtos. NUEVO análisis de esta pasada."],
        pop_txt=POP_197,
        ctitle="Delta Z-score: Puro (184) vs Mixto (13)")

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE II — Slide 20: Separador 2.3 SAT
# ─────────────────────────────────────────────────────────────────────────────
def gen_sep_23():
    items = [
        "Evolución SAT por Período: Formados vs Control  (6 períodos)  ★ hallazgo clave",
        "Evolución SAT por Período por Facultad                         ★ hallazgo clave",
    ]
    fig = _base()
    _sep_cascade(fig, "2.3   SAT y Recomendación — Formados vs Control", items, pop_txt=POP_816)
    return _save(fig, "30_sep23.png")

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE II — Slide 21: SAT doble (Nota + Recomendación)
# ─────────────────────────────────────────────────────────────────────────────
def gen_sat_doble():
    img_nota = _img("G6_control_vs_tratamiento_918.png") or os.path.join(ORIG_MEDIA, "media_image19.png")
    img_rec  = _img("G6.2_control_bin_918.png")          or os.path.join(ORIG_MEDIA, "media_image18.png")

    fig = _base()
    _title(fig, "Evolución SAT por Período: Formados vs Control (6 períodos)")
    _pop(fig, POP_816)

    # Dos gráficos lado a lado
    half_w = CHART_W / 2 - 0.01
    if img_nota and os.path.exists(img_nota):
        _embed_chart(fig, img_nota, cx=CHART_X, cw=half_w)
    if img_rec and os.path.exists(img_rec):
        _embed_chart(fig, img_rec, cx=CHART_X + half_w + 0.02, cw=half_w)

    fig.text(CHART_X, CHART_Y + CHART_H + 0.008,
             "SAT Nota (z-score)", ha="left", va="bottom",
             fontsize=9.5, color="white", transform=fig.transFigure, zorder=7)
    fig.text(CHART_X + half_w + 0.02, CHART_Y + CHART_H + 0.008,
             "SAT % Recomendación (z-score)", ha="left", va="bottom",
             fontsize=9.5, color="white", transform=fig.transFigure, zorder=7)

    _bullets(fig, [
        "HALLAZGO CLAVE: Los docentes formados superan al grupo control en SAT los 6 períodos "
        "sin excepción — el beneficio no es un evento de un solo semestre.",
        "Brecha promedio en nota: +0.18z, estadísticamente significativa en 3 de 6 períodos. "
        "En recomendación la brecha final es +0.155z.",
        "Grupo Control muestra tendencia plana a negativa — los formados consolidan su ventaja progresivamente."],
    )
    return _save(fig, "31_sat_doble.png")

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE II — Slide 22: Control SAT por facultad
# ─────────────────────────────────────────────────────────────────────────────
def gen_control_fac():
    img = _img("G1_linea_z_918.png") or os.path.join(ORIG_MEDIA, "media_image20.png")
    return gen_chart_slide("32_control_fac.png",
        "SAT Formados vs Control — Detalle por Facultad",
        img,
        ["El grupo de tratamiento está consistentemente sobre el promedio de su facultad (z positivo) "
         "en la mayoría de facultades.",
         "Esto no prueba causalidad. Los docentes que se capacitan pueden ser sistemáticamente distintos "
         "(más motivados, más activos). El diseño observacional no permite establecer causalidad.",
         "Sin embargo, la consistencia del patrón a través de 6 facultades distintas refuerza "
         "la interpretación de un efecto real de la formación."],
        pop_txt=POP_816,
        ctitle="SAT Formados vs Control por Facultad — z-score")

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE II — Slide 23: Separador 2.4
# ─────────────────────────────────────────────────────────────────────────────
def gen_sep_24():
    items = [
        "Antigüedad × Tipo de formacion — SAT Nota y Recomendacion     (129 con antiguedad)",
        "Evolución por Jerarquía — Talleres                             (subgrupo Taller)",
        "Evolución por Jerarquía — Diplomados                           (subgrupo Diplomado)",
    ]
    fig = _base()
    _sep_cascade(fig, "2.4   Comparación por Tipo, Jerarquía y Antigüedad", items, pop_txt=POP_197)
    return _save(fig, "33_sep24.png")

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE II — Slides 24-26: gráficas 2.4
# ─────────────────────────────────────────────────────────────────────────────
def gen_ant_tipo():
    img = _img("G3_heatmap_antiguedad_918.png") or os.path.join(ORIG_MEDIA, "media_image21.png")
    return gen_chart_slide("34_ant_tipo.png",
        "Antigüedad × Tipo — SAT Nota y Recomendación",
        img,
        ["129 de 197 aptos P3 tienen dato de antigüedad disponible. "
         "Universo: 917 → 197 aptos P3 → 129 con antigüedad informada.",
         "El heatmap permite identificar qué combinación antigüedad × modalidad produce "
         "los mejores resultados en SAT nota y SAT recomendación.",
         "NUEVO análisis de esta pasada. Permite la detección de nichos de alta efectividad formativa."],
        pop_txt=POP_197,
        ctitle="Antiguedad x Tipo — Heatmap SAT Nota y Recomendacion (n=129)")

def gen_evol_jer_taller():
    img = _img("G7_taller_jerarquia_918.png") or os.path.join(ORIG_MEDIA, "media_image22.png")
    return gen_chart_slide("35_evol_jer_taller.png",
        "Evolución por Jerarquía — Talleres",
        img,
        ["Asistentes Docentes (n=52) parten bajo la media y alcanzan su peak durante talleres. "
         "Universo: 917 → 197 aptos P3 → subgrupo Taller.",
         "Los rangos intermedios (Asociado) no logran reposicionarse con el taller. "
         "El efecto es más pronunciado en los rangos de entrada al escalafón.",
         "NUEVO hallazgo de esta pasada de análisis con universo 197 Aptos P3."],
        pop_txt=POP_197,
        ctitle="Evolucion por Jerarquia — Subgrupo Taller")

def gen_evol_jer_diplomado():
    img = _img("G7_diplomado_jerarquia_918.png") or os.path.join(ORIG_MEDIA, "media_image23.png")
    return gen_chart_slide("36_evol_jer_diplomado.png",
        "Evolución por Jerarquía — Diplomados",
        img,
        ["Instructores Docentes (n=18): caso de mayor éxito. "
         "Parten bajo la media, escalan durante el diplomado y consolidan al alza. "
         "Universo: 917 → 197 aptos P3 → subgrupo Diplomado.",
         "El Diplomado parece tener efecto más diferenciado entre jerarquías que el Taller.",
         "NUEVO análisis con universo 197 Aptos P3."],
        pop_txt=POP_197,
        ctitle="Evolucion por Jerarquia — Subgrupo Diplomado")

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE III — Slide 27: Separador
# ─────────────────────────────────────────────────────────────────────────────
def gen_b3_sep():
    items = [
        "Notas y Tasas de Aprobacion según SAT                         (834 con notas y SAT)",
        "  ├── Formados vs Control — Aprobacion (%)                    (834)",
        "  ├── Tasa de Aprobacion × Antiguedad Docente                 (834)",
        "  ├── Tasa de Aprobacion × Jerarquia Academica                (834)",
        "  └── ¿Es Acumulativo el Efecto de la Formacion?              (816 con SAT y notas)",
    ]
    fig = _base()
    _title(fig, "BLOQUE III — Resultados de Calificaciones", fontsize=18)
    _pop(fig, POP_834)

    ax = fig.add_axes([PIC_RECT[0]+0.01, PIC_RECT[1]+0.03,
                       PIC_RECT[2]-0.02, PIC_RECT[3]-0.05],
                      facecolor="none", zorder=5)
    ax.set_xlim(0,1); ax.set_ylim(0,1); ax.axis("off")
    ax.text(0, 0.96, "917 docentes jerarquizados", ha="left", va="top",
            fontsize=11, color="white", fontfamily="monospace", transform=ax.transAxes)
    ax.text(0, 0.84, "  │", ha="left", va="top",
            fontsize=11, color="white", fontfamily="monospace", transform=ax.transAxes)
    ax.text(0, 0.74, "  └── 834 con notas y SAT disponibles", ha="left", va="top",
            fontsize=11, color="white", fontfamily="monospace", transform=ax.transAxes)
    for i, item in enumerate(items):
        ax.text(0.06, 0.62 - i*0.12, item, ha="left", va="top",
                fontsize=10.5, color="white", fontfamily="monospace", transform=ax.transAxes)
    return _save(fig, "37_b3_sep.png")

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE III — Slides 28-32
# ─────────────────────────────────────────────────────────────────────────────
def gen_aprobacion1():
    img = _img("GN_aprobacion_918.png") or os.path.join(ORIG_MEDIA, "media_image24.jpeg")
    return gen_chart_slide("38_aprobacion1.png",
        "Aprobación: Formados vs Control",
        img,
        ["Los docentes formados logran +1.5 puntos porcentuales más de aprobación en los alumnos "
         "respecto al grupo control.",
         "A partir del período 2024-01 la brecha se consolida notablemente. "
         "El efecto es más claro en los semestres posteriores a la formación.",
         "HALLAZGO CLAVE: La formación docente tiene un efecto positivo en las calificaciones "
         "de los alumnos, más allá de la percepción docente."],
        pop_txt=POP_834,
        ctitle="Aprobacion (%) — Formados vs Control")

def gen_aprobacion2():
    img = os.path.join(ORIG_MEDIA, "media_image25.jpeg")
    if not os.path.exists(img):
        img = _img("GN_aprobacion_918.png")
    return gen_chart_slide("39_aprobacion2.png",
        "Aprobación: Tendencia y Brecha Consolidada",
        img,
        ["La brecha en aprobación a favor del grupo formado se consolida en el tiempo — "
         "el cierre del estudio (2025-01) muestra la mayor distancia entre grupos.",
         "Los formados logran ubicarse de manera más consistente en el cuadrante de "
         "'bien evaluados y con altas notas' según el análisis bidimensional.",
         "Los resultados son consistentes con los hallazgos SAT: la formación tiene efecto "
         "transversal tanto en evaluación docente como en desempeño estudiantil."],
        pop_txt=POP_834,
        ctitle="Tendencia de Aprobacion — Consolidacion del efecto formativo")

def gen_aprob_ant():
    img = _img("G_aprobacion_antiguedad_918.png") or os.path.join(ORIG_MEDIA, "media_image26.png")
    return gen_chart_slide("40_aprob_ant.png",
        "Tasa de Aprobación × Antigüedad Docente",
        img,
        ["La tasa de aprobación es relativamente estable entre tramos de antigüedad. "
         "No hay un efecto antigüedad marcado en las calificaciones de los alumnos.",
         "Universo: 917 → 834 docentes con notas y SAT disponibles. NUEVO análisis.",
         "La estabilidad sugiere que la antigüedad docente no es un predictor fuerte "
         "de los resultados de sus alumnos — la formación es más determinante que la experiencia."],
        pop_txt=POP_834,
        ctitle="Tasa de Aprobacion x Antiguedad (n=834)")

def gen_aprob_jer():
    img = _img("G_aprobacion_jerarquia_918.png") or os.path.join(ORIG_MEDIA, "media_image27.png")
    return gen_chart_slide("41_aprob_jer.png",
        "Tasa de Aprobación × Jerarquía Académica",
        img,
        ["Instructor Docente lidera con la mediana más alta de aprobación (94.6%). "
         "Las jerarquías superiores no muestran ventaja sistemática.",
         "Universo: 917 → 834 docentes con notas y SAT disponibles. NUEVO análisis.",
         "Los rangos de entrada al escalafón combinan alta energía formativa y buenos resultados "
         "estudiantiles — reforzando la prioridad de formación en este segmento."],
        pop_txt=POP_834,
        ctitle="Tasa de Aprobacion x Jerarquia (n=834)")

def gen_acumulativo():
    img = _img("G_acumulacion_formacion_918.png") or os.path.join(ORIG_MEDIA, "media_image28.png")
    return gen_chart_slide("42_acumulativo.png",
        "¿Es Acumulativo el Efecto de la Formación?",
        img,
        ["A partir de 3 iniciativas de formación la mediana de aprobación supera al grupo sin formación. "
         "Universo: 917 → 816 con SAT y notas.",
         "HALLAZGO CLAVE: El impacto de la formación se vuelve evidente recién a partir de 3 iniciativas — "
         "la participación única o doble no produce el mismo efecto.",
         "NUEVO análisis de esta pasada. Implica que las políticas de formación deben apuntar a la "
         "intensidad y no solo a la cobertura."],
        pop_txt=POP_834,
        ctitle="Efecto Acumulativo — Aprobacion por cantidad de iniciativas (n=816)")

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE IV — Slide 33: Separador
# ─────────────────────────────────────────────────────────────────────────────
def gen_b4_sep():
    items = [
        "Evolución EDD 2022-2025                                        (486 con EDD disponible)",
        "  ├── 134 formados con EDD",
        "  └── 352 sin formacion con EDD",
        "EDD — Poblacion Mixta vs Pura                                  (197 aptos P3 con EDD)",
    ]
    fig = _base()
    _title(fig, "BLOQUE IV — Evaluación de Desempeño Docente (EDD)", fontsize=18)
    _pop(fig, POP_486)

    ax = fig.add_axes([PIC_RECT[0]+0.01, PIC_RECT[1]+0.03,
                       PIC_RECT[2]-0.02, PIC_RECT[3]-0.05],
                      facecolor="none", zorder=5)
    ax.set_xlim(0,1); ax.set_ylim(0,1); ax.axis("off")
    ax.text(0, 0.96, "917 docentes jerarquizados", ha="left", va="top",
            fontsize=11, color="white", fontfamily="monospace", transform=ax.transAxes)
    ax.text(0, 0.84, "  │", ha="left", va="top",
            fontsize=11, color="white", fontfamily="monospace", transform=ax.transAxes)
    ax.text(0, 0.74, "  └── 486 con EDD disponible", ha="left", va="top",
            fontsize=11, color="white", fontfamily="monospace", transform=ax.transAxes)
    for i, item in enumerate(items):
        ax.text(0.06, 0.62 - i*0.14, item, ha="left", va="top",
                fontsize=10.5, color="white", fontfamily="monospace", transform=ax.transAxes)
    return _save(fig, "43_b4_sep.png")

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE IV — Slides 34-35
# ─────────────────────────────────────────────────────────────────────────────
def gen_edd_evol():
    img = _img("G9_edd_evolucion_918.png") or os.path.join(ORIG_MEDIA, "media_image29.jpeg")
    return gen_chart_slide("44_edd_evol.png",
        "EDD: Evolución Formados vs Sin Formación",
        img,
        ["La brecha de +0.260 en 2025 entre formados (0.864) y sin formación (0.604) sugiere "
         "que la formación actúa como escudo protector frente a la evaluación de jefaturas.",
         "El grupo sin formación experimenta una caída libre desde 2023 (0.872) a 2025 (0.604) — "
         "perdiendo un 25% de su valoración en 3 años.",
         "HALLAZGO CLAVE: La formación actúa como escudo protector frente a la evaluación de jefes: "
         "el efecto más grande y contundente del estudio se observa en EDD."],
        pop_txt=POP_486,
        ctitle="Evolucion EDD 2022-2025 — Formados vs Sin Formacion")

def gen_edd_puro():
    img = _img("G9_puro_edd_918.png") or os.path.join(ORIG_MEDIA, "media_image30.png")
    return gen_chart_slide("45_edd_puro.png",
        "EDD — Población Mixta vs Pura",
        img,
        ["Diplomado puro sube de 0.823 a 0.854 al excluir los mixtos — "
         "el efecto del Diplomado sobre EDD es más limpio en población pura.",
         "Universo: 917 → 197 aptos P3 → subgrupo con EDD disponible. NUEVO análisis.",
         "La consistencia del efecto positivo del Diplomado en EDD (tanto en población completa "
         "como en pura) refuerza la robustez del hallazgo."],
        pop_txt=POP_197,
        ctitle="EDD — Comparacion Mixta vs Pura (197 aptos P3)")

# ─────────────────────────────────────────────────────────────────────────────
# Conclusiones — Slide 36
# ─────────────────────────────────────────────────────────────────────────────
def gen_conclusiones():
    conclusiones = [
        ("SAT y Recomendación (197 aptos P3)",
         "Los formados superan al control en los 6 períodos sin excepción. "
         "Brecha +0.18z en nota, +0.155z en recomendación. El beneficio es persistente."),
        ("Aprobación de Alumnos (834 con notas)",
         "+1.5 pp de aprobación a favor de los formados. La brecha se consolida en 2024-2025. "
         "Efecto cruzado: la formación mejora tanto la percepción como los resultados concretos."),
        ("Efecto Acumulativo (816 con SAT y notas)",
         "El impacto emerge recién a partir de 3+ iniciativas. "
         "Dosis-respuesta clara: intensidad importa más que cobertura."),
        ("EDD (486 con EDD disponible)",
         "El efecto más contundente del estudio: +0.26 en EDD en 2025. "
         "La formación actúa como escudo protector frente a evaluación de jefaturas."),
        ("Brechas de cobertura (917 universo)",
         "6 de cada 10 docentes no han accedido a formación. "
         "La oportunidad de mejora más grande es ampliar la cobertura, especialmente en rangos superiores."),
    ]

    fig = _base()
    _title(fig, "Conclusiones Consolidadas", fontsize=20)
    _pop(fig, "Universo: 917 jerarquizados  ·  197 aptos P3  ·  834 con notas  ·  486 con EDD")

    ax = fig.add_axes([PIC_RECT[0]+0.01, PIC_RECT[1]+0.01,
                       PIC_RECT[2]-0.02, PIC_RECT[3]-0.02],
                      facecolor="none", zorder=5)
    ax.set_xlim(0,1); ax.set_ylim(0,1); ax.axis("off")

    top, step = 0.95, 0.185
    colors_h = ["#64B5F6","#80DEEA","#A5D6A7","#FFB74D","#F48FB1"]
    for i, (header, body) in enumerate(conclusiones):
        ax.text(0, top, f"  {i+1}.  {header}", ha="left", va="top",
                fontsize=11, fontweight="bold", color=colors_h[i],
                transform=ax.transAxes)
        for j, line in enumerate(textwrap.wrap(body, width=115)):
            ax.text(0.06, top - 0.055 - j*0.048, line, ha="left", va="top",
                    fontsize=9.5, color="white", transform=ax.transAxes, clip_on=True)
        top -= step

    return _save(fig, "46_conclusiones.png")

# ─────────────────────────────────────────────────────────────────────────────
# Ensamblar PPTX de BLOQUES 2-3-4
# ─────────────────────────────────────────────────────────────────────────────
def build_pptx_b234(paths):
    prs = Presentation()
    prs.slide_width  = Emu(12192000)
    prs.slide_height = Emu(6858000)
    blank = prs.slide_layouts[6]
    for p in paths:
        sl = prs.slides.add_slide(blank)
        sl.shapes.add_picture(p, Emu(0), Emu(0), prs.slide_width, prs.slide_height)
    prs.save(OUT_PPTX_B234)
    print(f"\n✓ PPTX BLOQUES 2-3-4: {OUT_PPTX_B234}")
    print(f"  {len(prs.slides)} diapositivas")

# ─────────────────────────────────────────────────────────────────────────────
# Concatenar con BLOQUE I
# ─────────────────────────────────────────────────────────────────────────────
def concatenar_final():
    """Concatena BLOQUE I (10 slides) + BLOQUES 234 (36 slides) = 46 slides totales."""
    from pptx.util import Emu
    import shutil, zipfile

    # Leer los dos PPTX
    prs1 = Presentation(OUT_PPTX_B1)
    prs2 = Presentation(OUT_PPTX_B234)

    # Crear presentación final copiando desde bloque1
    shutil.copy2(OUT_PPTX_B1, OUT_PPTX_FINAL)
    prs_final = Presentation(OUT_PPTX_FINAL)

    blank = prs_final.slide_layouts[6]

    # Agregar slides de B234 copiando las imágenes
    slides_b234 = sorted([f for f in os.listdir(OUT_DIR) if f.endswith(".png")],
                         key=lambda x: int(x.split("_")[0]))

    for png_name in slides_b234:
        png_path = os.path.join(OUT_DIR, png_name)
        sl = prs_final.slides.add_slide(blank)
        sl.shapes.add_picture(png_path, Emu(0), Emu(0),
                               prs_final.slide_width, prs_final.slide_height)

    prs_final.save(OUT_PPTX_FINAL)
    print(f"\n✓ PPTX FINAL COMPLETA: {OUT_PPTX_FINAL}")
    print(f"  {len(prs_final.slides)} diapositivas totales")

# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("Generando diapositivas dark BLOQUES II, III, IV + Conclusiones...")
    print()

    paths_b234 = [
        # BLOQUE II
        gen_b2_sep(),       # 11
        gen_mapa_b2(),      # 12
        gen_embudo(),       # 13
        gen_sep_21(),       # 14
        gen_evolucion(),    # 15
        gen_cobertura(),    # 16
        gen_jerarquia_p2(), # 17
        gen_antiguedad_p2(),# 18
        gen_modalidad(),    # 19
        gen_brechas(),      # 20
        gen_intensidad(),   # 21
        gen_combinaciones(),# 22
        gen_sep_22(),       # 23
        gen_venn_pob(),     # 24
        gen_contraste_normal(), # 25
        gen_sat_pura(),     # 26
        gen_evol_diplomado(),   # 27
        gen_evol_taller(),  # 28
        gen_delta_z(),      # 29
        gen_sep_23(),       # 30
        gen_sat_doble(),    # 31
        gen_control_fac(),  # 32
        gen_sep_24(),       # 33
        gen_ant_tipo(),     # 34
        gen_evol_jer_taller(),      # 35
        gen_evol_jer_diplomado(),   # 36
        # BLOQUE III
        gen_b3_sep(),       # 37
        gen_aprobacion1(),  # 38
        gen_aprobacion2(),  # 39
        gen_aprob_ant(),    # 40
        gen_aprob_jer(),    # 41
        gen_acumulativo(),  # 42
        # BLOQUE IV
        gen_b4_sep(),       # 43
        gen_edd_evol(),     # 44
        gen_edd_puro(),     # 45
        # Conclusiones
        gen_conclusiones(), # 46
    ]

    print(f"\nTotal slides B234: {len(paths_b234)}")
    build_pptx_b234(paths_b234)

    print("\nConcatenando con BLOQUE I...")
    concatenar_final()

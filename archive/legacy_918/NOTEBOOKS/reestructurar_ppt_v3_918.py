import sys; sys.stdout.reconfigure(encoding="utf-8")
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

BASE   = os.path.dirname(__file__)
SRC    = os.path.join(BASE, "..", "..", "PRESENTACION_final borrador1.pptx")
OUT    = os.path.join(BASE, "..", "..", "PRESENTACION_final_v3.pptx")
IMGDIR = BASE

prs = Presentation(SRC)
layouts = {l.name: l for l in prs.slide_layouts}
LY_SEPARADOR = layouts["Encabezado de sección"]
LY_BLANK     = layouts["Blank"]

# ── Geometría (idéntica a la v2 ya validada) ──────────────────────────────────
SLIDE_W, SLIDE_H = Emu(12192000), Emu(6858000)
MARGIN_X = Emu(700000)

TITLE_TOP, TITLE_H = Emu(330000), Emu(560000)
GAP_1 = Emu(160000)
PIC_TOP_MAX = TITLE_TOP + TITLE_H + GAP_1
PIC_H_MAX = Emu(4050000)
GAP_2 = Emu(170000)
BULLETS_TOP = PIC_TOP_MAX + PIC_H_MAX + GAP_2
BULLETS_H = SLIDE_H - BULLETS_TOP - Emu(220000)

TITLE_POS  = dict(left=MARGIN_X, top=TITLE_TOP, width=SLIDE_W - 2*MARGIN_X, height=TITLE_H)
PIC_POS    = dict(left=MARGIN_X, top=PIC_TOP_MAX, width=SLIDE_W - 2*MARGIN_X, height=PIC_H_MAX)
BULLET_POS = dict(left=MARGIN_X, top=BULLETS_TOP, width=SLIDE_W - 2*MARGIN_X, height=BULLETS_H)

MARKER_W, MARKER_H = Emu(1350000), Emu(360000)
MARKER_POS = dict(left=SLIDE_W - MARKER_W - Emu(140000), top=Emu(120000),
                   width=MARKER_W, height=MARKER_H)

TX1_COLOR = RGBColor(0x26, 0x26, 0x26)

def title_font_size(text):
    n = len(text)
    if n <= 52: return Pt(20)
    if n <= 65: return Pt(17)
    return Pt(15)

def set_title(slide, text):
    box = slide.shapes.add_textbox(**TITLE_POS)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = title_font_size(text)
    run.font.bold = True
    run.font.color.rgb = TX1_COLOR
    return box

def add_picture_scaled(slide, img_path):
    from PIL import Image as PILImage
    with PILImage.open(img_path) as im:
        ow, oh = im.size
    aspect = oh / ow
    max_w, max_h = PIC_POS["width"], PIC_POS["height"]
    w = max_w
    h = int(w * aspect)
    if h > max_h:
        h = max_h
        w = int(h / aspect)
    left = PIC_POS["left"] + (max_w - w) // 2
    top  = PIC_POS["top"] + (max_h - h) // 2
    slide.shapes.add_picture(img_path, left, top, w, h)

def add_bullets(slide, bullets):
    box = slide.shapes.add_textbox(**BULLET_POS)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"• {b}"
        run.font.size = Pt(12)
        run.font.color.rgb = TX1_COLOR
        p.space_after = Pt(4)
        p.line_spacing = 1.0

def add_new_marker(slide):
    box = slide.shapes.add_textbox(**MARKER_POS)
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "★ NUEVO"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xC6, 0x28, 0x28)

def new_content_slide(titulo, img_filename, bullets, marcar=True):
    slide = prs.slides.add_slide(LY_BLANK)
    set_title(slide, titulo)
    if img_filename:
        img_path = os.path.join(IMGDIR, img_filename)
        if os.path.exists(img_path):
            add_picture_scaled(slide, img_path)
        add_bullets(slide, bullets)
    else:
        text_top = TITLE_TOP + TITLE_H + Emu(700000)
        text_h   = SLIDE_H - text_top - Emu(700000)
        box = slide.shapes.add_textbox(MARGIN_X, text_top, SLIDE_W - 2*MARGIN_X, text_h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = f"• {b}"
            run.font.size = Pt(18)
            run.font.color.rgb = TX1_COLOR
            p.space_after = Pt(20)
    if marcar:
        add_new_marker(slide)
    return slide

def new_separator(titulo, subtitulo):
    slide = prs.slides.add_slide(LY_SEPARADOR)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = titulo
        elif ph.placeholder_format.idx == 1:
            ph.text = subtitulo
    return slide

def mini_separator(titulo):
    """Sub-separador liviano dentro de un bloque (solo texto, sin layout completo)."""
    slide = prs.slides.add_slide(LY_BLANK)
    text_top = Emu(2800000)
    box = slide.shapes.add_textbox(MARGIN_X, text_top, SLIDE_W - 2*MARGIN_X, Emu(1200000))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = titulo
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
    return slide

# ══════════════════════════════════════════════════════════════════════════════
# 1. Referencias a slides ORIGINALES
# ══════════════════════════════════════════════════════════════════════════════
orig_slides = list(prs.slides)
S = {i+1: orig_slides[i] for i in range(len(orig_slides))}

# ── Índice (slide 6) — limpiar TODO (incluye SmartArt oculto) y reescribir ───
slide6 = S[6]
for shape in list(slide6.shapes):
    is_title = shape.has_text_frame and "Índice" in shape.text_frame.text
    if not is_title:
        sp = shape._element
        sp.getparent().remove(sp)
box = slide6.shapes.add_textbox(MARGIN_X, Emu(1700000), SLIDE_W - 2*MARGIN_X, Emu(4200000))
tf = box.text_frame
tf.word_wrap = True
items = [
    "Bloque I — Clasificación del Cuerpo Académico (Universo 917)",
    "Bloque II — Evaluación Docente Antes y Después (917 → Formados → Aptos P3)",
    "    2.1 Participación en Formación   ·   2.2 Población Pura vs Mixta",
    "    2.3 SAT y Recomendación: Formados vs Control   ·   2.4 Comparación por Tipo/Jerarquía",
    "Bloque III — Resultados de Calificaciones (Notas según SAT)",
    "Bloque IV — Evaluación de Desempeño Docente (EDD)",
    "Bloque V — Jornada vs Honorario (Universo NOMINA 957)",
    "Conclusiones y Recomendaciones",
]
for i, it in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = it
    run.font.size = Pt(15) if it.startswith("    ") else Pt(17)
    run.font.bold = not it.startswith("    ")
    run.font.color.rgb = TX1_COLOR if not it.startswith("    ") else RGBColor(0x55,0x55,0x55)
    p.space_after = Pt(10)

# ── Recortar texto largo de slide 3 (Metodología) ─────────────────────────────
slide3 = S[3]
for shape in list(slide3.shapes):
    if shape.has_text_frame and len(shape.text_frame.text) > 200:
        tf = shape.text_frame
        full_text = tf.text
        if "Universo de Análisis" in full_text or "917 docentes jerarquizados" in full_text:
            tf.text = ("917 docentes jerarquizados · 493 con perfil completo + 424 solo nómina\n"
                       "Z-score = posición relativa del docente en su facultad y período")
        elif "Marco Metodológico" in full_text:
            tf.text = "z = (SAT docente − promedio facultad·período) / desviación estándar"
        elif "dimensiones de satisfacción" in full_text:
            tf.text = "SAT Nota (escala 1–7) y SAT % Recomendación · correlación r=0.893"
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.size = Pt(13)

# ── Completar subtítulo del separador "HONORARIO VS JORNADA" (S[22] = Bloque V)
slide22 = S[22]
for ph in slide22.placeholders:
    if ph.placeholder_format.idx == 1:
        ph.text = "Bloque V — Universo NOMINA 957\n(587 Jornada + 370 Honorario)"

print("Slides originales preparadas.")

# ══════════════════════════════════════════════════════════════════════════════
# 2. Slides nuevas — SOLO las que quedan (se excluyen AFO/MET/dimensiones/tipologías)
# ══════════════════════════════════════════════════════════════════════════════

# ── BLOQUE I — Clasificación del Cuerpo Académico ────────────────────────────
N1 = {}
N1["jerarquia"] = new_content_slide(
    "Distribución de Jerarquía Académica", "G_jerarquia_918.png",
    ["Instructor Docente es la más numerosa (37.3%), seguida de Asistente Docente (28.7%)",
     "Universo 917 → 915 con jerarquía informada (2 sin dato vía dotación)"])

N1["nivel_formacion"] = new_content_slide(
    "Nivel de Formación del Cuerpo Docente", "G_nivel_formacion_918.png",
    ["87.4% posee Magíster o Doctorado — fuerte orientación al posgrado",
     "Universo 917 → 493 perfil completo → 478 con nivel de formación informado"])

N1["jer_formacion"] = new_content_slide(
    "Jerarquía × Nivel de Formación", "G_jerarquia_formacion_918.png",
    ["A mayor jerarquía, mayor proporción de Doctores",
     "Universo 917 → 493 perfil completo → 477 con ambos datos informados"])

N1["gradorec"] = new_content_slide(
    "Grado Académico Reconocido (GRADOREC)", "G_gradorec_918.png",
    ["46.8% Magíster Profesional, 28.5% Doctor — perfil aplicado más que investigativo",
     "Universo 917 → 493 perfil completo → 477 con grado clasificado"])

N1["institucion"] = new_content_slide(
    "Institución de Obtención del Grado", "G_institucion_grado_918.png",
    ["La propia Universidad Central forma al 15.3% de su planta jerarquizada",
     "Universo 917 → 493 perfil completo → 478 con institución informada"])

N1["pais"] = new_content_slide(
    "País de Obtención del Grado", "G_pais_grado_918.png",
    ["70.7% de los grados se obtuvo en Chile; España lidera el destino internacional (14.0%)",
     "Universo 917 → 493 perfil completo → 478 con país informado"])

N1["carga"] = new_content_slide(
    "Distribución de la Carga Académica", "G_carga_academica_918.png",
    ["Docencia concentra el 49.5% de la carga declarada en dotación",
     "Universo 917 → 493 con cargo asignado en dotación"])

N1["anios_jerarquia"] = new_content_slide(
    "Años desde Ingreso hasta Jerarquización", "G_anios_hasta_jerarquia_918.png",
    ["Instructor Docente se jerarquiza casi de inmediato (mediana 0.7 años)",
     "Universo 917 → 433 con fechas de ingreso y jerarquización disponibles"])

N1["edad_jerarquizarse"] = new_content_slide(
    "Edad al Momento de Jerarquizarse", "G_edad_al_jerarquizarse_918.png",
    ["La edad mediana de jerarquización sube de 37 (Instructor) a 62 años (Titular)",
     "Universo 917 → 433 con fechas de ingreso y jerarquización disponibles"])

print(f"Bloque I: {len(N1)} slides nuevas.")

# ── BLOQUE II.1 — Participación en Formación ─────────────────────────────────
N21 = {}
N21["facultad"] = new_content_slide(
    "Participación en Formación por Facultad", "p2_g22_facultad.png",
    ["Brechas de hasta 30pp entre la facultad más y menos activa",
     "Universo 917 → 493 perfil completo (única base con dato de facultad)"])

N21["antiguedad"] = new_content_slide(
    "Participación en Formación por Antigüedad", "p2_g23_antiguedad.png",
    ["El tramo 0-4 años concentra la mayor tasa de participación (50%)",
     "Universo 917 → 493 perfil completo → 491 con antigüedad informada"])

N21["modalidad"] = new_content_slide(
    "Distribución por Modalidad de Formación", "p2_g31_modalidad.png",
    ["Taller concentra el 61% de las 615 instancias totales",
     "Universo 917 → 357 docentes formados → 615 actividades de formación"])

N21["brechas"] = new_content_slide(
    "Brechas: Perfil de No Participantes", "p2_g41_brechas.png",
    ["560 docentes sin ninguna instancia de formación (61.1% del universo)",
     "Universo 917 jerarquizados completo"])

N21["intensidad_global"] = new_content_slide(
    "Intensidad de Participación Global", "p2_g42_intensidad.png",
    ["61% de los formados participó en una sola instancia",
     "Universo 917 → 357 docentes con al menos 1 instancia de formación"])

N21["venn917"] = new_content_slide(
    "Combinaciones de Formación — Universo 917", "G_venn_formacion_917_918.png",
    ["357 docentes con formación: 76% hizo una sola modalidad",
     "Universo 917 → 357 docentes formados"])

N21["int_talleres"] = new_content_slide(
    "Intensidad de Participación en Talleres", "G_intensidad_talleres_918.png",
    ["59% de los docentes tomó un solo taller — experiencia puntual",
     "Universo 917 → 215 docentes con al menos 1 taller"])

N21["int_proyectos"] = new_content_slide(
    "Intensidad de Participación en Proyectos", "G_intensidad_proyectos_918.png",
    ["88% participó en un solo proyecto; máximo 4 (Juan Carlos Araya Vargas)",
     "Universo 917 → 32 docentes con al menos 1 proyecto"])

print(f"Bloque II.1 (Participación): {len(N21)} slides nuevas.")

# ── BLOQUE II.2 — Población Formada: Pura vs Mixta ───────────────────────────
N22 = {}
N22["venn_puro"] = new_content_slide(
    "Combinaciones de Modalidades — Población P3", "G_venn_poblaciones_918.png",
    ["93.4% de los 197 aptos P3 hizo una sola modalidad (población pura)",
     "Universo 917 → 197 aptos P3"])

N22["contraste1"] = new_content_slide(
    "Trayectoria SAT — Normal vs Población Pura", "G1_contraste_normal_puro_918.png",
    ["El Diplomado puro muestra mayor caída (Δ=-0.070) que con mixtos (Δ=-0.028)",
     "Universo 917 → 197 aptos P3 → 184 puros / 13 mixtos"])

N22["puro_nota_recom"] = new_content_slide(
    "SAT Nota y Recomendación — Población Pura", "G1_puro_nota_recom_918.png",
    ["Mismos patrones en ambas métricas: Proyecto cae más, Taller es estable",
     "Universo 917 → 197 aptos P3 → 184 puros"])

N22["delta_puro_mixto"] = new_content_slide(
    "Δ Z-score: Puro vs Mixto", "G_delta_z_puro_vs_mixto_918.png",
    ["Los 13 mixtos muestran mejor resultado (62% mejoró) que los puros (49%)",
     "Universo 917 → 197 aptos P3 → 184 puros / 13 mixtos"])

print(f"Bloque II.2 (Puro vs Mixto): {len(N22)} slides nuevas.")

# ── BLOQUE II.3 — SAT y Recomendación: Formados vs Control ───────────────────
N23 = {}
N23["t1"] = new_content_slide(
    "Validación Estadística — SAT Nota", "G11_pruebas_t_918.png",
    ["Diferencia significativa en 3 de 6 períodos (p<0.05)",
     "Universo 917 → 816 con SAT disponible · n fluctúa por período"])

N23["t2"] = new_content_slide(
    "Validación Estadística — % Recomendación", "G11.2_pruebas_t_bin_918.png",
    ["Significativo en 4 de 6 períodos",
     "Universo 917 → 816 con SAT disponible · n fluctúa por período"])

print(f"Bloque II.3 (SAT validación): {len(N23)} slides nuevas.")

# ── BLOQUE II.4 — Comparación por Tipo, Jerarquía y Antigüedad ───────────────
N24 = {}
N24["antig_par"] = new_content_slide(
    "Antigüedad × Tipo — SAT Nota y Recomendación", "G3_par_antiguedad_918.png",
    ["129 de 197 aptos P3 tienen dato de antigüedad disponible",
     "Universo 917 → 197 aptos P3 → 129 con antigüedad informada"])

N24["ev_taller"] = new_content_slide(
    "Evolución por Jerarquía — Talleres", "G7_taller_jerarquia_918.png",
    ["Asistentes Docentes parten bajo la media y alcanzan su peak durante talleres",
     "Universo 917 → 197 aptos P3 → subgrupo Taller"])

N24["ev_diplomado"] = new_content_slide(
    "Evolución por Jerarquía — Diplomados", "G7_diplomado_jerarquia_918.png",
    ["Instructores Docentes: caso de mayor éxito, parten bajo y suben",
     "Universo 917 → 197 aptos P3 → subgrupo Diplomado"])

N24["t4"] = new_content_slide(
    "T-test Jerarquía × Diplomado", "G13_ttest_jerarquia_diplomado_918.png",
    ["ANOVA no significativo (p=0.718): el Diplomado beneficia por igual a todos los rangos",
     "Universo 917 → 197 aptos P3 → subgrupo Diplomado"])

N24["t5"] = new_content_slide(
    "T-test Jerarquía × Taller", "G13.2_ttest_jerarquia_taller_918.png",
    ["Ningún par de jerarquías muestra diferencia significativa tras Bonferroni",
     "Universo 917 → 197 aptos P3 → subgrupo Taller"])

print(f"Bloque II.4 (Comparación): {len(N24)} slides nuevas.")

# ── BLOQUE III — Resultados de Calificaciones (Notas según SAT) ─────────────
N3 = {}
N3["aprob_antig"] = new_content_slide(
    "Tasa de Aprobación × Antigüedad Docente", "G_aprobacion_antiguedad_918.png",
    ["La tasa de aprobación es relativamente estable entre tramos de antigüedad",
     "Universo 917 → 834 docentes con notas y SAT disponibles"])

N3["aprob_jer"] = new_content_slide(
    "Tasa de Aprobación × Jerarquía Académica", "G_aprobacion_jerarquia_918.png",
    ["Instructor Docente lidera con la mediana más alta (94.6%)",
     "Universo 917 → 834 docentes con notas y SAT disponibles"])

N3["acumulativo"] = new_content_slide(
    "¿Es Acumulativo el Efecto de la Formación?", "G_acumulacion_formacion_918.png",
    ["A partir de 3 instancias la mediana de aprobación supera al grupo sin formación",
     "Universo 917 → 816 con SAT y notas → por intensidad de instancias"])

print(f"Bloque III: {len(N3)} slides nuevas.")

# ── BLOQUE IV — Evaluación de Desempeño Docente (EDD) ────────────────────────
N4 = {}
N4["t3"] = new_content_slide(
    "Validación Estadística — EDD", "G12_pruebas_t_edd_918.png",
    ["Significativo en 2024 (p=0.002) y 2025 (p<0.001, d=+0.955)",
     "Universo 917 → 486 con EDD disponible (134 formados + 352 sin formación)"])

N4["edd_contraste"] = new_content_slide(
    "EDD — Población Normal vs Pura", "G9_contraste_edd_918.png",
    ["Diplomado puro sube de 0.823 a 0.854 al excluir los mixtos",
     "Universo 917 → 197 aptos P3 → con EDD disponible"])

print(f"Bloque IV: {len(N4)} slides nuevas.")

# ── BLOQUE V — Jornada vs Honorario (adición) ────────────────────────────────
N5 = {}
N5["venn_contrato"] = new_content_slide(
    "Combinaciones de Formación — Jornada vs Honorario", "G_venn_formacion_contrato_918.png",
    ["Ninguna intersección de Honorario incluye Proyecto",
     "Universo NOMINA 957 → 587 Jornada + 370 Honorario"])

print(f"Bloque V: {len(N5)} slides nuevas.")

# ── Separadores de bloque y sub-bloque ────────────────────────────────────────
SEP_I   = new_separator("BLOQUE I", "Clasificación del Cuerpo Académico\nUniverso: 917 Docentes Jerarquizados")
SEP_II  = new_separator("BLOQUE II", "Evaluación Docente Antes y Después\nUniverso: 917 → Formados → Aptos P3")
SUB_21  = mini_separator("2.1   Participación en Formación")
SUB_22  = mini_separator("2.2   Población Formada: Pura vs Mixta")
SUB_23  = mini_separator("2.3   SAT y Recomendación — Formados vs Control")
SUB_24  = mini_separator("2.4   Comparación por Tipo, Jerarquía y Antigüedad")
SEP_III = new_separator("BLOQUE III", "Resultados de Calificaciones de Asignaturas\n(Notas y Tasas de Aprobación según SAT)")
SEP_IV  = new_separator("BLOQUE IV", "Evaluación de Desempeño Docente (EDD)\nUniverso: 917 → con EDD Disponible")

CIERRE1 = new_content_slide(
    "Conclusiones Consolidadas", None,
    ["La formación genera impacto medible en SAT, EDD y aprobación — efecto de selección + protección",
     "No hay diferencia Jornada/Honorario en calidad docente: el contrato no determina el desempeño",
     "La cobertura (39%) es la principal oportunidad: 6 de cada 10 nunca participaron",
     "El impacto se consolida a partir de 3+ instancias — la participación única no basta"],
    marcar=False)
CIERRE2 = new_content_slide(
    "Recomendaciones", None,
    ["Institucionalizar Diplomados como requisito de avance en la carrera académica",
     "Diseñar rutas formativas progresivas en lugar de instancias aisladas",
     "Ampliar cobertura a Honorarios: mismo retorno con menor acceso actual",
     "Crear mecanismo de evaluación de desempeño (EDD) para Honorarios"],
    marcar=False)

print("Separadores y cierre creados.")

# ══════════════════════════════════════════════════════════════════════════════
# 3. Orden final completo
# ══════════════════════════════════════════════════════════════════════════════
final_order = (
    [S[1], S[2], S[6]] +                                # Apertura + Índice
    [SEP_I] +
    [S[13], S[14], S[15], S[16]] +                       # Caracterización original (edad/sexo/facultad/antig/edad×jer)
    [N1["jerarquia"], N1["nivel_formacion"], N1["jer_formacion"], N1["gradorec"],
     N1["institucion"], N1["pais"], N1["carga"],
     N1["anios_jerarquia"], N1["edad_jerarquizarse"]] +
    [SEP_II] +
    [S[3], S[4], S[5]] +                                 # Universo / metodología / funnel P3
    [SUB_21] +
    [S[18], S[19], S[21], S[17]] +                       # P2 original (incl. tipo×jerarquía reubicado)
    [N21["facultad"], N21["antiguedad"], N21["modalidad"], N21["brechas"],
     N21["intensidad_global"], N21["venn917"], N21["int_talleres"], N21["int_proyectos"]] +
    [SUB_22] +
    [N22["venn_puro"], N22["contraste1"], N22["puro_nota_recom"], N22["delta_puro_mixto"]] +
    [SUB_23] +
    [S[7]] +                                             # G6/G6.2 evolución (2 img)
    [N23["t1"], N23["t2"]] +
    [SUB_24] +
    [S[11]] +                                            # G1 trayectoria por tipo original
    [N24["antig_par"], N24["ev_taller"], N24["ev_diplomado"], N24["t4"], N24["t5"]] +
    [SEP_III] +
    [S[8], S[9]] +                                       # scatter SAT-notas + aprobación GN
    [N3["aprob_antig"], N3["aprob_jer"], N3["acumulativo"]] +
    [SEP_IV] +
    [S[10]] +                                            # EDD evolución G9/G10
    [N4["t3"], N4["edd_contraste"]] +
    [S[22]] +                                            # Bloque V (separador ya existente)
    [S[23], S[24], S[25], S[26]] +
    [N5["venn_contrato"]] +
    [CIERRE1, CIERRE2]
)

# S[12] (heatmap antigüedad original, superado por N24["antig_par"] pareado) y
# S[20] (duplicado) quedan excluidos del orden final → eliminados del paquete.

sldIdLst = prs.slides._sldIdLst
current_slide_objs = list(prs.slides)
current_sldid_els  = list(sldIdLst)
slide_to_sldid = {id(s._element): sldid for s, sldid in zip(current_slide_objs, current_sldid_els)}

new_sldid_elements = [slide_to_sldid[id(s._element)] for s in final_order]
assert len(new_sldid_elements) == len(set(id(e) for e in new_sldid_elements)), "Slides repetidas en orden final"

ids_finales = {id(s._element) for s in final_order}
excluidas = [s for s in [S[12], S[20]] if id(s._element) not in ids_finales]
assert len(excluidas) == 2, "S[12] y S[20] deben quedar excluidas"

for el in list(sldIdLst):
    sldIdLst.remove(el)
for el in new_sldid_elements:
    sldIdLst.append(el)

# Eliminar relaciones de las slides excluidas (S12 redundante, S20 duplicada)
for s_excl in [S[12], S[20]]:
    rId = None
    for rel_id, rel in prs.part.rels.items():
        if rel.target_part == s_excl.part:
            rId = rel_id
            break
    if rId:
        prs.part.drop_rel(rId)

print(f"\nTotal slides en el deck final: {len(prs.slides)}")
prs.save(OUT)
print(f"Guardado: {OUT}")

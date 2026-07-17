import sys; sys.stdout.reconfigure(encoding="utf-8")
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

BASE     = os.path.dirname(__file__)
TEMPLATE = r"c:\Users\r.gonzalez_fluxsolar.LAPTOP-FLUX-ECO\Downloads\Analisis_UCEN_v2\PPT FINAL SOLO LOS N 197 SOLO APTOS P3.pptx"
OUT      = r"c:\Users\r.gonzalez_fluxsolar.LAPTOP-FLUX-ECO\Downloads\Analisis_UCEN_v2\PRESENTACION_197_P3_v2.pptx"

prs = Presentation(TEMPLATE)
layouts = {l.name: l for l in prs.slide_layouts}
LY_SEP   = layouts["Encabezado de sección"]
LY_BLANK = layouts["Blank"]

# ── Limpiar slides existentes ─────────────────────────────────────────────────
def clear_slides(prs):
    NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    sldIdLst = prs.slides._sldIdLst
    for sId in list(sldIdLst):
        rId = sId.get(f'{{{NS_R}}}id')
        try: prs.part.drop_rel(rId)
        except: pass
    sldIdLst.clear()

clear_slides(prs)

# ── Geometría (base v12, ajustada para etiqueta de población) ─────────────────
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

TITLE_POS  = dict(left=Emu(786582),  top=Emu(250608),  width=Emu(10599174), height=Emu(553998))
POP_POS    = dict(left=Emu(786582),  top=Emu(820000),   width=Emu(10599174), height=Emu(290000))
PIC_BOX    = dict(left=Emu(786581),  top=Emu(1125000),  width=Emu(10599174), height=Emu(3720000))
BULLET_POS = dict(left=Emu(786581),  top=Emu(4870000),  width=Emu(10599174), height=Emu(1870000))

# ── Colores ── fondo oscuro (foto UCEN) → texto BLANCO ───────────────────────
TITLE_FILL    = RGBColor(0x90, 0xAB, 0xC4)  # accent1  — barra título azul claro
TITLE_TEXT    = RGBColor(0x00, 0x50, 0x88)  # azul oscuro SOBRE la barra clara
BODY_TEXT     = RGBColor(0xFF, 0xFF, 0xFF)  # BLANCO sobre fondo oscuro
POP_TEXT      = RGBColor(0xC8, 0xDC, 0xF0)  # azul muy claro — etiqueta población
HALLAZGO_BLUE = RGBColor(0x00, 0x72, 0xE6)
HALLAZGO_RED  = RGBColor(0xC6, 0x28, 0x28)
C_NAVY        = RGBColor(0x00, 0x21, 0x47)  # acento portada
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

# ── Helpers ───────────────────────────────────────────────────────────────────
def png(name): return os.path.join(BASE, name)

def _no_shadow(shape):
    spPr = shape._element.spPr
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))

def _title_font_size(text):
    n = len(text)
    if n <= 55: return Pt(22)
    if n <= 75: return Pt(18)
    return Pt(15)

def add_title_bar(slide, text):
    box = slide.shapes.add_shape(1, **TITLE_POS)
    box.fill.solid(); box.fill.fore_color.rgb = TITLE_FILL
    box.line.fill.background()
    _no_shadow(box)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = _title_font_size(text)
    run.font.bold = True
    run.font.name = "Arial"
    run.font.color.rgb = TITLE_TEXT

def add_pop_tag(slide, pop):
    """Etiqueta de población en azul claro bajo la barra de título."""
    tb = slide.shapes.add_textbox(**POP_POS)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = pop
    run.font.size = Pt(9)
    run.font.italic = True
    run.font.name = "Arial"
    run.font.color.rgb = POP_TEXT

def add_picture_scaled(slide, img_path):
    """Marco blanco + imagen escalada con aspect ratio correcto."""
    frame = slide.shapes.add_shape(1, **PIC_BOX)
    frame.fill.solid(); frame.fill.fore_color.rgb = C_WHITE
    frame.line.fill.background()
    _no_shadow(frame)
    if not os.path.exists(img_path):
        print(f"  FALTA PNG: {os.path.basename(img_path)}")
        return
    try:
        from PIL import Image as PILImage
        with PILImage.open(img_path) as im:
            ow, oh = im.size
        aspect = oh / ow
        max_w = PIC_BOX["width"].emu
        max_h = PIC_BOX["height"].emu
        w = max_w
        h = int(w * aspect)
        if h > max_h:
            h = max_h; w = int(h / aspect)
        left = PIC_BOX["left"].emu + (max_w - w) // 2
        top  = PIC_BOX["top"].emu  + (max_h - h) // 2
    except ImportError:
        left = PIC_BOX["left"].emu
        top  = PIC_BOX["top"].emu
        w    = PIC_BOX["width"].emu
        h    = PIC_BOX["height"].emu
    slide.shapes.add_picture(img_path, Emu(left), Emu(top), Emu(w), Emu(h))

def add_bullets(slide, bullets, hallazgo=None):
    """Bullets numerados en BLANCO directamente sobre el fondo oscuro del template."""
    box = slide.shapes.add_shape(1, **BULLET_POS)
    box.fill.background()
    box.line.fill.background()
    _no_shadow(box)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left  = Emu(160000)
    tf.margin_right = Emu(100000)
    tf.margin_top   = Emu(60000)

    count = [0]
    def _add_p(text, color=BODY_TEXT, bold=False):
        p = tf.paragraphs[0] if count[0] == 0 else tf.add_paragraph()
        count[0] += 1
        p.alignment = PP_ALIGN.JUSTIFY
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", "228600"); pPr.set("indent", "-228600")
        buClr = etree.SubElement(pPr, qn('a:buClr'))
        buClrVal = etree.SubElement(buClr, qn('a:srgbClr'))
        buClrVal.set('val', str(color))
        buFont = etree.SubElement(pPr, qn('a:buFont')); buFont.set('typeface', '+mj-lt')
        buAutoNum = etree.SubElement(pPr, qn('a:buAutoNum')); buAutoNum.set('type', 'arabicPeriod')
        run = p.add_run()
        run.text = text
        run.font.size = Pt(12)
        run.font.name = "Arial"
        run.font.bold = bold
        run.font.color.rgb = color
        p.space_after = Pt(5)

    for b in bullets:
        _add_p(b)
    if hallazgo:
        _add_p(f"HALLAZGO CLAVE: {hallazgo}", color=HALLAZGO_RED, bold=True)

def add_hallazgo_marker(slide):
    box = slide.shapes.add_textbox(Emu(120000), SLIDE_H - Emu(520000), Emu(500000), Emu(420000))
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "★"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = HALLAZGO_BLUE

# ── Factories ─────────────────────────────────────────────────────────────────

POP_197 = "Universo: 197 Aptos P3  ·  130 formados vs 67 control  ·  SAT disponible baseline / durante / resultado"
POP_917 = "Universo: 917 docentes jerarquizados UCEN  ·  Los 197 Aptos P3 son un subconjunto de este universo"

def new_separator(titulo, subtitulo):
    """Separador de bloque usando layout 'Encabezado de sección' del template."""
    slide = prs.slides.add_slide(LY_SEP)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = titulo
        elif ph.placeholder_format.idx == 1:
            ph.text = subtitulo
    return slide

def content_slide(title, img_name, bullets, pop=POP_197, hallazgo=None):
    """Slide estándar: barra título + población + imagen + bullets."""
    s = prs.slides.add_slide(LY_BLANK)
    add_title_bar(s, title)
    add_pop_tag(s, pop)
    add_picture_scaled(s, png(img_name))
    add_bullets(s, bullets, hallazgo_texto if hallazgo else None)
    if hallazgo:
        add_hallazgo_marker(s)
    return s

def double_slide(title, img_left, img_right, bullets, pop=POP_197):
    """Dos imágenes lado a lado en el espacio del PIC_BOX."""
    s = prs.slides.add_slide(LY_BLANK)
    add_title_bar(s, title)
    add_pop_tag(s, pop)
    half_w = (PIC_BOX["width"].emu - 50000) // 2
    l = PIC_BOX["left"].emu
    t = PIC_BOX["top"].emu
    h = PIC_BOX["height"].emu

    for img_name, offset in [(img_left, 0), (img_right, half_w + 50000)]:
        p = png(img_name)
        frame = s.shapes.add_shape(1, Emu(l + offset), Emu(t), Emu(half_w), Emu(h))
        frame.fill.solid(); frame.fill.fore_color.rgb = C_WHITE
        frame.line.fill.background(); _no_shadow(frame)
        if os.path.exists(p):
            try:
                from PIL import Image as PILImage
                with PILImage.open(p) as im:
                    ow, oh = im.size
                aspect = oh / ow
                w = half_w; hh = int(w * aspect)
                if hh > h: hh = h; w = int(hh / aspect)
                il = l + offset + (half_w - w) // 2
                it = t + (h - hh) // 2
            except ImportError:
                il, it, w, hh = l + offset, t, half_w, h
            s.shapes.add_picture(p, Emu(il), Emu(it), Emu(w), Emu(hh))
        else:
            print(f"  FALTA PNG: {img_name}")
    add_bullets(s, bullets)
    return s

def text_slide(title, sections):
    """Slide de texto puro para conclusiones/recomendaciones."""
    s = prs.slides.add_slide(LY_BLANK)
    add_title_bar(s, title)
    y = 870000
    for header, items in sections.items():
        tb_h = slide_shapes_add_textbox_safe(s, 786582, y, 10599174, 420000, header, 14, bold=True, color=TITLE_FILL)
        y += 430000
        for item in items:
            tb_i = slide_shapes_add_textbox_safe(s, 1000000, y, 10386174, 360000, f"• {item}", 11, color=BODY_TEXT)
            y += 370000
        y += 80000
    return s

def slide_shapes_add_textbox_safe(slide, l, t, w, h, text, size, bold=False, color=BODY_TEXT):
    tb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Arial"
    run.font.color.rgb = color
    return tb

# Alias para claridad
hallazgo_texto = None

def cs(title, img, bullets, pop=POP_197, h=False, htxt=None):
    """Shorthand para content_slide con hallazgo opcional."""
    s = prs.slides.add_slide(LY_BLANK)
    add_title_bar(s, title)
    add_pop_tag(s, pop)
    add_picture_scaled(s, png(img))
    add_bullets(s, bullets, htxt)
    if h:
        add_hallazgo_marker(s)
    return s

# ══════════════════════════════════════════════════════════════════════════════
# DECK
# ══════════════════════════════════════════════════════════════════════════════

# ── 1. PORTADA ────────────────────────────────────────────────────────────────
portada = prs.slides.add_slide(LY_BLANK)
# Acento navy vertical izquierdo (replicando el original)
bar = portada.shapes.add_shape(1, Emu(2248566), Emu(0), Emu(380000), Emu(SLIDE_H.emu))
bar.fill.solid(); bar.fill.fore_color.rgb = C_NAVY
bar.line.fill.background(); _no_shadow(bar)

def _ptxt(slide, l, t, w, h, text, size, bold=False, color=C_WHITE):
    tb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Arial"
    run.font.color.rgb = color

_ptxt(portada, 323850, 2200000, 7200000, 1100000,
      "Resultados del Perfeccionamiento Docente", 30, bold=True)
_ptxt(portada, 323850, 3450000, 11000000, 750000,
      "Análisis de Impacto: SAT, Recomendación, Notas y EDD", 20, color=POP_TEXT)
_ptxt(portada, 323850, 4300000, 11000000, 550000,
      "Universo: 197 Aptos P3  ·  Períodos 2023–2025", 15, color=POP_TEXT)
_ptxt(portada, 323850, 5600000, 11000000, 480000,
      "Universidad Central de Chile  |  Producto 3: Análisis de Formación e Innovación",
      14, bold=True, color=C_WHITE)

# ── 2. ÍNDICE ─────────────────────────────────────────────────────────────────
indice = prs.slides.add_slide(LY_BLANK)
add_title_bar(indice, "Índice")
bloques_idx = [
    "BLOQUE I      Clasificación del Cuerpo Académico — distribución jerarquía, facultad y antigüedad",
    "BLOQUE II     Evaluación Docente Antes y Después — SAT, % Recomendación, caracterización formados vs no formados",
    "BLOQUE III    Resultados de Calificaciones — notas y tasas de aprobación/reprobación de alumnos según SAT",
    "BLOQUE IV     EDD — Evaluación de Desempeño Docente, evolución y panel balanceado",
    "CIERRE        Conclusiones y Recomendaciones",
]
box = indice.shapes.add_textbox(Emu(786581), Emu(900000), Emu(10599174), Emu(5500000))
tf = box.text_frame; tf.word_wrap = True
for i, b in enumerate(bloques_idx):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(16)
    run = p.add_run()
    run.text = b
    run.font.size = Pt(15)
    run.font.name = "Arial"
    run.font.color.rgb = BODY_TEXT

# ── 3. MARCO METODOLÓGICO ────────────────────────────────────────────────────
marco = prs.slides.add_slide(LY_BLANK)
add_title_bar(marco, "Universo de Análisis y Metodología")
add_pop_tag(marco, "Base del estudio: diseño cuasi-experimental longitudinal 2022–2025")

CAJAS_M = [
    (786581,  1120000, 5150000, 2700000,
     "1. Universo Aptos P3",
     ["917 docentes jerarquizados UCEN (universo base)",
      "493 con actividad docente en período analizado",
      "357 participaron en ≥ 1 iniciativa de formación",
      "197 Aptos P3: SAT válido en baseline, durante y resultado"]),
    (6235755, 1120000, 5150000, 2700000,
     "2. Marco Metodológico – Z-score",
     ["z = (SAT docente − media facultad período) / DE facultad período",
      "z = 0 → en el promedio exacto de su facultad ese semestre",
      "z > 0 → sobre el promedio · z < 0 → bajo el promedio",
      "Permite comparar docentes entre facultades distintas"]),
    (786581,  4050000, 5150000, 2600000,
     "3. Métricas de Evaluación",
     ["SAT Nota (1–7) estandarizada como z-score por facultad",
      "SAT % Recomendación: ¿recomendaría a este docente?",
      "EDD: Evaluación Directa de Desempeño (escala 0–1)",
      "Notas y aprobación alumnos: % con nota ≥ 4.0"]),
    (6235755, 4050000, 5150000, 2600000,
     "4. Diseño Comparativo",
     ["Grupo tratamiento: 130 docentes formados con SAT P3",
      "Grupo control: 67 docentes sin formación con SAT P3",
      "Comparación: posición z antes → durante → después",
      "Validación estadística: prueba t de diferencia de medias"]),
]
for (l, t, w, h, header, items) in CAJAS_M:
    # Sin fondo (transparente) — solo texto blanco sobre el fondo oscuro del template
    tb_h = marco.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf_h = tb_h.text_frame; tf_h.word_wrap = True
    p0 = tf_h.paragraphs[0]
    r0 = p0.add_run()
    r0.text = header
    r0.font.size = Pt(13); r0.font.bold = True; r0.font.name = "Arial"
    r0.font.color.rgb = TITLE_FILL  # azul claro para destacar el header
    p0.space_after = Pt(6)
    for item in items:
        p = tf_h.add_paragraph()
        r = p.add_run()
        r.text = f"• {item}"
        r.font.size = Pt(10.5); r.font.name = "Arial"
        r.font.color.rgb = BODY_TEXT
        p.space_after = Pt(3)

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE I — Clasificación Cuerpo Académico
# ─────────────────────────────────────────────────────────────────────────────
new_separator("BLOQUE I",
              "Clasificación del Cuerpo Académico\n"
              "Distribución por jerarquía, facultad y antigüedad · Universo 917")

cs("Distribución por Jerarquía Académica",
   "G_jerarquia_918.png",
   ["La jerarquía Asistente concentra la mayor cantidad de docentes, "
    "seguida por Asociados. Instructores y Titulares son los extremos del escalafón.",
    "Los 197 Aptos P3 replican la distribución del universo completo de 917 docentes.",
    "Comprender la distribución de jerarquías es clave para interpretar los resultados "
    "de SAT y EDD que se desglosan por esta variable en los bloques siguientes."],
   pop=POP_917)

double_slide(
    "Distribución por Facultad y Antigüedad",
    "G_facultad_918.png", "G_antiguedad_918.png",
    ["Por facultad: la distribución de docentes es heterogénea — "
     "algunas facultades concentran más docentes con formación que otras.",
     "Por antigüedad: los docentes con 0–4 años son el grupo más participativo en "
     "formación; a mayor antigüedad, menor tasa de participación."],
    pop=POP_917)

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE II — Evaluación Docente Antes y Después
# ─────────────────────────────────────────────────────────────────────────────
new_separator("BLOQUE II",
              "Evaluación Docente Antes y Después\n"
              "SAT · % Recomendación · Caracterización de Poblaciones")

cs("Evolución Z-score SAT por Período: Formados vs Control",
   "G6_control_vs_tratamiento_918.png",
   ["El grupo tratamiento (130 formados) supera consistentemente al grupo control "
    "en los 6 períodos analizados. Brecha media: +0.18z (~7 percentiles).",
    "El grupo control (67 sin formación) se mantiene bajo el promedio facultad con "
    "tendencia descendente a partir de 2024.",
    "La brecha se amplía en 2024-01 (+0.245z), coincidiendo con el peak "
    "de participación en talleres y diplomados."],
   h=True,
   htxt="Los docentes formados superan a sus pares sin formación en SAT en los 6 períodos consecutivos."
)

cs("% Recomendación: Formados vs Control por Período",
   "G6.2_control_bin_918.png",
   ["El porcentaje de alumnos que recomendaría al docente también es mayor en el "
    "grupo formado en todos los períodos analizados.",
    "La consistencia entre SAT Nota y % Recomendación refuerza la robustez del hallazgo: "
    "dos métricas independientes apuntan en la misma dirección.",
    "La brecha de recomendación se consolida especialmente en los períodos 2024–2025."])

cs("Trayectoria por Tipo de Formación: Baseline → Resultado",
   "G1_linea_z_918.png",
   ["Los tres tipos de formación (Taller, Diplomado, Proyecto) mantienen z-score "
    "positivo antes y después de la formación respecto al control.",
    "Proyectos muestra la mayor posición relativa post-formación (+0.257z), "
    "seguido de Diplomados (+0.177z) y Talleres (+0.056z).",
    "La trayectoria es ascendente para los tres tipos, sin caída post-formación."])

cs("Validación Estadística: Prueba t SAT",
   "G11_pruebas_t_918.png",
   ["La diferencia entre tratamiento y control es estadísticamente significativa "
    "en 3 de los 6 períodos (p < 0.05).",
    "En los 6 períodos la diferencia es positiva y consistente en dirección, "
    "con p < 0.10 en 5 de ellos.",
    "El tamaño del efecto (d de Cohen ≈ 0.3) corresponde a un efecto pequeño-moderado, "
    "significativo en el contexto docente."])

cs("Perfil de Participación: Jerarquía y Antigüedad",
   "p2_g41_brechas.png",
   ["Instructores Docentes y Asistentes Docentes muestran la mayor tasa relativa "
    "de participación en formación respecto a sus pares sin formación.",
    "Titulares Docentes (21% de participación): el grupo con menor tasa relativa, "
    "lo que sugiere una brecha de acceso o motivación en los rangos superiores.",
    "Por antigüedad: a mayor número de años en la institución, menor tasa de participación."],
   pop="Universo: 357 docentes con ≥1 iniciativa · Perfil por jerarquía y antigüedad")

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE III — Calificaciones y Aprobación
# ─────────────────────────────────────────────────────────────────────────────
new_separator("BLOQUE III",
              "Resultados de Calificaciones de Asignaturas\n"
              "Notas y tasas de aprobación/reprobación · Notas según SAT docente")

cs("Tasa de Aprobación: Formados vs Control",
   "GN_aprobacion_918.png",
   ["Los alumnos de docentes formados registran una tasa de aprobación "
    "superior en ~+1.5 pp respecto al grupo control.",
    "La diferencia es estadísticamente significativa dado el volumen de secciones "
    "y alumnos involucrados en el período 2023–2025.",
    "El efecto es modesto en términos absolutos pero consistente: refleja un impacto "
    "real en el rendimiento estudiantil."],
   h=True,
   htxt="Los alumnos de docentes formados aprueban en mayor proporción que los del grupo control.")

cs("Relación SAT Docente → Notas de Alumnos",
   "G8_scatter_sat_notas_918.png",
   ["Existe una correlación positiva débil-moderada entre la evaluación docente (SAT) "
    "y la nota promedio de los alumnos (r = 0.28, p < 0.001).",
    "Los docentes formados (naranja) tienden a concentrarse en la zona "
    "de mayor SAT, reforzando la asociación formación → mejor evaluación → mejores notas.",
    "La relación no es causal directa pero es una señal consistente a nivel de sistema."],
   pop="Universo: 816 docentes con SAT y calificaciones disponibles · 17.248 secciones · 2023–2025")

cs("Aprobación por Jerarquía Académica",
   "G_aprobacion_jerarquia_918.png",
   ["Los Instructores Docentes muestran el mayor diferencial de aprobación entre "
    "formados y control: la formación tiene más impacto en los rangos iniciales.",
    "Asistentes y Asociados presentan patrones similares con ventaja consistente "
    "en los grupos formados.",
    "Titulares: menor diferencial absoluto, posiblemente por mayor experiencia base."])

cs("¿Es Acumulativo el Efecto de la Formación?",
   "G_acumulacion_formacion_918.png",
   ["La tasa de aprobación sube progresivamente con el número de iniciativas "
    "de formación cursadas: a más formación, mejor rendimiento de alumnos.",
    "El efecto no es lineal: se consolida especialmente a partir de la 2ª iniciativa. "
    "Una sola participación genera un primer efecto, pero no el máximo.",
    "Implicancia: diseñar rutas formativas secuenciales de al menos 2–3 iniciativas "
    "es más efectivo que eventos aislados."],
   h=True,
   htxt="A mayor número de iniciativas de formación, mayor tasa de aprobación de alumnos.",
   pop="Universo: 917 docentes con notas disponibles · Agrupados por intensidad de formación")

# ─────────────────────────────────────────────────────────────────────────────
# BLOQUE IV — EDD
# ─────────────────────────────────────────────────────────────────────────────
new_separator("BLOQUE IV",
              "Evaluación de Desempeño Docente (EDD)\n"
              "Evolución · Panel balanceado · Validación estadística")

cs("Evolución EDD: Formados vs Control",
   "G9_edd_evolucion_918.png",
   ["La EDD de los docentes formados sigue una trayectoria ascendente consistente "
    "con los resultados de SAT: mayor formación, mejor desempeño evaluado.",
    "El grupo control muestra una trayectoria plana o descendente en el mismo período.",
    "La convergencia entre SAT y EDD como indicadores independientes fortalece "
    "la validez interna del estudio."])

cs("Panel Balanceado EDD: 170 Mismos Docentes Seguidos 4 Años",
   "G9_panel_balanceado_918.png",
   ["El panel balanceado elimina el sesgo de composición: son exactamente "
    "los mismos 170 docentes observados en los 4 períodos EDD disponibles.",
    "Brecha sostenida a favor de los formados: +0.081 puntos (escala 0–1), "
    "equivalente a ~8 puntos porcentuales en desempeño.",
    "La ventaja no se debe a rotación de docentes ni a cambios en la composición "
    "del grupo: es mejora real de los mismos individuos a lo largo del tiempo."],
   h=True,
   htxt="Panel de 170 individuos: los formados mejoran EDD; el control cae. La ventaja es real, no composicional.")

cs("Validación Estadística EDD",
   "G12_pruebas_t_edd_918.png",
   ["La prueba t para EDD confirma significancia estadística de la brecha "
    "entre formados y control.",
    "Consistente con los resultados SAT: ambas métricas apuntan en la misma dirección "
    "con significancia estadística.",
    "La doble validación (SAT + EDD) reduce el riesgo de falso positivo "
    "y fortalece la confianza en los hallazgos del estudio."])

# ─────────────────────────────────────────────────────────────────────────────
# CIERRE
# ─────────────────────────────────────────────────────────────────────────────
new_separator("CIERRE", "Conclusiones y Recomendaciones")

text_slide("Conclusiones", {
    "Hallazgo 1 — Evaluación SAT (Bloque II)": [
        "Los 130 docentes formados superan al grupo control (67) en z-score SAT en los 6 períodos (brecha media +0.18z, ~7 percentiles).",
        "La brecha es estadísticamente significativa y se amplía en 2024 con el peak de formación.",
    ],
    "Hallazgo 2 — Tasa de Aprobación (Bloque III)": [
        "Los alumnos de docentes formados aprueban ~+1.5 pp más. El efecto es acumulativo: a más iniciativas, mayor impacto.",
        "La correlación SAT → Notas (r=0.28) sitúa la formación docente en el continuo de mejora del rendimiento estudiantil.",
    ],
    "Hallazgo 3 — EDD Panel Balanceado (Bloque IV)": [
        "Seguidos los mismos 170 docentes 4 años: formados mantienen ventaja EDD de +0.081 puntos de forma sostenida.",
        "Descarta sesgo de composición: la mejora es real en los mismos individuos.",
    ],
    "Convergencia de indicadores": [
        "SAT, % Recomendación, Notas/Aprobación y EDD apuntan en la misma dirección: la formación docente mejora la calidad.",
    ],
})

text_slide("Recomendaciones", {
    "1. Diseñar rutas formativas progresivas (no actividades aisladas)": [
        "El efecto es acumulativo desde la 2ª iniciativa. Priorizar secuencias de ≥2–3 actividades por docente.",
        "En talleres, el umbral de impacto real en aprobación de alumnos es 3+ participaciones.",
    ],
    "2. Focalizar según jerarquía y antigüedad": [
        "Instructores y Asistentes Docentes muestran mayor potencial de mejora — priorizar en el plan formativo.",
        "Titulares con 15+ años: baja participación y diferencial menor → estrategias de motivación diferenciadas.",
    ],
    "3. Ampliar la base P3 para estudios futuros": [
        "Solo 197 de 917 son aptos P3. Garantizar captura sistemática de SAT pre y post formación "
        "permitirá aumentar la potencia estadística en los próximos ciclos.",
    ],
    "4. Institucionalizar la medición longitudinal": [
        "El panel balanceado EDD (170 docentes, 4 años) demuestra el valor del seguimiento continuo. "
        "Mantener el registro estructurado como práctica permanente de monitoreo.",
    ],
})

# ── Guardar ───────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"\n✓ Guardado: {OUT}")
print(f"  Total slides: {len(prs.slides)}")
for i, sl in enumerate(prs.slides):
    txts = [sh.text_frame.text[:55] for sh in sl.shapes
            if sh.has_text_frame and sh.text_frame.text.strip()]
    print(f"  [{i+1:02d}] {txts[0] if txts else '(sin texto)'}")

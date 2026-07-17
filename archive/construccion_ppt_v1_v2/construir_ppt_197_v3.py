import sys; sys.stdout.reconfigure(encoding="utf-8")
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

BASE     = os.path.dirname(__file__)
TEMPLATE = r"c:\Users\r.gonzalez_fluxsolar.LAPTOP-FLUX-ECO\Downloads\Analisis_UCEN_v2\PPT FINAL SOLO LOS N 197 SOLO APTOS P3.pptx"
OUT      = r"c:\Users\r.gonzalez_fluxsolar.LAPTOP-FLUX-ECO\Downloads\Analisis_UCEN_v2\PRESENTACION_197_P3_v3b.pptx"

prs = Presentation(TEMPLATE)
layouts = {l.name: l for l in prs.slide_layouts}
LY_SEP   = layouts["Encabezado de sección"]
LY_BLANK = layouts["Blank"]

# ── Limpiar slides existentes ─────────────────────────────────────────────────
def clear_slides(prs):
    NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    sldIdLst = prs.slides._sldIdLst
    for sId in list(sldIdLst):
        rId = sId.get(f'{{{NS_R}}}id')
        try: prs.part.drop_rel(rId)
        except: pass
    sldIdLst.clear()

clear_slides(prs)

# ── Geometría ─────────────────────────────────────────────────────────────────
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

TITLE_POS  = dict(left=Emu(786582),  top=Emu(250608),  width=Emu(10599174), height=Emu(553998))
POP_POS    = dict(left=Emu(786582),  top=Emu(820000),   width=Emu(10599174), height=Emu(290000))
PIC_BOX    = dict(left=Emu(786581),  top=Emu(1125000),  width=Emu(10599174), height=Emu(3720000))
BULLET_POS = dict(left=Emu(786581),  top=Emu(4870000),  width=Emu(10599174), height=Emu(1870000))

# ── Colores ───────────────────────────────────────────────────────────────────
TITLE_FILL    = RGBColor(0x90, 0xAB, 0xC4)  # barra título azul claro (accent1)
TITLE_TEXT    = RGBColor(0x00, 0x50, 0x88)  # azul oscuro sobre la barra
BODY_TEXT     = RGBColor(0xFF, 0xFF, 0xFF)  # BLANCO sobre fondo oscuro template
POP_TEXT      = RGBColor(0xC8, 0xDC, 0xF0)  # azul muy claro — etiqueta población
CASCADE_TXT   = RGBColor(0x00, 0x21, 0x47)  # navy oscuro — texto en caja blanca
HALLAZGO_BLUE = RGBColor(0x00, 0x72, 0xE6)
HALLAZGO_RED  = RGBColor(0xC6, 0x28, 0x28)
C_NAVY        = RGBColor(0x00, 0x21, 0x47)
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

# ── Helpers ───────────────────────────────────────────────────────────────────
def png(name): return os.path.join(BASE, name)

def _no_shadow(shape):
    spPr = shape._element.spPr
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))

def _title_font_size(text):
    n = len(text)
    if n <= 55: return Pt(22)
    if n <= 75: return Pt(18)
    return Pt(15)

def add_title_bar(slide, text):
    box = slide.shapes.add_shape(1, **TITLE_POS)
    box.fill.solid(); box.fill.fore_color.rgb = TITLE_FILL
    box.line.fill.background(); _no_shadow(box)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = _title_font_size(text)
    run.font.bold = True
    run.font.name = "Arial"
    run.font.color.rgb = TITLE_TEXT

def add_pop_tag(slide, pop):
    tb = slide.shapes.add_textbox(**POP_POS)
    tf = tb.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = pop
    run.font.size = Pt(9)
    run.font.italic = True
    run.font.name = "Arial"
    run.font.color.rgb = POP_TEXT

def add_picture_scaled(slide, img_path):
    if not os.path.exists(img_path):
        print(f"  FALTA PNG: {os.path.basename(img_path)}")
        return
    try:
        from PIL import Image as PILImage
        with PILImage.open(img_path) as im:
            ow, oh = im.size
        aspect = oh / ow
        max_w = PIC_BOX["width"].emu
        max_h = PIC_BOX["height"].emu
        w = max_w; h = int(w * aspect)
        if h > max_h: h = max_h; w = int(h / aspect)
        left = PIC_BOX["left"].emu + (max_w - w) // 2
        top  = PIC_BOX["top"].emu  + (max_h - h) // 2
    except ImportError:
        left, top, w, h = (PIC_BOX["left"].emu, PIC_BOX["top"].emu,
                           PIC_BOX["width"].emu, PIC_BOX["height"].emu)
    slide.shapes.add_picture(img_path, Emu(left), Emu(top), Emu(w), Emu(h))

def add_double_pic(slide, img_left, img_right):
    """Dos imágenes lado a lado en PIC_BOX (sin caja blanca)."""
    half_w = (PIC_BOX["width"].emu - 60000) // 2
    l = PIC_BOX["left"].emu; t = PIC_BOX["top"].emu; h = PIC_BOX["height"].emu
    for img_name, offset in [(img_left, 0), (img_right, half_w + 60000)]:
        p = png(img_name)
        if os.path.exists(p):
            try:
                from PIL import Image as PILImage
                with PILImage.open(p) as im:
                    ow, oh = im.size
                aspect = oh / ow
                w2 = half_w; hh = int(w2 * aspect)
                if hh > h: hh = h; w2 = int(hh / aspect)
                il = l + offset + (half_w - w2) // 2
                it = t + (h - hh) // 2
            except ImportError:
                il, it, w2, hh = l + offset, t, half_w, h
            slide.shapes.add_picture(p, Emu(il), Emu(it), Emu(w2), Emu(hh))
        else:
            print(f"  FALTA PNG: {img_name}")

def add_bullets(slide, bullets, hallazgo=None):
    box = slide.shapes.add_shape(1, **BULLET_POS)
    box.fill.background(); box.line.fill.background(); _no_shadow(box)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(160000); tf.margin_right = Emu(100000); tf.margin_top = Emu(60000)
    count = [0]
    def _add_p(text, color=BODY_TEXT, bold=False):
        p = tf.paragraphs[0] if count[0] == 0 else tf.add_paragraph()
        count[0] += 1; p.alignment = PP_ALIGN.JUSTIFY
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", "228600"); pPr.set("indent", "-228600")
        buClr = etree.SubElement(pPr, qn('a:buClr'))
        etree.SubElement(buClr, qn('a:srgbClr')).set('val', str(color))
        etree.SubElement(pPr, qn('a:buFont')).set('typeface', '+mj-lt')
        etree.SubElement(pPr, qn('a:buAutoNum')).set('type', 'arabicPeriod')
        run = p.add_run()
        run.text = text; run.font.size = Pt(12); run.font.name = "Arial"
        run.font.bold = bold; run.font.color.rgb = color; p.space_after = Pt(5)
    for b in bullets:
        _add_p(b)
    if hallazgo:
        _add_p(f"HALLAZGO CLAVE: {hallazgo}", color=HALLAZGO_RED, bold=True)

def add_hallazgo_marker(slide):
    box = slide.shapes.add_textbox(Emu(120000), SLIDE_H - Emu(520000), Emu(500000), Emu(420000))
    tf = box.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = "★"
    run.font.size = Pt(28); run.font.bold = True; run.font.color.rgb = HALLAZGO_BLUE

# ── Etiquetas de población ────────────────────────────────────────────────────
POP_197 = "Universo: 197 Aptos P3  ·  130 formados + 67 control  ·  SAT disponible baseline / durante / resultado"
POP_917 = "Universo: 917 docentes jerarquizados UCEN  ·  Los 197 Aptos P3 son un subconjunto de este universo"
POP_816 = "Universo: 816 docentes con SAT y calificaciones disponibles  ·  17.248 secciones  ·  2023–2025"
POP_357 = "Universo: 357 docentes con ≥1 iniciativa de formación  ·  Perfil por jerarquía y antigüedad"

# ── Factories ─────────────────────────────────────────────────────────────────
def cs(title, img, bullets, pop=POP_197, h=False, htxt=None):
    """Slide estándar: título + etiqueta población + imagen + bullets."""
    s = prs.slides.add_slide(LY_BLANK)
    add_title_bar(s, title)
    add_pop_tag(s, pop)
    add_picture_scaled(s, png(img))
    add_bullets(s, bullets, htxt)
    if h: add_hallazgo_marker(s)
    return s

def ds(title, img_l, img_r, bullets, pop=POP_197):
    """Slide con dos imágenes lado a lado."""
    s = prs.slides.add_slide(LY_BLANK)
    add_title_bar(s, title)
    add_pop_tag(s, pop)
    add_double_pic(s, img_l, img_r)
    add_bullets(s, bullets)
    return s

def new_separator(titulo, subtitulo):
    """Separador usando layout Encabezado de sección (para slides sin cascada)."""
    slide = prs.slides.add_slide(LY_SEP)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:   ph.text = titulo
        elif ph.placeholder_format.idx == 1: ph.text = subtitulo
    return slide

def _dotlead(text, width=84):
    """Replace multi-space gap between label and trailing (n=...) with dot leaders."""
    import re as _re
    m = _re.match(r'^(.*?)\s{3,}(\(.+\))\s*$', text.strip())
    if not m:
        return text
    left  = m.group(1).rstrip()
    right = m.group(2)
    n_dots = max(3, width - len(left) - len(right) - 2)
    return f"{left} {'.' * n_dots} {right}"

def mini_separator(titulo, items, pop=None):
    """Separador de bloque: cascada de texto directo sobre fondo oscuro, sin caja blanca."""
    s = prs.slides.add_slide(LY_BLANK)
    add_title_bar(s, titulo)
    if pop:
        add_pop_tag(s, pop)

    PAD_L = 200000
    PAD_T = 200000
    tb = s.shapes.add_textbox(
        Emu(PIC_BOX["left"].emu + PAD_L),
        Emu(PIC_BOX["top"].emu  + PAD_T),
        Emu(PIC_BOX["width"].emu  - 2 * PAD_L),
        Emu(PIC_BOX["height"].emu - 2 * PAD_T),
    )
    tf = tb.text_frame; tf.word_wrap = False

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        is_last = (i == len(items) - 1)
        prefix  = "  └──  " if is_last else "  ├──  "
        run = p.add_run()
        run.text = prefix + _dotlead(item)
        run.font.size = Pt(13)
        run.font.name = "Courier New"
        run.font.color.rgb = BODY_TEXT
        run.font.bold = False
        p.space_after = Pt(10)

    return s


# ══════════════════════════════════════════════════════════════════════════════
# DECK
# ══════════════════════════════════════════════════════════════════════════════

# ── 1. PORTADA ────────────────────────────────────────────────────────────────
portada = prs.slides.add_slide(LY_BLANK)
bar = portada.shapes.add_shape(1, Emu(2248566), Emu(0), Emu(380000), SLIDE_H)
bar.fill.solid(); bar.fill.fore_color.rgb = C_NAVY
bar.line.fill.background(); _no_shadow(bar)

def _ptxt(slide, l, t, w, h, text, size, bold=False, color=C_WHITE):
    tb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; run = p.add_run()
    run.text = text; run.font.size = Pt(size); run.font.bold = bold
    run.font.name = "Arial"; run.font.color.rgb = color

_ptxt(portada, 323850, 2200000, 7200000, 1100000,
      "Resultados del Perfeccionamiento Docente", 30, bold=True)
_ptxt(portada, 323850, 3450000, 11000000, 750000,
      "Análisis de Impacto: SAT, Recomendación, Notas y EDD", 20, color=POP_TEXT)
_ptxt(portada, 323850, 4300000, 11000000, 550000,
      "Universo: 197 Aptos P3  ·  Períodos 2023–2025", 15, color=POP_TEXT)
_ptxt(portada, 323850, 5600000, 11000000, 480000,
      "Universidad Central de Chile  |  Producto 3: Análisis de Formación e Innovación",
      14, bold=True, color=C_WHITE)

# ── 2. ÍNDICE ─────────────────────────────────────────────────────────────────
indice = prs.slides.add_slide(LY_BLANK)
add_title_bar(indice, "Índice")
bloques_idx = [
    "BLOQUE I      Clasificación del Cuerpo Académico — edad, sexo, facultad, antigüedad, jerarquía y grado  (n=197 Aptos P3)",
    "BLOQUE II     Evaluación Docente Antes y Después — SAT, % Recomendación, formados vs control  (n=197 Aptos P3)",
    "BLOQUE III    Calificaciones de Alumnos — notas y aprobación según SAT docente  (n=816 con SAT+notas)",
    "BLOQUE IV     EDD — Evaluación de Desempeño Docente, evolución y panel balanceado  (n=197 Aptos P3)",
    "CIERRE        Conclusiones y Recomendaciones",
]
box = indice.shapes.add_textbox(Emu(786581), Emu(900000), Emu(10599174), Emu(5500000))
tf = box.text_frame; tf.word_wrap = True
for i, b in enumerate(bloques_idx):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(18)
    run = p.add_run(); run.text = b
    run.font.size = Pt(14); run.font.name = "Arial"; run.font.color.rgb = BODY_TEXT

# ── 3. MARCO METODOLÓGICO ────────────────────────────────────────────────────
marco = prs.slides.add_slide(LY_BLANK)
add_title_bar(marco, "Universo de Análisis y Metodología")
add_pop_tag(marco, "Base del estudio: diseño cuasi-experimental longitudinal 2022–2025")
CAJAS_M = [
    (786581,  1120000, 5150000, 2700000,
     "1. Universo Aptos P3",
     ["917 docentes jerarquizados UCEN (universo base)",
      "493 con actividad docente en período analizado",
      "357 participaron en ≥ 1 iniciativa de formación",
      "197 Aptos P3: SAT válido en baseline, durante y resultado"]),
    (6235755, 1120000, 5150000, 2700000,
     "2. Marco Metodológico – Z-score",
     ["z = (SAT docente − media facultad período) / DE facultad período",
      "z = 0 → en el promedio exacto de su facultad ese semestre",
      "z > 0 → sobre el promedio · z < 0 → bajo el promedio",
      "Permite comparar docentes entre facultades distintas"]),
    (786581,  4050000, 5150000, 2600000,
     "3. Métricas de Evaluación",
     ["SAT Nota (1–7) estandarizada como z-score por facultad",
      "SAT % Recomendación: ¿recomendaría a este docente?",
      "EDD: Evaluación Directa de Desempeño (escala 0–1)",
      "Notas y aprobación alumnos: % con nota ≥ 4.0"]),
    (6235755, 4050000, 5150000, 2600000,
     "4. Diseño Comparativo",
     ["Grupo tratamiento: 130 docentes formados con SAT P3",
      "Grupo control: 67 docentes sin formación con SAT P3",
      "Comparación: posición z antes → durante → después",
      "Validación estadística: prueba t de diferencia de medias"]),
]
for (l, t, w, h, header, items) in CAJAS_M:
    tb_h = marco.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf_h = tb_h.text_frame; tf_h.word_wrap = True
    p0 = tf_h.paragraphs[0]; r0 = p0.add_run()
    r0.text = header; r0.font.size = Pt(13); r0.font.bold = True
    r0.font.name = "Arial"; r0.font.color.rgb = TITLE_FILL; p0.space_after = Pt(6)
    for item in items:
        p = tf_h.add_paragraph(); r = p.add_run()
        r.text = f"• {item}"; r.font.size = Pt(10.5); r.font.name = "Arial"
        r.font.color.rgb = BODY_TEXT; p.space_after = Pt(3)


# ═════════════════════════════════════════════════════════════════════════════
# BLOQUE I — Clasificación del Cuerpo Académico  (n=197 Aptos P3)
# ═════════════════════════════════════════════════════════════════════════════
mini_separator(
    titulo="BLOQUE I — Clasificación del Cuerpo Académico",
    pop=POP_197,
    items=[
        "Diapo 5:   Tramo de Edad y Sexo                           (n=129 con edad / 197 con sexo)",
        "Diapo 6:   Distribución por Facultad/Unidad               (n=129 con dato)",
        "Diapo 7:   Antigüedad en la Institución                   (n=129 con dato)",
        "Diapo 8:   Distribución de Jerarquía Académica            (n=197)",
        "Diapo 9:   Grado Académico Reconocido                     (n=129 con dotación)",
        "Diapo 10:  Institución de Obtención del Grado             (n=129 con dotación)",
    ]
)

# Diapo 5: Edad + Sexo (doble)
ds("Distribución por Tramo de Edad y Sexo",
   "G_edad_197.png", "G_sexo_197.png",
   ["El grupo etario dominante en los 197 Aptos P3 es el tramo 40–44 años (31 doc.), "
    "con un núcleo 35–54 que concentra el 74% de los docentes con dato disponible.",
    "La distribución por sexo muestra mayor presencia masculina, aunque el perfil "
    "está dentro de los rangos esperados para docencia universitaria en Chile.",
    "Nota: 65% de los 197 tienen dato de edad (129 doc.); el 35% restante corresponde "
    "a honorarios sin registro en dotación."],
   pop=POP_197)

# Diapo 6: Facultad
cs("Distribución por Unidad/Facultad",
   "G_facultad_197.png",
   ["La Facultad de Medicina y Ciencias de la Salud concentra el mayor número "
    "de Aptos P3 (47 doc., 36%), seguida por Derecho y Humanidades (22, 17%) "
    "e Ingeniería y Arquitectura (21, 16%).",
    "Las cinco facultades principales agrupan el 90% de los docentes Aptos P3. "
    "La distribución refleja la mayor concentración de actividad SAT en salud y derecho.",
    "Vicerrectoría de Investigación e Innovación aporta 11 docentes (8%), "
    "con un perfil más orientado a investigación que a docencia de pregrado."],
   pop=POP_197)

# Diapo 7: Antigüedad
cs("Antigüedad en la Institución",
   "G_antiguedad_197.png",
   ["El tramo de mayor frecuencia entre los Aptos P3 es 5–9 años (39 doc.), "
    "seguido por 0–4 años (35 doc.). Los docentes más jóvenes en la institución "
    "son los que más frecuentemente reúnen las condiciones para ser Aptos P3.",
    "A mayor antigüedad, menor presencia en el grupo de Aptos P3, lo que es "
    "consistente con que los docentes más experimentados tienden menos a "
    "participar en iniciativas de formación.",
    "Mediana estimada de antigüedad: 9–10 años en la institución."],
   pop=POP_197)

# Diapo 8: Jerarquía
cs("Distribución de Jerarquía Académica",
   "G_jerarquia_197.png",
   ["Los Aptos P3 se concentran en las jerarquías de entrada al escalafón: "
    "Asistente Docente (65 doc., 33%) e Instructor Docente (62 doc., 31%), "
    "en línea con el perfil de mayor participación en formación.",
    "Los Asociados Docentes aportan 45 doc. (23%), confirmando que la mitad "
    "del escalafón también está representada. Titulares son el grupo "
    "con menor representación (7 doc., 4%).",
    "El Cuerpo Docente (DOCENTE) representa el 87% vs el Cuerpo Regular (13%), "
    "lo que es consistente con la mayor propensión a participar en formación "
    "del perfil docente orientado a la enseñanza."],
   pop=POP_197)

# Diapo 11: Grado Académico
cs("Grado Académico Reconocido",
   "G_gradorec_197.png",
   ["De los 129 Aptos P3 con dotación completa, el grado más frecuente es "
    "Magíster (Profesional + Académico), seguido por Doctor. "
    "El perfil de posgrado predomina entre los docentes formados.",
    "Los Doctores y Post-Doctores representan el núcleo de mayor calificación "
    "académica. El segmento Profesional (sin posgrado) es menor, "
    "concentrado posiblemente en rangos de entrada al escalafón.",
    "Nota: 68 docentes Aptos P3 no tienen dato de grado (honorarios sin dotación). "
    "Los porcentajes son sobre n=129 con dato disponible."],
   pop=POP_197)

# Diapo 12: Institución del grado
cs("Institución de Obtención del Grado",
   "G_institucion_197.png",
   ["La Universidad Central de Chile aparece como principal fuente de grados "
    "académicos, seguida por universidades del CRUCH. "
    "El patrón es consistente con el universo completo de 917.",
    "Una proporción de docentes obtuvo su grado en el extranjero, "
    "lo que enriquece la diversidad de formación del cuerpo docente.",
    "Nota: dato disponible para 129 de los 197 Aptos P3 "
    "(solo docentes con dotación registrada en el sistema)."],
   pop=POP_197)


# ── Guardar ───────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"\n✓ Guardado: {OUT}")
print(f"  Total slides: {len(prs.slides)}")
for i, sl in enumerate(prs.slides):
    txts = [sh.text_frame.text[:50] for sh in sl.shapes
            if sh.has_text_frame and sh.text_frame.text.strip()]
    print(f"  [{i+1:02d}] {txts[0] if txts else '(sin texto)'}")

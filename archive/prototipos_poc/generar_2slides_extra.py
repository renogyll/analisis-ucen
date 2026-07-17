# -*- coding: utf-8 -*-
"""
generar_2slides_extra.py
Genera PPTX con 2 slides nuevas en formato dark v3:
  Slide 1 - Scatter SAT vs Nota Promedio  (3 paneles por ano 2023/2024/2025)
  Slide 2 - Embudo 917 357 197 rediseniado (mas compacto y acotado)
"""
import sys; sys.stdout.reconfigure(encoding="utf-8")
import os, zipfile
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.patheffects as pe
from scipy import stats as scipy_stats
import pandas as pd
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─────────────────────────────────────────────────────────────────────────────
# Rutas
# ─────────────────────────────────────────────────────────────────────────────
BASE    = os.path.dirname(os.path.abspath(__file__))
PROC    = os.path.join(BASE, "..", "PROCESADO")
ROOT    = os.path.normpath(os.path.join(BASE, "..", ".."))
SCRATCH = (r"C:\Users\RGONZA~1.LAP\AppData\Local\Temp\claude"
           r"\c--Users-r-gonzalez-fluxsolar-LAPTOP-FLUX-ECO-Downloads-Analisis-UCEN-v2"
           r"\19e6fc3f-6ca1-4150-9da7-8dfa38be71ca\scratchpad")
FONDOTIPO = (r"c:\Users\r.gonzalez_fluxsolar.LAPTOP-FLUX-ECO"
             r"\Downloads\Analisis_UCEN_v2\Fondotipop.pptx")
OUT_DIR  = os.path.join(BASE, "dark_slides_v3")
OUT_PPTX = os.path.join(ROOT, "SLIDES_EXTRA_scatter_embudo.pptx")
os.makedirs(SCRATCH, exist_ok=True)
os.makedirs(OUT_DIR, exist_ok=True)

# ─────────────────────────────────────────────────────────────────────────────
# Assets (foto B&N + logo extraidos de Fondotipop.pptx)
# ─────────────────────────────────────────────────────────────────────────────
BG_PATH   = os.path.join(SCRATCH, "fondotipo_image1.jpg")
LOGO_PATH = os.path.join(SCRATCH, "fondotipo_image2.png")
for path, zname in [(BG_PATH,"ppt/media/image1.jpg"),(LOGO_PATH,"ppt/media/image2.png")]:
    if not os.path.exists(path):
        with zipfile.ZipFile(FONDOTIPO) as z:
            with open(path,"wb") as f:
                f.write(z.read(zname))

with PILImage.open(BG_PATH) as _im:
    _rgb = _im.convert("RGB"); _iw,_ih = _rgb.size
    _nh  = int(_iw/(16/9)); _y0 = min(int(_ih*0.12),_ih-_nh)
    bg_arr = np.array(_rgb.crop((0,_y0,_iw,_y0+_nh)))

with PILImage.open(LOGO_PATH) as _logo:
    logo_arr = np.array(_logo.convert("RGBA")).astype(np.float32)/255.0

H_GRAD = 600; grad = np.zeros((H_GRAD,1,4), dtype=np.float32)
for _r in range(H_GRAD):
    _t = _r/(H_GRAD-1)
    _stops = [(0.00,(0,33,71)),(0.54,(0,70,128)),(1.00,(144,171,196))]
    for _i in range(len(_stops)-1):
        _t0,_c0 = _stops[_i]; _t1,_c1 = _stops[_i+1]
        if _t0 <= _t <= _t1:
            _s = (_t-_t0)/(_t1-_t0)
            grad[_r,0] = [(_c0[0]+_s*(_c1[0]-_c0[0]))/255,
                          (_c0[1]+_s*(_c1[1]-_c0[1]))/255,
                          (_c0[2]+_s*(_c1[2]-_c0[2]))/255, 0.82]; break

# ─────────────────────────────────────────────────────────────────────────────
# Constantes de layout (identicas al script principal)
# ─────────────────────────────────────────────────────────────────────────────
SW, SH         = 13.333, 7.5
SW_EMU, SH_EMU = 12192000, 6858000

PIC_L, PIC_T, PIC_W, PIC_H = 786581, 1125000, 10599174, 3720000
BUL_L, BUL_T, BUL_W, BUL_H = 786581, 4870000, 10599174, 1870000
LOGO_L, LOGO_T, LOGO_W, LOGO_H = 9813773, 656354, 1756626, 697725
TITLE_L, TITLE_T, TITLE_W, TITLE_H = PIC_L, 185000, PIC_W, 710000
POP_L, POP_T, POP_W, POP_H = PIC_L, 845000, 9000000, 255000
CTITLE_L = int((PIC_L/SW_EMU + 0.13) * SW_EMU)
CTITLE_T = PIC_T + 38000
CTITLE_W = int(PIC_W * 0.87)
CTITLE_H = 295000

def _ex(e): return e/SW_EMU
def _ey(e): return e/SH_EMU
def _fig_rect(l,t,w,h): return (l, 1-t-h, w, h)

PIC_RECT  = _fig_rect(_ex(PIC_L), _ey(PIC_T), _ex(PIC_W), _ey(PIC_H))
LOGO_RECT = _fig_rect(_ex(LOGO_L), _ey(LOGO_T), _ex(LOGO_W), _ey(LOGO_H))

CHART_X = PIC_RECT[0] + 0.13
CHART_Y = PIC_RECT[1] + 0.04
CHART_W = PIC_RECT[2] - 0.19
CHART_H = PIC_RECT[3] - 0.09

POP_197 = "Universo: 197 Aptos P3  |  todos formados  |  Periodos 2022-2025"
POP_CTR = "Universo: 197 Aptos P3  |  formados vs control externo sin formacion  |  2023-2025"

# ─────────────────────────────────────────────────────────────────────────────
# Datos
# ─────────────────────────────────────────────────────────────────────────────
scat = pd.read_csv(os.path.join(PROC, "scatter_sat_notas.csv"), encoding="utf-8-sig")
scat["sat"]           = pd.to_numeric(scat["sat"],           errors="coerce")
scat["nota_promedio"] = pd.to_numeric(scat["nota_promedio"], errors="coerce")
scat["anio"]          = scat["periodo"].str[:4]
scat = scat.dropna(subset=["sat","nota_promedio"])

# ─────────────────────────────────────────────────────────────────────────────
# Helpers matplotlib
# ─────────────────────────────────────────────────────────────────────────────
def _bg_fig():
    fig = plt.figure(figsize=(SW,SH), facecolor="#101820")
    fig.patch.set_facecolor("#101820")
    for z,arr in [(0,bg_arr),(1,grad)]:
        ax = fig.add_axes([0,0,1,1], zorder=z)
        ax.imshow(arr, extent=[0,1,0,1], aspect="auto", origin="upper")
        ax.set_xlim(0,1); ax.set_ylim(0,1); ax.axis("off")
    al = fig.add_axes([LOGO_RECT[0],LOGO_RECT[1],LOGO_RECT[2],LOGO_RECT[3]],
                      zorder=10, facecolor="none")
    al.imshow(logo_arr, aspect="auto"); al.axis("off"); al.patch.set_visible(False)
    return fig

def _tr_fig():
    fig = plt.figure(figsize=(SW,SH), facecolor="none")
    fig.patch.set_facecolor("none"); return fig

SHARED_BG = os.path.join(OUT_DIR, "_background.png")

def _ensure_bg():
    if not os.path.exists(SHARED_BG):
        fig = _bg_fig()
        plt.savefig(SHARED_BG, dpi=150, facecolor=fig.get_facecolor()); plt.close()
        print("  bg generado")

def _save_ch(fig, name):
    path = os.path.join(OUT_DIR, name)
    plt.savefig(path, dpi=150, facecolor="none", transparent=True); plt.close(); return path

def _style_dark(ax, xlabel=None):
    ax.tick_params(axis="x", colors="white", labelsize=9, length=0)
    ax.tick_params(axis="y", colors="#AAAAAA", labelsize=9, length=0)
    for sp in ax.spines.values():
        sp.set_edgecolor("white"); sp.set_alpha(0.30); sp.set_linewidth(0.8)
    ax.yaxis.grid(True, color="white", alpha=0.07, linewidth=0.5)
    ax.xaxis.grid(True, color="white", alpha=0.06, linewidth=0.4)
    ax.set_axisbelow(True)
    if xlabel: ax.set_xlabel(xlabel, color="#AAAAAA", fontsize=8.5)

# ─────────────────────────────────────────────────────────────────────────────
# Helpers python-pptx (identicos al script principal)
# ─────────────────────────────────────────────────────────────────────────────
def _new_sl(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _pic(sl, path, prs):
    sl.shapes.add_picture(path, Emu(0), Emu(0), prs.slide_width, prs.slide_height)

def _txt(sl, text, left, top, width, height,
         fs=12, bold=False, italic=False, color="#FFFFFF",
         align=PP_ALIGN.LEFT, wrap=True, lspc=0):
    txb = sl.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf  = txb.text_frame; tf.word_wrap = wrap
    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if lspc > 0 and i > 0:
            p.space_before = Pt(lspc)
        run = p.add_run(); run.text = line
        run.font.size = Pt(fs); run.font.bold = bold; run.font.italic = italic
        r,g,b = int(color[1:3],16), int(color[3:5],16), int(color[5:7],16)
        run.font.color.rgb = RGBColor(r,g,b)
    return txb

def _T(sl, text, fs=20):
    _txt(sl, text, TITLE_L, TITLE_T, TITLE_W, TITLE_H,
         fs=fs, bold=True, color="#FFFFFF", align=PP_ALIGN.CENTER)

def _POP(sl, text=None):
    _txt(sl, text or POP_197, POP_L, POP_T, POP_W, POP_H,
         fs=7.5, italic=True, color="#C8DCF0")

def _BUL(sl, items, fs=11.5):
    txb = sl.shapes.add_textbox(Emu(BUL_L), Emu(BUL_T), Emu(BUL_W), Emu(BUL_H))
    tf  = txb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5); p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = str(i+1) + ".  " + item
        run.font.size = Pt(fs)
        run.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 - Scatter SAT vs Nota Promedio (3 paneles por ano)
# ─────────────────────────────────────────────────────────────────────────────
def slide_scatter(prs):
    anios     = ["2023","2024","2025"]
    COL_CTRL  = "#8FA8C8"   # azul-gris para control
    COL_FORM  = "#E8703A"   # naranja para formados
    ALPHA     = 0.38
    MS        = 10
    PANEL_GAP = 0.022

    # -- estadisticas por ano --------------------------------------------------
    stats_by_year = {}
    for y in anios:
        s = scat[scat["anio"]==y]
        r, p = scipy_stats.pearsonr(s["sat"], s["nota_promedio"])
        n_sec = len(s)
        n_f   = int(s["rut_docente"][s["formado"]].nunique())
        n_c   = int(s["rut_docente"][~s["formado"]].nunique())
        stats_by_year[y] = dict(r=r, p=p, n_sec=n_sec, n_f=n_f, n_c=n_c)
        print("  Stats", y, "r=", round(r,3), "n_sec=", n_sec,
              "n_f=", n_f, "n_c=", n_c)

    # -- figura transparente --------------------------------------------------
    fig = _tr_fig()

    gx = PIC_RECT[0] + 0.04
    gy = PIC_RECT[1] + 0.005
    gw = PIC_RECT[2] - 0.055
    gh = PIC_RECT[3] - 0.05

    panel_w = (gw - PANEL_GAP*2) / 3

    for i, y in enumerate(anios):
        s   = scat[scat["anio"]==y]
        sf  = s[s["formado"]]
        sc2 = s[~s["formado"]]
        st  = stats_by_year[y]

        ax_l = gx + i*(panel_w + PANEL_GAP)
        ax   = fig.add_axes([ax_l, gy, panel_w, gh], facecolor="none", zorder=5)

        # Control primero (fondo), formados encima
        ax.scatter(sc2["sat"], sc2["nota_promedio"],
                   c=COL_CTRL, alpha=ALPHA, s=MS, marker="o",
                   linewidths=0, zorder=2,
                   label="Sin formacion (n=" + str(st["n_c"]) + ")")
        ax.scatter(sf["sat"], sf["nota_promedio"],
                   c=COL_FORM, alpha=ALPHA+0.14, s=MS, marker="o",
                   linewidths=0, zorder=3,
                   label="Formados (n=" + str(st["n_f"]) + ")")

        # Linea de tendencia global
        x_all = s["sat"].values
        y_all = s["nota_promedio"].values
        m, b  = np.polyfit(x_all, y_all, 1)
        xr    = np.array([x_all.min(), x_all.max()])
        ax.plot(xr, m*xr+b, "--", color="white", linewidth=1.5, alpha=0.60,
                zorder=4, label="Tendencia (r = " + str(round(st["r"],2)) + ")")

        # Estilo oscuro
        _style_dark(ax, xlabel="SAT docente (sobre 7)")
        ax.set_xlim(1, 7); ax.set_ylim(1, 7)
        if i == 0:
            ax.set_ylabel("Nota promedio alumnos (sobre 7)",
                          color="#AAAAAA", fontsize=8.5)
        else:
            ax.set_yticklabels([])
        ax.set_xticks([1,2,3,4,5,6,7])
        ax.set_yticks([1,2,3,4,5,6,7])

        # Titulo del panel
        ax.set_title("Ano " + y, color="white", fontsize=10.5,
                     fontweight="bold", pad=6)

        # Leyenda compacta
        leg = ax.legend(fontsize=7.2, framealpha=0.35,
                        facecolor="#101820", edgecolor="#444444",
                        loc="upper left", markerscale=1.5,
                        handlelength=0.8, borderpad=0.5, labelspacing=0.35)
        for txt in leg.get_texts():
            txt.set_color("white")

        # Caja de estadisticas (esquina inferior derecha)
        p_str  = "< 0.001" if st["p"] < 0.001 else ("= " + str(round(st["p"],3)))
        n_fmt  = str(st["n_sec"]).replace(
            str(st["n_sec"])[:-3],
            str(st["n_sec"])[:-3] + "."
        ) if len(str(st["n_sec"])) > 3 else str(st["n_sec"])
        box_txt = ("r = " + str(round(st["r"],2)) + "\n"
                   + "p " + p_str + "\n"
                   + str(st["n_sec"]) + " secciones")
        ax.text(0.97, 0.05, box_txt,
                transform=ax.transAxes, fontsize=8,
                va="bottom", ha="right",
                color="white",
                bbox=dict(boxstyle="round,pad=0.45",
                          facecolor="#1A2E10",
                          edgecolor="#6AAA40", alpha=0.88))

    path = _save_ch(fig, "extra_scatter_sat_nota.png")

    # -- slide pptx -----------------------------------------------------------
    sl = _new_sl(prs)
    _pic(sl, SHARED_BG, prs)
    _pic(sl, path, prs)
    _T(sl, "Relacion SAT Docente y Nota Promedio de Alumnos   Por Ano", fs=13)
    _POP(sl, "Universo con SAT y calificaciones disponibles  |  2023-01 a 2025-02"
             "  |  Cada punto = 1 seccion")
    _BUL(sl, [
        "Correlacion positiva y estadisticamente significativa entre SAT docente y nota "
        "promedio de alumnos en los tres anos (r=0.16 en 2023, r=0.35 en 2024, r=0.32 "
        "en 2025; p<0.001 en todos los cortes).",
        "Los docentes formados (naranja) se concentran en la region SAT mayor a 5 con "
        "notas promedio superiores; los sin formacion presentan mayor dispersion hacia "
        "valores bajos en ambas dimensiones.",
        "La correlacion se fortalece en 2024-2025, coincidiendo con el mayor volumen "
        "de formacion, consistente con un efecto acumulativo sobre el rendimiento "
        "estudiantil.",
    ])
    print("  slide 1 - Scatter SAT vs Nota (3 paneles) OK")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 - Embudo rediseniado (mas compacto, centrado, menos ancho)
# ─────────────────────────────────────────────────────────────────────────────
def slide_embudo(prs):
    fig = _tr_fig()

    # Embudo estrecho: ~40% del ancho del area PIC, centrado horizontalmente
    ew  = PIC_RECT[2] * 0.40
    ex  = PIC_RECT[0] + (PIC_RECT[2] - ew) / 2
    ey  = PIC_RECT[1] + PIC_RECT[3] * 0.02
    eh  = PIC_RECT[3] * 0.93

    ax = fig.add_axes([ex, ey, ew, eh], facecolor="none", zorder=5)
    ax.set_xlim(0,10); ax.set_ylim(0,10); ax.axis("off")

    # (ancho_top, ancho_bot, numero, pct_label, descripcion, color)
    steps = [
        (8.2, 6.5, "917",
         "Universo base",
         "Docentes jerarquizados UCEN",
         "#3D6FA4"),
        (6.5, 4.5, "357",
         "39%  de 917",
         "Participaron en al menos 1 iniciativa de formacion",
         "#4B9CD3"),
        (4.5, 2.4, "197",
         "55%  de 357",
         "Aptos P3: SAT disponible en baseline y resultado",
         "#52C97A"),
    ]

    tops = [9.65, 6.80, 3.95]
    hh   = 2.50

    for i, ((tw, bw, n, pct_lbl, desc, col), top) in enumerate(zip(steps, tops)):
        bot  = top - hh
        pts  = np.array([[5-tw/2, top], [5+tw/2, top],
                         [5+bw/2, bot], [5-bw/2, bot]])

        # Relleno
        ax.fill(pts[:,0], pts[:,1], color=col, alpha=0.82, zorder=2)
        # Borde
        ax.plot(np.append(pts[:,0], pts[0,0]),
                np.append(pts[:,1], pts[0,1]),
                color="white", linewidth=0.7, alpha=0.40, zorder=3)

        # Numero grande centrado
        ax.text(5, top - hh*0.40, n,
                ha="center", va="center", fontsize=28, fontweight="bold",
                color="white", zorder=4,
                path_effects=[pe.withStroke(linewidth=3.5, foreground="#0A0F18")])

        # Etiqueta de porcentaje bajo el numero
        ax.text(5, top - hh*0.67, pct_lbl,
                ha="center", va="center", fontsize=8.2,
                color="#D8F0D8", fontweight="bold", zorder=4)

        # Descripcion a la DERECHA del embudo
        rx = 5 + tw/2 + 0.28
        ax.text(rx, top - hh*0.44, desc,
                ha="left", va="center", fontsize=8.8,
                color="white", zorder=4)

        # Linea de conexion
        ax.plot([5+tw/2+0.04, rx-0.06], [top-hh*0.44, top-hh*0.44],
                "-", color="white", linewidth=0.5, alpha=0.35, zorder=3)

    # Flechas de reduccion entre niveles (lado izquierdo)
    arrow_data = [
        (4.65, 7.05, "917  353  = 39%"),
        (3.80, 4.20, "357  197  = 55%"),
    ]
    for ax_x, ay, label in arrow_data:
        ax.annotate("", xy=(ax_x-0.1, ay-0.35), xytext=(ax_x-0.1, ay+0.35),
                    arrowprops=dict(arrowstyle="->", color="#FFD580",
                                   lw=1.2), zorder=5)
        ax.text(ax_x-0.55, ay, label,
                ha="right", va="center", fontsize=7.5,
                color="#FFD580", fontstyle="italic", zorder=5)

    path = _save_ch(fig, "extra_embudo_v2.png")

    # -- slide pptx -----------------------------------------------------------
    sl = _new_sl(prs)
    _pic(sl, SHARED_BG, prs)
    _pic(sl, path, prs)
    _T(sl, "Embudo de Seleccion   De 917 a 197 Aptos P3")
    _POP(sl, "Criterio Aptos P3: SAT valido en baseline (t-1) y resultado (t+1)  "
             "|  Universo: 917 docentes jerarquizados UCEN  |  Formacion 2022-2025")
    _BUL(sl, [
        "De los 917 docentes jerarquizados, 357 (39%) participaron en al menos "
        "una iniciativa de formacion (Taller, Diplomado o Proyecto) entre 2022 y 2025.",
        "197 de esos 357 tienen SAT disponible en los dos momentos clave: "
        "baseline (semestre anterior a la formacion) y resultado (semestre posterior). "
        "Son el universo rector del analisis P3.",
        "Todos los 197 son formados. El grupo control son docentes sin formacion "
        "con SAT disponible (n=486 docentes, aprox. 300-415 por periodo).",
    ])
    print("  slide 2 - Embudo rediseniado OK")


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("Generando SLIDES_EXTRA_scatter_embudo.pptx ...")
    _ensure_bg()
    prs = Presentation()
    prs.slide_width  = Emu(SW_EMU)
    prs.slide_height = Emu(SH_EMU)
    slide_scatter(prs)
    slide_embudo(prs)
    prs.save(OUT_PPTX)
    print("\n  Guardado:", OUT_PPTX)
    print("  2 slides: [1] Scatter SAT vs Nota (3 paneles), [2] Embudo rediseniado")

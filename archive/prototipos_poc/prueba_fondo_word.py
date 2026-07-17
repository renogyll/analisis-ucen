"""
prueba_fondo_word.py — Prueba de concepto: 3 fondos para gráficos en informe Word
Genera una PPTX con 3 slides mostrando el mismo gráfico en cada color de fondo propuesto.
"""
import sys; sys.stdout.reconfigure(encoding="utf-8")
import os, io
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE   = os.path.dirname(os.path.abspath(__file__))
SLIDES = os.path.join(BASE, "dark_slides_v3")
ROOT   = os.path.normpath(os.path.join(BASE, "..", ".."))

CHART_FILE = "22_chart.png"   # SAT z-score 6 periodos — buen ejemplo visual

FONDOS = [
    ("#54668E", ( 84, 102, 142), "Opción A — #54668E  (azul-gris medio)"),
    ("#363955", ( 54,  57,  85), "Opción B — #363955  (índigo oscuro)"),
    ("#879EC6", (135, 158, 198), "Opción C — #879EC6  (azul-lavanda claro)"),
    ("#59598E", ( 89,  89, 142), "Opción D — #59598E  (violeta-azul medio)"),
]

SCRATCH = (r"C:\Users\RGONZA~1.LAP\AppData\Local\Temp\claude"
           r"\c--Users-r-gonzalez-fluxsolar-LAPTOP-FLUX-ECO-Downloads-Analisis-UCEN-v2"
           r"\19e6fc3f-6ca1-4150-9da7-8dfa38be71ca\scratchpad")
os.makedirs(SCRATCH, exist_ok=True)

SW = 12192000   # EMU ancho slide
SH =  6858000   # EMU alto slide

# ─── Recortar márgenes transparentes del PNG del gráfico ──────────────────────
def _autocrop(img_rgba):
    """Elimina el exceso de espacio transparente alrededor del gráfico."""
    r, g, b, a = img_rgba.split()
    bbox = a.getbbox()   # bounding box del área opaca
    if bbox:
        pad = 30          # píxeles de margen proporcional en todos los lados
        W, H = img_rgba.size
        l = max(0, bbox[0] - pad)
        t = max(0, bbox[1] - pad)
        r2 = min(W, bbox[2] + pad)
        b2 = min(H, bbox[3] + pad)
        return img_rgba.crop((l, t, r2, b2))
    return img_rgba

# ─── Componer gráfico sobre fondo sólido ──────────────────────────────────────
def _compose(chart_rgba_cropped, rgb_fondo, out_size=(1920, 1080)):
    bg = PILImage.new("RGB", out_size, rgb_fondo)
    ch = chart_rgba_cropped.resize(out_size, PILImage.LANCZOS)
    bg.paste(ch, (0, 0), ch)
    return bg

# ─── Construcción de la PPTX de prueba ────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Emu(SW)
prs.slide_height = Emu(SH)

chart_path = os.path.join(SLIDES, CHART_FILE)
if not os.path.exists(chart_path):
    print(f"ERROR: no se encontró {chart_path}")
    sys.exit(1)

# Cargar y recortar el chart una sola vez
chart_orig = PILImage.open(chart_path).convert("RGBA")
chart_crop = _autocrop(chart_orig)
print(f"Chart original: {chart_orig.size}  →  recortado: {chart_crop.size}")

for hex_color, rgb, label in FONDOS:
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

    # Fondo sólido del slide
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(*rgb)

    # Etiqueta en parte superior
    tx = slide.shapes.add_textbox(Emu(400000), Emu(80000),
                                   Emu(SW - 800000), Emu(380000))
    tf = tx.text_frame
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = label
    run.font.size  = Pt(20)
    run.font.bold  = True
    run.font.color.rgb = RGBColor(0, 33, 71)

    # Componer imagen con el fondo correspondiente
    composed = _compose(chart_crop, rgb, out_size=(1920, 1080))
    tmp = os.path.join(SCRATCH, f"poc_{hex_color[1:]}.png")
    composed.save(tmp, "PNG")

    # Insertar imagen con margen proporcional uniforme (4% por lado)
    margin_x = int(SW * 0.035)
    margin_y = int(SH * 0.14)    # mayor margen arriba para la etiqueta
    img_w    = SW - 2 * margin_x
    img_h    = SH - margin_y - int(SH * 0.04)   # 4% de margen abajo

    slide.shapes.add_picture(tmp,
                              Emu(margin_x), Emu(margin_y),
                              Emu(img_w),    Emu(img_h))

out = os.path.join(ROOT, "PRUEBA_fondos_informe_2.pptx")
prs.save(out)
print(f"\n✓ Guardado: {out}")
print(f"  3 slides → fondos: {', '.join(h for h,_,_ in FONDOS)}")
